# -*- coding: utf-8 -*-
"""書式順守の最終監査"""
from pptx import Presentation
from pptx.oxml.ns import qn
from collections import Counter
prs = Presentation("../屋根ビジネス研究会_法則で学ぶ学習資料.pptx")
W, H = prs.slide_width, prs.slide_height
print("スライドサイズ:", W, "x", H, "→", round(W/914400,2), "x", round(H/914400,2), "inch",
      "/ 比率", round(W/H,4), "(4:3 =", round(4/3,4), ")")
allowed = {'002060','0070C0','FFC000','FF0000','C00000','000000','FFFFFF'}
colors, fonts, sizes = Counter(), Counter(), Counter()
oob, nofoot, nopage = [], [], []
for i, s in enumerate(prs.slides, 1):
    txts = []
    for sh in s.shapes:
        if sh.left is None: continue
        if sh.left < -1000 or sh.top < -1000 or sh.left+sh.width > W+1000 or sh.top+sh.height > H+1000:
            oob.append((i, sh.shape_type, sh.left, sh.top, sh.width, sh.height))
        try:
            if sh.fill.type == 1: colors[str(sh.fill.fore_color.rgb)] += 1
        except Exception: pass
        try:
            if sh.line.fill.type == 1: colors[str(sh.line.color.rgb)] += 1
        except Exception: pass
        if sh.has_text_frame:
            txts.append(sh.text_frame.text)
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size: sizes[r.font.size.pt] += 1
                    rPr = r._r.find(qn('a:rPr'))
                    if rPr is not None:
                        for tag in ('a:latin','a:ea'):
                            e = rPr.find(qn(tag))
                            if e is not None: fonts[e.get('typeface')] += 1
                    try:
                        if r.font.color.rgb: colors[str(r.font.color.rgb)] += 1
                    except Exception: pass
        if sh.has_table:
            for row in sh.table.rows:
                for c in row.cells:
                    for p in c.text_frame.paragraphs:
                        for r in p.runs:
                            if r.font.size: sizes[r.font.size.pt] += 1
                            rPr = r._r.find(qn('a:rPr'))
                            if rPr is not None:
                                for tag in ('a:latin','a:ea'):
                                    e = rPr.find(qn(tag))
                                    if e is not None: fonts[e.get('typeface')] += 1
                    try:
                        if c.fill.type == 1: colors[str(c.fill.fore_color.rgb)] += 1
                    except Exception: pass
    joined = "\n".join(txts)
    if "Copyright©2020 Funai Consulting Inc." not in joined: nofoot.append(i)
    if not any(t.strip() == str(i) for t in txts): nopage.append(i)

print("\nフォント:", dict(fonts))
print("使用色:", sorted(colors.keys()))
bad = [c for c in colors if c not in allowed]
print("パレット外の色:", bad if bad else "なし（規定5色＋白黒のみ）")
print("\n文字サイズ分布:", sorted(sizes.items(), key=lambda x:-x[1]))
print("\nスライド枠外にはみ出す図形:", oob if oob else "なし")
print("フッター欠落スライド:", nofoot if nofoot else "なし")
print("ページ番号欠落スライド:", nopage, "（P1表紙・P35まとめは意図的に番号なし）")
print("総スライド数:", len(prs.slides._sldIdLst))
