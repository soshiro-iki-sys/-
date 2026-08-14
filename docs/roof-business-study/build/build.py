# -*- coding: utf-8 -*-
"""屋根ビジネス研究会_法則で学ぶ学習資料.pptx ビルダー"""
from deck_lib import *
from content import *
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = new_deck()
P = 0
def nxt():
    global P
    P += 1
    return P

CH0 = "第0部 本資料の使い方"
CH1 = "第1部 全体像"
CH2 = "第2部 8つの法則"
CH3 = "第3部 現場別クイックリファレンス"
CH4 = "第4部 定着チェック"


# ══════════════════════════ 第0部 ══════════════════════════
# P1 表紙（章扉と同じ濃紺ベタ）
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 1.0, 2.35, 8.0, 0.035, fill=ORANGE)
textbox(s, 1.0, 1.35, 8.0, 0.5, [('Roof Reform-Manual', 16, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 1.0, 1.75, 8.0, 0.62, [('屋根ビジネス研究会　船井流マニュアル', 20, ORANGE, True)], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 1.0, 2.60, 8.0, 1.5, [('法則で学ぶ', 40, WHITE, True), ('学習資料', 40, WHITE, True)],
        anchor=MSO_ANCHOR.TOP, spacing=1.05)
textbox(s, 1.0, 4.35, 8.0, 1.2,
        [('全4編・102ページを　1つの根本原理 ＋ 8つの法則 に再構成', 14, WHITE, False),
         ('対象：新任営業マン・店長／想定学習時間：1〜2時間', 12, WHITE, False)], spacing=1.5)
textbox(s, 1.0, 6.55, 8.0, 0.3, [('Copyright©2020 Funai Consulting Inc. All rights reserved.', 10, WHITE, False)])
nxt()

# P2 本資料の使い方
s = base_slide(prs, CH0, "本資料の読み方 ─ これは要約ではなく「法則の地図」です", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.46,
      "覚えるのは8個ではありません。根本原理1つを覚え、そこから8法則を導いてください。",
      14, WHITE, True, fill=BLUE)
# 圧縮マップ
srcs = [("Ⅰ.事業コンセプト", "17頁"), ("Ⅲ.集客", "44頁"), ("Ⅲ.店舗", "14頁"), ("Ⅳ.営業", "27頁")]
for i, (nm, pg) in enumerate(srcs):
    y = 1.78 + i * 0.60
    label(s, MARGIN, y, 2.3, 0.48, nm + "　" + pg, 12, WHITE, True, fill=NAVY)
arrow(s, 2.72, 2.55, 0.55, 0.7, BLUE)
label(s, 3.40, 1.78, 2.75, 2.28,
      ["根本原理", "", "屋根は、客が", "一生見ない商品である"],
      15, WHITE, True, fill=DKRED, spacing=1.25)
arrow(s, 6.30, 2.55, 0.55, 0.7, BLUE)
label(s, 6.98, 1.78, 2.72, 2.28,
      ["8つの法則", "", "①一番化　②可視化", "③掛け算　④3フェーズ", "⑤階段　⑥標準化", "⑦ホーム　⑧先回り"],
      13, WHITE, True, fill=BLUE, spacing=1.2)
# 構成
label(s, MARGIN, 4.28, CW, 0.36, "本資料の構成", 12, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
table(s, MARGIN, 4.68, CW,
      [["部", "内容", "使い方"],
       ["第1部", "全体像 ─ 屋根ビジネスの1本の線", "最初に通しで読む"],
       ["第2部", "根本原理と8つの法則　★本資料の核", "暗記せず“導ける”ようにする"],
       ["第3部", "現場別クイックリファレンス", "実務中に“引く”ための頁"],
       ["第4部", "定着チェック20問＋CASE＋7日間ロードマップ", "1週間後に自己採点する"]],
      col_w=[1.1, 4.4, 3.3], row_h=0.33, size=11, first_col_fill=BLUE)
source(s, "Ⅰ.事業コンセプト／Ⅲ.集客／Ⅲ.店舗／Ⅳ.営業　全4編")


# ══════════════════════════ 第1部 ══════════════════════════
# P3 モデル比較
s = base_slide(prs, CH1, "船井流屋根ビジネスモデルとは ─ 他モデルとの決定的な違い", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.44,
      "“圧倒的一番化を目指すための、戸建て住宅向け元請け屋根リフォームモデル”",
      14, WHITE, True, fill=BLUE)
table(s, MARGIN, 1.70, CW,
      [["", "一般的リフォームモデル", "一般的塗装モデル", "船井流屋根ビジネスモデル"],
       ["業態", "総合リフォーム", "塗装のみの取り扱い", "屋根・外壁リフォーム専門特化"],
       ["商品", "各種オーダーメイド商品", "外壁塗装・屋根塗装", "屋根カバー・屋根葺き替え"],
       ["集客", "訪問販売やOB客メイン", "訪問販売やOB・紹介がメイン", "屋根専門チラシ・WEB・店舗"],
       ["店舗", "事務所", "事務所", "屋根リフォーム体感型SR"],
       ["営業", "属人的営業", "属人的営業", "標準化マニュアル営業"],
       ["収益", "粗利率30％前後", "粗利率35〜37％", "粗利率40％"],
       ["商圏", "商圏発想なし", "商圏発想なし", "1拠点50万人"]],
      col_w=[1.0, 2.6, 2.7, 3.1], row_h=0.36, size=11, first_col_fill=BLUE)
label(s, MARGIN, 4.55, CW, 0.42,
      "⇒ 他リフォーム会社・他塗装店と違い、“屋根”に特化して集客・営業を行う",
      13, WHITE, True, fill=NAVY)
label(s, MARGIN, 5.15, 4.6, 1.5,
      ["【事業コンセプト】", "地域の住宅の屋根を守り、", "安心して暮らせる住環境を提供する"],
      12, NAVY, True, fill=WHITE, line=NAVY, spacing=1.35)
label(s, 5.10, 5.15, 4.6, 1.5,
      ["【ライフサイクル】", "元請け屋根ビジネスは“成長期”。", "競合が乱立する前に地域一番店シェアを取る"],
      12, NAVY, True, fill=WHITE, line=NAVY, spacing=1.35)
source(s, "Ⅰ.事業コンセプト P.4,7,10")

# P4 数値の背骨
s = base_slide(prs, CH1, "数値の背骨 ─ 2.7億円はどこから来るのか", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.42,
      "同じ2.7億円を、市場側からと現場側から、2通りの掛け算で出せるようにする",
      13, WHITE, True, fill=BLUE)
# 上段：市場側
label(s, MARGIN, 1.70, 1.5, 0.34, "① 市場から", 11, WHITE, True, fill=NAVY)
seq1 = [("MS", "2,090円"), ("商圏人口", "50万人"), ("シェア率", "26％")]
x = MARGIN
for i, (k, v) in enumerate(seq1):
    label(s, x, 2.12, 1.95, 0.85, [(k, 11, NAVY, False), (v, 17, NAVY, True)], fill=WHITE, line=BLUE, spacing=1.15)
    x += 1.95
    if i < 2:
        label(s, x, 2.12, 0.42, 0.85, "×", 18, BLUE, True)
        x += 0.42
label(s, x, 2.12, 0.42, 0.85, "＝", 18, BLUE, True)
label(s, x + 0.42, 2.12, 2.24, 0.85, [("受注売上", 11, WHITE, False), ("2.7億円", 19, WHITE, True)],
      fill=ORANGE, line=None, spacing=1.15)
# 下段：現場側
label(s, MARGIN, 3.20, 1.5, 0.34, "② 現場から", 11, WHITE, True, fill=NAVY)
seq2 = [("現場調査数", "年250件"), ("契約率", "60％"), ("客単価", "180万円")]
x = MARGIN
for i, (k, v) in enumerate(seq2):
    label(s, x, 3.62, 1.95, 0.85, [(k, 11, NAVY, False), (v, 17, NAVY, True)], fill=WHITE, line=BLUE, spacing=1.15)
    x += 1.95
    if i < 2:
        label(s, x, 3.62, 0.42, 0.85, "×", 18, BLUE, True)
        x += 0.42
label(s, x, 3.62, 0.42, 0.85, "＝", 18, BLUE, True)
label(s, x + 0.42, 3.62, 2.24, 0.85, [("受注売上", 11, WHITE, False), ("2.7億円", 19, WHITE, True)],
      fill=ORANGE, line=None, spacing=1.15)
# 月次
label(s, MARGIN, 4.70, CW, 0.36, "月次に落とすと", 11, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
label(s, MARGIN, 5.12, CW, 0.52,
      "現場調査 21件　×　契約率 60％　×　客単価 180万円　＝　月間受注売上 2,250万円",
      15, NAVY, True, fill=WHITE, line=ORANGE)
label(s, MARGIN, 5.78, CW, 0.92,
      ["【契約率が60％と50％の2つ出てくる理由】",
       "事業計画のKPIは契約率60％。一方、営業マン個人の目標は現調7件×50％＝4件で設定されている。",
       "営業3名×7件＝現調21件、契約4件×3名＝12件 ≒ 事業計画12.5件。個人目標は安全側に置いている。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=NAVY, spacing=1.3)
source(s, "Ⅰ.事業コンセプト P.5,9,13／Ⅳ.営業 P.3")

# P5 全体フロー
s = base_slide(prs, CH1, "全体フロー ─ 集客から契約までの1本の線", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.40,
      "リフォーム事業のKPIは「新規の現場調査数」。全ての販促はここに集約される",
      13, WHITE, True, fill=BLUE)
# 集客5媒体
label(s, MARGIN, 1.66, 1.45, 0.34, "集客", 11, WHITE, True, fill=NAVY)
media = [("チラシ", "35％"), ("WEB", "35％"), ("イベント", "10％"), ("近隣挨拶", "15％"), ("その他", "5％")]
x = MARGIN
for nm, pc in media:
    label(s, x, 2.06, 1.80, 0.80, [(nm, 11, NAVY, True), (pc, 16, BLUE, True)], fill=WHITE, line=BLUE, spacing=1.15)
    x += 1.88
arrow(s, 4.75, 2.96, 0.5, 0.34, BLUE, 'down')
# フェーズ帯
ph = [("現場調査", "月21件", "劣化写真は最低30枚"),
      ("初回訪問", "ここで8割決まる", "7つの阻害要因を解消し仮クロージング"),
      ("クロージング", "契約率60％", "ショールームで見積提出・決断を迫る"),
      ("契約", "客単価180万円", "月間受注 2,250万円")]
y = 3.38
for i, (nm, kv, note) in enumerate(ph):
    fill = ORANGE if i == 1 else NAVY
    tcol = NAVY if i == 1 else WHITE
    label(s, MARGIN, y, 2.25, 0.62, nm, 14, tcol, True, fill=fill)
    label(s, 2.62, y, 1.85, 0.62, kv, 12, NAVY, True, fill=WHITE, line=BLUE)
    label(s, 4.60, y, 5.10, 0.62, note, 12, BLACK, False, PP_ALIGN.LEFT, fill=WHITE, line=BLUE)
    y += 0.72
label(s, MARGIN, 6.28, CW, 0.44,
      "初回訪問は“ただ実測に行く場”ではない。ここでの差が、そのまま契約率の差になる。",
      12, WHITE, True, fill=DKRED)
source(s, "Ⅲ.集客 P.2,4／Ⅳ.営業 P.3,11,23")

# P6 KPIツリー
s = base_slide(prs, CH1, "KPI因数分解ツリー ─ 不振の原因を1つに特定する", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.40,
      "「売上が足りない」で止めない。必ず因数まで割ってから打ち手を決める",
      13, WHITE, True, fill=BLUE)
label(s, 3.30, 1.66, 3.4, 0.52, "受注売上", 16, WHITE, True, fill=NAVY)
arrow(s, 4.80, 2.24, 0.4, 0.28, BLUE, 'down')
tri = [("現場調査数", ["チラシ：枚数×反響率1/10,000", "WEB：セッション×CVR", "イベント：年3〜4回", "近隣挨拶：やり切れるか"]),
       ("契約率", ["7つの阻害要因の解消度", "初回訪問での仮クロージング", "クロージングは来店で行う", "地雷3種を仕掛けたか"]),
       ("客単価", ["カバー・葺き替え220万円", "小工事50万円", "屋根材5種以上の品揃え", "安い順に紹介し上へ誘導"])]
x = MARGIN
for nm, items in tri:
    label(s, x, 2.60, 3.02, 0.46, nm, 14, WHITE, True, fill=BLUE)
    label(s, x, 3.12, 3.02, 2.05, bullets(items), 11, BLACK, False,
          PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.5)
    x += 3.19
label(s, MARGIN, 5.35, CW, 0.36, "この分解ができると、打ち手はこう変わる", 12, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
table(s, MARGIN, 5.75, CW,
      [["低い因数", "やってはいけない指示", "本来の打ち手"],
       ["現場調査数", "「もっと訪問しろ」", "エリア別反響率を出し、悪いエリアを削る"],
       ["契約率", "「気合いを入れろ」", "阻害要因のどれで落ちたかを案件ごとに記録する"],
       ["客単価", "「高いものを売れ」", "屋根材の品揃えと施工事例の価格帯を増やす"]],
      col_w=[1.9, 3.0, 4.5], row_h=0.31, size=11, first_col_fill=BLUE)
source(s, "Ⅰ.事業コンセプト P.5,13／Ⅲ.集客 P.11,12／Ⅳ.営業 P.19,27")


# ══════════════════════════ 第2部 ══════════════════════════
# P7 根本原理
s = base_slide(prs, CH2, "根本原理 ─ すべてはこの一点から始まる", nxt(), 20)
label(s, MARGIN, 1.20, CW, 1.05, ["屋根は、", "客が一生見ない商品である"], 26, WHITE, True, fill=DKRED, spacing=1.2)
label(s, MARGIN, 2.42, CW, 1.25,
      ["「お客様は普段屋根の劣化や破損を目にすることはなく、メンテナンスしなければ",
       "　いけないということを認識していない」　　　　　　　　　　　　（Ⅳ.営業マニュアル P.2）",
       "「屋根は普段から目にする機会が少なく、専門的な知識が多いため、",
       "　一般的な方が知識を持っていない商材です」　　　　　　　　　　（Ⅲ.店舗マニュアル P.10）"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=NAVY, spacing=1.45)
label(s, MARGIN, 3.85, CW, 0.38, "この1点から、屋根ビジネスの“おかしな常識”がすべて説明できる", 12, WHITE, True, fill=NAVY)
cons = [("客は必要性を\n認識できない", "だから、こちらから\n見せに行く必要がある"),
        ("客は判断基準を\n持てない", "だから、基準を\nこちらから与える"),
        ("客は不安になる\n（高額×初体験）", "だから、一段ずつ\n安心を積み上げる"),
        ("市場全体が\n未認識ニーズの塊", "だから、確率と仕組みで\n刈り取れる")]
x = MARGIN
for a, b in cons:
    label(s, x, 4.32, 2.24, 0.78, a.split('\n'), 12, WHITE, True, fill=BLUE, spacing=1.25)
    arrow(s, x + 0.97, 5.16, 0.3, 0.26, NAVY, 'down')
    label(s, x, 5.50, 2.24, 0.80, b.split('\n'), 11, NAVY, False, fill=WHITE, line=BLUE, spacing=1.3)
    x += 2.36
label(s, MARGIN, 6.42, CW, 0.34,
      "＝ 値引きではなく【価格/価値】を上げる。これが屋根セールスの本体である。",
      12, WHITE, True, fill=ORANGE)
source(s, "Ⅳ.営業 P.2／Ⅲ.店舗 P.10")

# P8 導出ツリー
s = base_slide(prs, CH2, "導出ツリー ─ 8法則は暗記せず「導く」", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.38, "根本原理から枝を辿れば、8つの法則は自分で再現できる", 13, WHITE, True, fill=BLUE)
label(s, 2.55, 1.62, 4.9, 0.50, "屋根は、客が一生見ない商品である", 15, WHITE, True, fill=DKRED)
arrow(s, 4.80, 2.18, 0.4, 0.26, NAVY, 'down')
branch = [("必要性を\n認識できない", ["法則2 可視化"], BLUE),
          ("判断基準を\n持てない", ["法則1 一番化", "法則8 先回り"], BLUE),
          ("不安になる", ["法則5 階段設計", "法則7 ホームで戦う"], BLUE),
          ("市場が未認識\nニーズの塊", ["法則3 掛け算分解", "法則4 3フェーズ"], BLUE)]
x = MARGIN
for nm, laws, col in branch:
    label(s, x, 2.52, 2.24, 0.66, nm.split('\n'), 12, WHITE, True, fill=NAVY, spacing=1.2)
    y = 3.30
    for lw in laws:
        label(s, x, y, 2.24, 0.46, lw, 12, WHITE, True, fill=col)
        y += 0.54
    x += 2.36
label(s, 2.55, 4.60, 4.9, 0.46, "説明が長く複雑になる", 13, WHITE, True, fill=NAVY)
arrow(s, 4.80, 5.12, 0.4, 0.26, NAVY, 'down')
label(s, 2.55, 5.44, 4.9, 0.46, "法則6 標準化", 13, WHITE, True, fill=BLUE)
label(s, MARGIN, 6.10, CW, 0.60,
      ["8法則はすべて4冊すべてに現れることを確認済み（＝1冊だけの記述は法則に昇格させていない）。",
       "「未来への投資」「業態発想」「やり切り」は2冊以下だったため、他の法則に統合した。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=NAVY, spacing=1.35)
source(s, "全4編の横断検証による")

# P9〜P16 法則1〜8
for law in LAWS:
    s = base_slide(prs, CH2, "法則%d　%s ─ %s" % (law['no'], law['name'], law['sub']), nxt(), 20)
    label(s, MARGIN, 1.12, CW, 0.42, law['define'], 14, WHITE, True, fill=BLUE)
    # 4象限
    keys = ["事業", "集客", "店舗", "営業"]
    pos = [(MARGIN, 1.68), (5.10, 1.68), (MARGIN, 3.30), (5.10, 3.30)]
    for k, (bx, by) in zip(keys, pos):
        rect(s, bx, by, 4.60, 1.52, fill=None, line=BLUE)
        label(s, bx, by, 4.60, 0.32, k, 11, WHITE, True, fill=BLUE)
        textbox(s, bx + 0.06, by + 0.36, 4.48, 1.12, bullets(law['cells'][k]),
                10, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, spacing=1.35)
    # 下段3ボックス
    # NG・明日の行動は1文なので結合し、自然折り返しに任せる
    trio = [("数値・基準", law['nums'], ORANGE, NAVY),
            ("やりがちなNG", ["".join(law['ng'])], DKRED, WHITE),
            ("明日の行動", ["".join(law['act'])], NAVY, WHITE)]
    bx = MARGIN
    for ttl, items, col, tcol in trio:
        rect(s, bx, 4.94, 3.05, 1.32, fill=None, line=col)
        label(s, bx, 4.94, 3.05, 0.30, ttl, 10, tcol, True, fill=col)
        textbox(s, bx + 0.06, 5.28, 2.93, 0.94, items, 10, BLACK, False,
                PP_ALIGN.LEFT, MSO_ANCHOR.TOP, spacing=1.3)
        bx += 3.18
    source(s, law['src'])

# P17 7つの阻害要因
s = base_slide(prs, CH2, "法則8の実装 ─ 7つの阻害要因と、その打ち手", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.40,
      "初回訪問でこの7つを解消できれば仮クロージングに進める。1つ残れば契約は落ちる",
      13, WHITE, True, fill=BLUE)
rows = [["阻害要因", "顧客の中で起きていること", "原典に書かれた打ち手"]]
for a, b, c in OBSTACLES:
    rows.append([a.replace('\n', ''), b, c])
table(s, MARGIN, 1.66, CW, rows, col_w=[1.55, 3.0, 4.85], row_h=0.62, size=10,
      header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 6.10, CW, 0.58,
      ["アプローチブックを“順に読み上げる”のは、この7つを漏れなく潰すための設計である。",
       "渡すだけでは読まれず、要因が残ったまま商談が終わる。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=DKRED, spacing=1.35)
source(s, "Ⅳ.営業 P.3,10（7要因の内訳は貴社ご提供資料による）")


# ══════════════════════════ 第3部 ══════════════════════════
# P18 集客の全体像
s = base_slide(prs, CH3, "集客① 全体像 ─ 入口戦略と受け皿戦略に分けて考える", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.40,
      "集客は「認知を取る入口」と「問い合わせを受け止める受け皿」の2層でできている",
      13, WHITE, True, fill=BLUE)
label(s, MARGIN, 1.66, 4.60, 0.36, "入口戦略（認知を取る）", 12, WHITE, True, fill=NAVY)
label(s, MARGIN, 2.06, 4.60, 1.72,
      bullets(["WEB：PPC広告・自然検索・その他流入", "アナログ：チラシ・新聞広告・近隣営業",
               "ニュースレター・DM（OB・失注客の掘り起こし）", "認知獲得：野立て看板・バス広告・CM"]),
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.55)
label(s, 5.10, 1.66, 4.60, 0.36, "受け皿戦略（問い合わせに変える）", 12, WHITE, True, fill=NAVY)
label(s, 5.10, 2.06, 4.60, 1.72,
      bullets(["ホームページ（導線・更新・メールフォーム）", "ショールーム店舗（自然来店・来店予約）",
               "足場近隣・OB・過去客・失注客", "＝ 法則7「ホームで戦う」の入口"]),
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.55)
label(s, MARGIN, 3.92, CW, 0.36, "媒体比率と、月間現調21件の内訳", 12, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
table(s, MARGIN, 4.32, CW,
      [["媒体", "比率", "月間現調の目安", "反響の基準値"],
       ["チラシ", "35％", "約7件", "配布枚数×1/10,000"],
       ["WEB", "35％", "約7件", "CV＝セッション×CVR"],
       ["イベント", "10％", "約2件", "イベントチラシ×1/7,000"],
       ["近隣挨拶", "15％", "約3件", "やり切れるかが最大の分岐点"],
       ["その他", "5％", "約1件", "タウンメール×1/100・DM×10％"]],
      col_w=[1.7, 1.0, 2.0, 4.7], row_h=0.36, size=11, first_col_fill=BLUE)
label(s, MARGIN, 6.55, CW, 0.30,
      "※ リフォーム事業のKPIは新規の現場調査数。販促に意欲的に取り組むことが前提である",
      10, NAVY, False, PP_ALIGN.LEFT)
source(s, "Ⅲ.集客 P.2,3,4,34")

# P19 チラシ3フェーズ
s = base_slide(prs, CH3, "集客② チラシ ─ 作る／当てる／回すの3フェーズ", nxt(), 20)
phase_tabs(s, ["１．作成フェーズ", "２．エリア選定フェーズ", "３．運用フェーズ"], -1)
cols = [("１．作成", ["見出しは上段中央に大きく", "「なんのお店か」を一目で", "工事リストを最上部に記載",
                    "左上にSR写真など視覚訴求", "最下限金額を出し裾を広げる", "電話番号・QRは大きく目立たせる"]),
        ("２．エリア選定", ["一次商圏＝車で5〜10分圏", "一次商圏は月2〜3回配布", "二次商圏は月1〜2回配布",
                        "戸建比率が高い・OBが多いエリアも", "　複数回配布してよい", "折込表は指定枚数×0.85で記載"]),
        ("３．運用", ["エリア別に反響件数・枚数を入力", "エリアごとの反響率を抽出する", "悪いエリアは回数減または中止",
                    "良いエリアの枚数を増やす", "反響率が維持できるなら月4回も可", "繁忙期3〜5月・9〜11月は増やす"])]
x = MARGIN
for nm, items in cols:
    label(s, x, 1.56, 3.05, 0.42, nm, 13, WHITE, True, fill=NAVY)
    label(s, x, 2.04, 3.05, 2.30, bullets(items), 11, BLACK, False,
          PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.5)
    x += 3.18
label(s, MARGIN, 4.55, CW, 0.36, "紙面レイアウトの原則", 12, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
label(s, MARGIN, 4.95, 4.60, 1.45,
      ["【表面】", "上段中央：何屋かがわかる見出し",
       "最上部：自社の工事リスト", "左上：ショールーム写真・資格・集合写真", "下部：最下限の商品パック／電話番号・QR"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.4)
label(s, 5.10, 4.95, 4.60, 1.45,
      ["【裏面・Zの法則】", "人は 左上→右上→左下→右下 の順に見る",
       "→ 一番右下にCV導線（店舗MAP・無料診断）", "必ず入れる：施工事例／スタッフ顔出し"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.4)
label(s, MARGIN, 6.48, CW, 0.32, "【図：チラシ表面・裏面サンプル ※原典 Ⅲ.集客 P.8,9 の画像を差し込む】",
      10, BLUE, False)
source(s, "Ⅲ.集客 P.5,6,7,10,11,40")

# P20 WEB3フェーズ
s = base_slide(prs, CH3, "集客③ WEB ─ CV＝セッション×CVR に尽きる", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.44,
      "反響数(CV) ＝ セッション数（流入） × コンバージョン率（反響率）",
      15, WHITE, True, fill=ORANGE)
cols = [("１．HP作成", ["見出しで「何屋か」がわかる", "自社のUSP（独自の強み）を掲載",
                     "導線誘導ページ(LP)／更新ページ／", "　メールフォームページの3種を持つ",
                     "ユーザーの15の疑問に答える", "　コンテンツを揃える"]),
        ("２．広告配信", ["Google広告とYahoo広告を作成", "配信エリアごとにキャンペーン",
                      "1広告グループのKWは15個以下", "広告ランク＝入札価格×広告の品質",
                      "→ 品質を上げれば費用は下がる", "狙うのは「顕在×緊急度高」の層"]),
        ("３．運用", ["2週間ごとに数値入力し反省と対策", "予算通り・目標CPA以下なら月1回に",
                    "SEOはロングテールKWで数を打つ", "EAT基準（専門性・権威性・信頼性）",
                    "Search Consoleで表示回数が多く", "　順位が低いKWから記事化する"])]
x = MARGIN
for nm, items in cols:
    label(s, x, 1.70, 3.05, 0.42, nm, 13, WHITE, True, fill=NAVY)
    label(s, x, 2.18, 3.05, 2.32, bullets(items), 11, BLACK, False,
          PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.5)
    x += 3.18
label(s, MARGIN, 4.70, 4.60, 0.36, "失敗パターン", 12, WHITE, True, fill=DKRED)
label(s, MARGIN, 5.10, 4.60, 1.35,
      ["セッションはあるのに反響がない", "＝ 流入顧客層が不適正",
       "「とりあえず安いところ」「一括見積り」", "「A社で検討中」ばかりが流入している"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=DKRED, spacing=1.45)
label(s, 5.10, 4.70, 4.60, 0.36, "成功パターン", 12, WHITE, True, fill=BLUE)
label(s, 5.10, 5.10, 4.60, 1.35,
      ["流入顧客属性の最適化", "「屋根のことを真剣に考えている」",
       "「優良施工業者を探している」層を集める", "同業他商圏の事例・入札KW・検索KWを使う"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.45)
label(s, MARGIN, 6.55, CW, 0.30,
      "※ 目標から逆算する：売上目標→必要契約数→必要現調数→媒体別集客数→CV数→セッション数×CVR",
      10, NAVY, False, PP_ALIGN.LEFT)
source(s, "Ⅲ.集客 P.12〜32")

# P21 イベント・近隣
s = base_slide(prs, CH3, "集客④ イベントと近隣集客 ─ 2段階オープンとやり切り", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.40,
      "プレオープンは“オープン前でしかできない販促”。後の販促効率UPに繋がる",
      13, WHITE, True, fill=BLUE)
table(s, MARGIN, 1.62, CW,
      [["段階", "目的", "主な準備物", "目標"],
       ["①ドアコール\n＋OB向けDM", "開店の挨拶と\n強力告知", "ドアコール計画表／ご招待状／挨拶文\n鉛筆くじ／不在宅用メッセージ／DM",
        "近隣約3,000件を訪問\nOBは全員に配布"],
       ["②プレオープン\n（特別ご招待会）", "来場者を\nDMリスト化", "鉛筆くじ景品／当日来場特典\nお子様向け特典",
        "ドアコール件数×1％の来場\nDMは配布数の3〜5％"],
       ["③グランド\nオープン", "新規顧客へ\nインパクト", "グランドオープンチラシ／抽選会景品\n来場・見積・成約特典",
        "チラシ配布枚数×1/7,000\n総来場数×10〜20％が現調"],
       ["④通常販促\n（安定）", "店舗認知と\n来店・現調獲得", "定番チラシの作成・折込\nOB向けDM",
        "チラシ×1/10,000\nタウンメール×1/100"]],
      col_w=[1.9, 1.6, 3.6, 2.7], row_h=0.72, size=10, header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 4.62, 4.60, 0.36, "イベント運営の要点", 12, WHITE, True, fill=NAVY)
label(s, MARGIN, 5.02, 4.60, 1.62,
      bullets(["イベントは年3〜4回（多いと飽きられる）", "チラシはB3サイズ・表面に商品は載せない",
               "特典は来店→見積→成約の3段構成", "最低3人（駐車場1・受付1・商談1）",
               "来場者はA/B/Cランクで首掛け名札を色分け"]),
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.45)
label(s, 5.10, 4.62, 4.60, 0.36, "近隣集客（ドローン活用）", 12, WHITE, True, fill=NAVY)
label(s, 5.10, 5.02, 4.60, 1.62,
      bullets(["あくまで近隣挨拶。訪問販売とは違う", "ドローンは思いっきり目立たせる（ベスト・ノボリ）",
               "「よろしければ○○さんの屋根でも飛ばしますよ」", "地図上で挨拶済み・温度感を可視化する",
               "最大のポイントは、やり切れるかどうか"]),
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.45)
source(s, "Ⅲ.集客 P.33〜44")

# P22 店舗：立地
s = base_slide(prs, CH3, "店舗① 立地 ─ 反響の“係数”は契約前に決まっている", nxt(), 20)
phase_tabs(s, ["１．立地", "２．外観", "３．内観"], 0)
label(s, MARGIN, 1.56, CW, 0.40,
      "立地は幹線道路沿いの認知性の高い好立地。地域査定と物件査定の2段階で見る",
      13, WHITE, True, fill=BLUE)
label(s, MARGIN, 2.06, 4.60, 0.34, "地域査定要素", 12, WHITE, True, fill=NAVY)
table(s, MARGIN, 2.42, 4.60,
      [["項目", "基準"],
       ["交通量", "昼間12時間交通量 8,000〜12,000台"],
       ["視認性", "200m手前からお店（看板）が見える"],
       ["住宅密度", "自動車20分圏で6万世帯"],
       ["持家比率", "持家比率60％以上"],
       ["所得指数", "所得指数1.0以上"],
       ["人口増減", "人口・世帯数が増加している"],
       ["道路状況", "中央分離帯がない（理想）"],
       ["近隣施設", "スーパー等の商業施設周辺"]],
      col_w=[1.4, 3.2], row_h=0.335, size=10, header_size=10, first_col_fill=BLUE)
label(s, 5.10, 2.06, 4.60, 0.34, "物件査定要素", 12, WHITE, True, fill=NAVY)
table(s, 5.10, 2.42, 4.60,
      [["項目", "基準"],
       ["看板", "店頭に縦幅1m以上設置可能"],
       ["", "道路から見える自立看板がある"],
       ["過去履歴", "過去履歴の出入頻度が少ない"],
       ["面積", "店舗面積10坪以上"],
       ["駐車場", "駐車場3台以上"],
       ["間口", "間口に8m以上取れる"],
       ["設備", "空調、水廻りが整っている"],
       ["家賃", "家賃が月20万円前後で借りられる"]],
      col_w=[1.4, 3.2], row_h=0.335, size=10, header_size=10, first_col_fill=BLUE)
label(s, MARGIN, 5.42, CW, 0.36, "来店型店舗の3つのメリット（＝なぜ店舗を持つのか）", 12, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
trio = [("１. 信頼度UP", "店舗を有することで\n“逃げない企業”とアピールできる"),
        ("２. 集客数UP", "認知性向上による自然来店客の増加\nチラシ・WEB反響の増加"),
        ("３. 契約率UP", "目的来店客の増加による契約率向上\n店舗営業のマニュアル化")]
x = MARGIN
for nm, txt in trio:
    label(s, x, 5.82, 3.05, 0.34, nm, 11, WHITE, True, fill=BLUE)
    label(s, x, 6.18, 3.05, 0.62, txt.split('\n'), 10, BLACK, False,
          PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.3)
    x += 3.18
source(s, "Ⅲ.店舗 P.2,3")

# P23 店舗：外観
s = base_slide(prs, CH3, "店舗② 外観 ─ 走行中の車から3秒で伝える", nxt(), 20)
phase_tabs(s, ["１．立地", "２．外観", "３．内観"], 1)
label(s, MARGIN, 1.56, CW, 0.40, "目的は「気軽に外壁・屋根工事の相談ができる場所づくり」", 13, WHITE, True, fill=BLUE)
label(s, MARGIN, 2.06, 4.60, 0.36, "外観ラフ・企画の4ポイント", 12, WHITE, True, fill=NAVY)
label(s, MARGIN, 2.46, 4.60, 1.95,
      ["① 自社の業態と社名を看板に記載する",
       "② ロゴやキャラクターでイメージの刷り込みを行う",
       "③ 自社の取り扱い商品を明示する",
       "④ ノボリを立て、店舗の視認性を高める"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.85)
label(s, 5.10, 2.06, 4.60, 0.36, "ノボリの3ポイント", 12, WHITE, True, fill=NAVY)
label(s, 5.10, 2.46, 4.60, 1.95,
      ["① 3m〜4m間隔で配置する",
       "② 色は看板と同じ色で統一する（逆転も可）",
       "③ ノボリの種類は限定する",
       "　（OPEN／業種／相談会実施中）"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.85)
label(s, MARGIN, 4.58, CW, 0.40,
      "3つの外観条件：入店しやすく何屋かわかる外観／屋根工事がわかりやすい空間／視認性の高い立地",
      12, WHITE, True, fill=NAVY)
label(s, MARGIN, 5.10, CW, 0.78,
      ["ノボリの種類を絞るのは、法則1（一番化）の実装である。",
       "「屋根リフォーム」「雨漏り・防水」など、伝える言葉を絞るほど記憶に残る。"],
      12, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.4)
label(s, MARGIN, 6.05, CW, 0.60,
      ["【図：外観デザインラフ・ノボリ作成事例】",
       "※原典 Ⅲ.店舗マニュアル P.4,5,6 の画像を差し込む"],
      10, BLUE, False, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, spacing=1.4)
source(s, "Ⅲ.店舗 P.2,4,5,6")

# P24 店舗：内観
s = base_slide(prs, CH3, "店舗③ 内観 ─ 全動線が「疑問・不安の解消」でできている", nxt(), 20)
phase_tabs(s, ["１．立地", "２．外観", "３．内観"], 2)
label(s, MARGIN, 1.56, CW, 0.38, "必要なもの8点セット", 12, WHITE, True, fill=BLUE)
items8 = ["① エントランス（ウェルカムボード）", "② 屋根模型", "③ 会社紹介POP", "④ 屋根知識POP",
          "⑤ 劣化診断POP", "⑥ 屋根材・ルーフィング材の見本や資料", "⑦ 商談机", "⑧ お茶出し用の給湯施設"]
for i, it in enumerate(items8):
    label(s, MARGIN + (i % 2) * 4.78, 2.02 + (i // 2) * 0.42, 4.62, 0.36, it, 11, BLACK, False,
          PP_ALIGN.LEFT, fill=WHITE, line=BLUE)
label(s, MARGIN, 3.98, CW, 0.36, "3種類のPOPと、それぞれが潰す不安", 12, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
table(s, MARGIN, 4.38, CW,
      [["POPの種類", "役割", "潰している阻害要因"],
       ["①知識系POP", "屋根材や劣化症状を知ってもらい、会社選びの\nポイントと判断基準を伝える", "プラン・時期"],
       ["②会社紹介系POP", "高額商材ゆえ、安心できる会社であることを訴求\n（会社紹介・強み・スタッフ紹介）", "会社・自分"],
       ["③施工事例＆\n　お客様の声", "家ごとに違う完成像を伝える。量・質ともに重要\n完工アンケートで口コミも掲載する", "金額・競合"]],
      col_w=[1.9, 5.4, 2.1], row_h=0.60, size=10, header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 6.42, CW, 0.34,
      "家模型の目的：家の構造を理解し基礎知識を身につけてもらう（屋根裏の釘・雨漏り箇所を赤で示す）",
      11, WHITE, True, fill=ORANGE)
source(s, "Ⅲ.店舗 P.7,9,10,11,13,14")

# P25 営業：全体フロー
s = base_slide(prs, CH3, "営業① 全体フロー ─ 初回訪問で8割が決まる", nxt(), 20)
phase_tabs(s, ["１．事前ヒアリング", "２．初回訪問", "３．クロージング"], -1)
label(s, MARGIN, 1.56, CW, 0.44,
      "屋根セールスコンセプト：屋根のメンテナンスの必要性を訴求し、お客様が安心して暮らせる住宅にする",
      12, WHITE, True, fill=BLUE)
flow = [("１．事前ヒアリング", ["電話反響受付票に沿って応対", "15分以内に担当から連絡する",
                          "お客様アンケートを上から順に", "「1〜2分よろしいでしょうか」と許可を取る",
                          "日程は3〜4候補から選ばせる", "キーマン（ご夫婦）の同席を依頼"]),
        ("２．初回訪問　★核心", ["ご挨拶・自己開示シートを渡す", "会社開示（アプローチブック冒頭）",
                            "現場調査（写真は最低30枚）", "診断結果の報告・VR・サーモ",
                            "ルーフィングの説明／工法と価格の明示", "仮クロージング（逆算スケジュール記入）"]),
        ("３．クロージング", ["前回商談の振り返り（お悩み相談シート）", "ショールームで模型・屋根材を体感",
                        "見積は3パターン提出", "「ご契約という形でよろしいでしょうか？」", "決まらなければ必ずその場で次アポ",
                        "宙ぶらりんの状態にはしない"])]
x = MARGIN
for nm, items in flow:
    fill = ORANGE if "★" in nm else NAVY
    tcol = NAVY if "★" in nm else WHITE
    label(s, x, 2.10, 3.05, 0.44, nm, 12, tcol, True, fill=fill)
    label(s, x, 2.60, 3.05, 2.42, bullets(items), 10, BLACK, False,
          PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.5)
    x += 3.18
label(s, MARGIN, 5.20, CW, 0.36, "営業において誤認しがちな3つの事実", 12, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
label(s, MARGIN, 5.60, CW, 1.10,
      ["① お客様自身も住まいの悩みを正確に認識していない → こちらから積極的にアドバイスするべき",
       "② 「価格が高いから」というのは自身のプレゼン力不足 → 【価格/価値】をどう出すかがポイント",
       "③ 営業力は才能に依存しない → トークマニュアル・営業フローの反復で能力は向上する"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=NAVY, spacing=1.5)
source(s, "Ⅳ.営業 P.2,3,4,5,6")

# P26 営業：初回訪問の型
s = base_slide(prs, CH3, "営業② 初回訪問の型 ─ 言う順番が決まっている", nxt(), 20)
phase_tabs(s, ["１．事前ヒアリング", "２．初回訪問", "３．クロージング"], 1)
rows = [["順", "やること", "なぜそうするのか（原典のポイント）"],
        ["1", "ご挨拶・上がり込み", "立ち話を避け、話を聞くマインドセットと安心感を得る"],
        ["2", "自己開示シートを渡す", "営業マンではなく1人の人間として自分をアピールする"],
        ["3", "会社開示", "数百万円の工事。会社を理解し信頼した上で契約いただく"],
        ["4", "アプローチブックを読み上げる", "渡すだけでは読まれない。7つの阻害要因を順に潰す"],
        ["5", "現調結果の報告（劣化写真）", "まとめシートで順に説明。iPadで見やすくお見せする"],
        ["6", "VR・サーモグラフィー", "普段見られない劣化を直接見せ、他社と差別化する"],
        ["7", "ルーフィングの説明", "塗装ではカバーできない点。“雨漏り対策の最後の砦”"],
        ["8", "工法の比較（カバー／葺き替え）", "買うか否かではなく、どちらを行うかに選択肢を切り替える"],
        ["9", "屋根材と価格の明示", "安い素材から順に紹介。実際の施工現場を例に価格を示す"],
        ["10", "仮クロージング", "お引渡し日から逆算して契約日を一緒に書き込む"]]
table(s, MARGIN, 1.56, CW, rows, col_w=[0.55, 3.15, 5.7], row_h=0.415, size=10, header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 6.15, CW, 0.55,
      ["見積提出は“ただ見積を出す場”ではなく“判断をする場”と認識いただく。",
       "そのために初回訪問で興味・関心を持っていただく必要がある。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.35)
source(s, "Ⅳ.営業 P.8〜23")

# P27 営業：クロージングと地雷
s = base_slide(prs, CH3, "営業③ クロージングと“地雷” ─ 競合を先回りで無力化する", nxt(), 20)
phase_tabs(s, ["１．事前ヒアリング", "２．初回訪問", "３．クロージング"], 2)
label(s, MARGIN, 1.56, 4.60, 0.36, "クロージングの手順", 12, WHITE, True, fill=NAVY)
label(s, MARGIN, 1.96, 4.60, 2.00,
      bullets(["お悩み相談シートで前回を振り返る", "アプローチブックを再度順に読み上げる",
               "SRの屋根材・家模型を見て・聞いて・体験", "見積は3パターン＋ルーフィング材の説明",
               "「一番しっくりくる材料はございますか？」", "「こちらでご契約という形でよろしいでしょうか？」"]),
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.5)
label(s, 5.10, 1.56, 4.60, 0.36, "決まらなかった場合", 12, WHITE, True, fill=DKRED)
label(s, 5.10, 1.96, 4.60, 2.00,
      ["待っていてもお客様から依頼が来ることはない",
       "と考えるべきである。",
       "",
       "「○○社のお見積りが出た後にまたお打ち合わせを。",
       "　○日と×日ではどちらがご都合よろしいですか？」",
       "",
       "→ 必ずその場で次アポを取得する"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=DKRED, spacing=1.3)
label(s, MARGIN, 4.12, CW, 0.36, "相見積もり対策の“地雷” ─ 3種類", 12, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
table(s, MARGIN, 4.52, CW,
      [["種類", "内容", "参考トーク"],
       ["見積地雷", "見積書の書き方・一式表記などの\n注意点を先に告知しておく", "（自社の見積の見方を先に説明しておく）"],
       ["工法地雷", "その屋根に適さない工法を\n先に指摘しておく", "「この屋根材は塗装では正しく施工できません。\n他社が塗装で提案してきたら気を付けてください」"],
       ["システム地雷", "保証・アフター・施工体制などの\n仕組み面の注意点を告知", "（自社のアフター体制を基準として示しておく）"]],
      col_w=[1.5, 3.5, 4.4], row_h=0.60, size=10, header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 6.42, CW, 0.34,
      "禁止事項：(1) 自爆しない　(2) あからさまな他社批判は慎む　(3) 「こう聞いて」と他社に言わせない",
      11, WHITE, True, fill=DKRED)
source(s, "Ⅳ.営業 P.13,24,25,26,27")

# P28 数値一覧
s = base_slide(prs, CH3, "数値基準 一覧 ─ この1枚を暗記する", nxt(), 20)
half = 14
rows_l = [["区分", "基準値"]] + [[a, b] for a, b, c in NUMBERS[:half]]
rows_r = [["区分", "基準値"]] + [[a, b] for a, b, c in NUMBERS[half:]]
table(s, MARGIN, 1.14, 4.62, rows_l, col_w=[1.35, 3.27], row_h=0.375, size=9, header_size=10, first_col_fill=BLUE)
table(s, 5.08, 1.14, 4.62, rows_r, col_w=[1.35, 3.27], row_h=0.375, size=9, header_size=10, first_col_fill=BLUE)
source(s, "Ⅰ.事業コンセプト／Ⅲ.集客／Ⅲ.店舗／Ⅳ.営業　全4編")


# ══════════════════════════ 第4部 ══════════════════════════
def quiz_slide(title, qs, page, answer=False):
    s = base_slide(prs, CH4, title, page, 20)
    if answer:
        rows = [["No", "設問", "解答", "根拠"]]
        for no, q, a, src in qs:
            rows.append([no, q, a, src])
        table(s, MARGIN, 1.16, CW, rows, col_w=[0.5, 3.1, 3.9, 1.9], row_h=0.50,
              size=9, header_size=10, first_col_fill=BLUE)
    else:
        rows = [["No", "設問", "解答欄"]]
        for no, q, a, src in qs:
            rows.append([no, q, ""])
        table(s, MARGIN, 1.16, CW, rows, col_w=[0.5, 5.0, 3.9], row_h=0.50,
              size=10, header_size=10, first_col_fill=BLUE)
    return s

# P29-P30 設問
s = quiz_slide("理解度チェック（1/2） ─ Q1〜Q10", QUIZ[:10], nxt())
label(s, MARGIN, 6.42, CW, 0.34, "目標：18問／20問以上。誤答は必ず原典の該当ページに戻って確認する",
      11, WHITE, True, fill=ORANGE)
source(s, "全4編")
s = quiz_slide("理解度チェック（2/2） ─ Q11〜Q20", QUIZ[10:], nxt())
label(s, MARGIN, 6.42, CW, 0.34, "Q13が答えられない場合は、第2部 P.17（7つの阻害要因）に戻る",
      11, WHITE, True, fill=ORANGE)
source(s, "全4編")

# P31-P32 解答
s = quiz_slide("理解度チェック 解答（1/2） ─ Q1〜Q10", QUIZ[:10], nxt(), answer=True)
source(s, "全4編")
s = quiz_slide("理解度チェック 解答（2/2） ─ Q11〜Q20", QUIZ[10:], nxt(), answer=True)
source(s, "全4編")

# P33 CASE
s = base_slide(prs, CH4, "CASE演習 ─ どの法則が破られているか", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.38,
      "現場で起きた出来事を読み、破られている法則を特定する。これができれば定着している",
      12, WHITE, True, fill=BLUE)
rows = [["", "起きたこと", "破られている法則", "打ち手"]]
for no, body, law, fix in CASES:
    rows.append([no, body.replace('\n', ''), law, fix])
table(s, MARGIN, 1.62, CW, rows, col_w=[0.85, 3.5, 2.15, 3.4], row_h=0.72,
      size=10, header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 6.30, CW, 0.42,
      "失注が出たら、必ず「7つの阻害要因のどれで落ちたか」を1つ特定して記録する。",
      12, WHITE, True, fill=DKRED)
source(s, "全4編（法則との対応は本資料による整理）")

# P34 ロードマップ
s = base_slide(prs, CH4, "7日間 学習ロードマップ", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.38, "1日30分×7日。読むだけで終わらせず、毎日1つ手を動かす", 12, WHITE, True, fill=BLUE)
rows = [["", "テーマ", "読む範囲", "その日のゴール"]]
for d, th, rg, goal in ROADMAP:
    rows.append([d, th, rg, goal])
table(s, MARGIN, 1.62, CW, rows, col_w=[0.85, 2.5, 3.3, 3.25], row_h=0.62,
      size=10, header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 6.10, CW, 0.60,
      ["Day7で18問未満だった場合は、間違えた設問の該当法則（第2部）だけを読み直す。",
       "全部を読み返す必要はない ─ それが法則で学ぶことの利点である。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=NAVY, spacing=1.35)
source(s, "本資料の構成による")

# P35 まとめ
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, fill=NAVY)
textbox(s, MARGIN, 0.45, CW, 0.4, [('Roof Reform-Manual　まとめ', 12, WHITE, False)])
textbox(s, MARGIN, 0.95, CW, 0.55, [('この1枚だけ、机に貼ってください', 22, ORANGE, True)])
label(s, 1.60, 1.75, 6.8, 0.72, "屋根は、客が一生見ない商品である", 21, WHITE, True, fill=DKRED)
arrow(s, 4.85, 2.56, 0.3, 0.26, ORANGE, 'down')
grid = [("法則1　一番化", "何屋かを一目で分からせる"),
        ("法則2　可視化", "見せた瞬間にニーズが生まれる"),
        ("法則3　掛け算分解", "どの因数が低いかで診断する"),
        ("法則4　3フェーズ", "作る→当てる→回す。回すを省かない"),
        ("法則5　階段設計", "ハードルは一段ずつしか下がらない"),
        ("法則6　標準化", "才能ではなくツールと反復で決まる"),
        ("法則7　ホームで戦う", "決めるのは必ず自社の土俵"),
        ("法則8　先回り", "先に言った側が信用される")]
for i, (nm, txt) in enumerate(grid):
    bx = MARGIN + (i % 2) * 4.78
    by = 2.95 + (i // 2) * 0.82
    rect(s, bx, by, 4.62, 0.70, fill=None, line=WHITE)
    textbox(s, bx + 0.10, by + 0.04, 4.42, 0.30, [(nm, 12, ORANGE, True)], align=PP_ALIGN.LEFT)
    textbox(s, bx + 0.10, by + 0.34, 4.42, 0.30, [(txt, 11, WHITE, False)], align=PP_ALIGN.LEFT)
textbox(s, MARGIN, 6.42, CW, 0.34,
        [('8つ覚えるのではない。1つ覚えて、8つ導く。', 13, WHITE, True)], align=PP_ALIGN.CENTER)
textbox(s, MARGIN, 7.05, CW, 0.3,
        [('Copyright©2020 Funai Consulting Inc. All rights reserved.', 9, WHITE, False)], align=PP_ALIGN.LEFT)
nxt()

OUT = "../屋根ビジネス研究会_法則で学ぶ学習資料.pptx"
prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst), "pages counted:", P)
