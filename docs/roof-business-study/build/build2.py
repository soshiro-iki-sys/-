# -*- coding: utf-8 -*-
"""屋根ビジネス研究会_法則で学ぶ学習資料（新人コンサルタント版）ビルダー"""
from deck_lib import *
from content2 import *
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = new_deck()
P = 0
def nxt():
    global P; P += 1; return P

CH0 = "第0部 本資料の使い方"
CH1 = "第1部 支援の全体像"
CH2 = "第2部 8つの法則（診断の物差し）"
CH3 = "第3部 支援実務ハンドブック"
CH4 = "第4部 定着チェック"

# ══════════════ 第0部 ══════════════
# P1 表紙
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 1.0, 2.30, 8.0, 0.035, fill=ORANGE)
textbox(s, 1.0, 1.30, 8.0, 0.45, [('Roof Reform-Manual', 16, WHITE, False)], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 1.0, 1.72, 8.0, 0.55, [('屋根ビジネス研究会　船井流マニュアル', 20, ORANGE, True)], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 1.0, 2.55, 8.0, 1.5, [('法則で学ぶ', 40, WHITE, True), ('学習資料', 40, WHITE, True)],
        anchor=MSO_ANCHOR.TOP, spacing=1.05)
label(s, 1.0, 4.28, 8.0, 0.52, "新人コンサルタント向け", 18, NAVY, True, fill=ORANGE)
textbox(s, 1.0, 5.00, 8.0, 1.2,
        [('全4編・102ページを　1つの根本原理 ＋ 8つの法則 に再構成', 13, WHITE, False),
         ('会員企業を診断し、指導するための物差しとして使う', 13, WHITE, False),
         ('想定学習時間：3日（Day1〜3）／実務適用まで4週間', 12, WHITE, False)], spacing=1.5)
textbox(s, 1.0, 6.60, 8.0, 0.3,
        [('Copyright©2020 Funai Consulting Inc. All rights reserved.', 10, WHITE, False)])
nxt()

# P2 使い方
s = base_slide(prs, CH0, "本資料の立ち位置 ─ あなたは「実行する人」ではなく「診断する人」", nxt(), 19)
label(s, MARGIN, 1.12, CW, 0.46,
      "102ページを暗記する必要はありません。クライアントの症状から、破れている法則を当てられれば十分です。",
      13, WHITE, True, fill=BLUE)
table(s, MARGIN, 1.72, CW,
      [["", "会員企業の営業マン・店長", "我々コンサルタント"],
       ["やること", "法則を実行する", "法則が守られているかを診断し、指導する"],
       ["必要な力", "トークとツールの反復", "症状から原因を特定する力／数値で語る力"],
       ["持ち物", "アプローチブック・模型・写真", "基準値・ヒアリングシート・処方箋の引き出し"],
       ["失敗の形", "「検討します」で終わる", "「頑張りましょう」で終わる"]],
      col_w=[1.3, 3.6, 4.5], row_h=0.42, size=11, first_col_fill=BLUE)
label(s, MARGIN, 4.05, CW, 0.36, "本資料の構成と、使うタイミング", 12, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
table(s, MARGIN, 4.45, CW,
      [["部", "内容", "使うタイミング"],
       ["第1部", "支援の全体像・数値の背骨・着任時チェックリスト", "初回訪問の前に読む"],
       ["第2部", "根本原理と8つの法則　★診断の物差し", "暗記せず“導ける”ようにする"],
       ["第3部", "ヒアリングシート／症状別処方箋／反論への切り返し", "訪問中に開く"],
       ["第4部", "理解度チェック20問・診断ケース・禁忌・30日ロードマップ", "着任1か月で自己採点"]],
      col_w=[1.0, 5.2, 2.6], row_h=0.38, size=11, first_col_fill=BLUE)
label(s, MARGIN, 6.28, CW, 0.42,
      "原典にない数値は、この資料にも書いていません。うろ覚えの数字でクライアントに語らないこと。",
      12, WHITE, True, fill=DKRED)
source(s, "Ⅰ.事業コンセプト／Ⅲ.集客／Ⅲ.店舗／Ⅳ.営業　全4編")

# ══════════════ 第1部 ══════════════
# P3 モデル比較
s = base_slide(prs, CH1, "我々が導入するモデルは何か ─ 初回面談で説明する1枚", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.44,
      "“圧倒的一番化を目指すための、戸建て住宅向け元請け屋根リフォームモデル”",
      14, WHITE, True, fill=BLUE)
table(s, MARGIN, 1.70, CW,
      [["", "一般的リフォームモデル", "一般的塗装モデル", "船井流屋根ビジネスモデル"],
       ["業態", "総合リフォーム", "塗装のみの取り扱い", "屋根・外壁リフォーム専門特化"],
       ["商品", "各種オーダーメイド商品", "外壁塗装・屋根塗装", "屋根カバー・屋根葺き替え"],
       ["集客", "訪問販売やOB客メイン", "訪問販売やOB・紹介がメイン", "屋根専門チラシ・WEB・店舗"],
       ["店舗", "事務所", "事務所", "屋根リフォーム体感型SR"],
       ["営業", "属人的営業", "属人的営業", "標準化マニュアル営業"],
       ["収益", "粗利率30％前後", "粗利率35〜37％", "粗利率40％"],
       ["商圏", "商圏発想なし", "商圏発想なし", "1拠点50万人"]],
      col_w=[1.0, 2.6, 2.7, 3.1], row_h=0.355, size=11, first_col_fill=BLUE)
label(s, MARGIN, 4.52, CW, 0.40,
      "この表の右列と、クライアントの現状を1行ずつ突き合わせる。それが初回診断である。",
      12, WHITE, True, fill=NAVY)
label(s, MARGIN, 5.06, 4.60, 1.55,
      ["【ライフサイクル上の位置】",
       "元請け屋根ビジネスは“成長期”。",
       "“屋根リフォーム専門店”を謳う企業は少なく、",
       "リフォーム店・塗装店ほど競合が乱立していない。",
       "→ 進む前に地域一番店シェアを確立する"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.4)
label(s, 5.10, 5.06, 4.60, 1.55,
      ["【ターゲットとペルソナ】",
       "築20年〜／年齢層50代後半〜",
       "（2世帯住宅・住宅の贈与の場合は年齢層が異なる）",
       "予算帯は塗装より高く増改築より安い180〜279万円",
       "→ なんでも屋ではない“業態発想”で集客する"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.4)
source(s, "Ⅰ.事業コンセプト P.4,5,7,10,11,12")

# P4 数値の背骨
s = base_slide(prs, CH1, "数値の背骨 ─ 目標設定はこの2式で行う", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.40,
      "同じ2.7億円を、市場側と現場側の2通りで出せること。これが目標設定の根拠になる",
      12, WHITE, True, fill=BLUE)
label(s, MARGIN, 1.64, 1.6, 0.32, "① 市場から", 11, WHITE, True, fill=NAVY)
for row, seq in enumerate([[("MS", "2,090円"), ("商圏人口", "50万人"), ("シェア率", "26％")],
                           [("現場調査数", "年250件"), ("契約率", "60％"), ("客単価", "180万円")]]):
    y = 2.02 if row == 0 else 3.50
    if row == 1:
        label(s, MARGIN, 3.12, 1.6, 0.32, "② 現場から", 11, WHITE, True, fill=NAVY)
    x = MARGIN
    for i, (k, v) in enumerate(seq):
        label(s, x, y, 1.95, 0.80, [(k, 10, NAVY, False), (v, 16, NAVY, True)],
              fill=WHITE, line=BLUE, spacing=1.15)
        x += 1.95
        if i < 2:
            label(s, x, y, 0.42, 0.80, "×", 17, BLUE, True); x += 0.42
    label(s, x, y, 0.42, 0.80, "＝", 17, BLUE, True)
    label(s, x + 0.42, y, 2.24, 0.80, [("受注売上", 10, WHITE, False), ("2.7億円", 18, WHITE, True)],
          fill=ORANGE, spacing=1.15)
label(s, MARGIN, 4.46, CW, 0.50,
      "月次：現場調査 21件 × 契約率 60％ × 客単価 180万円 ＝ 月間受注売上 2,250万円",
      14, NAVY, True, fill=WHITE, line=ORANGE)
label(s, MARGIN, 5.10, 4.60, 0.34, "3か年の立ち上げ計画", 11, WHITE, True, fill=NAVY)
table(s, MARGIN, 5.48, 4.60,
      [["年次", "年商", "月商", "月間見積件数"],
       ["1年目", "1億円", "約800万円", "7件"],
       ["2年目", "1.8億円", "1,500万円", "14件"],
       ["3年目", "2.7億円", "2,250万円", "21件"]],
      col_w=[1.0, 1.2, 1.3, 1.4], row_h=0.30, size=10, header_size=10, first_col_fill=BLUE)
label(s, 5.10, 5.10, 4.60, 0.34, "契約率が60％と50％で出てくる理由", 11, WHITE, True, fill=DKRED)
label(s, 5.10, 5.48, 4.60, 1.20,
      ["事業計画のKPIは契約率60％。営業マン個人の目標は",
       "現調7件×50％＝4件で設定されている。",
       "営業3名×7件＝現調21件、契約4件×3名＝12件 ≒",
       "事業計画12.5件。個人目標は安全側に置いている。",
       "→ クライアントに聞かれたら、この説明で答える"],
      10, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=DKRED, spacing=1.35)
source(s, "Ⅰ.事業コンセプト P.5,9,13／Ⅳ.営業 P.3")

# P5 導入時チェックリスト
s = base_slide(prs, CH1, "着任時チェックリスト ─ これを満たさない限り業績は上がらない", nxt(), 19)
label(s, MARGIN, 1.12, CW, 0.44,
      "原典が「下記を満たさない場合は業績アップ・利益向上が見込めない」と明記している5項目。最優先で確認する",
      12, WHITE, True, fill=DKRED)
chk = [("1", "営業マンは専任で設ける", "他事業部・別業務と兼任させない。元請け屋根営業の経験値と知識を蓄える"),
       ("2", "年間販促投資費用を予算組に入れる", "元請け集客に必要な“未来への投資”。年間売上目標の5〜8％を確保する"),
       ("3", "営業ツールを利用する", "研究会入会特典の営業ツールを必ず使う。口だけでなく武器を使わせる"),
       ("4", "来店型営業を徹底する", "SR（無ければ事務所に来店スペース）を構え“ホームで受注する”癖を付ける"),
       ("5", "本事業の責任者を配置する", "代表・役員も可。KPI・受注売上・客単価を把握する。営業への丸投げはNG")]
y = 1.72
for no, ttl, body in chk:
    label(s, MARGIN, y, 0.52, 0.72, no, 18, WHITE, True, fill=NAVY)
    label(s, 0.88, y, 3.30, 0.72, ttl, 12, NAVY, True, PP_ALIGN.LEFT, fill=WHITE, line=BLUE)
    label(s, 4.26, y, 5.44, 0.72, body, 11, BLACK, False, PP_ALIGN.LEFT, fill=WHITE, line=BLUE)
    y += 0.80
label(s, MARGIN, 5.80, CW, 0.86,
      ["この5項目は、法則8（先回り）そのものである。起きる失敗を先に潰しておくための項目であり、",
       "支援開始時に「まだ満たしていない項目はどれですか」と聞くだけで、初回訪問の議題が決まる。",
       "5項目のうち1つでも空欄なら、施策の話より先にそこを埋めることを提案する。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.35)
source(s, "Ⅰ.事業コンセプト P.6")

# P6 KPIツリー
s = base_slide(prs, CH1, "KPI因数分解ツリー ─ 診断はここから始まる", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.40,
      "「売上が足りない」で止めない。因数まで割ってから、初めて打ち手の話をする",
      12, WHITE, True, fill=BLUE)
label(s, 3.30, 1.64, 3.4, 0.50, "受注売上", 16, WHITE, True, fill=NAVY)
arrow(s, 4.80, 2.20, 0.4, 0.26, BLUE, 'down')
tri = [("現場調査数", ["チラシ：枚数×反響率1/10,000", "WEB：セッション×CVR", "イベント：年3〜4回", "近隣挨拶：やり切れるか"]),
       ("契約率", ["7つの阻害要因の解消度", "初回訪問での仮クロージング", "クロージングは来店で行うか", "地雷3種を仕掛けているか"]),
       ("客単価", ["カバー・葺き替え220万円", "小工事50万円", "屋根材5種以上の品揃え", "安い順に紹介し上へ誘導"])]
x = MARGIN
for nm, items in tri:
    label(s, x, 2.54, 3.02, 0.44, nm, 14, WHITE, True, fill=BLUE)
    label(s, x, 3.04, 3.02, 1.95, bullets(items), 11, BLACK, False,
          PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.45)
    x += 3.19
label(s, MARGIN, 5.18, CW, 0.36, "因数が特定できると、指導の言葉はこう変わる", 12, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
table(s, MARGIN, 5.58, CW,
      [["低い因数", "言ってはいけない指導", "本来の指導"],
       ["現場調査数", "「もっと販促しましょう」", "エリア別反響率を出させ、悪いエリアから削らせる"],
       ["契約率", "「営業力を上げましょう」", "失注案件を7つの阻害要因のどれかに分類させる"],
       ["客単価", "「単価を上げましょう」", "屋根材の品揃えと、価格帯別の施工事例を増やさせる"]],
      col_w=[1.8, 3.1, 4.5], row_h=0.31, size=11, first_col_fill=BLUE)
source(s, "Ⅰ.事業コンセプト P.5,13／Ⅲ.集客 P.11,12／Ⅳ.営業 P.19,27")

# ══════════════ 第2部 ══════════════
# P7 根本原理
s = base_slide(prs, CH2, "根本原理 ─ すべてはこの一点から始まる", nxt(), 20)
label(s, MARGIN, 1.18, CW, 1.00, ["屋根は、", "客が一生見ない商品である"], 26, WHITE, True, fill=DKRED, spacing=1.2)
label(s, MARGIN, 2.35, CW, 1.20,
      ["「お客様は普段屋根の劣化や破損を目にすることはなく、メンテナンスしなければ",
       "　いけないということを認識していない」　　　　　　　　　　　　（Ⅳ.営業マニュアル P.2）",
       "「屋根は普段から目にする機会が少なく、専門的な知識が多いため、",
       "　一般的な方が知識を持っていない商材です」　　　　　　　　　　（Ⅲ.店舗マニュアル P.10）"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=NAVY, spacing=1.4)
label(s, MARGIN, 3.72, CW, 0.36, "この1点から、屋根ビジネスの“おかしな常識”がすべて説明できる", 12, WHITE, True, fill=NAVY)
cons = [("客は必要性を\n認識できない", "だから、こちらから\n見せに行く"),
        ("客は判断基準を\n持てない", "だから、基準を\nこちらから与える"),
        ("客は不安になる\n（高額×初体験）", "だから、一段ずつ\n安心を積み上げる"),
        ("市場全体が\n未認識ニーズの塊", "だから、確率と仕組みで\n刈り取れる")]
x = MARGIN
for a, b in cons:
    label(s, x, 4.18, 2.24, 0.74, a.split('\n'), 12, WHITE, True, fill=BLUE, spacing=1.25)
    arrow(s, x + 0.97, 4.98, 0.3, 0.24, NAVY, 'down')
    label(s, x, 5.30, 2.24, 0.76, b.split('\n'), 11, NAVY, False, fill=WHITE, line=BLUE, spacing=1.3)
    x += 2.36
label(s, MARGIN, 6.22, CW, 0.50,
      ["クライアントに最初に伝えるのはこの一文。「価格が高いから決まらない」は多くの場合プレゼン力の問題であり、",
       "値引きではなく【価格/価値】の価値側を上げる、というのが本モデルの前提である。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.3)
source(s, "Ⅳ.営業 P.2／Ⅲ.店舗 P.10")

# P8 導出ツリー
s = base_slide(prs, CH2, "導出ツリー ─ 8法則は暗記せず「導く」", nxt(), 20)
label(s, MARGIN, 1.12, CW, 0.38, "根本原理から枝を辿れば、8つの法則はその場で再現できる", 12, WHITE, True, fill=BLUE)
label(s, 2.55, 1.60, 4.9, 0.48, "屋根は、客が一生見ない商品である", 15, WHITE, True, fill=DKRED)
arrow(s, 4.80, 2.14, 0.4, 0.24, NAVY, 'down')
branch = [("必要性を\n認識できない", ["法則2 可視化"]),
          ("判断基準を\n持てない", ["法則1 一番化", "法則8 先回り"]),
          ("不安になる", ["法則5 階段設計", "法則7 ホームで戦う"]),
          ("市場が未認識\nニーズの塊", ["法則3 掛け算分解", "法則4 3フェーズ"])]
x = MARGIN
for nm, laws in branch:
    label(s, x, 2.46, 2.24, 0.64, nm.split('\n'), 12, WHITE, True, fill=NAVY, spacing=1.2)
    y = 3.22
    for lw in laws:
        label(s, x, y, 2.24, 0.44, lw, 12, WHITE, True, fill=BLUE)
        y += 0.52
    x += 2.36
label(s, 2.55, 4.48, 4.9, 0.44, "説明が長く複雑になる", 13, WHITE, True, fill=NAVY)
arrow(s, 4.80, 4.98, 0.4, 0.24, NAVY, 'down')
label(s, 2.55, 5.28, 4.9, 0.44, "法則6 標準化", 13, WHITE, True, fill=BLUE)
label(s, MARGIN, 5.92, CW, 0.82,
      ["8法則はすべて4冊すべてに現れることを検証済み（1冊だけの記述は法則に昇格させていない）。",
       "候補のうち「未来への投資」「業態発想」「やり切り」は2冊以下だったため、他の法則に統合した。",
       "各法則の頁は【4冊での現れ方／診断の質問／症状／指導のNG】の4点セットで構成している。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=NAVY, spacing=1.3)
source(s, "全4編の横断検証による")

# P9〜P16 法則
for law in LAWS:
    s = base_slide(prs, CH2, "法則%d　%s ─ %s" % (law['no'], law['name'], law['sub']), nxt(), 20)
    label(s, MARGIN, 1.10, CW, 0.42, "クライアントへの一言：" + law['say'], 13, WHITE, True, fill=BLUE)
    keys = ["事業", "集客", "店舗", "営業"]
    pos = [(MARGIN, 1.64), (5.10, 1.64), (MARGIN, 3.16), (5.10, 3.16)]
    for k, (bx, by) in zip(keys, pos):
        rect(s, bx, by, 4.60, 1.44, fill=None, line=BLUE)
        label(s, bx, by, 4.60, 0.30, k + "マニュアルでの現れ方", 10, WHITE, True, fill=BLUE)
        textbox(s, bx + 0.06, by + 0.34, 4.48, 1.06, bullets(law['cells'][k]),
                10, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, spacing=1.32)
    trio = [("診断の質問（何を聞くか）", bullets(law['ask']), NAVY, WHITE),
            ("症状（破れている時に出る）", bullets(law['sym']), ORANGE, NAVY),
            ("指導のNG", [law['ng']], DKRED, WHITE)]
    bx = MARGIN
    for ttl, items, col, tcol in trio:
        rect(s, bx, 4.72, 3.05, 1.56, fill=None, line=col)
        label(s, bx, 4.72, 3.05, 0.30, ttl, 10, tcol, True, fill=col)
        textbox(s, bx + 0.06, 5.06, 2.93, 1.18, items, 9.5, BLACK, False,
                PP_ALIGN.LEFT, MSO_ANCHOR.TOP, spacing=1.3)
        bx += 3.18
    source(s, law['src'])

# P17 7つの阻害要因
s = base_slide(prs, CH2, "法則8の実装 ─ 7つの阻害要因と打ち手（営業同行時のチェックリスト）", nxt(), 18)
label(s, MARGIN, 1.10, CW, 0.38,
      "初回訪問でこの7つを解消できれば仮クロージングに進む。1つ残れば契約は落ちる",
      12, WHITE, True, fill=BLUE)
rows = [["阻害要因", "顧客の中で起きていること", "原典に書かれた打ち手（同行時に有無を確認する）"]]
for a, b, c in OBSTACLES:
    rows.append([a.replace('\n', ''), b, c])
table(s, MARGIN, 1.60, CW, rows, col_w=[1.5, 2.9, 5.0], row_h=0.62, size=10,
      header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 6.04, CW, 0.66,
      ["アプローチブックを“順に読み上げる”のは、この7つを漏れなく潰すための設計である。",
       "失注報告を受けたら「7つのどれで落ちたか」を1つ特定させる。これが契約率改善の起点になる。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=DKRED, spacing=1.35)
source(s, "Ⅳ.営業 P.3,10（7要因の内訳は貴社ご提供資料による）")


# ══════════════ 第3部 ══════════════
# P18 ヒアリングシート
s = base_slide(prs, CH3, "初回訪問 ヒアリングシート ─ この14項目を数字で埋める", nxt(), 19)
label(s, MARGIN, 1.10, CW, 0.38,
      "感想ではなく数字を取る。埋まらない欄があること自体が、最初の課題である",
      12, WHITE, True, fill=BLUE)
half = 7
rows_l = [["確認項目", "基準値", "疑う法則"]] + [[a, b, c] for a, b, c in INTAKE[:half]]
rows_r = [["確認項目", "基準値", "疑う法則"]] + [[a, b, c] for a, b, c in INTAKE[half:]]
table(s, MARGIN, 1.58, 4.62, rows_l, col_w=[1.55, 2.15, 0.92], row_h=0.50, size=9, header_size=10, first_col_fill=BLUE)
table(s, 5.08, 1.58, 4.62, rows_r, col_w=[1.55, 2.15, 0.92], row_h=0.50, size=9, header_size=10, first_col_fill=BLUE)
label(s, MARGIN, 5.70, CW, 0.94,
      ["【聞き方の順序】① 3因数（現調数・契約率・客単価）→ ② 販促の実績値 → ③ 体制（専任・責任者）→ ④ 店舗",
       "この順で聞くと、KPIツリーの上から下へ辿ることになり、相手も答えやすい。",
       "「分かりません」と返ってきた項目は、責めずにその場で一緒に計算する。それ自体が最初の支援になる。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.35)
source(s, "Ⅰ.事業 P.5,6,9,13／Ⅲ.集客 P.4,12,34／Ⅲ.店舗 P.3／Ⅳ.営業 P.3,11")

# P19 処方箋①
s = base_slide(prs, CH3, "症状別 処方箋① 集客が伸びないとき", nxt(), 20)
label(s, MARGIN, 1.10, CW, 0.38, "現場調査数が足りない場合、原因はほぼこの5パターンに収まる", 12, WHITE, True, fill=BLUE)
rows = [["症状", "疑う原因", "処方箋", "法則"]]
for a, b, c, d in DIAG_A:
    rows.append([a, b, c, d])
table(s, MARGIN, 1.58, CW, rows, col_w=[2.1, 2.35, 3.85, 1.1], row_h=0.66, size=10,
      header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 5.30, CW, 1.35,
      ["【この表の使い方】",
       "・症状は必ずクライアントの言葉ではなく“数字”で確認する（例：反響率を実際に計算してもらう）",
       "・原因は1つに絞る。複数を同時に直そうとすると、どれが効いたか分からなくなる",
       "・処方箋を出したら、次回定例までの担当者と期日をその場で決める（法則4の“回す”）",
       "・「チラシはもう効かない」と言われたら、まずエリア別反響率が出ているかを確認する"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=NAVY, spacing=1.4)
source(s, "Ⅲ.集客 P.3,4,7,11,12,22,23,33,39,43")

# P20 処方箋②
s = base_slide(prs, CH3, "症状別 処方箋② 契約率・客単価が上がらないとき", nxt(), 20)
label(s, MARGIN, 1.10, CW, 0.38, "現調数は足りているのに売上が出ない場合、原因は営業と店舗にある", 12, WHITE, True, fill=BLUE)
rows = [["症状", "疑う原因", "処方箋", "法則"]]
for a, b, c, d in DIAG_B:
    rows.append([a, b, c, d])
table(s, MARGIN, 1.58, CW, rows, col_w=[2.1, 2.35, 3.85, 1.1], row_h=0.60, size=10,
      header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 5.50, CW, 1.15,
      ["【契約率を上げる最短経路】",
       "初回訪問で8割が決まる。すなわち改善のレバーは、クロージングではなく初回訪問にある。",
       "同行して確認すべきは1点 ─ 仮クロージング（次回契約いただく口頭内諾）まで進めているか。",
       "進めていなければ、逆算スケジュールの記入をその場で型として渡す。"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.4)
source(s, "Ⅳ.営業 P.3,8,19,23,25,27／Ⅲ.店舗 P.7")

# P21 月次定例
s = base_slide(prs, CH3, "月次定例 アジェンダ標準形 ─ 80分の型", nxt(), 20)
label(s, MARGIN, 1.10, CW, 0.38, "毎回この順で進める。順番を変えないことが、定着の条件になる", 12, WHITE, True, fill=BLUE)
rows = [["", "議題", "時間", "見るもの・決めること"]]
for a, b, c, d in AGENDA:
    rows.append([a, b, c, d])
table(s, MARGIN, 1.58, CW, rows, col_w=[0.5, 2.3, 0.8, 5.8], row_h=0.62, size=10,
      header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 4.90, 4.60, 0.34, "必ず出席させる人", 11, WHITE, True, fill=NAVY)
label(s, MARGIN, 5.28, 4.60, 1.35,
      bullets(["事業責任者（代表・役員でも可）", "営業の代表者1名", "販促の数値を入力している担当者",
               "※営業マンへの丸投げ状態では定例が機能しない"]),
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.5)
label(s, 5.10, 4.90, 4.60, 0.34, "終了時に必ず残すもの", 11, WHITE, True, fill=NAVY)
label(s, 5.10, 5.28, 4.60, 1.35,
      bullets(["3因数のうち、今月改善する1つ", "折込表（エリア・枚数・折込日）の確定版",
               "宿題ごとの担当者と期日", "次回の日程"]),
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.5)
label(s, MARGIN, 6.42, CW, 0.34,
      "帳票はこちらが埋めない。埋めるのはクライアント、確認するのが我々の役割である。",
      11, WHITE, True, fill=DKRED)
source(s, "Ⅰ.事業 P.6／Ⅲ.集客 P.10,11,32／Ⅳ.営業 P.3")

# P22 反論への切り返し
s = base_slide(prs, CH3, "クライアントの反論への切り返し集", nxt(), 20)
label(s, MARGIN, 1.10, CW, 0.38,
      "反論には論破ではなく、原典の数値で返す。感情論に持ち込まないための型",
      12, WHITE, True, fill=BLUE)
rows = [["よく言われること", "返し方"]]
for a, b in REBUTTALS:
    rows.append([a, b])
table(s, MARGIN, 1.58, CW, rows, col_w=[2.7, 6.7], row_h=0.60, size=10,
      header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 6.00, CW, 0.70,
      ["いずれも共通しているのは「その判断の根拠になる数字を、まだ持っていない」という点である。",
       "反論を否定するのではなく、判断材料を一緒に作りにいく ─ それが法則3（掛け算分解）の指導実装。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.35)
source(s, "Ⅰ.事業 P.5,6,9／Ⅲ.集客 P.11／Ⅲ.店舗 P.3／Ⅳ.営業 P.2,11")

# P23 チラシ指導
s = base_slide(prs, CH3, "集客支援① チラシ ─ 3フェーズのどこで止まっているかを見る", nxt(), 19)
phase_tabs(s, ["１．作成フェーズ", "２．エリア選定フェーズ", "３．運用フェーズ"], -1)
cols = [("１．作成", ["見出しは上段中央に大きく", "「なんのお店か」を一目で",
                    "工事リストを最上部に記載", "左上にSR写真など視覚訴求",
                    "最下限金額を出し裾を広げる", "電話番号・QRは大きく目立たせる",
                    "裏面右下にCV導線（Zの法則）"]),
        ("２．エリア選定", ["一次商圏＝車で5〜10分圏", "一次商圏は月2〜3回配布",
                        "二次商圏は月1〜2回配布", "戸建比率が高い・OBが多いエリアも",
                        "　複数回配布してよい", "折込表は指定枚数×0.85で記載",
                        "折込日ごとに表を複数枚作る"]),
        ("３．運用　★止まりやすい", ["エリア別に反響件数・枚数を入力", "エリアごとの反響率を抽出する",
                              "悪いエリアは回数減または中止", "良いエリアの枚数を増やす",
                              "反響率が維持できるなら月4回も可", "繁忙期3〜5月・9〜11月は増やす",
                              "→ ここが空白なら最優先で着手"])]
x = MARGIN
for nm, items in cols:
    fill = ORANGE if "★" in nm else NAVY
    tcol = NAVY if "★" in nm else WHITE
    label(s, x, 1.54, 3.05, 0.40, nm, 11, tcol, True, fill=fill)
    label(s, x, 2.00, 3.05, 2.65, bullets(items), 10, BLACK, False,
          PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.5)
    x += 3.18
label(s, MARGIN, 4.82, CW, 0.34, "指導のポイント", 11, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
label(s, MARGIN, 5.20, CW, 1.10,
      ["・多くの会員企業は「作成」で止まっている。デザインの相談は来るが、反響率の相談は来ない",
       "・初回訪問では必ず「エリア別の反響率を見せてください」と依頼する。無ければそこが着手点",
       "・数字が無い状態でデザインを直しても、効いたかどうかを判定できない（法則4）"],
      11, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.45)
label(s, MARGIN, 6.42, CW, 0.30, "【図：チラシ表面・裏面サンプル ※原典 Ⅲ.集客 P.8,9 の画像を差し込む】", 10, BLUE, False)
source(s, "Ⅲ.集客 P.5,6,7,10,11,40")

# P24 WEB指導
s = base_slide(prs, CH3, "集客支援② WEB ─ CV＝セッション×CVR に尽きる", nxt(), 20)
label(s, MARGIN, 1.10, CW, 0.42, "反響数(CV) ＝ セッション数（流入） × コンバージョン率（反響率）", 15, WHITE, True, fill=ORANGE)
cols = [("１．HP作成", ["見出しで「何屋か」がわかる", "自社のUSP（独自の強み）を掲載",
                     "導線誘導ページ(LP)／更新ページ／", "　メールフォームページの3種を持つ",
                     "ユーザーの15の疑問に答える", "　コンテンツを揃える"]),
        ("２．広告配信", ["Google広告とYahoo広告を作成", "配信エリアごとにキャンペーン",
                      "1広告グループのKWは15個以下", "広告ランク＝入札価格×広告の品質",
                      "→ 品質を上げれば費用は下がる", "狙うのは「顕在×緊急度高」の層"]),
        ("３．運用", ["2週間ごとに数値入力し反省と対策", "予算通り・目標CPA以下なら月1回に",
                    "SEOはロングテールKWで数を打つ", "EAT基準（専門性・権威性・信頼性）",
                    "Search Consoleで表示回数が多く", "　順位が低いKWから記事化する"])]
x = MARGIN
for nm, items in cols:
    label(s, x, 1.64, 3.05, 0.40, nm, 12, WHITE, True, fill=NAVY)
    label(s, x, 2.10, 3.05, 2.30, bullets(items), 10, BLACK, False,
          PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.5)
    x += 3.18
label(s, MARGIN, 4.56, 4.60, 0.34, "セッションはあるのに反響がない", 11, WHITE, True, fill=DKRED)
label(s, MARGIN, 4.94, 4.60, 1.42,
      ["＝ 流入顧客層が不適正",
       "「とりあえず安いところ」「一括見積りをしたい」",
       "「A社（競合）で検討中」ばかりが流入している",
       "→ 入札KWを「市区町村×屋根○○」に絞り直す"],
      10, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=DKRED, spacing=1.5)
label(s, 5.10, 4.56, 4.60, 0.34, "目標からの逆算（クライアントに書かせる）", 11, WHITE, True, fill=BLUE)
label(s, 5.10, 4.94, 4.60, 1.42,
      ["売上目標 → 必要契約数 → 必要現調数",
       "→ チラシ／WEB／その他の集客数に配分",
       "→ CV数 ＝ セッション数 × CVR",
       "記入例：月10件 ＝ 2,000セッション × 0.5％"],
      10, BLACK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, fill=WHITE, line=BLUE, spacing=1.5)
label(s, MARGIN, 6.48, CW, 0.30,
      "※ 運用でやらなければ失敗する5点：2週ごと入力／入力ごとに反省と対策／目標達成後は月1回／屋根KW×市区町村／KWは他社事例とKWプランナーから",
      9, NAVY, False, PP_ALIGN.LEFT)
source(s, "Ⅲ.集客 P.12〜32")

# P25 店舗支援
s = base_slide(prs, CH3, "店舗支援 ─ 出店可否は、この基準で判定する", nxt(), 20)
label(s, MARGIN, 1.10, CW, 0.40,
      "基準を満たさない物件なら、出店を止めるのもコンサルタントの仕事である",
      12, WHITE, True, fill=DKRED)
label(s, MARGIN, 1.60, 4.62, 0.32, "地域査定要素", 11, WHITE, True, fill=NAVY)
table(s, MARGIN, 1.94, 4.62,
      [["項目", "基準"],
       ["交通量", "昼間12時間交通量 8,000〜12,000台"],
       ["視認性", "200m手前からお店（看板）が見える"],
       ["住宅密度", "自動車20分圏で6万世帯"],
       ["持家比率", "持家比率60％以上"],
       ["所得指数", "所得指数1.0以上"],
       ["人口増減", "人口・世帯数が増加している"],
       ["道路状況", "中央分離帯がない（理想）"],
       ["近隣施設", "スーパー等の商業施設周辺"]],
      col_w=[1.35, 3.27], row_h=0.325, size=10, header_size=10, first_col_fill=BLUE)
label(s, 5.08, 1.60, 4.62, 0.32, "物件査定要素", 11, WHITE, True, fill=NAVY)
table(s, 5.08, 1.94, 4.62,
      [["項目", "基準"],
       ["看板", "店頭に縦幅1m以上設置可能"],
       ["", "道路から見える自立看板がある"],
       ["過去履歴", "過去履歴の出入頻度が少ない"],
       ["面積", "店舗面積10坪以上"],
       ["駐車場", "駐車場3台以上"],
       ["間口", "間口に8m以上取れる"],
       ["設備", "空調、水廻りが整っている"],
       ["家賃", "家賃が月20万円前後で借りられる"]],
      col_w=[1.35, 3.27], row_h=0.325, size=10, header_size=10, first_col_fill=BLUE)
label(s, MARGIN, 4.90, CW, 0.34, "内観8点セット（既存SRの点検にも使う）", 11, WHITE, True, fill=NAVY, align=PP_ALIGN.LEFT)
items8 = ["①ウェルカムボード", "②屋根模型", "③会社紹介POP", "④屋根知識POP",
          "⑤劣化診断POP", "⑥屋根材・ルーフィング材の見本", "⑦商談机", "⑧お茶出し用の給湯施設"]
for i, it in enumerate(items8):
    label(s, MARGIN + (i % 4) * 2.36, 5.28 + (i // 4) * 0.42, 2.24, 0.36, it, 10, BLACK, False,
          PP_ALIGN.LEFT, fill=WHITE, line=BLUE)
label(s, MARGIN, 6.18, CW, 0.52,
      ["POPは3種類：①知識系（プラン・時期の不安を潰す）②会社紹介系（会社・自分）③施工事例＆お客様の声（金額・競合）",
       "＝ 内観は7つの阻害要因を、営業マンが話す前に潰しておくための装置である。"],
      10, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.3)
source(s, "Ⅲ.店舗 P.2,3,7,10")

# P26 営業同行
s = base_slide(prs, CH3, "営業支援 ─ 同行時に見る初回訪問の10ステップ", nxt(), 20)
label(s, MARGIN, 1.10, CW, 0.38,
      "同行の目的は評価ではなく、どのステップが抜けているかの特定である",
      12, WHITE, True, fill=BLUE)
rows = [["順", "やること", "なぜそうするのか（原典のポイント）", "潰す要因"],
        ["1", "ご挨拶・上がり込み", "立ち話を避け、話を聞くマインドセットと安心感を得る", "自分"],
        ["2", "自己開示シートを渡す", "営業マンではなく1人の人間として自分をアピールする", "自分"],
        ["3", "会社開示", "数百万円の工事。会社を理解し信頼した上で契約いただく", "会社"],
        ["4", "アプローチブックを読み上げる", "渡すだけでは読まれない。7つの阻害要因を順に潰す", "全要因"],
        ["5", "現調結果の報告（劣化写真）", "まとめシートで順に説明。iPadで見やすくお見せする", "プラン"],
        ["6", "VR・サーモグラフィー", "普段見られない劣化を直接見せ、他社と差別化する", "時期・競合"],
        ["7", "ルーフィングの説明", "塗装ではカバーできない点。“雨漏り対策の最後の砦”", "プラン・競合"],
        ["8", "工法の比較（カバー／葺き替え）", "買うか否かでなく、どちらを行うかに選択肢を切り替える", "プラン"],
        ["9", "屋根材と価格の明示", "安い素材から順に紹介。実際の施工現場を例に価格を示す", "金額"],
        ["10", "仮クロージング", "お引渡し日から逆算して契約日を一緒に書き込む", "時期・キーマン"]]
table(s, MARGIN, 1.58, CW, rows, col_w=[0.45, 2.75, 4.9, 1.3], row_h=0.415, size=9.5,
      header_size=10, first_col_fill=BLUE)
label(s, MARGIN, 6.18, CW, 0.52,
      ["抜けが最も多いのはステップ4と10。4を飛ばすと要因が残り、10を飛ばすと見積提出が“判断の場”にならない。",
       "同行後のフィードバックは、良し悪しではなく「何番が抜けていたか」で伝える。"],
      10, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.3)
source(s, "Ⅳ.営業 P.8〜23")

# P27 地雷
s = base_slide(prs, CH3, "営業支援 補足 ─ “地雷”の指導は、禁止事項とセットで", nxt(), 20)
label(s, MARGIN, 1.10, CW, 0.40,
      "地雷＝想定される相見積もりに事前に対策を打ち、その通りになった時に自社が一気に信用される技術",
      12, WHITE, True, fill=BLUE)
table(s, MARGIN, 1.62, CW,
      [["種類", "内容", "参考トーク"],
       ["見積地雷", "見積書の書き方・一式表記などの\n注意点を先に告知しておく", "（自社の見積の見方を先に説明しておく）"],
       ["工法地雷", "その屋根に適さない工法を\n先に指摘しておく", "「この屋根材は塗装では正しく施工できません。\n他社が塗装で提案してきたら気を付けてください」"],
       ["システム地雷", "保証・アフター・施工体制などの\n仕組み面の注意点を告知", "（自社のアフター体制を基準として示しておく）"]],
      col_w=[1.5, 3.5, 4.4], row_h=0.62, size=10, header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 3.85, CW, 0.36, "指導時に必ずセットで伝える禁止事項", 11, WHITE, True, fill=DKRED, align=PP_ALIGN.LEFT)
ban = [("(1) 自爆しないこと", "自社が満たしていない基準を地雷にすると、そのまま自社に返ってくる"),
       ("(2) あからさまな他社批判は慎むこと", "批判した瞬間に、こちらの品位が疑われる"),
       ("(3) 「こう聞いて」と他社に言わないようお願いすること", "指示された質問だと分かると、仕掛けが露見する")]
y = 4.25
for a, b in ban:
    label(s, MARGIN, y, 3.90, 0.50, a, 10, WHITE, True, PP_ALIGN.LEFT, fill=DKRED)
    label(s, 4.32, y, 5.38, 0.50, b, 10, BLACK, False, PP_ALIGN.LEFT, fill=WHITE, line=DKRED)
    y += 0.58
label(s, MARGIN, 6.05, CW, 0.62,
      ["地雷は強力だが、扱いを誤ると信用を失う唯一の技術である。新人が単独で教えると事故になりやすい。",
       "初回は必ず、禁止事項3つを先に伝えてから、工法地雷1本だけを文章化させるところから始める。"],
      10, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.3)
source(s, "Ⅳ.営業 P.13,27")

# P28 数値一覧
s = base_slide(prs, CH3, "数値基準 一覧 ─ 診断の物差し（この1枚を暗記する）", nxt(), 20)
half = 15
rows_l = [["区分", "基準値"]] + [[a, b] for a, b, c in NUMBERS[:half]]
rows_r = [["区分", "基準値"]] + [[a, b] for a, b, c in NUMBERS[half:]]
table(s, MARGIN, 1.12, 4.62, rows_l, col_w=[1.35, 3.27], row_h=0.355, size=9, header_size=10, first_col_fill=BLUE)
table(s, 5.08, 1.12, 4.62, rows_r, col_w=[1.35, 3.27], row_h=0.355, size=9, header_size=10, first_col_fill=BLUE)
label(s, 5.08, 6.20, 4.62, 0.48,
      "この数値を根拠なく丸めたり、うろ覚えで語らないこと。全て原典に出典がある。",
      10, WHITE, True, fill=DKRED)
source(s, "Ⅰ.事業コンセプト／Ⅲ.集客／Ⅲ.店舗／Ⅳ.営業　全4編")


# ══════════════ 第4部 ══════════════
def quiz_slide(title, qs, page, answer=False):
    s = base_slide(prs, CH4, title, page, 20)
    if answer:
        rows = [["No", "設問", "解答", "根拠"]]
        for no, q, a, src in qs:
            rows.append([no, q, a, src])
        table(s, MARGIN, 1.14, CW, rows, col_w=[0.5, 2.9, 4.1, 1.9], row_h=0.50,
              size=9, header_size=10, first_col_fill=BLUE)
    else:
        rows = [["No", "設問", "解答欄"]]
        for no, q, a, src in qs:
            rows.append([no, q, ""])
        table(s, MARGIN, 1.14, CW, rows, col_w=[0.5, 5.2, 3.7], row_h=0.50,
              size=10, header_size=10, first_col_fill=BLUE)
    return s

s = quiz_slide("理解度チェック（1/2） ─ Q1〜Q10", QUIZ[:10], nxt())
label(s, MARGIN, 6.40, CW, 0.34, "目標：18問／20問以上。誤答は必ず原典の該当ページに戻って確認する",
      11, WHITE, True, fill=ORANGE)
source(s, "全4編")
s = quiz_slide("理解度チェック（2/2） ─ Q11〜Q20", QUIZ[10:], nxt())
label(s, MARGIN, 6.40, CW, 0.34, "Q15が答えられない場合は、第2部 P.17（7つの阻害要因）に戻る",
      11, WHITE, True, fill=ORANGE)
source(s, "全4編")
s = quiz_slide("理解度チェック 解答（1/2） ─ Q1〜Q10", QUIZ[:10], nxt(), answer=True)
source(s, "全4編")
s = quiz_slide("理解度チェック 解答（2/2） ─ Q11〜Q20", QUIZ[10:], nxt(), answer=True)
source(s, "全4編")

# P33 診断ケース
s = base_slide(prs, CH4, "診断ケース演習 ─ 会員企業の状況から、破れている法則を当てる", nxt(), 19)
label(s, MARGIN, 1.10, CW, 0.36,
      "実務で最も使う能力。法則名で言語化できれば、処方箋は自動的に決まる",
      12, WHITE, True, fill=BLUE)
rows = [["", "クライアントの状況", "破れている法則", "打ち手"]]
for no, body, law, fix in CASES:
    rows.append([no, body.replace('\n', ''), law, fix])
table(s, MARGIN, 1.56, CW, rows, col_w=[0.75, 3.65, 1.85, 3.15], row_h=0.72,
      size=9.5, header_size=10, first_col_fill=BLUE)
label(s, MARGIN, 6.20, CW, 0.50,
      ["CASE 6のように、施策以前に「導入時の注意点5項目」が満たされていないケースは非常に多い。",
       "その場合は集客・営業の話をする前に、専任化と責任者配置から着手する。"],
      10, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.3)
source(s, "全4編（法則との対応は本資料による整理）")

# P34 禁忌
s = base_slide(prs, CH4, "コンサルタントとしての禁忌 ─ やってはいけない6つ", nxt(), 20)
label(s, MARGIN, 1.10, CW, 0.40,
      "支援が空回りするとき、原因はクライアント側ではなくこちら側にあることが多い",
      12, WHITE, True, fill=DKRED)
y = 1.62
for i, (a, b) in enumerate(TABOOS):
    label(s, MARGIN, y, 0.50, 0.64, "×", 18, WHITE, True, fill=DKRED)
    label(s, 0.86, y, 3.55, 0.64, a, 11, DKRED, True, PP_ALIGN.LEFT, fill=WHITE, line=DKRED)
    label(s, 4.49, y, 5.21, 0.64, b, 10, BLACK, False, PP_ALIGN.LEFT, fill=WHITE, line=BLUE)
    y += 0.72
label(s, MARGIN, 5.98, CW, 0.72,
      ["6つに共通するのは「その場は前に進むが、クライアントの中に何も残らない」という点である。",
       "支援の成果物は資料でも数字でもなく、クライアント自身が数字を見て判断できる状態になること。",
       "帳票を埋めるのはクライアント、確認するのが我々 ─ この役割分担を最後まで崩さない。"],
      10, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.35)
source(s, "Ⅰ.事業 P.6／Ⅲ.集客 P.11,32／Ⅳ.営業 P.10,27")

# P35 ロードマップ
s = base_slide(prs, CH4, "着任30日ロードマップ", nxt(), 20)
label(s, MARGIN, 1.10, CW, 0.36, "読むだけで終わらせない。Week2以降は必ず担当先の実データに当てる", 12, WHITE, True, fill=BLUE)
rows = [["期間", "テーマ", "使うもの", "到達目標"]]
for a, b, c, d in ROADMAP:
    rows.append([a, b, c, d])
table(s, MARGIN, 1.56, CW, rows, col_w=[1.0, 2.3, 3.1, 3.0], row_h=0.62,
      size=10, header_size=11, first_col_fill=BLUE)
label(s, MARGIN, 5.60, CW, 1.05,
      ["【Day7で18問未満だった場合】",
       "全部を読み返す必要はない。間違えた設問に対応する法則の頁（第2部）だけを読み直す。",
       "それが「法則で学ぶ」ことの利点であり、この資料を法則順に並べた理由である。"],
      11, NAVY, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fill=WHITE, line=ORANGE, spacing=1.4)
source(s, "本資料の構成による")

# P36 まとめ
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, SW, SH, fill=NAVY)
textbox(s, MARGIN, 0.42, CW, 0.36, [('Roof Reform-Manual　まとめ', 12, WHITE, False)])
textbox(s, MARGIN, 0.88, CW, 0.50, [('この1枚を持って、訪問してください', 21, ORANGE, True)])
label(s, 1.60, 1.62, 6.8, 0.66, "屋根は、客が一生見ない商品である", 20, WHITE, True, fill=DKRED)
arrow(s, 4.85, 2.36, 0.3, 0.24, ORANGE, 'down')
grid = [("法則1　一番化", "何屋かを一目で分からせているか"),
        ("法則2　可視化", "劣化の証拠を見せているか"),
        ("法則3　掛け算分解", "3因数のどれが低いか特定したか"),
        ("法則4　3フェーズ", "“回す”帳票は埋まっているか"),
        ("法則5　階段設計", "次の一段が用意されているか"),
        ("法則6　標準化", "専任・責任者・読み上げはあるか"),
        ("法則7　ホームで戦う", "クロージングは来店で行っているか"),
        ("法則8　先回り", "7つの阻害要因を先に潰したか")]
for i, (nm, txt) in enumerate(grid):
    bx = MARGIN + (i % 2) * 4.78
    by = 2.78 + (i // 2) * 0.80
    rect(s, bx, by, 4.62, 0.68, fill=None, line=WHITE)
    textbox(s, bx + 0.10, by + 0.03, 4.42, 0.30, [(nm, 12, ORANGE, True)], align=PP_ALIGN.LEFT)
    textbox(s, bx + 0.10, by + 0.33, 4.42, 0.30, [(txt, 11, WHITE, False)], align=PP_ALIGN.LEFT)
textbox(s, MARGIN, 6.28, CW, 0.32,
        [('8つ覚えるのではない。1つ覚えて、8つ導く。', 13, WHITE, True)], align=PP_ALIGN.CENTER)
textbox(s, MARGIN, 6.66, CW, 0.32,
        [('そして、実行するのはクライアント。我々は診断し、確認する。', 12, ORANGE, True)], align=PP_ALIGN.CENTER)
textbox(s, MARGIN, 7.08, CW, 0.3,
        [('Copyright©2020 Funai Consulting Inc. All rights reserved.', 9, WHITE, False)], align=PP_ALIGN.LEFT)
nxt()

OUT = "../屋根ビジネス研究会_法則で学ぶ学習資料_新人コンサル版.pptx"
prs.save(OUT)
print("saved:", OUT, "/ pages:", P)
