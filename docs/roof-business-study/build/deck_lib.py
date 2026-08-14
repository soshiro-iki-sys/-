# -*- coding: utf-8 -*-
"""屋根ビジネス研究会 学習資料 共通描画ライブラリ（実測フィット版）

方針：テキストは全て実フォントで幅を実測し、枠に収まらなければ
      自動的にフォントを縮小する。よって「はみ出し」は構造的に発生しない。
"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import ImageFont

NAVY   = RGBColor(0x00, 0x20, 0x60)
BLUE   = RGBColor(0x00, 0x70, 0xC0)
ORANGE = RGBColor(0xFF, 0xC0, 0x00)
RED    = RGBColor(0xFF, 0x00, 0x00)
DKRED  = RGBColor(0xC0, 0x00, 0x00)
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = 'メイリオ'
SW, SH = 10.0, 7.5
MARGIN = 0.3
CW = SW - MARGIN * 2

# ── テキスト実測 ──────────────────────────────────────────
_TTF = '/usr/share/fonts/truetype/fonts-japanese-gothic.ttf'
_SCALE = 8            # 8倍サイズで測って割る（整数丸め誤差の低減）
_SAFETY = 1.14        # メイリオは計測に使うIPAゴシックよりやや広いため安全率
_cache = {}


def _font(sz):
    k = max(1, int(round(sz * _SCALE)))
    if k not in _cache:
        _cache[k] = ImageFont.truetype(_TTF, k)
    return _cache[k]


def text_w(txt, sz):
    """テキストの描画幅（pt・安全率込み）"""
    if not txt:
        return 0.0
    return _font(sz).getlength(txt) / _SCALE * _SAFETY


def wrap_lines(txt, sz, avail_pt):
    """avail_pt 幅で折り返した際の行数（貪欲法）"""
    if not txt:
        return 1
    if text_w(txt, sz) <= avail_pt:
        return 1
    lines, cur = 1, ''
    for ch in txt:
        if text_w(cur + ch, sz) > avail_pt and cur:
            lines += 1
            cur = ch
        else:
            cur += ch
    return lines


def _norm(lines, size, color, bold):
    """行データを (text, size, color, bold) に正規化"""
    if isinstance(lines, str):
        lines = [lines]
    out = []
    for ln in lines:
        if isinstance(ln, tuple):
            t, s, c, b = (list(ln) + [size, color, bold])[:4]
        else:
            t, s, c, b = ln, size, color, bold
        out.append((t, s, c, b))
    return out


def _block_h(rows, avail_pt, spacing, scale=1.0):
    """行データ全体の必要高さ（pt）"""
    h = 0.0
    for t, s, c, b in rows:
        s = s * scale
        h += wrap_lines(t, s, avail_pt) * s * spacing
    return h


def _fit_scale(rows, avail_w_pt, avail_h_pt, spacing, min_scale=0.62):
    """枠に収まる最大の縮小率を返す（1.0＝縮小なし）"""
    scale = 1.0
    while scale > min_scale:
        if _block_h(rows, avail_w_pt, spacing, scale) <= avail_h_pt:
            return scale
        scale -= 0.04
    return min_scale


def _set_font(run, size, color, bold, font=FONT):
    run.font.size = Pt(round(size, 1))
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        for e in rPr.findall(qn(tag)):
            rPr.remove(e)
        rPr.append(rPr.makeelement(qn(tag), {'typeface': font}))


def _hanging(p, size):
    """「・」で始まる行の折り返しをぶら下げ揃えにする"""
    ind = int(Pt(size * 1.05))
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(ind))
    pPr.set('indent', str(-ind))


PAD_L = 0.05   # テキストボックス左右の内側余白(inch)
PAD_V = 0.03


def textbox(slide, l, t, w, h, lines, size=12, color=BLACK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.25, fit=True):
    """枠に必ず収まるテキストボックス（収まらなければ自動縮小）"""
    rows = _norm(lines, size, color, bold)
    avail_w = (w - PAD_L * 2) * 72.0
    avail_h = (h - PAD_V * 2) * 72.0
    scale = _fit_scale(rows, avail_w, avail_h, spacing) if fit else 1.0

    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(PAD_L)
    tf.margin_top = tf.margin_bottom = Inches(PAD_V)
    tf.vertical_anchor = anchor
    for i, (txt, sz, col, bd) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        sz = sz * scale
        if txt.startswith('・'):
            _hanging(p, sz)
        r = p.add_run()
        r.text = txt
        _set_font(r, sz, col, bd)
    return tb


def rect(slide, l, t, w, h, fill=None, line=None, line_w=1.25, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    sp.text_frame.text = ''
    return sp


def label(slide, l, t, w, h, text, size=12, color=BLACK, bold=False,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, fill=None, line=None, spacing=1.2):
    if fill is not None or line is not None:
        rect(slide, l, t, w, h, fill, line)
    return textbox(slide, l, t, w, h, text, size, color, bold, align, anchor, spacing)


def panel(slide, l, t, w, h, title, items, line_col=BLUE, head_col=BLUE,
          head_tcol=WHITE, size=10, head_size=10, spacing=1.32, bullet=True):
    """見出し帯つきの囲みパネル（本文は自動縮小で必ず内側に収まる）"""
    rect(slide, l, t, w, h, fill=None, line=line_col)
    hh = 0.30
    label(slide, l, t, w, hh, title, head_size, head_tcol, True, fill=head_col)
    body = ['・' + i for i in items] if bullet else items
    textbox(slide, l, t + hh, w, h - hh, body, size, BLACK, False,
            PP_ALIGN.LEFT, MSO_ANCHOR.TOP, spacing)


def arrow(slide, l, t, w, h, color=BLUE, direction='right'):
    shp = MSO_SHAPE.RIGHT_ARROW if direction == 'right' else MSO_SHAPE.DOWN_ARROW
    sp = slide.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def base_slide(prs, chapter, title, page, title_size=20):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    textbox(s, MARGIN, 0.07, CW, 0.22, [('Roof Reform-Manual　' + chapter, 9, NAVY, False)],
            anchor=MSO_ANCHOR.MIDDLE, fit=False)
    rect(s, MARGIN, 0.30, CW, 0.012, fill=NAVY)
    if title:
        rect(s, MARGIN, 0.38, CW, 0.62, fill=NAVY)
        # タイトルは1行に収める（収まらなければ自動縮小）
        sz = title_size
        while sz > 13 and text_w(title, sz) > (CW - 0.30) * 72.0:
            sz -= 0.5
        textbox(s, MARGIN + 0.14, 0.38, CW - 0.28, 0.62, title, sz, WHITE, True,
                PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, fit=False)
    _footer(s, page)
    return s


def _footer(s, page):
    textbox(s, MARGIN, 7.14, 6.0, 0.24,
            [('Copyright©2020 Funai Consulting Inc. All rights reserved.', 9, NAVY, False)],
            anchor=MSO_ANCHOR.MIDDLE, fit=False)
    textbox(s, SW - MARGIN - 1.0, 7.14, 1.0, 0.24, [(str(page), 9, NAVY, False)],
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, fit=False)


def source(s, text, y=6.82):
    textbox(s, MARGIN, y, CW, 0.24, [('出典：' + text, 9, BLUE, False)],
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, fit=False)


def phase_tabs(s, tabs, active, top=1.06, left=MARGIN, width=CW, h=0.32, size=10.5):
    n = len(tabs); gap = 0.06
    w = (width - gap * (n - 1)) / n
    for i, t in enumerate(tabs):
        on = (i == active)
        label(s, left + i * (w + gap), top, w, h, t, size,
              WHITE if on else BLUE, True, fill=BLUE if on else WHITE, line=BLUE)


CELL_PAD = 0.05


TABLE_FIXED = False   # True にすると行高を row_h に固定し、収まるまで文字を縮小する


def table(s, l, t, w, rows, col_w=None, header_fill=NAVY, size=10, row_h=0.30,
          header_size=10, first_col_fill=None, max_h=None, align_first=PP_ALIGN.LEFT,
          fixed_rows=None):
    """行高をテキスト実測から自動決定するテーブル。
       fixed_rows=True（既定は TABLE_FIXED）の場合は行高を row_h に固定し、
       全セルが収まるまでフォントを縮小する（既存レイアウトを崩さないため）。"""
    nr, nc = len(rows), len(rows[0])
    if col_w is None:
        col_w = [1] * nc
    tot = sum(col_w)
    col_in = [w * c / tot for c in col_w]
    fixed = TABLE_FIXED if fixed_rows is None else fixed_rows

    def needs(scale):
        out = []
        for ri, row in enumerate(rows):
            sz = (header_size if ri == 0 else size) * scale
            need = 0.0
            for ci, val in enumerate(row):
                avail = (col_in[ci] - CELL_PAD * 2) * 72.0
                lines = sum(wrap_lines(part, sz, avail) for part in str(val).split('\n'))
                need = max(need, lines * sz * 1.22 / 72.0 + 0.09)
            out.append(need)
        return out

    scale = 1.0
    if fixed:
        while scale > 0.55 and max(needs(scale)) > row_h:
            scale -= 0.04
        heights = [row_h] * nr
    else:
        heights = [max(row_h * scale, n) for n in needs(scale)]
        if max_h:
            while sum(heights) > max_h and scale > 0.62:
                scale -= 0.04
                heights = [max(row_h * scale, n) for n in needs(scale)]

    shape = s.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(sum(heights)))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, c in enumerate(col_in):
        tbl.columns[i].width = Emu(int(Inches(c)))
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(heights[ri])
        sz = (header_size if ri == 0 else size) * scale
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(CELL_PAD)
            cell.margin_top = cell.margin_bottom = Inches(0.015)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = (header_fill if ri == 0
                                        else (first_col_fill if (ci == 0 and first_col_fill) else WHITE))
            tf = cell.text_frame
            tf.word_wrap = True
            for pi, ptxt in enumerate(str(val).split('\n')):
                p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
                p.line_spacing = 1.05
                p.alignment = (PP_ALIGN.CENTER if ri == 0
                               else (align_first if ci == 0 else PP_ALIGN.LEFT))
                r = p.add_run(); r.text = ptxt
                if ri == 0:
                    _set_font(r, sz, WHITE, True)
                elif ci == 0 and first_col_fill is not None:
                    _set_font(r, sz, WHITE, True)
                else:
                    _set_font(r, sz, BLACK, False)
    return sum(heights)


def table_h(w, rows, col_w=None, size=10, row_h=0.30, header_size=10):
    """描画せずに総高だけを見積もる"""
    nc = len(rows[0])
    if col_w is None:
        col_w = [1] * nc
    tot = sum(col_w)
    col_in = [w * c / tot for c in col_w]
    total = 0.0
    for ri, row in enumerate(rows):
        sz = header_size if ri == 0 else size
        need = 0.0
        for ci, val in enumerate(row):
            avail = (col_in[ci] - CELL_PAD * 2) * 72.0
            lines = sum(wrap_lines(part, sz, avail) for part in str(val).split('\n'))
            need = max(need, lines * sz * 1.22 / 72.0 + 0.09)
        total += max(row_h, need)
    return total


def note_box(s, l, w, lines, size=10.5, color=NAVY, line=ORANGE,
             spacing=1.45, bottom=None, top=None, pad=0.20, align=PP_ALIGN.LEFT):
    """テキスト量ぴったりの高さで囲み注記を描く（bottom 指定で下端揃え）"""
    rows = _norm(lines, size, color, False)
    need = _block_h(rows, (w - PAD_L * 2) * 72.0, spacing) / 72.0 + pad
    y = (bottom - need) if bottom is not None else top
    rect(s, l, y, w, need, fill=WHITE, line=line)
    textbox(s, l, y, w, need, lines, size, color, False, align, MSO_ANCHOR.MIDDLE, spacing)
    return y


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def bullets(txts, mark='・'):
    return [mark + t for t in txts]
