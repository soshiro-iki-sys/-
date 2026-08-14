# -*- coding: utf-8 -*-
"""図形の重なり・枠外を幾何的に検出する"""
from pptx import Presentation
import sys
PATH = "../屋根ビジネス研究会_法則で学ぶ学習資料_現場向け.pptx"
prs = Presentation(PATH)
E = 914400.0
W, H = prs.slide_width / E, prs.slide_height / E
TOL = 0.02   # inch

def box(sh):
    return (sh.left / E, sh.top / E, (sh.left + sh.width) / E, (sh.top + sh.height) / E)

def inter(a, b):
    x = min(a[2], b[2]) - max(a[0], b[0])
    y = min(a[3], b[3]) - max(a[1], b[1])
    return (x, y) if x > TOL and y > TOL else None

issues = []
for si, s in enumerate(prs.slides, 1):
    txts = []
    for sh in s.shapes:
        if sh.left is None:
            continue
        b = box(sh)
        if b[0] < -TOL or b[1] < -TOL or b[2] > W + TOL or b[3] > H + TOL:
            issues.append((si, "枠外", type(sh).__name__, [round(v, 2) for v in b]))
        has_text = sh.has_text_frame and sh.text_frame.text.strip()
        if has_text:
            txts.append((sh, b, sh.text_frame.text.strip()[:26]))
        if sh.has_table:
            tb = box(sh)
            txts.append((sh, tb, "[TABLE]"))
    for i in range(len(txts)):
        for j in range(i + 1, len(txts)):
            ov = inter(txts[i][1], txts[j][1])
            if ov:
                issues.append((si, "重なり", f"{txts[i][2]} × {txts[j][2]}",
                               f"重畳 {ov[0]:.2f}x{ov[1]:.2f}in"))
print("=== 重なり・枠外の検出 ===")
if not issues:
    print("  なし")
for it in issues:
    print(" ", it)
print("slides:", len(prs.slides._sldIdLst))
