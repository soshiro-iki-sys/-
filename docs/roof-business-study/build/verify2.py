# -*- coding: utf-8 -*-
"""全スライドのテキストはみ出し検証（保守的な推定）"""
import math, unicodedata, sys
from pptx import Presentation
from pptx.util import Pt

PATH = "../屋根ビジネス研究会_法則で学ぶ学習資料_新人コンサル版.pptx"
prs = Presentation(PATH)

def width_units(t):
    """全角=1.0, 半角=0.55 でテキスト幅を em 単位換算"""
    u = 0.0
    for ch in t:
        u += 1.0 if unicodedata.east_asian_width(ch) in ('W', 'F', 'A') else 0.55
    return u

def check_tf(tf, w_in, h_in, ctx, issues, pad=0.08, lsp_default=1.15):
    inner_w = (w_in - pad) * 72.0
    inner_h = (h_in - 0.04) * 72.0
    total = 0.0
    max_line_over = None
    for p in tf.paragraphs:
        txt = "".join(r.text for r in p.runs)
        sizes = [r.font.size.pt for r in p.runs if r.font.size] or [12.0]
        sz = max(sizes)
        lsp = p.line_spacing if isinstance(p.line_spacing, float) else lsp_default
        if not txt:
            total += sz * lsp
            continue
        cpl = inner_w / sz
        lines = max(1, math.ceil(width_units(txt) / cpl))
        if lines > 1 and width_units(txt) / cpl > 1.0:
            ratio = width_units(txt) / cpl
            if max_line_over is None or ratio > max_line_over[0]:
                max_line_over = (ratio, txt[:34], sz, lines)
        total += lines * sz * lsp
    if total > inner_h + 0.5:
        issues.append((ctx, "縦オーバー", round(total, 1), round(inner_h, 1),
                       max_line_over[1] if max_line_over else ""))
    return total

issues = []
wraps = []
for si, slide in enumerate(prs.slides, 1):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            check_tf(sh.text_frame, sh.width / 914400, sh.height / 914400,
                     f"P{si} shape", issues)
            # 1行想定のラベルが折り返していないか
            for p in sh.text_frame.paragraphs:
                txt = "".join(r.text for r in p.runs)
                if not txt:
                    continue
                sizes = [r.font.size.pt for r in p.runs if r.font.size] or [12.0]
                sz = max(sizes)
                inner_w = (sh.width / 914400 - 0.08) * 72.0
                if width_units(txt) * sz > inner_w:
                    wraps.append((si, round(width_units(txt) * sz / inner_w, 2), sz, txt[:40]))
        if sh.has_table:
            tbl = sh.table
            for ri, row in enumerate(tbl.rows):
                for ci, cell in enumerate(row.cells):
                    if not cell.text.strip():
                        continue
                    cw = tbl.columns[ci].width / 914400
                    rh = row.height / 914400
                    check_tf(cell.text_frame, cw - 0.06, rh, f"P{si} table[{ri},{ci}]", issues, pad=0.06, lsp_default=1.0)

print("=== 縦オーバー（枠からはみ出す可能性） ===")
if not issues:
    print("  なし")
for i in issues:
    print(" ", i)
print("\n=== 横方向の折り返し（1行に収まらない箇所・上位20） ===")
wraps.sort(key=lambda x: -x[1])
for w in wraps[:20]:
    print(f"  P{w[0]} 幅比{w[1]} {w[2]}pt : {w[3]}")
print("\nslides:", len(prs.slides._sldIdLst))
