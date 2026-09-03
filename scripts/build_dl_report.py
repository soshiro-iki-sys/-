# -*- coding: utf-8 -*-
"""DLレポート第1号「工務店のための 太陽光・蓄電池 提携モデル解説」たたき台
企画書 P20「見本｜DLレポート」の目次どおり、A4縦12ページで生成する。
"""
import os
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

# ---------- パレット（株式会社ワン・ミニットの資料から採取） ----------
CHAR  = RGBColor(0x2D, 0x2C, 0x2A)   # 濃色｜見出し帯・強い文字（チャコール）
INK   = RGBColor(0x2B, 0x2A, 0x28)   # 本文
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORNG  = RGBColor(0xE7, 0x75, 0x2F)   # アクセント｜数値・強調（オレンジ）
AMBER = RGBColor(0xFA, 0xA4, 0x13)   # 濃色帯の上の強調（山吹）
PALE  = RGBColor(0xF7, 0xF3, 0xEE)   # 淡い地（ウォームグレー）
GRAY  = RGBColor(0xB8, 0xB4, 0xAE)
SOFT  = RGBColor(0xF3, 0xF1, 0xEE)

FONT = 'Meiryo UI'
ISSUER = '株式会社ワン・ミニット'
COPY   = '© 株式会社ワン・ミニット　All Rights Reserved.'

W, H = 21.0, 29.7                      # A4縦
M, G  = 1.60, 0.30
CW    = W - M * 2                      # 17.80
RE    = M + CW                         # 19.40
CT, CB = 4.60, 27.10

T_TITLE, T_LEAD, T_LEADEM = 20, 13, 15
T_HEAD, T_BODY, T_SUB = 13, 11, 9
T_NUM_L, T_NUM_M = 20, 15


def cols(n, gutter=G):
    w = (CW - gutter * (n - 1)) / n
    return [M + i * (w + gutter) for i in range(n)], w


X2, W2 = cols(2)
X3, W3 = cols(3)

prs = Presentation()
prs.slide_width  = Cm(W)
prs.slide_height = Cm(H)
BLANK = prs.slide_layouts[6]

# ---------- 自動フィット ----------
_PFONT = '/usr/share/fonts/opentype/ipafont-gothic/ipagp.ttf'
_PPC, _PT2CM, _LH = 50.0, 2.54 / 72.0, 1.22
_fc = {}
_probe = ImageDraw.Draw(Image.new('RGB', (8, 8)))
FITWARN = []


def _f(pt):
    px = max(6, int(round(pt * _PT2CM * _PPC)))
    if px not in _fc:
        _fc[px] = ImageFont.truetype(_PFONT, px)
    return _fc[px]


def _wrap_count(text, pt, avail_cm):
    if not text:
        return 1
    fnt = _f(pt); avail = avail_cm * _PPC
    n, cur = 1, ''
    for ch in text:
        if cur and _probe.textlength(cur + ch, font=fnt) > avail:
            n += 1; cur = ch
        else:
            cur += ch
    return n


def _need(body, w_cm, spc):
    tot = 0.0
    for parts in body:
        txt = ''.join(p[0] for p in parts)
        sz = max(p[1] for p in parts)
        tot += _wrap_count(txt, sz, w_cm) * sz * _PT2CM * spc * _LH
    return tot


def fit_spc(body, w, h, spc, mx=0.14, my=0.05, tag=''):
    aw, ah = w - mx * 2, h - my * 2
    sp = spc
    while sp > 1.0 and _need(body, aw, sp) > ah:
        sp = round(sp - 0.05, 2)
    if _need(body, aw, sp) > ah:
        FITWARN.append('%s: %.2f > %.2f' % (tag or 'box', _need(body, aw, sp), ah))
    return sp


# ---------- 低レベル ----------
def style(run, size, bold, color, font=FONT):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set('typeface', font)


def para(tf, parts, align=PP_ALIGN.LEFT, spc=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if spc:
        p.line_spacing = spc
    for t in parts:
        text, size, color = t[0], t[1], t[2]
        bold = t[3] if len(t) > 3 else True
        font = t[4] if len(t) > 4 else FONT
        r = p.add_run(); r.text = text
        style(r, size, bold, color, font)
    return p


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE):
    tb = s.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    return tf


def shape(s, x, y, w, h, fill=None, line=None, lw=1.0, kind=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(kind, Cm(x), Cm(y), Cm(w), Cm(h))
    st = sh.element.find(qn('p:style'))
    if st is not None:
        sh.element.remove(st)
    sh.shadow.inherit = False
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.14)
    tf.margin_top = tf.margin_bottom = Cm(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def arrow(s, x, y, w, h, kind, fill=ORNG):
    return shape(s, x, y, w, h, fill=fill, line=None, kind=kind)


def plain(s, x, y, w, h, lines, fill=None, line=CHAR, lw=1.0,
          align=PP_ALIGN.CENTER, spc=1.25):
    sh = shape(s, x, y, w, h, fill=fill, line=line, lw=lw)
    sp = fit_spc(lines, w, h, spc, tag='plain')
    for i, parts in enumerate(lines):
        para(sh.text_frame, parts, align, sp, first=(i == 0))
    return sh


def card(s, x, y, w, hh, bh, head, body, body_fill=PALE, head_size=T_HEAD,
         body_size=T_BODY, body_align=PP_ALIGN.LEFT, body_spc=1.5, lw=1.0):
    sh = shape(s, x, y, w, hh, fill=CHAR, line=CHAR, lw=lw)
    para(sh.text_frame, head, PP_ALIGN.CENTER, first=True)
    sb = shape(s, x, y + hh, w, bh, fill=body_fill, line=CHAR, lw=lw)
    sp = fit_spc(body, w, bh, body_spc, tag='card')
    for i, ln in enumerate(body):
        para(sb.text_frame, ln, body_align, sp, first=(i == 0))
    return sh, sb


def photo(s, x, y, w, h, label='図 版 ・ 写 真'):
    sh = shape(s, x, y, w, h, fill=SOFT, line=CHAR, lw=1.0)
    sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    para(sh.text_frame, [(label, T_BODY, GRAY)], PP_ALIGN.CENTER, first=True)
    return sh


SLIDES = []


def page(title, lead_parts=None, strip=True):
    s = prs.slides.add_slide(BLANK)
    SLIDES.append({'s': s, 'strip': strip})
    if title is not None:
        tf = textbox(s, M, 1.55, 12.60, 1.10)
        para(tf, [(title, T_TITLE, INK)], PP_ALIGN.LEFT, first=True)
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(0), Cm(2.90), Cm(W), Cm(2.90))
        c.line.color.rgb = CHAR; c.line.width = Pt(1.75)
        _lg = shape(s, RE - 2.40, 1.53, 2.40, 1.23, fill=SOFT, line=GRAY, lw=1.0)
        _lg.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        para(_lg.text_frame, [('ロ ゴ', T_SUB, GRAY)], PP_ALIGN.CENTER, first=True)
    if lead_parts:
        tf = textbox(s, M, 3.35, CW, 1.00, MSO_ANCHOR.TOP)
        para(tf, lead_parts, PP_ALIGN.LEFT, 1.15, first=True)
    return s


def footer(s, n, strip=True):
    if strip:
        shape(s, 0, 28.60, W, 0.50, fill=CHAR)
    tf = textbox(s, (W - 4.00) / 2, 28.45, 4.00, 0.80)
    para(tf, [(str(n), 9, WHITE, False)], PP_ALIGN.CENTER, first=True)
    tf = textbox(s, RE - 8.00, 28.68, 8.00, 0.40)
    para(tf, [(COPY, 7, WHITE, False)], PP_ALIGN.RIGHT, first=True)


# ================= 本編（A4縦・全12ページ） =================

# --- P1 表紙 ---
s = page(None, strip=False)
shape(s, 0, 0, W, 0.95, fill=ORNG)
tf = textbox(s, M, 0.08, CW, 0.80)
para(tf, [('D O W N L O A D   R E P O R T', 11, WHITE, True, 'Arial')],
     PP_ALIGN.CENTER, first=True)
tf = textbox(s, M, 5.20, CW, 1.20, MSO_ANCHOR.TOP)
para(tf, [('工務店のための', 16, CHAR)], PP_ALIGN.LEFT, first=True)
tf = textbox(s, M, 6.70, CW, 5.20, MSO_ANCHOR.TOP)
para(tf, [('太陽光・蓄電池', 34, CHAR)], PP_ALIGN.LEFT, 1.25, first=True)
para(tf, [('提携モデル 解説', 34, CHAR)], PP_ALIGN.LEFT, 1.25)
shape(s, M, 12.60, 4.20, 0.14, fill=ORNG)
tf = textbox(s, M, 13.30, CW, 2.20, MSO_ANCHOR.TOP)
para(tf, [('― 工事を持たずに、売上をつくる3つの型 ―', 15, ORNG)], PP_ALIGN.LEFT, 1.4, first=True)
para(tf, [('新築1棟あたりの粗利を、施工体制を持たずに上げる方法', 12, INK, False)],
     PP_ALIGN.LEFT, 1.4)
photo(s, M, 16.20, CW, 8.40, '表 紙 ビ ジ ュ ア ル （ 施 工 写 真 な ど ）')
shape(s, M, 25.40, CW, 2.30, fill=PALE, line=GRAY)
_lg = shape(s, M + 0.40, 25.75, 3.40, 1.60, fill=SOFT, line=GRAY, lw=1.0)
_lg.line.dash_style = MSO_LINE_DASH_STYLE.DASH
para(_lg.text_frame, [('ロ　ゴ', T_BODY, GRAY)], PP_ALIGN.CENTER, first=True)
tf = textbox(s, M + 4.30, 25.70, CW - 4.70, 1.85)
para(tf, [(ISSUER, 13, INK)], PP_ALIGN.RIGHT, 1.3, first=True)
para(tf, [('【　　　　　　部　／　担当　　　　　　】', 10, CHAR, False)], PP_ALIGN.RIGHT, 1.3)
para(tf, [('A4縦 ／ 全12ページ ／ 2026年◯月', 9, CHAR, False)], PP_ALIGN.RIGHT, 1.3)

# --- P2 なぜいま太陽光・蓄電池なのか ---
s = page('なぜいま、太陽光・蓄電池なのか',
         [('■　', T_LEAD, INK), ('制度・電気代・施主の関心', T_LEADEM, ORNG),
          ('の3つが、同じ方向に動いています', T_LEAD, INK)])
_y = CT
for _hd, _ls in [
    ('① 制　度｜省エネ基準への適合が前提になった',
     ['・2025年4月から、原則すべての新築住宅で省エネ基準への適合が義務化【要確認】',
      '・自治体によっては新築への太陽光設置を求める制度も始まっている【要確認】',
      '・「載せる／載せない」を施主に説明できることが、実務上の前提になりつつある']),
    ('② 電気代｜施主が自分ごととして気にしている',
     ['・電気料金の水準は【　　】年比で【　　】％上昇【出典を入れる】',
      '・光熱費は入居後30年ぶん続く支出。建物価格より効く場面がある',
      '・「毎月いくら下がるか」を出せると、商談の温度が上がる']),
    ('③ 施　主｜比較の土俵が「性能」に移った',
     ['・住宅ローン減税・補助金の要件に省エネ性能が絡む【要確認】',
      '・他社が提案していれば、提案しない工務店は候補から外れる',
      '・逆に言えば、提案できるだけで差がつく状態がまだ続いている'])]:
    card(s, M, _y, CW, 1.05, 2.85, [(_hd, T_HEAD, WHITE)],
         [[(t, T_BODY, INK)] for t in _ls])
    _y += 4.30
plain(s, M, 17.50, CW, 2.20,
      [[('とはいえ、多くの工務店が', T_LEADEM, INK)],
       [('「提案したいが、体制がない」', T_LEADEM, ORNG), ('で止まっています。', T_LEADEM, INK)]],
      fill=PALE, spc=1.4)
photo(s, M, 20.10, CW, 5.60, '市 場 デ ー タ の グ ラ フ （ 電 気 代 推 移 な ど ）')
tf = textbox(s, M, 25.90, CW, 0.90, MSO_ANCHOR.TOP)
para(tf, [('※ 制度・数値は執筆時点の情報です。最新の公表資料でご確認ください。', T_SUB, GRAY, False)],
     PP_ALIGN.LEFT, first=True)

# --- P3 取りこぼしの構造 ---
s = page('工務店が取りこぼしている売上　①',
         [('■　提案しない理由は、やる気ではなく', T_LEAD, INK),
          ('体制', T_LEADEM, ORNG), ('にあります', T_LEAD, INK)])
for x, hd, ls in zip(X3,
        ['答えられない', '施工できない', '責任を持てない'],
        [['・機器の選定基準がわからない', '・発電量シミュレーションを', '　出せない',
          '・見積の内訳を説明できない'],
         ['・電気工事の職人がいない', '・系統連系の申請が煩雑', '・屋根の納まりに不安が残る'],
         ['・メーカー保証の窓口が不明', '・雨漏りの責任範囲が曖昧', '・10年後の点検を誰がやるか']]):
    card(s, x, CT, W3, 1.05, 3.60, [(hd, T_HEAD, WHITE)],
         [[(t, T_BODY, INK)] for t in ls])
arrow(s, (W - 1.60) / 2, 9.55, 1.60, 1.00, MSO_SHAPE.DOWN_ARROW, ORNG)
plain(s, M, 10.95, CW, 2.60,
      [[('結果、', T_LEADEM, INK), ('「うちではできません」', T_LEADEM, ORNG),
        ('と答えることになり、', T_LEADEM, INK)],
       [('施主は太陽光の専門業者へ流れていきます。', T_LEADEM, INK)]],
      fill=PALE, spc=1.4)
card(s, M, 14.10, CW, 1.05, 5.20, [('そのとき、失っているもの', T_HEAD, WHITE)],
     [[('・搭載工事の粗利（1棟あたり【　　】万円）', T_BODY, INK)],
      [('・専門業者が入ることによる、引き渡し前後の主導権', T_BODY, INK)],
      [('・「あの工務店は詳しい」という評判と、そこからの紹介', T_BODY, INK)],
      [('・入居後のメンテナンス・リフォームへの接点', T_BODY, INK)]], body_spc=1.7)
photo(s, M, 20.70, CW, 5.90, '取 り こ ぼ し の 構 造 図')

# --- P4 いくら取りこぼしているか ---
s = page('工務店が取りこぼしている売上　②',
         [('■　年20棟の工務店なら、', T_LEAD, INK), ('年間【　　】万円', T_LEADEM, ORNG),
          ('の機会損失', T_LEAD, INK)])
card(s, M, CT, CW, 1.05, 2.20, [('試算の式', T_HEAD, AMBER)],
     [[('年間棟数 × 搭載率 × 1棟あたりの粗利 ＝ 取りこぼしている粗利', T_LEADEM, INK)]],
     body_fill=WHITE, body_align=PP_ALIGN.CENTER, lw=1.5)
plain(s, M, 8.35, CW, 0.85,
      [[('年間棟数と搭載率でみた、1年あたりの粗利（1棟あたり37万円で試算）', T_BODY, WHITE)]],
      fill=CHAR, line=CHAR)
cwT = [4.45, 4.45, 4.45, 4.45]
cxT = [M]
for w in cwT[:-1]:
    cxT.append(cxT[-1] + w)
for x, w, t in zip(cxT, cwT, ['年間棟数', '搭載率 20％', '搭載率 30％', '搭載率 40％']):
    plain(s, x, 9.20, w, 0.85, [[(t, T_BODY, WHITE)]], fill=CHAR, line=CHAR)
_y = 10.05
for _n, _v in [('10棟', ('74万円', '111万円', '148万円')),
               ('20棟', ('148万円', '222万円', '296万円')),
               ('30棟', ('222万円', '333万円', '444万円'))]:
    plain(s, cxT[0], _y, cwT[0], 1.00, [[(_n, T_BODY, CHAR)]], fill=PALE, line=CHAR)
    for x, w, v in zip(cxT[1:], cwT[1:], _v):
        plain(s, x, _y, w, 1.00, [[(v, T_LEAD, ORNG)]], fill=WHITE, line=CHAR)
    _y += 1.00
tf = textbox(s, M, 13.20, CW, 0.80, MSO_ANCHOR.TOP)
para(tf, [('※ 1棟あたりの粗利は「販売代行型」（平均単価250万円 × 粗利率15％）で試算しています。'
           '型によって変わります。', T_SUB, GRAY, False)], PP_ALIGN.LEFT, first=True)
card(s, M, 14.40, CW, 1.05, 4.60, [('この数字の読み方', T_HEAD, WHITE)],
     [[('・搭載率は、提案しなければ0％。提案すれば0％にはなりません', T_BODY, INK)],
      [('・年20棟・搭載率30％なら、6棟ぶんの粗利が毎年出ていっています', T_BODY, INK)],
      [('・棟数を増やさずに粗利を上げられる、数少ない打ち手です', T_BODY, INK)]], body_spc=1.7)
photo(s, M, 20.40, CW, 6.20, '機 会 損 失 の グ ラ フ')


def model_page(no, name, sub, roles, merit, demerit, income, fit):
    """提携モデルのページ。役割分担表＋メリット／注意点＋1棟あたりの収益。"""
    s = page('提携モデル%s　%s' % (no, name),
             [('■　', T_LEAD, INK), (sub, T_LEADEM, ORNG)])
    plain(s, M, CT, CW, 0.85, [[('役割分担', T_BODY, WHITE)]], fill=CHAR, line=CHAR)
    cwR = [4.60, 6.60, 6.60]
    cxR = [M]
    for w in cwR[:-1]:
        cxR.append(cxR[-1] + w)
    for x, w, t in zip(cxR, cwR, ['工　程', '貴社（工務店）', '当社（ワン・ミニット）']):
        plain(s, x, 5.45, w, 0.85, [[(t, T_BODY, WHITE)]], fill=CHAR, line=CHAR)
    _y = 6.30
    for _p, _a, _b in roles:
        plain(s, cxR[0], _y, cwR[0], 0.95, [[(_p, T_BODY, CHAR)]], fill=PALE, line=CHAR)
        plain(s, cxR[1], _y, cwR[1], 0.95, [[(_a, T_BODY, INK)]], fill=WHITE, line=CHAR)
        plain(s, cxR[2], _y, cwR[2], 0.95, [[(_b, T_BODY, INK)]], fill=WHITE, line=CHAR)
        _y += 0.95
    _y += 0.55
    card(s, X2[0], _y, W2, 1.00, 3.60, [('向いている工務店', T_HEAD, WHITE)],
         [[(t, T_BODY, INK)] for t in fit], body_spc=1.6)
    card(s, X2[1], _y, W2, 1.00, 3.60, [('注意しておくこと', T_HEAD, WHITE)],
         [[(t, T_BODY, INK)] for t in demerit], body_spc=1.6)
    _y += 5.00
    card(s, M, _y, CW, 1.00, 2.30, [('1棟あたりの収益イメージ', T_HEAD, AMBER)],
         [[(income, T_LEADEM, INK)]], body_fill=WHITE, body_align=PP_ALIGN.CENTER, lw=1.5)
    _y += 3.60
    plain(s, M, _y, CW, CB - _y - 1.10,
          [[('ひとことで言うと', T_BODY, CHAR)], [(merit, T_LEADEM, INK)]],
          fill=PALE, spc=1.5)
    tf = textbox(s, M, CB - 0.95, CW, 0.80, MSO_ANCHOR.TOP)
    para(tf, [('※ 収益は平均単価250万円・粗利率15％で試算した目安です。実際の条件は個別にご相談ください。',
               T_SUB, GRAY, False)], PP_ALIGN.LEFT, first=True)
    return s


# --- P5 紹介型 ---
model_page('①', '紹介型', '施主を紹介するだけ。手を動かさずに手数料を受け取る',
           [('提　案', '施主に当社を紹介', '訪問・提案・見積'),
            ('機器調達', '―', '当社が手配'),
            ('施　工', '―', '当社が施工'),
            ('申　請', '―', '当社が代行'),
            ('保　証', '―', '当社が窓口'),
            ('収　益', '紹介手数料', '工事売上')],
           '手間はほぼゼロ。まずここから始めて、感触をつかむ型です。',
           ['・工事の粗利は取れない', '・施主との接点が当社に移る',
            '・「丸投げ」に見えない伝え方が要る'],
           '紹介手数料　250万円 × 5％ ＝ 約 12万円 ／ 棟',
           ['・まず1棟やってみたい', '・電気工事の職人がいない',
            '・引き渡し前で手が回らない'])

# --- P6 販売代行型 ---
model_page('②', '販売代行型', '提案と受注は貴社。施工だけを当社が引き受ける',
           [('提　案', '貴社が提案・見積', '資料と試算を提供'),
            ('機器調達', '―', '当社が手配'),
            ('施　工', '―', '当社が施工'),
            ('申　請', '―', '当社が代行'),
            ('保　証', '一次窓口', '当社がバックアップ'),
            ('収　益', '工事粗利', '仕入・施工')],
           '施主との関係を持ったまま、粗利を取りにいく型です。',
           ['・提案できる人を1人つくる必要がある', '・見積・試算の作成に時間がかかる',
            '・一次対応は貴社が受ける'],
           '工事粗利　250万円 × 15％ ＝ 約 37万円 ／ 棟',
           ['・自社で提案までやりたい', '・施主との関係を手放したくない',
            '・年10棟以上を手がけている'])

# --- P7 共同施工型 ---
model_page('③', '共同施工型', '電気工事など一部を貴社が担当し、工事売上も取る',
           [('提　案', '貴社が提案・見積', '資料と試算を提供'),
            ('機器調達', '―', '当社が手配'),
            ('施　工', '電気工事・屋根まわり', '架台・パネル・蓄電池'),
            ('申　請', '―', '当社が代行'),
            ('保　証', '一次窓口', '当社がバックアップ'),
            ('収　益', '工事粗利 ＋ 施工売上', '仕入・施工')],
           '3つの型でいちばん残るぶん、体制の準備も必要になります。',
           ['・電気工事士の確保が前提', '・工程の擦り合わせが要る',
            '・初回は当社が同行して手順を共有'],
           '工事粗利 37万円 ＋ 施工分 約 9万円 ＝ 約 46万円 ／ 棟',
           ['・電気工事の職人が社内にいる', '・年20棟以上を手がけている',
            '・将来は自社完結を目指したい'])


# --- P8 導入事例① ---
s = page('導入事例　T社＠山梨県　①',
         [('■　', T_LEAD, INK), ('「売り方がわからない」', T_LEADEM, ORNG),
          ('ところから始めた工務店の話', T_LEAD, INK)])
plain(s, M, CT, CW, 0.85, [[('会社概要', T_BODY, WHITE)]], fill=CHAR, line=CHAR)
cwP = [4.60, 13.20]
_y = 5.45
for _k, _v in [('所 在 地', '山梨県【　　　市】'), ('年間棟数', '【　　】棟'),
               ('従業員数', '【　　】名'), ('主力商品', '【　　　　　　　　　　　】'),
               ('提携時期', '【　　】年【　　】月')]:
    plain(s, M, _y, cwP[0], 0.90, [[(_k, T_BODY, CHAR)]], fill=PALE, line=CHAR)
    plain(s, M + cwP[0], _y, cwP[1], 0.90, [[(_v, T_BODY, INK)]], fill=WHITE, line=CHAR,
          align=PP_ALIGN.LEFT)
    _y += 0.90
card(s, M, 10.30, CW, 1.05, 4.20, [('提携前の課題', T_HEAD, WHITE)],
     [[('・施主から太陽光の相談を受けても、その場で答えられなかった', T_BODY, INK)],
      [('・見積を出すたびに他社に比較され、価格の話だけになっていた', T_BODY, INK)],
      [('・電気工事の職人がおらず、受けても回せる見込みがなかった', T_BODY, INK)]], body_spc=1.7)
card(s, M, 15.85, CW, 1.05, 5.20, [('取り組んだこと', T_HEAD, WHITE)],
     [[('・まず「紹介型」で1棟だけ試し、施工品質と対応を自分の目で確認した', T_BODY, INK)],
      [('・2棟目から「販売代行型」に切り替え、提案と見積を自社で持った', T_BODY, INK)],
      [('・当社の試算シートを商談で使い、毎月の光熱費で説明するようにした', T_BODY, INK)],
      [('・引き渡し後の点検も自社の名前で案内し、接点を残した', T_BODY, INK)]], body_spc=1.7)
photo(s, M, 22.40, CW, 4.20, '担 当 者 の 写 真 ／ 現 場 写 真')

# --- P9 導入事例② ---
s = page('導入事例　T社＠山梨県　②',
         [('■　', T_LEAD, INK), ('提案するようになっただけ', T_LEADEM, ORNG),
          ('で、数字はここまで動きました', T_LEAD, INK)])
plain(s, M, CT, CW, 0.85, [[('提携前 → 提携後（1年）', T_BODY, WHITE)]], fill=CHAR, line=CHAR)
cwK = [7.00, 3.60, 3.60, 3.60]
cxK = [M]
for w in cwK[:-1]:
    cxK.append(cxK[-1] + w)
for x, w, t in zip(cxK, cwK, ['項　目', '提携前', '提携後', '増　減']):
    plain(s, x, 5.45, w, 0.85, [[(t, T_BODY, WHITE)]], fill=CHAR, line=CHAR)
_y = 6.30
for _k, _a, _b, _c in [('年間棟数', '【　】棟', '【　】棟', '【　】棟'),
                       ('太陽光の搭載率', '【　】％', '【　】％', '＋【　】pt'),
                       ('搭載した棟数', '【　】棟', '【　】棟', '＋【　】棟'),
                       ('太陽光の粗利（年）', '【　】万円', '【　】万円', '＋【　】万円'),
                       ('1棟あたり粗利', '【　】万円', '【　】万円', '＋【　】万円')]:
    plain(s, cxK[0], _y, cwK[0], 1.00, [[(_k, T_BODY, CHAR)]], fill=PALE, line=CHAR,
          align=PP_ALIGN.LEFT)
    for x, w, v, col in zip(cxK[1:], cwK[1:], (_a, _b, _c), (INK, ORNG, ORNG)):
        plain(s, x, _y, w, 1.00, [[(v, T_BODY, col)]], fill=WHITE, line=CHAR)
    _y += 1.00
tf = textbox(s, M, 11.45, CW, 0.80, MSO_ANCHOR.TOP)
para(tf, [('※ 実績値はヒアリング後に確定させます。', T_SUB, GRAY, False)], PP_ALIGN.LEFT, first=True)
card(s, M, 12.60, CW, 1.05, 5.20, [('うまくいった理由', T_HEAD, WHITE)],
     [[('・いきなり全部やらず、紹介型で1棟だけ試したこと', T_BODY, INK)],
      [('・「毎月いくら下がるか」で説明したこと。機器の性能では話さない', T_BODY, INK)],
      [('・断られた施主にも試算だけは渡し、着工前にもう一度声をかけたこと', T_BODY, INK)],
      [('・引き渡し後の点検を自社の名前で案内し、紹介が生まれたこと', T_BODY, INK)]], body_spc=1.7)
plain(s, M, 18.90, CW, 2.60,
      [[('「うちではできません」を', T_LEADEM, INK), ('やめただけ', T_LEADEM, ORNG),
        ('です。', T_LEADEM, INK)],
       [('新しく雇った人も、増やした棟数もありません。', T_LEADEM, INK)]],
      fill=PALE, spc=1.4)
photo(s, M, 22.10, CW, 4.50, '成 果 の グ ラ フ ／ 施 主 の 声')

# --- P10 収益シミュレーション ---
s = page('収益シミュレーション',
         [('■　年20棟・搭載率30％（年6棟）で、', T_LEAD, INK),
          ('型ごとにいくら残るか', T_LEADEM, ORNG)])
plain(s, M, CT, CW, 0.85, [[('3つの型の比較（年6棟で試算）', T_BODY, WHITE)]],
      fill=CHAR, line=CHAR)
cwS = [4.40, 4.40, 4.50, 4.50]
cxS = [M]
for w in cwS[:-1]:
    cxS.append(cxS[-1] + w)
for x, w, t in zip(cxS, cwS, ['', '① 紹介型', '② 販売代行型', '③ 共同施工型']):
    plain(s, x, 5.45, w, 0.85, [[(t, T_BODY, WHITE)]], fill=CHAR, line=CHAR)
_y = 6.30
for _k, _v, _big in [('1棟あたりの収益', ('12万円', '37万円', '46万円'), False),
                     ('年6棟での収益', ('75万円', '225万円', '279万円'), True),
                     ('貴社の手間', ('紹介のみ', '提案・見積', '提案・見積＋工事'), False),
                     ('必要な体制', ('なし', '提案できる人1名', '提案1名＋電気工事士'), False)]:
    plain(s, cxS[0], _y, cwS[0], 1.05, [[(_k, T_BODY, CHAR)]], fill=PALE, line=CHAR,
          align=PP_ALIGN.LEFT)
    for x, w, v in zip(cxS[1:], cwS[1:], _v):
        plain(s, x, _y, w, 1.05, [[(v, T_NUM_M if _big else T_BODY, ORNG if _big else INK)]],
              fill=PALE if _big else WHITE, line=CHAR)
    _y += 1.05
tf = textbox(s, M, 10.70, CW, 0.80, MSO_ANCHOR.TOP)
para(tf, [('※ 平均単価250万円・粗利率15％で試算した目安です。地域・機器・屋根条件で変わります。',
           T_SUB, GRAY, False)], PP_ALIGN.LEFT, first=True)
card(s, M, 11.90, CW, 1.05, 4.60, [('どこから始めるか', T_HEAD, WHITE)],
     [[('・1年目は ① 紹介型 で1〜2棟。施工品質と対応を自分の目で確かめる', T_BODY, INK)],
      [('・感触がよければ ② 販売代行型 へ。ここが収益と手間のバランス点です', T_BODY, INK)],
      [('・電気工事士がいる会社だけ、③ 共同施工型 まで進める価値があります', T_BODY, INK)]],
     body_spc=1.7)
plain(s, M, 17.80, CW, 2.60,
      [[('②から始めれば、年20棟の工務店で', T_LEADEM, INK)],
       [('年間 225万円', T_NUM_L, ORNG), ('。棟数は増やさずに、です。', T_LEADEM, INK)]],
      fill=PALE, spc=1.4)
photo(s, M, 21.00, CW, 5.60, '3 つ の 型 の 比 較 グ ラ フ')


# --- P11 導入の流れ ---
s = page('導入の流れ',
         [('■　初回のご相談から1棟目の引き渡しまで、', T_LEAD, INK),
          ('およそ3〜4ヶ月', T_LEADEM, ORNG)])
_y = CT
for _no, _hd, _ls in [
    ('STEP 1', '個別相談（オンライン・60分）',
     ['・現在の棟数・搭載率・体制をうかがいます',
      '・3つの型のうち、どこから始めるのが合うかをその場でご提案します']),
    ('STEP 2', '試算のご提示（1週間以内）',
     ['・貴社の棟数に合わせた収益シミュレーションをお出しします',
      '・想定される手間と、必要な体制も具体的にお伝えします']),
    ('STEP 3', '提携内容のすり合わせ（2〜3週間）',
     ['・役割分担・収益条件・保証の範囲を書面で確定させます',
      '・引き渡し後のメンテナンス体制もここで決めます']),
    ('STEP 4', '1棟目の実施（1〜2ヶ月）',
     ['・初回は当社が同行し、提案から施工まで一緒に進めます',
      '・試算シート・提案資料・チラシ類は当社がご用意します']),
    ('STEP 5', '振り返りと本格稼働',
     ['・1棟目の結果を見て、型を続けるか変えるかを判断します',
      '・以降は貴社主体で回し、当社は施工とバックアップに回ります'])]:
    plain(s, M, _y, 3.30, 1.00, [[(_no, T_BODY, WHITE)]], fill=ORNG, line=ORNG)
    plain(s, M + 3.30, _y, CW - 3.30, 1.00, [[(_hd, T_HEAD, WHITE)]],
          fill=CHAR, line=CHAR, align=PP_ALIGN.LEFT)
    plain(s, M, _y + 1.00, CW, 2.30, [[(t, T_BODY, INK)] for t in _ls],
          fill=WHITE, line=CHAR, align=PP_ALIGN.LEFT, spc=1.6)
    _y += 3.75
plain(s, M, 23.55, CW, 2.10,
      [[('1棟目までは、当社が横について進めます。', T_LEADEM, INK)],
       [('いきなり自社だけで回す必要はありません。', T_LEADEM, ORNG)]],
      fill=PALE, spc=1.4)

# --- P12 個別相談のご案内 ---
s = page('個別相談のご案内',
         [('■　', T_LEAD, INK), ('無料・オンライン60分', T_LEADEM, ORNG),
          ('。まず現状をうかがうところから', T_LEAD, INK)])
card(s, M, CT, CW, 1.05, 4.80, [('こんな方に向いています', T_HEAD, WHITE)],
     [[('・施主から太陽光の相談を受けたが、答えられなかったことがある', T_BODY, INK)],
      [('・提案したいが、施工体制がないので踏み出せていない', T_BODY, INK)],
      [('・他社に太陽光だけ持っていかれた経験がある', T_BODY, INK)],
      [('・棟数を増やさずに、1棟あたりの粗利を上げたい', T_BODY, INK)]], body_spc=1.7)
card(s, M, 10.75, CW, 1.05, 4.20, [('当日おうかがいすること', T_HEAD, WHITE)],
     [[('・年間棟数と、直近の太陽光の搭載率', T_BODY, INK)],
      [('・社内に電気工事士がいるかどうか', T_BODY, INK)],
      [('・いつまでに何棟で試したいか', T_BODY, INK)]], body_spc=1.7)
plain(s, M, 16.15, CW, 1.05, [[('お申し込み', T_HEAD, WHITE)]], fill=CHAR, line=CHAR)
_y = 17.20
for _k, _v in [('お 電 話', '【　　-　　　　-　　　　】（平日9:00〜18:00）'),
               ('メ ー ル', '【　　　　　　　　　　＠　　　　　　　　　】'),
               ('W E B', '【　　　　　　　　　　　　　　　　　　　　】'),
               ('担　　当', '【　　　　　　部】　【　　　　　　　　】')]:
    plain(s, M, _y, 4.20, 1.05, [[(_k, T_BODY, CHAR)]], fill=PALE, line=CHAR)
    plain(s, M + 4.20, _y, CW - 4.20, 1.05, [[(_v, T_BODY, INK)]], fill=WHITE, line=CHAR,
          align=PP_ALIGN.LEFT)
    _y += 1.05
photo(s, M, 21.90, 8.60, 3.20, '二 次 元 コ ー ド')
plain(s, M + 8.90, 21.90, CW - 8.90, 3.20,
      [[('この二次元コードから', T_BODY, INK)], [('そのままご予約いただけます', T_BODY, INK)]],
      fill=PALE, spc=1.5)
plain(s, M, 25.50, CW, 1.60,
      [[('まずは60分。合わなければ、それで終わりで構いません。', T_LEADEM, INK)]],
      fill=CHAR, line=CHAR)
s.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = AMBER


# ================= 保存 =================
for i, r in enumerate(SLIDES, start=1):
    footer(r['s'], i, strip=r['strip'])

os.makedirs('output', exist_ok=True)
out = 'output/DLレポート第1号_工務店のための太陽光蓄電池提携モデル解説_たたき台.pptx'
prs.save(out)
for w in FITWARN:
    print('FIT WARNING:', w)
print('saved:', out, '/ pages:', len(prs.slides._sldIdLst), '/ warnings:', len(FITWARN))
