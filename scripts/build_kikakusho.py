# -*- coding: utf-8 -*-
"""企画書 空パッケージ ジェネレータ
prompts/企画書_空パッケージ_生成プロンプト.md の仕様に完全準拠。
レイアウトは §グリッド（M / G / 列幅）と §タイプスケール のみを使う。
"""
import os
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

# ---------- パレット ----------
NAVY  = RGBColor(0x00, 0x20, 0x60)
INK   = RGBColor(0x1F, 0x1F, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED   = RGBColor(0xFF, 0x00, 0x00)
YEL   = RGBColor(0xFF, 0xFF, 0x00)
PALE  = RGBColor(0xE9, 0xEE, 0xF6)
GRAY  = RGBColor(0xBF, 0xBF, 0xBF)
SOFT  = RGBColor(0xF2, 0xF2, 0xF2)

FONT = 'Meiryo UI'
LOGO = 'templates/assets/funai_logo.png'
CASE_HOUMON = 'templates/assets/case_houmon.png'
CASE_MAIL   = 'templates/assets/case_mailmaga.png'
COPY = 'Copyright Funai Consulting Inc. All Rights Reserved.'

# ---------- キャンバス（A4横） ----------
W, H = 27.52, 19.05

# ---------- グリッド ----------
M   = 0.30                    # 左右マージン（全要素の左端・右端はここに揃える）
G   = 0.32                    # 列間ガター
CW  = W - M * 2               # 26.92 本文の全幅
CT  = 4.05                    # 本文エリア上端
CB  = 15.20                   # 本文エリア下端（全ページ共通）
RE  = M + CW                  # 27.22 本文の右端


def cols(n, gutter=G):
    """n分割したときの (x座標のリスト, 列幅) を返す。"""
    w = (CW - gutter * (n - 1)) / n
    return [M + i * (w + gutter) for i in range(n)], w


X2, W2 = cols(2)              # 13.30
X3, W3 = cols(3)              # 8.76
X4, W4 = cols(4)              # 6.49

# ---------- タイプスケール ----------
# 54/44 表紙 ／ 32 大書 ／ 28 ページ見出し・大数値 ／ 26 キーメッセージ
# 24 リード強調 ／ 22 まとめ帯・中数値 ／ 20 リード・公式 ／ 18 カード見出し
# 16 章タグ・表ヘッダー ／ 14 本文 ／ 12 補助 ／ 9-8 フッター
T_TITLE, T_LEAD, T_LEADEM = 28, 20, 24
T_TAG, T_HEAD, T_BODY, T_SUB = 16, 18, 14, 12
T_SUM, T_KEY = 22, 26
T_NUM_L, T_NUM_M = 28, 22

prs = Presentation()
prs.slide_width  = Emu(9906000)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]


# ---------- 低レベルヘルパ ----------
def style(run, size=T_BODY, bold=True, color=INK, font=FONT):
    f = run.font
    f.size = Pt(size); f.bold = bold
    f.color.rgb = color
    f.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set('typeface', font)


def para(tf, parts, align=PP_ALIGN.LEFT, spc=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if spc:
        p.line_spacing = spc
    for t in parts:
        text, size, color = t[0], t[1], t[2]
        bold = t[3] if len(t) > 3 else True
        font = t[4] if len(t) > 4 else FONT
        r = p.add_run(); r.text = text
        style(r, size, bold, color, font)
    return p


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE):
    tb = s.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    return tf


def shape(s, x, y, w, h, fill=None, line=None, lw=1.0, kind=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(kind, Cm(x), Cm(y), Cm(w), Cm(h))
    st = sh.element.find(qn('p:style'))
    if st is not None:
        sh.element.remove(st)
    sh.shadow.inherit = False
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.16)
    tf.margin_top = tf.margin_bottom = Cm(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def arrow(s, x, y, w, h, kind, fill=RED):
    return shape(s, x, y, w, h, fill=fill, line=None, kind=kind)



# ---------- 自動フィット（箱に対して行間→字数を測って収める） ----------
from PIL import Image, ImageDraw, ImageFont

_PFONT = '/usr/share/fonts/opentype/ipafont-gothic/ipagp.ttf'
_PPC, _PT2CM, _LH = 50.0, 2.54 / 72.0, 1.22
_fcache = {}
_probe = ImageDraw.Draw(Image.new('RGB', (8, 8)))
FITWARN = []


def _f(pt):
    px = max(6, int(round(pt * _PT2CM * _PPC)))
    if px not in _fcache:
        _fcache[px] = ImageFont.truetype(_PFONT, px)
    return _fcache[px]


def _wrap_count(text, pt, avail_cm):
    if not text:
        return 1
    fnt = _f(pt); avail = avail_cm * _PPC
    n, cur = 1, ''
    for ch in text:
        if cur and _probe.textlength(cur + ch, font=fnt) > avail:
            n += 1; cur = ch
        else:
            cur += ch
    return n


def _need(body, w_cm, spc):
    tot = 0.0
    for parts in body:
        txt = ''.join(p[0] for p in parts)
        sz = max(p[1] for p in parts)
        tot += _wrap_count(txt, sz, w_cm) * sz * _PT2CM * spc * _LH
    return tot


def fit_spc(body, w, h, spc, mx=0.16, my=0.06, tag=''):
    """箱に収まる最大の行間を返す。1.0でも収まらなければ警告を残す。"""
    aw, ah = w - mx * 2, h - my * 2
    sp = spc
    while sp > 1.0 and _need(body, aw, sp) > ah:
        sp = round(sp - 0.05, 2)
    if _need(body, aw, sp) > ah:
        FITWARN.append('%s: %.2fcm > %.2fcm' % (tag or 'box', _need(body, aw, sp), ah))
    return sp


# ---------- 共通 chrome ----------
def footer(s, page, strip=True):
    if strip:
        shape(s, 0, 18.55, W, 0.50, fill=NAVY)
    tf = textbox(s, (W - 6.42) / 2, 18.31, 6.42, 1.01)
    para(tf, [(str(page), 9, WHITE, False)], PP_ALIGN.CENTER, first=True)
    tf = textbox(s, RE - 9.00, 18.47, 9.00, 0.60)
    para(tf, [(COPY, 8, WHITE, False, 'Arial')], PP_ALIGN.RIGHT, first=True)


def header(s, title, tag):
    tf = textbox(s, M, 0.30, 12.78, 1.30)
    para(tf, [(title, T_TITLE, INK)], PP_ALIGN.LEFT, first=True)
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(0), Cm(1.75), Cm(W), Cm(1.75))
    c.line.color.rgb = NAVY; c.line.width = Pt(2.25)
    s.shapes.add_picture(LOGO, Cm(RE - 2.72), Cm(0.28), Cm(2.72), Cm(1.40))
    sh = shape(s, M, 1.95, 7.20, 0.85, fill=NAVY)
    para(sh.text_frame, [(tag, T_TAG, WHITE)], PP_ALIGN.CENTER, first=True)


def lead(s, parts):
    tf = textbox(s, M, 2.90, CW, 1.15)
    para(tf, parts, PP_ALIGN.LEFT, 1.0, first=True)


def band(s, parts):
    shape(s, 0, 15.93, W, 2.40, fill=NAVY)
    tf = textbox(s, M, 15.93, CW, 2.40)
    para(tf, parts, PP_ALIGN.CENTER, 0.9, first=True)


def keymsg(s, msg, size=T_KEY):
    band(s, [(msg, size, YEL)])


PENDING = []


def rec(owner, pending):
    PENDING.append((len(SLIDES), owner, pending))


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ---------- 高レベル部品 ----------
def card(s, x, y, w, hh, bh, head, body, body_fill=PALE, head_size=T_HEAD,
         body_size=T_BODY, body_align=PP_ALIGN.LEFT, body_spc=1.45, lw=1.5):
    sh = shape(s, x, y, w, hh, fill=NAVY, line=NAVY, lw=lw)
    para(sh.text_frame, head, PP_ALIGN.CENTER, first=True)
    sb = shape(s, x, y + hh, w, bh, fill=body_fill, line=NAVY, lw=lw)
    sp = fit_spc(body, w, bh, body_spc, tag='card body')
    for i, line_parts in enumerate(body):
        para(sb.text_frame, line_parts, body_align, sp, first=(i == 0))
    return sh, sb


def plain(s, x, y, w, h, lines, fill=None, line=NAVY, lw=1.5,
          align=PP_ALIGN.CENTER, spc=1.25):
    sh = shape(s, x, y, w, h, fill=fill, line=line, lw=lw)
    sp = fit_spc(lines, w, h, spc, tag='plain')
    for i, parts in enumerate(lines):
        para(sh.text_frame, parts, align, sp, first=(i == 0))
    return sh


def photo(s, x, y, w, h, label='写 真'):
    sh = shape(s, x, y, w, h, fill=SOFT, line=NAVY, lw=1.0)
    sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    para(sh.text_frame, [(label, T_BODY, GRAY)], PP_ALIGN.CENTER, first=True)
    return sh


def bullets(s, x, y, w, h, lines, size=T_BODY, spc=1.5, fill=WHITE):
    """カード本体の箇条書き。行間と余白をここで一元管理する。"""
    sb = shape(s, x, y, w, h, fill=fill, line=NAVY, lw=1.5)
    body = [t if isinstance(t, list) else [(t, size, INK)] for t in lines]
    sp = fit_spc(body, w, h, spc, tag='bullets')
    for i, parts in enumerate(body):
        para(sb.text_frame, parts, PP_ALIGN.LEFT, sp, first=(i == 0))
    return sb


# ---------- スライド種別 ----------
SLIDES = []


def slide_blank():
    s = prs.slides.add_slide(BLANK)
    SLIDES.append({'s': s, 'strip': True})
    return s


def cover():
    s = slide_blank()
    shape(s, 0, 0, W, H, fill=NAVY)
    tf = textbox(s, 1.40, 4.64, 18.97, 2.00)
    para(tf, [('[ 企画書 ]　', 44, WHITE)], PP_ALIGN.LEFT, first=True)
    tf = textbox(s, 1.40, 6.53, 19.12, 5.40)
    para(tf, [('アライアンス提携による', 54, WHITE)], PP_ALIGN.LEFT, 1.15, first=True)
    para(tf, [('名簿獲得戦略', 54, WHITE)], PP_ALIGN.LEFT, 1.15)
    shape(s, M, 15.61, CW, 2.50, fill=WHITE)
    s.shapes.add_picture(LOGO, Cm(0.90), Cm(16.25), Cm(3.19), Cm(1.64))
    tf = textbox(s, 9.31, 16.07, RE - 9.31, 1.58)
    para(tf, [('株式会社 船井総合研究所', T_TAG, INK)], PP_ALIGN.RIGHT, 1.0, first=True)
    para(tf, [('専門工事支援部　住宅エネルギーチーム　佐野、伊木、岸野', T_TAG, INK)],
         PP_ALIGN.RIGHT, 1.0)
    SLIDES[-1]['strip'] = False
    notes(s, '表紙。会議名・開催日を入れる場合はタイトル下に追記する。')
    return s


def toc():
    s = slide_blank()
    header(s, '目次', '本日の検討範囲')
    items = [
        ('①', 'Ⅰ. 結論', '決めていただきたいこと', 'P3'),
        ('②', 'Ⅱ. 戦略', 'センターピンはどこか', 'P4'),
        ('③', 'Ⅲ. 施策', 'どう獲り、育て、刈るか', 'P8'),
        ('④', 'Ⅳ. 数値', 'いくら投じ、いくら返るか', 'P14'),
        ('⑤', 'Ⅴ. 実行', '誰が、いつまでに', 'P16'),
        ('⑥', 'Ⅵ. 別軸', '本編とは切り離して見る2件', 'P25'),
    ]
    y = 4.42
    for num, label, sub, pg in items:
        tf = textbox(s, 1.70, y, 1.94, 1.40)
        para(tf, [(num, T_NUM_L, INK)], PP_ALIGN.CENTER, first=True)
        tf = textbox(s, 3.90, y, 7.50, 1.40)
        para(tf, [(label, T_TITLE, INK)], PP_ALIGN.LEFT, first=True)
        tf = textbox(s, 11.60, y, 11.00, 1.40)
        para(tf, [('～ ' + sub + ' ～', T_HEAD, NAVY)], PP_ALIGN.LEFT, first=True)
        tf = textbox(s, 23.00, y, RE - 23.00, 1.40)
        para(tf, [(pg, T_LEAD, NAVY)], PP_ALIGN.RIGHT, first=True)
        y += 1.80
    notes(s, '各章の掲載ページを示す。')
    return s


def case_page(tag, lead_parts, point, results, owner, img=None, img_rect=None,
              note_text=None, pending=None):
    """事例ページ。POINT欄（4.05〜5.33）と成果帯、その間に図版を置く。"""
    s = slide_blank(); header(s, 'Ⅲ. 施策', tag)
    lead(s, lead_parts)
    sh = shape(s, M, CT, 3.30, 1.28, fill=RED, line=RED, lw=1.0)
    para(sh.text_frame, [('P O I N T', T_HEAD, WHITE)], PP_ALIGN.CENTER, first=True)
    plain(s, 3.94, CT, RE - 3.94, 1.28, [[(point, T_HEAD, INK)]], fill=WHITE, line=NAVY,
          lw=1.0, align=PP_ALIGN.LEFT)
    if img:
        x, y, w, h = img_rect
        s.shapes.add_picture(img, Cm(x), Cm(y), Cm(w), Cm(h))
    band(s, [('成果｜', T_LEADEM, WHITE), (results, T_NUM_L, YEL)])
    if pending:
        rec(owner, pending)
    notes(s, note_text or ('事例紹介ページ。主担当：%s。' % owner))
    return s


def exec_page(tag, lead_parts, items, target, message, owner, pending=None):
    """実行ページ。items は (実行すること, 担当, 開始時期, 月次目標)。'／' で改行。"""
    s = slide_blank(); header(s, 'Ⅴ. 実行', tag)
    lead(s, lead_parts)
    cw = [1.80, 13.40, 4.00, 3.52, 4.20]
    cx = [M]
    for w in cw[:-1]:
        cx.append(cx[-1] + w)
    for x, w, t in zip(cx, cw, ['No.', '実行すること', '担　当', '開始時期', '月次目標']):
        plain(s, x, CT, w, 1.00, [[(t, T_TAG, WHITE)]], fill=NAVY, line=NAVY, lw=1.0)
    y = CT + 1.00
    rh, ts = (1.60, T_TAG) if len(items) <= 5 else (1.45, T_SUB)
    for i, row in enumerate(items, start=1):
        it, vals = row[0], list(row[1:]) + [''] * (4 - len(row))
        plain(s, cx[0], y, cw[0], rh, [[(str(i), T_TAG, NAVY)]], fill=PALE, line=NAVY, lw=1.0)
        plain(s, cx[1], y, cw[1], rh, [[(it, ts, INK)]], fill=WHITE, line=NAVY, lw=1.0,
              align=PP_ALIGN.LEFT)
        for x, w, v in zip(cx[2:], cw[2:], vals):
            plain(s, x, y, w, rh, [[(t, T_SUB, NAVY)] for t in v.split('／')],
                  fill=WHITE, line=NAVY, lw=1.0, spc=1.2)
        y += rh
    plain(s, M, y + 0.35, CW, CB - (y + 0.35),
          [[('この施策の月次目標｜', T_LEAD, INK), (target, T_LEAD, RED)]], fill=PALE, lw=1.5)
    if pending:
        rec(owner, pending)
    keymsg(s, message)
    notes(s, '施策ごとの実行計画。主担当：%s。' % owner)
    return s


# ================= 本編（打合せ資料・26ページ） =================

# --- 1 表紙 ---
cover()

# --- 2 目次 ---
toc()

# --- 3 エグゼクティブサマリー ---
s = slide_blank(); header(s, 'Ⅰ. 結論', 'エグゼクティブサマリー')
lead(s, [('■　提携26社 × 1社あたり名簿68件 → 粗利', T_LEAD, INK),
         ('3,320万円', T_LEADEM, RED)])
for x, hd, val in zip(X3, ['① 提携社数', '② 1社あたり名簿数', '③ 獲得名簿数'],
                      ['26社', '68件', '1,755件']):
    card(s, x, CT, W3, 1.30, 2.60, [(hd, T_HEAD, WHITE)], [[(val, T_NUM_L, RED)]],
         body_align=PP_ALIGN.CENTER)
card(s, M, 8.35, CW, 1.20, 2.40, [('KGI｜粗利額', T_LEAD, YEL)],
     [[('粗利 ＝ 獲得名簿数 1,755件 × アポ率 12％ × 商談化率 70％', T_LEAD, INK)],
      [('× 成約率 30％ × 平均単価 250万円 × 粗利率 30％ ＝ 3,320万円', T_LEAD, INK)]],
     body_fill=WHITE, body_align=PP_ALIGN.CENTER, head_size=T_LEAD, body_spc=1.3)
plain(s, M, 12.35, CW, CB - 12.35,
      [[('この企画で狙うのは粗利', T_SUM, INK), ('3,320万円', T_SUM, RED), ('。', T_SUM, INK)],
       [('その全ては「獲得名簿数」という1つの数字に依存する。', T_SUM, INK)]],
      fill=PALE, spc=1.35)
rec('起案者', '—（全数値を反映済み）')
keymsg(s, '初年度は提携26社・名簿1,755件で、粗利3,320万円を狙う。')
notes(s, 'KGI＝粗利額。全数値が確定してから最後に執筆する。主担当：起案者。')

# --- 4 背景と課題 ---
s = slide_blank(); header(s, 'Ⅱ. 戦略', '背景と課題')
lead(s, [('■　個人向け集客はもう伸びない。だから、', T_LEAD, INK),
         ('名簿を持つ企業と組む', T_LEADEM, RED)])
card(s, X2[0], CT, W2, 1.40, 4.40, [('こ れ ま で', T_LEADEM, WHITE)],
     [[('・訪問、セミナーによる属人的な集客', T_HEAD, INK)],
      [('・訪問の出口がセミナーなのか提携なのか', T_HEAD, INK)],
      [('　不明確で、効率が悪い', T_HEAD, INK)]], body_size=T_HEAD, body_spc=1.6)
card(s, X2[1], CT, W2, 1.40, 4.40, [('こ れ か ら', T_LEADEM, WHITE)],
     [[('・DLやメルマガも活用し、より数多くの', T_HEAD, INK)],
      [('　企業にアプローチする', T_HEAD, INK)],
      [('・訪問の出口を「提携」に一本化する', T_HEAD, INK)]], body_size=T_HEAD, body_spc=1.6)
plain(s, M, 10.25, CW, CB - 10.25,
      [[('1,755件', 32, RED), ('を獲りにいく', 32, INK)]], fill=PALE)
rec('経営企画', '直近CPA 40,000円の裏取り')
keymsg(s, '1件ずつ獲る時代は終わった。名簿ごと獲りにいく。')
notes(s, '現状のリード獲得単価を自社実績から出す。CPA高騰の根拠データを別紙添付。主担当：経営企画。')

# --- 5 センターピンの定義 ---
s = slide_blank(); header(s, 'Ⅱ. 戦略', 'センターピンの定義')
lead(s, [('■　センターピン＝', T_LEAD, INK), ('獲得名簿数', T_LEADEM, RED)])
plain(s, M, CT, CW, 1.40, [[('セ ン タ ー ピ ン', T_LEADEM, YEL)]],
      fill=NAVY, line=NAVY, lw=1.5)
plain(s, M, 5.75, CW, 3.30,
      [[('獲得名簿数 ＝ 提携社数 ', T_NUM_L, INK), ('26社', T_NUM_L, RED)],
       [('× 1社あたり名簿数 ', T_NUM_L, INK), ('68件', T_NUM_L, RED),
        ('　＝　', T_NUM_L, INK), ('1,755件', T_NUM_L, RED)]],
      fill=WHITE, lw=2.5, spc=1.2)
for x, hd, lines in zip(X3,
                        ['① 目的を一本化', '② 提携後は必ず回収', '③ 他のKPIは追わない'],
                        [['訪問もDLもメルマガも、', '目的は提携社数'],
                         ['提携＝ゴールではない。', '名簿の回収までが1件'],
                         ['名簿数に効かない施策は、', 'やらないと決める']]):
    card(s, x, 9.45, W3, 1.30, 3.10, [(hd, T_HEAD, WHITE)],
         [[(t, T_TAG, INK)] for t in lines], body_align=PP_ALIGN.CENTER, body_spc=1.5)
plain(s, M, 14.20, CW, CB - 14.20,
      [[('この式の左辺を最大化すること以外は、やらない。', T_LEAD, INK)]], fill=PALE, lw=1.0)
rec('起案者', '1社あたり保有名簿300件・名簿共有率35％の妥当性')
keymsg(s, '追う数字は「獲得名簿数」ただ1つ。')
notes(s, '獲得名簿数 = 提携社数 × 1社あたり名簿数。未確定：目標提携社数、1社あたり想定名簿数。主担当：起案者。')

# --- 6 KPIツリー ---
s = slide_blank(); header(s, 'Ⅱ. 戦略', 'KPIツリー')
lead(s, [('■　粗利（KGI）は、名簿数（中間KPI）と施策KPIに', T_LEAD, INK),
         ('分解できる', T_LEADEM, RED)])
plain(s, M, CT, CW, 1.40,
      [[('KGI｜粗利　', T_LEADEM, WHITE), ('3,320万円', T_LEADEM, YEL)]],
      fill=NAVY, line=NAVY, lw=1.5)
arrow(s, (W - 1.00) / 2, 5.65, 1.00, 0.80, MSO_SHAPE.DOWN_ARROW, NAVY)
plain(s, M, 6.65, CW, 1.40,
      [[('中間KPI｜獲得名簿数　', T_LEADEM, WHITE), ('1,755件', T_LEADEM, YEL)]],
      fill=NAVY, line=NAVY, lw=1.5)
arrow(s, (W - 1.00) / 2, 8.25, 1.00, 0.80, MSO_SHAPE.DOWN_ARROW, NAVY)
for x, hd, lines in zip(X4,
        ['施策KPI ①\n訪問', '施策KPI ②\nフォーム・メルマガ・DL',
         '施策KPI ③\nセミナー（年4回）', '施策KPI ④\n名簿共有'],
        [['提携社数', '26社'], ['DL獲得数', '460件'],
         ['実参加社数', '28社'], ['名簿共有率', '35％']]):
    sh = shape(s, x, 9.35, W4, 1.60, fill=NAVY, line=NAVY, lw=1.0)
    for i, t in enumerate(hd.split('\n')):
        para(sh.text_frame, [(t, T_SUB, WHITE)], PP_ALIGN.CENTER, 1.2, first=(i == 0))
    sb = shape(s, x, 10.95, W4, 2.60, fill=PALE, line=NAVY, lw=1.0)
    para(sb.text_frame, [(lines[0], T_BODY, INK)], PP_ALIGN.CENTER, 1.35, first=True)
    para(sb.text_frame, [(lines[1], T_NUM_M, RED)], PP_ALIGN.CENTER, 1.35)
plain(s, M, 13.90, CW, CB - 13.90,
      [[('下の4つを動かせば、上の2つが動く。', T_LEAD, INK)]], fill=PALE, lw=1.0)
rec('経営企画', '—（KPI前提と事例ページに集約）')
keymsg(s, '粗利は、下段の施策KPIを動かさない限り動かない。')
notes(s, 'KGI（粗利）→ 中間KPI（獲得名簿数）→ 施策KPI の3階層。'
         '各率は「KPI前提と事例」ページと一致させる。主担当：経営企画。')

# --- 7 施策別のやることとKPI ---
s = slide_blank(); header(s, 'Ⅱ. 戦略', '施策別のやることとKPI')
lead(s, [('■　施策ごとに、', T_LEAD, INK), ('やること', T_LEADEM, RED), ('と', T_LEAD, INK),
         ('追うKPI', T_LEADEM, RED), ('を並べる', T_LEAD, INK)])
for x, hd, todo, kpi in zip(X3,
        ['① 訪問・飛び込み', '② フォーム・メルマガ・DL', '③ セミナー（年4回）'],
        [['・3名で月30社を訪問する', '・断られても名刺を獲る',
          '・獲った名刺はメルマガへ', '・商談化した先に提携提案',
          '・年1,080社を回りきる'],
         ['・DL資料を2ヶ月に1本作る', '・フォームで3,000社に配信',
          '・FB広告でDLを訴求する', '・DL者へ週1メルマガ配信',
          '・反応者を営業へ渡す'],
         ['・オンラインで年4回開催', '・メルマガ・DMで告知する',
          '・架電で参加を後押しする', '・当日に収益試算を提示',
          '・個別面談で提携を打診']],
        [[[('訪問 1,080社', T_BODY, INK)], [('→　提携 ', T_BODY, INK), ('16社', T_NUM_M, RED)]],
         [[('DL獲得 460件', T_BODY, INK)], [('→　提携 ', T_BODY, INK), ('4社', T_NUM_M, RED)]],
         [[('実参加 28社', T_BODY, INK)], [('→　提携 ', T_BODY, INK), ('6社', T_NUM_M, RED)]]]):
    sh = shape(s, x, CT, W3, 1.30, fill=NAVY, line=NAVY, lw=1.5)
    para(sh.text_frame, [(hd, T_HEAD, WHITE)], PP_ALIGN.CENTER, first=True)
    bullets(s, x, 5.35, W3, 4.30, todo, size=T_BODY, spc=1.6)
    plain(s, x, 9.65, W3, 0.85, [[('こ の 施 策 の K P I', T_SUB, YEL)]],
          fill=NAVY, line=NAVY, lw=1.0)
    plain(s, x, 10.50, W3, 2.00, kpi, fill=PALE, line=NAVY, lw=1.0,
          align=PP_ALIGN.CENTER, spc=1.35)
X4b, W4b = cols(4, 0.75)
for x, lb, vl in zip(X4b, ['提携社数', '名簿共有率', '獲得名簿数', '粗利（KGI）'],
                     ['26社', '35％', '1,755件', '3,320万円']):
    plain(s, x, 12.90, W4b, CB - 12.90, [[(lb, T_BODY, INK)], [(vl, T_NUM_M, RED)]],
          fill=PALE, line=NAVY, lw=1.5, spc=1.3)
for x in X4b[:3]:
    arrow(s, x + W4b + 0.03, 13.75, 0.69, 0.56, MSO_SHAPE.RIGHT_ARROW, RED)
rec('起案者', '—（KPI前提と事例ページに集約）')
keymsg(s, '3つの施策は、最後は「提携社数」の1本に合流する。')
notes(s, '施策ごとに「やること」と「追うKPI」を並べ、下段でそれが提携社数 → 名簿 → 粗利に'
         'つながることを示すページ。訪問16社＋フォーム/DL 4社＋セミナー6社＝提携26社。'
         '各KPIの数値は「KPI前提と事例」ページと一致させる。主担当：起案者。')

# --- 8 訪問・飛び込み営業 ---
s = slide_blank(); header(s, 'Ⅲ. 施策', '訪問・飛び込み営業')
lead(s, [('■　断られても名刺は必ず獲る。それがそのまま', T_LEAD, INK),
         ('メルマガの母数', T_LEADEM, RED), ('になる', T_LEAD, INK)])
card(s, M, CT, CW, 1.30, 2.40, [('提携社数の公式', T_LEAD, YEL)],
     [[('提携社数 ＝ 訪問数 ', T_LEAD, INK), ('1,080社', T_LEAD, RED),
       (' × 商談率 ', T_LEAD, INK), ('7.5％', T_LEAD, RED),
       (' × 提携率 ', T_LEAD, INK), ('20％', T_LEAD, RED), ('　＝　', T_LEAD, INK),
       ('16社', T_LEAD, RED)]],
     body_fill=WHITE, body_align=PP_ALIGN.CENTER, head_size=T_LEAD, lw=2.5)
card(s, X2[0], 8.15, W2, 1.40, 3.40, [('即効獲得｜経営層に直接提示', T_HEAD, WHITE)],
     [[('・販売利益（フロー収益）', T_HEAD, INK)],
      [('・メンテ手数料（ストック収益）', T_HEAD, INK)],
      [('・3名 × 月30社 × 12ヶ月で1,080社', T_HEAD, INK)]], body_size=T_HEAD, body_spc=1.55)
card(s, X2[1], 8.15, W2, 1.40, 3.40, [('リスト補給｜断られた後の動線', T_HEAD, WHITE)],
     [[('・名刺は必ず獲得する', T_HEAD, INK)],
      [('・即日メルマガへ合流させる', T_HEAD, INK)],
      [('・DL・セミナーで再接触する', T_HEAD, INK)]], body_size=T_HEAD, body_spc=1.55)
plain(s, M, 13.30, CW, CB - 13.30,
      [[('訪問の成果はゼロにならない。', T_SUM, INK), ('断り ＝ リード獲得', T_SUM, RED),
        ('と定義する。', T_SUM, INK)]], fill=PALE)
rec('営業', '訪問先リスト1,080社の確保')
keymsg(s, '訪問1,080社が、提携16社と名刺864件を同時に生む。')
notes(s, '提携社数 = 訪問数 × 商談率 × 提携率／未確定：訪問先リストの確保、商談率。主担当：営業。')

# --- 9 事例｜訪問 ---
case_page('事例｜訪問（セミナー経由）',
          [('■　訪問で成果を出している', T_LEAD, INK), ('T社＠山梨県', T_LEADEM, RED),
           ('の取り組み', T_LEAD, INK)],
          '営業の機会損失が発生している（売り方がわからない）工務店にアプローチ',
          '訪問 → 提携　2.5％', '営業',
          img=CASE_HOUMON, img_rect=((W - 21.79) / 2, 5.63, 21.79, 9.57),
          note_text='T社＠山梨県の実績。訪問120社 → セミナー参加15社（12.5％）→ 締結3社（20％）。'
                    '集客は遅くとも当日1ヶ月半前から動き、事務員・担当者と関係をつくる。'
                    '締結は社長同席のもと、先方の決裁者に提携の詳細を伝えて締結書をもらう。主担当：営業。')

# --- 10 フォーム・メルマガ・DL ---
s = slide_blank(); header(s, 'Ⅲ. 施策', 'フォーム・メルマガ・DL')
lead(s, [('■　同じDL資料を、', T_LEAD, INK), ('押して届ける', T_LEADEM, RED),
         ('と', T_LEAD, INK), ('広告で取りにいく', T_LEADEM, RED),
         ('の2経路で配る', T_LEAD, INK)])
card(s, M, CT, CW, 1.30, 2.40, [('DL獲得数の公式', T_LEAD, YEL)],
     [[('① フォーム送信 3,000社 × 年6回（約17,500通）→ 累積DL率 ', T_LEAD, INK),
       ('7.0％', T_LEAD, RED), ('　→　', T_LEAD, INK), ('210件', T_LEAD, RED)],
      [('② FB広告（DL訴求）月30万円 × 12ヶ月 → ', T_LEAD, INK), ('250件', T_LEAD, RED),
       ('　＝　DL獲得 ', T_LEAD, INK), ('460件', T_LEAD, RED)]],
     body_fill=WHITE, body_align=PP_ALIGN.CENTER, head_size=T_LEAD, body_spc=1.3, lw=2.5)
for x, hd, ls in zip(X3,
        ['① プッシュ｜フォーム送信', '② プル｜FB広告', '③ 育成｜メルマガ'],
        [['・商圏3,000社を全社抽出', '・2ヶ月に1本のDLを年6回',
          '・月平均1,460通で配信', '・累積DL率7.0％ → 210件'],
         ['・目的はDLのダウンロード', '・DL申込LPへ誘導する',
          '・月30万円 × 12ヶ月', '・年250件（1件 約1,440円）'],
         ['・DL者と訪問名刺を積み上げ', '・期中平均662件へ週1配信',
          '・月2,900通（枠の19％）', '・セミナー集客の母数になる']]):
    card(s, x, 8.15, W3, 1.30, 3.50, [(hd, T_TAG, WHITE)],
         [[(t, T_BODY, INK)] for t in ls], body_size=T_BODY, body_spc=1.6)
plain(s, M, 13.30, CW, CB - 13.30,
      [[('資料は1つ。', T_SUM, INK), ('届け方を2つ', T_SUM, RED),
        ('持つから、DLが積み上がる。', T_SUM, INK)]], fill=PALE)
rec('マーケ', 'フォーム送信のDL率、FB広告のDL単価、ツール割引の適用期間')
keymsg(s, '同じDL資料を、押して届け、広告で取りにいく。')
notes(s, 'DL獲得数 = フォーム送信210件 ＋ FB広告250件 ＝ 460件。'
         'FB広告はDL資料のダウンロードを促す目的で使い、DL申込LPへ誘導する。'
         'フォームの母数は商圏3,000社の全社（訪問先も除外しない）。主担当：マーケ。')

# --- 11 事例｜育成・刈り取り ---
s = case_page('事例｜育成・刈り取り',
              [('■　メルマガで成果を出している', T_LEAD, INK), ('船井総研', T_LEADEM, RED),
               ('の実績', T_LEAD, INK)],
              '住宅用太陽光・蓄電池の販売店に7,700部撒いて10社がクリックしている',
              '開封率16.9％　→　クリック率0.8％　→　DL率70％', 'マーケ',
              img=CASE_MAIL, img_rect=(10.22, 6.73, 17.00, 7.38),
              note_text='船井総研のメルマガ実績。配布7,775件 → 開封1,313件（16.9％）→ '
                        'URLクリック10件（開封の0.8％）→ ダウンロード7件（クリックの70％）。'
                        'アライアンスとは別文脈の数値だが、開封率の裏づけとして使う。主担当：マーケ。')
plain(s, M, 5.63, 9.60, 1.20, [[('メルマガKPI', T_NUM_M, WHITE)]],
      fill=NAVY, line=NAVY, lw=1.5)
_y = 6.83
for _lb, _vl in [('配布部数', '7,775件'), ('↓　開封率 16.9％', ''),
                 ('開封数', '1,313件'), ('↓　クリック率 0.8％', ''),
                 ('URLクリック数', '10件'), ('↓　ダウンロード率 70％', ''),
                 ('ダウンロード数', '7件')]:
    if _vl:
        plain(s, M, _y, 5.80, 1.30, [[(_lb, T_BODY, INK)]], fill=WHITE, line=NAVY,
              lw=1.0, align=PP_ALIGN.LEFT)
        plain(s, 6.10, _y, 3.80, 1.30, [[(_vl, T_HEAD, RED)]], fill=WHITE, line=NAVY, lw=1.0)
        _y += 1.30
    else:
        plain(s, M, _y, 9.60, 1.06, [[(_lb, T_SUB, NAVY)]], fill=PALE, line=NAVY, lw=1.0)
        _y += 1.06

# --- 12 セミナー ---
s = slide_blank(); header(s, 'Ⅲ. 施策', 'セミナー')
lead(s, [('■　年4回のセミナーで、', T_LEAD, INK), ('提携を決める', T_LEADEM, RED)])
XS, WS = cols(3, 0.85)
for x, hd, lines in zip(XS, ['STEP 1｜告　知', 'STEP 2｜送　客', 'STEP 3｜当　日'],
        [['メルマガとDMで', 'セミナーを案内し', '申込を集める'],
         ['反応者へ架電し', '参加を後押しする', '（テレアポは上積み）'],
         ['個別の収益試算を', '提示し、その場で', '提携を打診する']]):
    card(s, x, CT, WS, 1.40, 3.20, [(hd, T_HEAD, WHITE)],
         [[(t, T_TAG, INK)] for t in lines], body_align=PP_ALIGN.CENTER, body_spc=1.5)
for x in XS[:2]:
    arrow(s, x + WS + 0.10, 5.90, 0.65, 0.90, MSO_SHAPE.RIGHT_ARROW, RED)
card(s, M, 9.05, CW, 1.30, 2.40, [('提携社数の公式', T_LEAD, YEL)],
     [[('期中平均リスト ', T_LEAD, INK), ('662件', T_LEAD, RED),
       (' × 申込率 ', T_LEAD, INK), ('2％', T_LEAD, RED),
       (' × 参加率 ', T_LEAD, INK), ('70％', T_LEAD, RED),
       (' × 年4回', T_LEAD, INK)],
      [('→ 実参加 ', T_LEAD, INK), ('28社', T_LEAD, RED),
       ('（重複25％を控除）× 提携率 ', T_LEAD, INK), ('20％', T_LEAD, RED),
       ('　＝　', T_LEAD, INK), ('6社', T_LEAD, RED)]],
     body_fill=WHITE, body_align=PP_ALIGN.CENTER, head_size=T_LEAD, body_spc=1.3, lw=2.5)
plain(s, M, 13.10, CW, CB - 13.10,
      [[('前提：セミナーは', T_SUM, INK),
        ('メルマガ・DM単独でも9社集客できる企画', T_SUM, RED), ('にする。', T_SUM, INK)],
       [('テレアポはその上積み。集客をテレアポに依存させない。', T_SUM, INK)]],
      fill=PALE, spc=1.3)
rec('マーケ／営業', 'セミナー申込率2％の実測、開催頻度')
keymsg(s, '年4回のセミナーが、提携を決める場になる。')
notes(s, 'フローはメルマガ（育成）→ テレアポ（送客）→ セミナー（刈り取り）。'
         'セミナーはDM単独でも成立する企画設計が前提で、テレアポは集客の上積み。主担当：マーケ／営業。')

# --- 13 事例｜セミナー ---
case_page('事例｜セミナー（訪問込み）',
          [('■　セミナーで成果を出している', T_LEAD, INK), ('T社＠山梨県', T_LEADEM, RED),
           ('の取り組み', T_LEAD, INK)],
          '訪問で接点をつくった先を勉強会に集め、社長同席の場で提携を決めている',
          'セミナー参加率12.5％　→　提携率20％', 'マーケ／営業',
          img=CASE_HOUMON, img_rect=((W - 21.79) / 2, 5.63, 21.79, 9.57),
          note_text='同じT社の事例をセミナー視点で見たもの。訪問120社のうち15社が勉強会に参加し'
                    '（12.5％）、そこから3社と締結（20％）。集客期間は当日1ヶ月半前から。'
                    '当日は社長同席で決裁者に提携の詳細を伝え、その場で締結書をもらう。'
                    '主担当：マーケ／営業。')

# --- 14 KPI前提と事例 ---
s = slide_blank(); header(s, 'Ⅳ. 数値', 'KPI前提と事例')
lead(s, [('■　各施策のKPIと、その', T_LEAD, INK), ('根拠となる事例', T_LEADEM, RED),
         ('まとめ', T_LEAD, INK)])
cw14 = [5.90, 6.90, 4.60, 9.52]
cx14 = [M]
for w in cw14[:-1]:
    cx14.append(cx14[-1] + w)
for x, w, t in zip(cx14, cw14, ['施　策', 'K P I', '仮説値', '事例｜実績値']):
    plain(s, x, CT, w, 1.00, [[(t, T_TAG, WHITE)]], fill=NAVY, line=NAVY, lw=1.0)
groups = [
    ('① 訪問・飛び込み',
     [('訪問数（年）', '1,080社', '3名 × 月30社 × 12ヶ月'),
      ('商談率', '7.5％', '目標値（現状1.0％・T社2.5％の中間）'),
      ('提携率', '20％', 'T社＠山梨県：商談15件 → 提携3社')]),
    ('② フォーム・メルマガ・DL',
     [('フォーム配信母数', '3,000社', '商圏の全社（訪問先も除外しない）'),
      ('累積DL率（年6回）', '7.0％', '210件がDL。月1,460通＝プラン枠の49％'),
      ('メルマガ開封率', '16.9％', '船井総研：7,775通 → 1,313件'),
      ('FB広告からのDL', '250件', '月30万円 × 12ヶ月。DL申込LPへ誘導')]),
    ('③ セミナー',
     [('セミナー申込率', '2％', '期中平均リスト662件 → 参加9社/回'),
      ('参加者からの提携率', '20％', '訪問の提携率と同水準で設定')]),
]
y = CT + 1.00
rh14 = (CB - y) / sum(len(g[1]) for g in groups)
for name, kpis in groups:
    plain(s, cx14[0], y, cw14[0], rh14 * len(kpis), [[(name, T_BODY, NAVY)]],
          fill=PALE, line=NAVY, lw=1.0, spc=1.25)
    for kpi, val, basis in kpis:
        plain(s, cx14[1], y, cw14[1], rh14, [[(kpi, T_BODY, INK)]], fill=WHITE, line=NAVY,
              lw=1.0, align=PP_ALIGN.LEFT)
        plain(s, cx14[2], y, cw14[2], rh14, [[(val, T_TAG, RED)]], fill=WHITE, line=NAVY, lw=1.0)
        plain(s, cx14[3], y, cw14[3], rh14, [[(basis, T_SUB, NAVY)]],
              fill=WHITE, line=NAVY, lw=1.0, align=PP_ALIGN.LEFT)
        y += rh14
rec('経営企画', '商談率7.5％の目標設定（現状1.0％からの引き上げ幅）')
keymsg(s, '仮説値はすべて、事例の実績値から置いている。')
notes(s, '施策ごとにKPIと事例をひとまとめにした前提条件表。仮説値は事例の実績値から置く。主担当：経営企画。')

# --- 15 収支シミュレーション ---
s = slide_blank(); header(s, 'Ⅳ. 数値', '収支シミュレーション')
lead(s, [('■　この1年間での', T_LEAD, INK), ('目標数値', T_LEADEM, RED),
         ('です。投資は年452万円', T_LEAD, INK)])
card(s, M, CT, CW, 1.20, 1.90, [('売上の公式', T_HEAD, YEL)],
     [[('売上 ＝ 獲得名簿数 × アポ率 × 商談化率 × 成約率 × 平均単価', T_LEAD, INK)]],
     body_fill=WHITE, body_align=PP_ALIGN.CENTER, head_size=T_HEAD, lw=2.5)
cw15 = [8.92, 6.00, 6.00, 6.00]
cx15 = [M]
for w in cw15[:-1]:
    cx15.append(cx15[-1] + w)
for x, w, t in zip(cx15, cw15, ['項　目', '悲観値', '標準値', '楽観値']):
    plain(s, x, 7.55, w, 1.00, [[(t, T_TAG, WHITE)]], fill=NAVY, line=NAVY, lw=1.0)
rows15 = [('提携社数', ('17社', '26社', '35社')),
          ('獲得名簿数', ('1,150件', '1,755件', '2,360件')),
          ('アポ率', ('9％', '12％', '15％')),
          ('成約率', ('22％', '30％', '35％')),
          ('契約件数', ('16件', '44件', '87件')),
          ('売上', ('4,000万円', '1億1,100万円', '2億1,700万円')),
          ('粗利（KGI）', ('1,200万円', '3,320万円', '6,500万円'))]
y = 8.55
rh15 = (CB - y) / len(rows15)
for a, vals in rows15:
    plain(s, cx15[0], y, cw15[0], rh15, [[(a, T_TAG, NAVY)]], fill=PALE, line=NAVY, lw=1.0,
          align=PP_ALIGN.LEFT)
    for x, w, v in zip(cx15[1:], cw15[1:], vals):
        plain(s, x, y, w, rh15, [[(v, T_TAG, RED)]], fill=WHITE, line=NAVY, lw=1.0)
    y += rh15
rec('経営企画', '採用ケースの決定（標準値を計画値とするか）')
keymsg(s, '投資452万円に対し、標準ケースで粗利3,320万円。')
notes(s, '売上 = 獲得名簿数 × アポ率 × 商談化率 × 成約率 × 平均単価。'
         '悲観値・標準値・楽観値の3ケースで置く。主担当：経営企画。')

# --- 16 ロードマップ ---
s = slide_blank(); header(s, 'Ⅴ. 実行', 'ロードマップ')
lead(s, [('■　', T_LEAD, INK), ('3ヶ月', T_LEADEM, RED),
         ('で、提携が生まれる状態をつくる', T_LEAD, INK)])
XR, WR = cols(3, 0.85)
for x, hd, lines in zip(XR, ['1ヶ月目', '2ヶ月目', '3ヶ月目'],
        [['訪問開始', 'DL第1号を制作'], ['メルマガ配信開始', 'DL広告を出稿'],
         ['初回セミナー開催', '直後に個別提案']]):
    card(s, x, CT, WR, 1.50, 4.60, [(hd, T_LEAD, WHITE)],
         [[(t, T_TAG, INK)] for t in lines], body_align=PP_ALIGN.CENTER, body_spc=1.6)
for x in XR[:2]:
    arrow(s, x + WR + 0.10, 7.15, 0.65, 0.90, MSO_SHAPE.RIGHT_ARROW, RED)
for x, hd, val in zip(X3, ['訪問数', 'DL獲得', '提携社数'],
                      ['90社 /月', '38件 /月', '2社 /月']):
    plain(s, x, 10.55, W3, 1.05, [[(hd, T_TAG, WHITE)]], fill=NAVY, line=NAVY, lw=1.0)
    plain(s, x, 11.60, W3, 1.50, [[(val, T_LEAD, RED)]], fill=PALE, line=NAVY, lw=1.0)
plain(s, M, 13.45, CW, CB - 13.45,
      [[('毎月、この', T_SUM, INK), ('3つの数字', T_SUM, RED),
        ('だけを見て軌道修正する。', T_SUM, INK)]], fill=PALE)
rec('起案者', '月次KPIの各月への配分')
keymsg(s, '3ヶ月で、提携が生まれる状態をつくる。')
notes(s, '3ヶ月で提携が生まれる状態までを引く。月次KPIは「KPI前提と事例」ページから逆算する。主担当：起案者。')

# --- 17 実行｜訪問 ---
exec_page('実行｜訪問',
          [('■　訪問を、', T_LEAD, INK), ('いつ・誰が・何件やるか', T_LEADEM, RED),
           ('まで落とす', T_LEAD, INK)],
          [('対象リストの作成（工務店・不動産会社を抽出）', '―', '済', '―'),
           ('訪問トークスクリプトの作成', '―', '済', '―'),
           ('訪問の実施（1名あたり 週8社 ／ 月30社）',
            '坂元様、島様／福島様', '随時', '30社／1人'),
           ('獲得名刺のメルマガ登録', '坂元様、島様／福島様', '9月', '83社／1人'),
           ('商談化した先への提携提案・契約締結', '山﨑様', '9月', '商談7件／提携1.3社')],
          '訪問 90社　→　商談 7件　→　提携 1.3社（年16社）',
          '訪問は月90社。ここが、すべての起点になる。', '営業')

# --- 18 実行｜フォーム・DL ---
exec_page('実行｜フォーム・DL',
          [('■　DL資料を、', T_LEAD, INK), ('何を作り、どう配るか', T_LEADEM, RED),
           ('まで落とす', T_LEAD, INK)],
          [('リストDLで商圏3,000社を抽出し、業種で精査する', '貴社', '9月', '―'),
           ('DL資料の年間ラインナップを決める（年6本）', '船井', '9月', '―'),
           ('第1号を制作する（工務店向け・提携モデルの解説）', '船井', '10月', '2ヶ月に1本'),
           ('DL申込LPとサンクスメールを用意する（広告の受け皿）', '貴社', '10月', '―'),
           ('フォーム送信を配信する（月平均1,460通）', '貴社', '10月', '18件DL'),
           ('FB広告でDLを訴求する（月30万円・LPへ誘導）', '船井', '10月', '20件DL')],
          'フォーム 18件　＋　FB広告 20件　＝　DL 38件 /月',
          'DLは2ヶ月に1本。押す経路と拾う経路の両方で配る。', 'マーケ')

# --- 19 見本｜DLレポート ---
s = slide_blank(); header(s, 'Ⅴ. 実行', '見本｜DLレポート')
lead(s, [('■　第1号は、', T_LEAD, INK), ('この形', T_LEADEM, RED),
         ('で作る。表紙・目次・制作要件のたたき台', T_LEAD, INK)])
CVW, CVH = 7.89, CB - CT           # 表紙モック（A4縦比）
shape(s, M, CT, CVW, CVH, fill=SOFT, line=NAVY, lw=1.5)
plain(s, M + 0.35, CT + 0.35, CVW - 0.70, 0.70,
      [[('D O W N L O A D   R E P O R T', T_SUB, WHITE, False, 'Arial')]],
      fill=RED, line=RED, lw=1.0)
tf = textbox(s, M + 0.35, 5.30, CVW - 0.70, 3.30, MSO_ANCHOR.TOP)
para(tf, [('工務店のための', T_BODY, NAVY)], PP_ALIGN.LEFT, 1.4, first=True)
para(tf, [('太陽光・蓄電池', T_NUM_M, NAVY)], PP_ALIGN.LEFT, 1.3)
para(tf, [('提携モデル 解説', T_NUM_M, NAVY)], PP_ALIGN.LEFT, 1.3)
tf = textbox(s, M + 0.35, 8.70, CVW - 0.70, 1.70, MSO_ANCHOR.TOP)
para(tf, [('― 工事を持たずに、', T_BODY, INK)], PP_ALIGN.LEFT, 1.35, first=True)
para(tf, [('　売上をつくる3つの型 ―', T_BODY, INK)], PP_ALIGN.LEFT, 1.35)
photo(s, M + 0.35, 10.30, CVW - 0.70, 2.60, '図 版 ・ 写 真')
tf = textbox(s, M + 0.35, 13.20, CVW - 0.70, 1.55, MSO_ANCHOR.TOP)
para(tf, [('株式会社 船井総合研究所', T_SUB, INK)], PP_ALIGN.LEFT, 1.45, first=True)
para(tf, [('A4縦 ／ 全12ページ ／ PDF', T_SUB, NAVY)], PP_ALIGN.LEFT, 1.45)
TX, TW = M + CVW + 0.34, 11.40     # 目次
plain(s, TX, CT, TW, 1.10, [[('目　次（全12ページ）', T_TAG, WHITE)]],
      fill=NAVY, line=NAVY, lw=1.0)
_y = CT + 1.10
_rows = [('P2', 'なぜいま太陽光・蓄電池なのか'),
         ('P3-4', '工務店が取りこぼしている売上'),
         ('P5', '提携モデル①　紹介型'),
         ('P6', '提携モデル②　販売代行型'),
         ('P7', '提携モデル③　共同施工型'),
         ('P8-9', '導入事例（T社＠山梨県）'),
         ('P10', '収益シミュレーション'),
         ('P11-12', '導入の流れ ／ 個別相談のご案内')]
_rh = (CB - _y) / len(_rows)
for _pg, _tt in _rows:
    plain(s, TX, _y, 2.40, _rh, [[(_pg, T_SUB, NAVY)]], fill=PALE, line=NAVY, lw=1.0)
    plain(s, TX + 2.40, _y, TW - 2.40, _rh, [[(_tt, T_BODY, INK)]], fill=WHITE,
          line=NAVY, lw=1.0, align=PP_ALIGN.LEFT)
    _y += _rh
RX, RW = TX + TW + G, RE - (TX + TW + G)
card(s, RX, CT, RW, 1.10, 4.30, [('制作要件', T_TAG, WHITE)],
     [[('・A4縦・全12ページ・PDF', T_SUB, INK)],
      [('・DL申込LP経由で提供する', T_SUB, INK)],
      [('・2ヶ月に1本・年6本', T_SUB, INK)],
      [('・第1号の公開は【　月】', T_SUB, INK)],
      [('・制作は船井が担当', T_SUB, INK)]], head_size=T_TAG, body_size=T_SUB, body_spc=1.6)
card(s, RX, 9.85, RW, 1.10, CB - 10.95, [('次号以降のテーマ案', T_TAG, WHITE)],
     [[('・第2号 補助金の使い方', T_SUB, INK)],
      [('・第3号 蓄電池リプレイス', T_SUB, INK)],
      [('・第4号 提携先の収益事例集', T_SUB, INK)],
      [('・第5号 訪問トーク集', T_SUB, INK)],
      [('・第6号 年間の総括レポート', T_SUB, INK)]],
     head_size=T_TAG, body_size=T_SUB, body_spc=1.6)
rec('マーケ', '第1号のテーマ確定と公開時期、掲載する事例')
keymsg(s, 'レポートは1本作れば、2つの経路で何度も配れる。')
notes(s, 'DLレポートのたたき台。左が表紙イメージ、中央が目次案、右が制作要件と次号以降の'
         'テーマ案。第1号は「工務店のための太陽光・蓄電池 提携モデル解説」を想定し、'
         'A4縦12ページのPDFとしてDL申込LP経由で配布する。主担当：マーケ。')

# --- 20 実行｜メルマガ ---
exec_page('実行｜メルマガ',
          [('■　メルマガを、', T_LEAD, INK), ('誰に・いつ・何を送るか', T_LEADEM, RED),
           ('まで落とす', T_LEAD, INK)],
          [('配信リストをCRMに一本化する（訪問名刺・DL者）', '', '', '―'),
           ('年間の配信カレンダーを作る（週1回・曜日固定）', '', '', '―'),
           ('配信フォーマットを決める（事例＋DL＋告知）', '', '', '―'),
           ('事例ストックを作る（営業から月2本吸い上げる）', '', '', '2本／月'),
           ('毎週配信し、開封・クリックを記録する（月2,900通）', '', '', '4本／月'),
           ('クリックした先を週次で営業へ渡す', '', '', '4件／月')],
          '配信 2,900通　→　開封 490件（16.9％）　→　クリック 4件',
          'メルマガは、訪問とDLで拾った先を温め続ける場。', 'マーケ',
          pending='配信の担当と開始時期')

# --- 21 見本｜メルマガ ---
s = slide_blank(); header(s, 'Ⅴ. 実行', '見本｜メルマガ')
lead(s, [('■　毎週の1通は、', T_LEAD, INK), ('3ブロック固定', T_LEADEM, RED),
         ('。事例 → DL → 勉強会の順に並べる', T_LEAD, INK)])
MW = W2
shape(s, M, CT, MW, CB - CT, fill=SOFT, line=NAVY, lw=1.5)
IX, IW = M + 0.30, MW - 0.60
plain(s, IX, CT + 0.30, IW, 0.85,
      [[('件名：【事例】工事を持たずに、太陽光で年間粗利◯◯万円', T_SUB, INK)]],
      fill=WHITE, line=NAVY, lw=1.0, align=PP_ALIGN.LEFT)
tf = textbox(s, IX, 5.30, IW, 0.55)
para(tf, [('差出人：船井総研 住宅エネルギーチーム　／　2026年◯月◯日', T_SUB, NAVY)],
     PP_ALIGN.LEFT, first=True)
_y = 6.00
for _hd, _bd in [
        ('① 今週の事例',
         ['山梨県のT社は、施工体制を持たないまま太陽光の提案を始め、',
          '半年で3社と締結。きっかけは1件の飛び込み訪問でした。']),
        ('② お役立ち資料（DL）のご案内',
         ['新レポート「工務店のための太陽光・蓄電池 提携モデル解説」を',
          '公開しました。全12ページ・無料でDLいただけます。'])]:
    plain(s, IX, _y, IW, 0.78, [[(_hd, T_SUB, WHITE)]], fill=NAVY, line=NAVY, lw=1.0,
          align=PP_ALIGN.LEFT)
    tf = textbox(s, IX + 0.20, _y + 0.88, IW - 0.40, 1.45, MSO_ANCHOR.TOP)
    for _i, _t in enumerate(_bd):
        para(tf, [(_t, T_SUB, INK)], PP_ALIGN.LEFT, 1.4, first=(_i == 0))
    _y += 2.55
plain(s, M + (MW - 6.70) / 2, 11.05, 6.70, 0.85,
      [[('▶　資料をダウンロードする', T_SUB, WHITE)]], fill=RED, line=RED, lw=1.0)
plain(s, IX, 12.10, IW, 0.78, [[('③ 勉強会のご案内', T_SUB, WHITE)]],
      fill=NAVY, line=NAVY, lw=1.0, align=PP_ALIGN.LEFT)
tf = textbox(s, IX + 0.20, 12.98, IW - 0.40, 1.45, MSO_ANCHOR.TOP)
para(tf, [('◯月◯日（◯）オンラインで「太陽光・蓄電池 収益化セミナー」を', T_SUB, INK)],
     PP_ALIGN.LEFT, 1.4, first=True)
para(tf, [('開催します。定員30社・参加費無料です。', T_SUB, INK)], PP_ALIGN.LEFT, 1.4)
tf = textbox(s, IX + 0.20, 14.50, IW - 0.40, 0.50)
para(tf, [('配信元：株式会社船井総合研究所　／　配信停止はこちら', T_SUB, GRAY)],
     PP_ALIGN.LEFT, first=True)
card(s, X2[1], CT, W2, 1.10, 3.60, [('配信ルール', T_TAG, WHITE)],
     [[('・毎週◯曜◯時に固定して配信する', T_BODY, INK)],
      [('・1通＝3ブロック構成で固定する', T_BODY, INK)],
      [('・件名は「【事例】」から始める', T_BODY, INK)],
      [('・本文は3分で読み切れる長さに収める', T_BODY, INK)]],
     head_size=T_TAG, body_size=T_BODY, body_spc=1.6)
card(s, X2[1], 9.15, W2, 1.10, CB - 10.25, [('3ブロックの狙い', T_TAG, WHITE)],
     [[('・① 事例｜提携先の成果を見せて信用をつくる', T_BODY, INK)],
      [('・② DL｜新しいレポートへ誘導し、DLを増やす', T_BODY, INK)],
      [('・③ 告知｜勉強会へ送客し、提携の場に連れていく', T_BODY, INK)],
      [('・②③のクリック者は、週次で営業へ渡す', T_BODY, INK)],
      [('・事例は営業から月2本吸い上げてストックする', T_BODY, INK)]],
     head_size=T_TAG, body_size=T_BODY, body_spc=1.6)
rec('マーケ', '配信曜日・時刻、初回配信日、差出人名義')
keymsg(s, '形を固定するから、毎週続けられる。')
notes(s, 'メルマガのたたき台。件名・差出人・3ブロック（事例／DL告知／勉強会告知）・CTA・'
         'フッターまでの誌面イメージ。右は配信ルールと各ブロックの狙い。'
         '形を固定することで制作負荷を下げ、週1配信を継続できるようにする。主担当：マーケ。')

# --- 22 実行｜セミナー ---
exec_page('実行｜セミナー',
          [('■　年4回のセミナーを、', T_LEAD, INK), ('実行単位', T_LEADEM, RED),
           ('まで落とす', T_LEAD, INK)],
          [('年間開催計画を確定する（年4回・日程を先に押さえる）', '船井', '9月', '―'),
           ('企画書と当日コンテンツを作る（開催2ヶ月前まで）', '船井', '9月', '―'),
           ('告知を開始する（1.5ヶ月前・メルマガ週1＋DM）', '貴社', '9月', '―'),
           ('申込者・反応者へ架電し、参加を後押しする', '貴社', '9月', '申込13社／回'),
           ('当日運営（個別試算を提示し、決裁者同席で打診）', '貴社', '11月', '参加9社／回'),
           ('翌営業日にお礼と個別面談の日程調整を送る', '貴社', '11月', '提携1.5社／回')],
          '1回あたり 申込 13社　→　参加 9社　→　提携 1.5社（年6社）',
          'セミナーは四半期に1回。開催日から逆算して動く。', 'マーケ／営業')

# --- 23 勉強会テーマ案 ---
s = slide_blank(); header(s, 'Ⅴ. 実行', '勉強会テーマ案')
lead(s, [('■　', T_LEAD, INK), ('相手が変われば、刺さるテーマも変わる', T_LEADEM, RED),
         ('。3パターンで用意する', T_LEAD, INK)])
for x, hd, lines in zip(X3,
        ['案A｜工務店・ビルダー', '案B｜不動産会社', '案C｜電気・設備工事'],
        [[('【対　象】', NAVY), ('　新築を手がける工務店', INK),
          ('【テーマ】', NAVY), ('　太陽光・蓄電池 収益化セミナー', INK),
          ('【訴　求】', NAVY), ('　工事を持たずに1棟の粗利を上げる', INK),
          ('【提携の形】', NAVY), ('　施主を紹介 → 当社が施工', INK),
          ('【名簿の質】', NAVY), ('　施主名簿。着工前で反応が早い', INK)],
         [('【対　象】', NAVY), ('　売買仲介・賃貸管理会社', INK),
          ('【テーマ】', NAVY), ('　所有物件の価値を上げる 太陽光活用', INK),
          ('【訴　求】', NAVY), ('　管理手数料以外の収益源をつくる', INK),
          ('【提携の形】', NAVY), ('　オーナー名簿へ案内 → 当社が提案', INK),
          ('【名簿の質】', NAVY), ('　オーナー名簿。件数がとにかく多い', INK)],
         [('【対　象】', NAVY), ('　電気工事・設備工事会社', INK),
          ('【テーマ】', NAVY), ('　既存客の蓄電池リプレイス提案', INK),
          ('【訴　求】', NAVY), ('　過去の顧客名簿がそのまま売上になる', INK),
          ('【提携の形】', NAVY), ('　既存客名簿へ共同で提案', INK),
          ('【名簿の質】', NAVY), ('　施工済み顧客。更新期が読める', INK)]]):
    sh = shape(s, x, CT, W3, 1.30, fill=NAVY, line=NAVY, lw=1.5)
    para(sh.text_frame, [(hd, T_HEAD, WHITE)], PP_ALIGN.CENTER, first=True)
    bullets(s, x, 5.35, W3, 7.60, [[(t, T_SUB, c)] for t, c in lines], spc=1.5)
plain(s, M, 13.30, CW, CB - 13.30,
      [[('まずは', T_LEAD, INK), ('案A（工務店）', T_LEAD, RED),
        ('で1回目を実施する。', T_LEAD, INK)],
       [('反応を見て、名簿が最も多い', T_LEAD, INK), ('案B（不動産）', T_LEAD, RED),
        ('へ広げる。', T_LEAD, INK)]], fill=PALE, lw=1.5, spc=1.3)
rec('マーケ／営業', '2回目以降にどのパターンを当てるか')
keymsg(s, '相手ごとにテーマを替えれば、同じ企画が3回使える。')
notes(s, '勉強会のテーマ案を対象業種ごとに3パターン用意したページ。'
         '案Aは工務店・ビルダー（施主名簿）、案Bは不動産会社（オーナー名簿・件数が多い）、'
         '案Cは電気・設備工事会社（施工済み顧客名簿・更新期が読める）。'
         '次回は案Aで実施し、反応を見て案Bへ広げる想定。主担当：マーケ／営業。')

# --- 24 次回勉強会の企画 ---
s = slide_blank(); header(s, 'Ⅴ. 実行', '次回勉強会の企画')
lead(s, [('■　テーマは', T_LEAD, INK), ('太陽光・蓄電池に特化', T_LEADEM, RED),
         ('。オンラインで集め、出口は', T_LEAD, INK), ('提携', T_LEADEM, RED)])
card(s, M, CT, CW, 1.10, 1.90,
     [('テ ー マ （ 案 ）　～ 提携手数料で売上をつくる ～', T_TAG, YEL)],
     [[('太陽光・蓄電池 収益化セミナー', 32, RED)]],
     body_fill=WHITE, head_size=T_TAG, body_align=PP_ALIGN.CENTER, lw=2.5)
for x, hd, lines in zip(X3,
        ['開催形式｜オンライン', '当日プログラム（90分）', '出　口｜提　携'],
        [['・Zoomウェビナー（90分）', '・定員 30社（1社2名まで）',
          '・参加費 無料', '・開催日【　　月　　日】'],
         ['・第1部 市場と補助金の最新', '・第2部 提携モデルの解説',
          '・第3部 収益試算のデモ', '・第4部 質疑・個別相談の案内'],
         ['・終了時に個別面談を案内', '・後日個別で試算を提示',
          '・決裁者同席で提携を打診', '・参加9社 → 提携1.5社が目標']]):
    card(s, x, 7.45, W3, 1.30, 2.90, [(hd, T_HEAD, WHITE)],
         [[(t, T_SUB, INK)] for t in lines], body_size=T_SUB, body_spc=1.65)
plain(s, M, 12.00, CW, 0.90,
      [[('集客の逆算スケジュール　―　申込 13社 → 参加 9社 → 提携 1.5社 が目標', T_TAG, WHITE)]],
      fill=NAVY, line=NAVY, lw=1.0)
X5, W5 = cols(5, 0.28)
for x, hd, lines in zip(X5,
        ['1.5ヶ月前｜告知', '1ヶ月前｜追込み', '2週間前｜確定',
         '前日｜最終確認', '当日｜刈り取り'],
        [['・メルマガで初告知', '・DMを発送する'],
         ['・反応者へ架電', '・メルマガ2回目'],
         ['・受講URLを案内', '・欠席防止の連絡'],
         ['・参加者リスト確定', '・個別提案先を選ぶ'],
         ['・収益試算を提示', '・個別面談へ誘導']]):
    plain(s, x, 12.90, W5, 0.80, [[(hd, T_SUB, WHITE)]], fill=NAVY, line=NAVY, lw=1.0)
    plain(s, x, 13.70, W5, CB - 13.70, [[(t, T_SUB, INK)] for t in lines],
          fill=PALE, line=NAVY, lw=1.0, align=PP_ALIGN.LEFT, spc=1.35)
rec('マーケ／営業', 'テーマの確定、開催日時、登壇者、DM発送先')
keymsg(s, '太陽光・蓄電池に絞り、オンラインで集め、提携で締める。')
notes(s, '次回勉強会の企画シート。テーマは太陽光・蓄電池に特化し、'
         'オンライン（Zoomウェビナー90分）で開催、出口は提携に置く。'
         'オンラインのため商圏を越えて集客できる一方、その場での締結はできないので、'
         '終了時に個別面談を案内し、後日個別試算を提示して決裁者同席で打診する。'
         'T社＠山梨県の事例のとおり、集客は遅くとも1ヶ月半前から動かす。主担当：マーケ／営業。')

# --- 25 別軸①｜TRENDE社 ---
s = slide_blank(); header(s, 'Ⅵ. 別軸', '別軸①｜TRENDE社')
lead(s, [('■　本編とは', T_LEAD, INK), ('別軸', T_LEADEM, RED),
         ('。伊藤忠商事の持分法適用会社との業務提携', T_LEAD, INK)])
card(s, X2[0], CT, W2, 1.30, 3.85, [('会社概要', T_HEAD, WHITE)],
     [[('・TRENDE株式会社（2017年設立・資本金14億円）', T_SUB, INK)],
      [('・東京電力グループ発。2023年に伊藤忠商事が', T_SUB, INK)],
      [('　子会社化し、現在は同社の持分法適用会社', T_SUB, INK)],
      [('・東京センチュリー・全農・東芝・東急不動産も出資', T_SUB, INK)],
      [('・小売電気／太陽光・蓄電池のリース・PPA', T_SUB, INK)]],
     body_size=T_SUB, body_spc=1.6)
plain(s, X2[1], CT, W2, 1.30, [[('顧客基盤（同社資料より）', T_HEAD, WHITE)]],
      fill=NAVY, line=NAVY, lw=1.5)
_y = CT + 1.30
_rows = [('蓄電池レンタルサービス', '約9,000件'),
         ('テラリス（蓄電池・PVリース）', '約1,500件'),
         ('ほっとでんき（PPA）', '約2,800件'),
         ('じぶん電力・ひだまりでんき', '約500件'),
         ('合　計（全国・家庭）', '約13,800件')]
_rh = 3.85 / len(_rows)
for _nm, _cnt in _rows:
    _last = (_cnt == '約13,800件')
    plain(s, X2[1], _y, 8.30, _rh, [[(_nm, T_SUB, NAVY if _last else INK)]],
          fill=PALE if _last else WHITE, line=NAVY, lw=1.0, align=PP_ALIGN.LEFT)
    plain(s, X2[1] + 8.30, _y, W2 - 8.30, _rh, [[(_cnt, T_SUB, RED)]],
          fill=PALE if _last else WHITE, line=NAVY, lw=1.0)
    _y += _rh
for x, hd, lines in zip(X3,
        ['案①｜名簿を受け取る', '案②｜案件を受け取る', '案③｜電力会社と組む'],
        [['・NEC蓄電池リースの', '　設置先名簿を受領する', '・当社が点検を提案する',
          '・最も名簿数に直結する'],
         ['・名簿は受領しない', '・同社の直販部門と組み', '　リフォーム案件のAPを受領',
          '・名簿ではなく案件が入る'],
         ['・同社が提携する電力会社', '　との協業に乗る', '・伊藤忠グループ案件も対象',
          '・中長期の広がりが大きい']]):
    card(s, x, 9.60, W3, 1.30, CB - 10.90, [(hd, T_HEAD, WHITE)],
         [[(t, T_BODY, INK)] for t in lines], body_size=T_BODY, body_spc=1.6)
rec('営業', 'TRENDEとの提携形態（案①／②／③）と、受領できる名簿の件数・範囲')
keymsg(s, '1社との提携で、名簿が桁で動く可能性がある。')
notes(s, '本編の施策（訪問／フォーム・メルマガ・DL／セミナー）とは切り離した別軸の案件。'
         'TRENDE社（伊藤忠商事の持分法適用会社）との業務提携の検討ページ。'
         '同社は2024年に買収した蓄電池レンタル（約9,000件）のリプレイス提案を推進中で、'
         'テラリス（リース）は2025年度実績約900件。提携形態は名簿受領・案件受領・電力会社協業の'
         '3案。案①なら獲得名簿数が一気に増えるが、名簿の範囲と件数は未確認。主担当：営業。')

# --- 26 別軸②｜イー・スマイル ---
s = slide_blank(); header(s, 'Ⅵ. 別軸', '別軸②｜イー・スマイル')
lead(s, [('■　こちらも', T_LEAD, INK), ('別軸', T_LEADEM, RED),
         ('。高橋顧問が持つ東海圏 5,000社のパイプ', T_LEAD, INK)])
card(s, X2[0], CT, W2, 1.30, 3.85, [('経　緯', T_HEAD, WHITE)],
     [[('・窓口を石原顧問から高橋顧問（イー・スマイル）へ', T_SUB, INK)],
      [('　変更。春日井市在住、43歳', T_SUB, INK)],
      [('・東海圏の工務店・設備に 5,000社 のパイプを持つ', T_SUB, INK)],
      [('・訪問は春日井市の事務所を拠点に、坂元・田中・島の', T_SUB, INK)],
      [('　3名で、金曜を除く平日16:00〜稼働', T_SUB, INK)]],
     body_size=T_SUB, body_spc=1.6)
card(s, X2[1], CT, W2, 1.30, 3.85, [('この軸で狙うこと', T_HEAD, WHITE)],
     [[('・5,000社から点検先を営業していく', T_SUB, INK)],
      [('・提携ではなく「点検案件」が先に立つ流れ', T_SUB, INK)],
      [('・本編の名簿獲得モデル（提携 → 名簿共有）とは', T_SUB, INK)],
      [('　入口も出口も別のルート', T_SUB, INK)],
      [('・進み方しだいで、本編とは別に数字が立つ', T_SUB, INK)]],
     body_size=T_SUB, body_spc=1.6)
for x, hd, lines in zip(X3,
        ['本編との違い①｜母数', '本編との違い②｜出口', '本編との違い③｜管理'],
        [['・本編の訪問1,080社は', '　自社リストから積んだ数字',
          '・5,000社はこことは別枠', '・重複の有無は未確認'],
         ['・本編は「提携 → 名簿共有」', '・こちらは「点検案件の獲得」',
          '・KGIへの効き方が違うため', '　同じ表では管理しない'],
         ['・本編の26社には含めない', '・件数・確度が見えた段階で',
          '　改めて数字を置く', '・当面は月次で進捗のみ共有']]):
    card(s, x, 9.60, W3, 1.30, CB - 10.90, [(hd, T_HEAD, WHITE)],
         [[(t, T_BODY, INK)] for t in lines], body_size=T_BODY, body_spc=1.6)
rec('営業', '5,000社の内訳と、実際に訪問できる社数')
keymsg(s, '5,000社のパイプは、本編とは別に数字を立てる。')
notes(s, '本編の施策とは切り離した別軸。イー・スマイル高橋顧問（春日井市）が持つ'
         '東海圏の工務店・設備5,000社のパイプから点検先を営業していくルート。'
         '本編の訪問1,080社とは母数も出口も別なので、同じKPI表では管理しない。'
         '件数と確度が見えた段階で改めて数字を置く。主担当：営業。')


# ================= ページ番号の付与と保存 =================
for i, rec_ in enumerate(SLIDES, start=1):
    footer(rec_['s'], i, strip=rec_['strip'])

os.makedirs('output', exist_ok=True)
out = 'output/企画書_アライアンス提携_空パッケージ.pptx'
prs.save(out)

lines = ['# 確定待ち事項リスト（企画書 空パッケージ・打合せ資料版）', '',
         '各ページで確定を待っている項目の一覧です。', '',
         '| P | 主担当 | 確定待ち事項 |', '|---|---|---|']
for pg, owner, pending in PENDING:
    lines.append('| %d | %s | %s |' % (pg, owner, pending))
open('output/確定待ち事項リスト.md', 'w', encoding='utf-8').write('\n'.join(lines) + '\n')

for wmsg in FITWARN:
    print('FIT WARNING:', wmsg)
print('saved:', out, '/ slides:', len(prs.slides._sldIdLst),
      '/ fit warnings:', len(FITWARN))
