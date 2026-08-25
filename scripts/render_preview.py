# -*- coding: utf-8 -*-
"""pptx の図形ジオメトリから プレビュー画像を生成し、はみ出しを検出する。"""
import sys, os, math
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SRC = sys.argv[1]
OUT = sys.argv[2] if len(sys.argv) > 2 else 'preview'
PPC = 50.0                      # px per cm
FONTP = '/usr/share/fonts/opentype/ipafont-gothic/ipagp.ttf'
PT2CM = 2.54 / 72.0
_cache = {}


def font(size_pt):
    px = max(6, int(round(size_pt * PT2CM * PPC)))
    if px not in _cache:
        _cache[px] = ImageFont.truetype(FONTP, px)
    return _cache[px]


def cm(v):
    return Emu(v).cm


def solid(sh):
    try:
        f = sh.fill
        if f.type is not None and 'SOLID' in str(f.type):
            return '#' + str(f.fore_color.rgb)
    except Exception:
        pass
    return None


def linecol(sh):
    try:
        ln = sh.line
        if ln.fill.type is not None and 'SOLID' in str(ln.fill.type):
            return '#' + str(ln.color.rgb), max(1, int(round((ln.width or 12700) / 12700 * 1.6)))
    except Exception:
        pass
    return None, 0


def wrap(text, fnt, maxw, draw):
    if not text:
        return ['']
    out, cur = [], ''
    for ch in text:
        t = cur + ch
        if draw.textlength(t, font=fnt) > maxw and cur:
            out.append(cur); cur = ch
        else:
            cur = t
    out.append(cur)
    return out


issues = []
prs = Presentation(SRC)
SW, SH = cm(prs.slide_width), cm(prs.slide_height)
W, H = int(SW * PPC), int(SH * PPC)

for idx, s in enumerate(prs.slides, 1):
    img = Image.new('RGB', (W, H), 'white')
    d = ImageDraw.Draw(img)
    for sh in s.shapes:
        x, y = cm(sh.left) * PPC, cm(sh.top) * PPC
        w, h = cm(sh.width) * PPC, cm(sh.height) * PPC
        st = str(sh.shape_type)
        if 'PICTURE' in st:
            d.rectangle([x, y, x + w, y + h], outline='#999999')
            d.text((x + 4, y + h / 2 - 6), 'LOGO', font=font(9), fill='#999999')
            continue
        if 'LINE' in st or 'CONNECT' in st:
            c, lw = linecol(sh)
            d.line([x, y, x + w, y + h], fill=c or '#000000', width=max(2, lw))
            continue
        fill = solid(sh)
        lc, lw = linecol(sh)
        name = sh.name or ''
        if 'Arrow' in name or 'ARROW' in st:
            if fill:
                d.polygon([(x + w * .5, y), (x + w, y + h * .5), (x + w * .75, y + h * .5),
                           (x + w * .75, y + h), (x + w * .25, y + h),
                           (x + w * .25, y + h * .5), (x, y + h * .5)], fill=fill)
            continue
        if fill or lc:
            d.rectangle([x, y, x + w, y + h], fill=fill,
                        outline=lc if lc else None, width=lw or 1)
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        ml = cm(tf.margin_left or 0) * PPC; mr = cm(tf.margin_right or 0) * PPC
        mt = cm(tf.margin_top or 0) * PPC; mb = cm(tf.margin_bottom or 0) * PPC
        avail = w - ml - mr
        lines = []
        for p in tf.paragraphs:
            runs = []
            for r in p.runs:
                try:
                    col = '#' + str(r.font.color.rgb)
                except Exception:
                    col = '#1F1F1F'
                runs.append((r.text, r.font.size.pt if r.font.size else 18, col))
            if not runs:
                continue
            size = max(r[1] for r in runs)
            spc = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
            txt = ''.join(r[0] for r in runs)
            fnt = font(size)
            for wl in wrap(txt, fnt, avail, d):
                lines.append((wl, size, runs[0][2], spc, p.alignment))
        if not lines:
            continue
        total = sum(sz * PT2CM * PPC * spc * 1.22 for _, sz, _, spc, _ in lines)
        anchor = tf.vertical_anchor
        if anchor == MSO_ANCHOR.TOP:
            cy = y + mt
        else:
            cy = y + (h - total) / 2
        over = total > (h - mt - mb) + 1.5
        for wl, sz, col, spc, align in lines:
            fnt = font(sz)
            lh = sz * PT2CM * PPC * spc * 1.22
            tw = d.textlength(wl, font=fnt)
            if align == PP_ALIGN.CENTER:
                tx = x + ml + (avail - tw) / 2
            elif align == PP_ALIGN.RIGHT:
                tx = x + ml + avail - tw
            else:
                tx = x + ml
            d.text((tx, cy + (lh - sz * PT2CM * PPC) / 2), wl, font=fnt, fill=col)
            cy += lh
        if over:
            issues.append('p%d %s: text %.2fcm > box %.2fcm' %
                          (idx, name, total / PPC, (h - mt - mb) / PPC))
            d.rectangle([x, y, x + w, y + h], outline='#FF00FF', width=4)
    img.save('%s-%02d.png' % (OUT, idx))

print('rendered %d slides -> %s-NN.png' % (len(prs.slides._sldIdLst), OUT))
print('OVERFLOW ISSUES: %d' % len(issues))
for i in issues:
    print('  ', i)
