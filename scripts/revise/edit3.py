# -*- coding: utf-8 -*-
"""3次修正：レイアウト装飾の補完と、はみ出しの解消"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
NAVY=RGBColor(0x00,0x20,0x60)
NS='{http://schemas.openxmlformats.org/drawingml/2006/main}'
prs=Presentation('final.pptx'); S=prs.slides; log=[]

def flat(sh,o):
    if sh.shape_type==6:
        for c in sh.shapes: flat(c,o)
    else: o.append(sh)
def shapes(i):
    o=[]
    for sh in S[i-1].shapes: flat(sh,o)
    return o
def norm(t): return t.replace('\x0b','').replace('\n','').replace(' ','')
def find(i,needle,nth=0):
    h=[s for s in shapes(i) if s.has_text_frame and norm(needle) in norm(s.text_frame.text)]
    if len(h)<=nth: raise SystemExit(f"NOT FOUND P{i}: {needle!r}")
    return h[nth]
def st(i,needle,new,label,nth=0):
    sh=find(i,needle,nth); tf=sh.text_frame
    paras=list(tf.paragraphs)
    for k,pa in enumerate(paras):
        want=new if k==0 else ""
        for br in pa._p.findall(NS+'br'): pa._p.remove(br)
        if pa.runs:
            pa.runs[0].text=want
            for r in pa.runs[1:]: r._r.getparent().remove(r._r)
    log.append(f"P{i} {label}: {new}")

def add_deco(i, rule=True, bar=True):
    sl=S[i-1]; tree=sl.shapes._spTree; added=[]
    if rule:
        s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.61), Inches(11.69), Inches(0.07))
        added.append(s)
    if bar:
        s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(8.05), Inches(11.69), Inches(0.22))
        added.append(s)
    for s in added:
        s.fill.solid(); s.fill.fore_color.rgb=NAVY
        s.line.fill.background(); s.shadow.inherit=False
        el=s._element; el.getparent().remove(el); tree.insert(2, el)   # 背面へ
    log.append(f"P{i} 装飾: 青線{'○' if rule else '−'} 下帯{'○' if bar else '−'}")

# ---- レイアウト装飾の補完 ----
add_deco(10); add_deco(11)
add_deco(15, rule=False); add_deco(16, rule=False)

# ---- P11 章タイトル幅・表ヘッダ ----
t=find(11,"1.長岡市のエネルギー事情")
t.width=Inches(8.95); log.append("P11 章タイトル幅を8.95inに")
tbl=[s for s in shapes(11) if s.has_table][0].table
for c,new in ((2,"発生件数（約）"),(3,"復旧日数（約）")):
    cell=tbl.rows[0].cells[c]; pa=cell.text_frame.paragraphs[0]
    for br in pa._p.findall(NS+'br'): pa._p.remove(br)
    if pa.runs:
        pa.runs[0].text=new
        for r in pa.runs[1:]: r._r.getparent().remove(r._r)
for row in tbl.rows:
    for cell in row.cells:
        for pa in cell.text_frame.paragraphs:
            for r in pa.runs: r.font.size=Pt(12)
log.append("P11 表：ヘッダを1行に短縮／全セル12ptに縮小")

# ---- P25 ----
st(25,"申請受付は5月18日から","①　申請受付は5/18から","①見出し")
st(25,"交付決定前の着工は対象外","②　交付決定前の着工はNG","②見出し")
st(25,"予算上限に達し次第、終了","③　予算上限で受付終了","③見出し")
st(25,"「交付決定→工事→補助金交付」の順","補助金は早い者勝ち。「決まってから工事」が鉄則です","結論バー")

# ---- P20 ----
st(20,"雪下ろしを前提としない設計","雪下ろしを前提としない","積雪対策②")

# ---- P41 ----
st(41,"参考：ヤシロは工務店として屋根の下地まで確認し",
   "参考：ヤシロは屋根の下地まで確認し、向かない場合は正直にお伝えします。","下部")

# ---- P34 ----
sh=find(34,"定期的なメンテナンスなど")
paras=list(sh.text_frame.paragraphs)
for k,want in enumerate(["定期メンテナンスなど","充実のアフターフォロー"]):
    pa=paras[k]
    for br in pa._p.findall(NS+'br'): pa._p.remove(br)
    if pa.runs:
        pa.runs[0].text=want
        for r in pa.runs[1:]: r._r.getparent().remove(r._r)
log.append("P34 対応①: 定期メンテナンスなど / 充実のアフターフォロー")

# ---- P14 補足テキストをカード内に収める ----
for sh in shapes(14):
    if sh.has_text_frame and abs(sh.top/914400-5.34)<0.02 and abs(sh.width/914400-1.88)<0.02:
        sh.top=Inches(5.05); sh.height=Inches(1.5)
log.append("P14 補足5枚をカード内に移動（y5.05 / h1.5）")

prs.save('final2.pptx')
print("\n".join(log)); print(f"\n--- {len(log)} 件 / final2.pptx")
