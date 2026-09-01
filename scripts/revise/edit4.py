# -*- coding: utf-8 -*-
"""4次修正：残るはみ出しの解消（すべて元ファイル由来の文字量オーバー）"""
from pptx import Presentation
from pptx.util import Inches
NS='{http://schemas.openxmlformats.org/drawingml/2006/main}'
prs=Presentation('final2.pptx'); S=prs.slides; log=[]
def flat(sh,o):
    if sh.shape_type==6:
        for c in sh.shapes: flat(c,o)
    else: o.append(sh)
def shapes(i):
    o=[]
    for sh in S[i-1].shapes: flat(sh,o)
    return o
def norm(t): return t.replace('\x0b','').replace('\n','').replace(' ','')
def find(i,needle,nth=0):
    h=[s for s in shapes(i) if s.has_text_frame and norm(needle) in norm(s.text_frame.text)]
    if len(h)<=nth: raise SystemExit(f"NOT FOUND P{i}: {needle!r}")
    return h[nth]
def setp(sh,lines):
    paras=list(sh.text_frame.paragraphs)
    for k,pa in enumerate(paras):
        want=lines[k] if k<len(lines) else ""
        for br in pa._p.findall(NS+'br'): pa._p.remove(br)
        if pa.runs:
            pa.runs[0].text=want
            for r in pa.runs[1:]: r._r.getparent().remove(r._r)
        elif want: pa.text=want
def st(i,needle,lines,label,nth=0):
    if isinstance(lines,str): lines=[lines]
    setp(find(i,needle,nth),lines); log.append(f"P{i} {label}: {' / '.join(lines)[:52]}")

# ---- P14 補足を元の位置に戻し、2行に収まる長さへ ----
for sh in shapes(14):
    if sh.has_text_frame and abs(sh.top/914400-5.05)<0.02 and abs(sh.width/914400-1.88)<0.02:
        sh.top=Inches(5.36); sh.height=Inches(1.24)
for old,new in [("エアコンなどが停止し","エアコンが止まり室温が下がる"),
                ("お湯が出ない","お湯が出ず入浴もできない"),
                ("ポンプ停止と配管凍結","ポンプ停止と凍結で水が止まる"),
                ("電動洗浄・ポンプ停止","電動洗浄が止まり衛生環境が悪化"),
                ("Wi-Fiも充電もなく","Wi-Fiも充電も不可。情報が届かない")]:
    st(14,old,new,"補足")
log.append("P14 補足の位置をy5.36・h1.24に調整")

# ---- P34 ----
st(34,"定期メンテナンスなど",["定期メンテナンスなど","充実のアフター対応"],"対応①")
st(34,"長く使っていくためには",["長く使うには定期的な","メンテナンスが必要",""],"デメリット②")
st(34,"メーカーの保証だけでなく",["メーカー保証に加え","弊社独自の工事保証つき",""],"対応③")

# ---- P41 ----
st(41,"④現地調査でどこを見るか","④現地調査の項目","チップ")
st(41,"参考：ヤシロは屋根の下地まで確認し","参考：ヤシロは屋根の下地まで確認します。","下部")

prs.save('done.pptx')
print("\n".join(log)); print(f"\n--- {len(log)} 件 / done.pptx")
