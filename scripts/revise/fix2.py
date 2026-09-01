# -*- coding: utf-8 -*-
"""2巡目：レンダリング再検証で残った折り返し・重なりを是正する"""
from pptx import Presentation
from pptx.util import Pt, Emu
from measure import flat, em_len, seg_runs, EMU, NS

prs = Presentation('polished.pptx')
S = prs.slides
log = []

def shapes(n):
    o = []
    for sh in S[n-1].shapes: flat(sh, o)
    return {sh.shape_id: sh for sh in o}

def set_sz(sh, pt):
    for pa in sh.text_frame.paragraphs:
        for r in pa.runs: r.font.size = Pt(pt)

def scale_sz(sh, f):
    for pa in sh.text_frame.paragraphs:
        for r in pa.runs:
            if r.font.size: r.font.size = Pt(round(r.font.size.pt*f*2)/2)

def box(sh, x=None, y=None, w=None, h=None):
    if x is not None: sh.left = Emu(int(x*EMU))
    if y is not None: sh.top = Emu(int(y*EMU))
    if w is not None: sh.width = Emu(int(w*EMU))
    if h is not None: sh.height = Emu(int(h*EMU))

def widest(sh):
    """最も幅の広い1行の実寸(in)と、そのときの最大フォントpt"""
    mx = 0.0; mpt = 0.0
    for pa in sh.text_frame.paragraphs:
        base = max([r.font.size.pt for r in pa.runs if r.font.size] or [18])
        for seg in seg_runs(pa):
            w = sum(em_len(t)*((sz or base)/72.0) for t, sz in seg)
            if w > mx: mx = w; mpt = base
    return mx, mpt

def fit_line(sh, pad=0.20):
    """1行に収まるところまで全ランを一律縮小"""
    avail = sh.width/EMU - pad
    w, _ = widest(sh)
    if w <= avail: return False
    scale_sz(sh, avail/w * 0.99)
    return True

# ── P8 「年間 約25,000円 の負担増」を1行に ────────────────────
d = shapes(8)
if fit_line(d[24]):
    log.append("P8 年間負担増の行 36/24pt→枠内1行に収まるサイズへ縮小")

# ── P9 左カード見出し・本文サイズを整理／右カード本文のあふれ解消 ────
d = shapes(9)
for sid in (13, 15, 17):
    box(d[sid], w=5.14); set_sz(d[sid], 20)
for sid in (14, 16, 18):
    set_sz(d[sid], 18); box(d[sid], h=0.78)
set_sz(d[28], 18); box(d[28], h=1.90)
for pa in d[28].text_frame.paragraphs:
    for br in pa._p.findall(NS+'br'): pa._p.remove(br)
log.append("P9 見出し24→20pt・本文20→18ptに統一／右カード本文のあふれと出典への重なりを解消")

# ── P13 「停電＋低気温＋除雪の遅れ」ラベル幅を2.60inに拡張して折り返し解消 ──
d = shapes(13)
pos = {22: 1.085, 23: 3.685, 24: 4.545, 25: 7.145, 26: 8.005}
for sid, x in pos.items():
    box(d[sid], x=x)
    if sid in (22, 24, 26): box(d[sid], w=2.60)
log.append("P13 3要素ラベルの枠幅を2.60inに拡張・行全体を再センタリング（「除雪の遅れ」の折り返し解消）")

# ── P18 見出し48→40pt／Q見出し28→24pt・枠幅を揃える ──────────
d = shapes(18)
set_sz(d[8], 40)
for sid, cx in ((13, 0.20), (19, 4.07), (25, 7.94)):
    set_sz(d[sid], 24); box(d[sid], x=cx+0.08, w=3.39)
log.append("P18 大見出し48→40pt（折り返し解消）／Qカード見出し28→24pt・枠幅3.39inに統一")

# ── P39 説明文6件 20→18pt（1行項目の折り返し解消） ───────────
d = shapes(39)
for sid in (16, 19, 22, 25, 28, 31): set_sz(d[sid], 18)
log.append("P39 説明文6件 20→18ptに再統一（折り返し解消）")

# ── P41 項目見出し6件 20→19pt・枠幅2.85in ──────────────────
d = shapes(41)
for sid in (14, 19, 24, 29, 34, 39):
    set_sz(d[sid], 19); box(d[sid], w=2.85)
log.append("P41 項目見出し6件 20→19pt・枠幅2.85inに拡張（「下地・野地板の健全性」の折り返し解消）")

# ── P42 まとめ5項目 26→22pt（全項目を1行に揃える） ─────────────
d = shapes(42)
for sid in (11, 15, 19, 23, 27): set_sz(d[sid], 22)
log.append("P42 まとめ5項目 26→22ptに統一（全項目を1行に揃える）")

prs.save('polished2.pptx')
print("\n".join("・"+x for x in log))
print(f"\n--- {len(log)}件 → polished2.pptx")
