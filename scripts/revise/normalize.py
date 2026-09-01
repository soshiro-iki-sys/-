# -*- coding: utf-8 -*-
"""全スライドの共通グリッドを規定値に揃える"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import unicodedata
EMU=914400
NS='{http://schemas.openxmlformats.org/drawingml/2006/main}'

# ---- 規定値（学習したフォーマット）----
TITLE  = dict(x=0.16, y=0.02, w=8.95, h=0.55, sz=30)
CHIP   = dict(x=0.16, y=0.73,          h=0.42, sz=20)   # 幅は文字数から算出
LEAD   = dict(x=0.16, y=1.12, w=11.37,          sz=22)  # 高さは行数から
SRC    = dict(x=0.16, y=7.04, w=11.37, h=0.26, sz=10.5)
BAR    = dict(x=0.00, y=7.36, w=11.69, h=0.91)
BARTX  = dict(x=0.35, y=7.36, w=10.99, h=0.91, maxsz=32)
DIV    = dict(x=0.80, y=1.20, w=10.09, h=5.87, sz=44)   # 章扉

prs=Presentation('v2.pptx'); S=prs.slides; log=[]

def wl(s):
    n=0
    for c in s:
        n += 1 if unicodedata.east_asian_width(c) in ('W','F','A') else 0.55
    return n
def flat(sh,o):
    if sh.shape_type==6:
        for c in sh.shapes: flat(c,o)
    else: o.append(sh)
def shapes(i):
    o=[]
    for sh in S[i-1].shapes: flat(sh,o)
    return o
def inch(v): return None if v is None else v/EMU
def size_of(sh):
    s=[r.font.size.pt for pa in sh.text_frame.paragraphs for r in pa.runs if r.font.size]
    return s[0] if s else None
def set_size(sh,pt):
    for pa in sh.text_frame.paragraphs:
        for r in pa.runs: r.font.size=Pt(pt)
def place(sh,x=None,y=None,w=None,h=None):
    moved=[]
    for key,val,attr in (('x',x,'left'),('y',y,'top'),('w',w,'width'),('h',h,'height')):
        if val is None: continue
        cur=inch(getattr(sh,attr))
        if cur is None or abs(cur-val)>0.0008:
            setattr(sh,attr,Inches(val)); moved.append(f"{key}{cur}→{val}")
    return moved

for n in range(1,len(S._sldIdLst)+1):
    o=shapes(n)
    title=chip=lead=src=bar=bartx=None
    divs=[]
    for sh in o:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            L,T,W,H=inch(sh.left),inch(sh.top),inch(sh.width),inch(sh.height)
            if L is not None and T is not None and W and abs(T-7.36)<0.35 and W>11: bar=sh
            continue
        L,T,W,H=inch(sh.left),inch(sh.top),inch(sh.width),inch(sh.height)
        if L is None: continue
        s=size_of(sh); txt=sh.text_frame.text.strip()
        if T<0.2 and W and W>4 and s and 28<=s<=32 and L<0.4: title=sh
        elif 0.60<T<0.95 and L<0.4 and s and 18<=s<=22 and W and W<8: chip=sh
        elif 1.0<T<1.30 and L<0.4 and W and W>10.5 and s and s==22: lead=sh
        elif 6.85<T<7.25 and s and s<=12: src=sh
        elif 7.0<T<7.6 and W and W>9 and s and s>=24: bartx=sh
        elif 1.0<T<1.5 and W and W>9.8 and s and s>=44 and L>0.5: divs.append(sh)

    msgs=[]
    if title:
        m=place(title, TITLE['x'], TITLE['y'], TITLE['w'], TITLE['h'])
        if size_of(title)!=TITLE['sz']: set_size(title,TITLE['sz']); m.append(f"sz→{TITLE['sz']}")
        if m: msgs.append("章タイトル "+" ".join(m))
    if chip:
        w=round(wl(chip.text_frame.text.strip())*0.31+0.36, 2)
        m=place(chip, CHIP['x'], CHIP['y'], w, CHIP['h'])
        if size_of(chip)!=CHIP['sz']: set_size(chip,CHIP['sz']); m.append(f"sz→{CHIP['sz']}")
        # チップ背景の矩形があれば同じ位置へ
        for sh in o:
            if sh is chip or sh.has_text_frame and sh.text_frame.text.strip(): continue
            L,T,W,H=inch(sh.left),inch(sh.top),inch(sh.width),inch(sh.height)
            if L is not None and abs(L-0.156)<0.06 and T and 0.6<T<0.95 and W and W<8:
                place(sh, CHIP['x'], CHIP['y'], w, CHIP['h']); m.append("背景矩形も整列")
        if m: msgs.append("チップ "+" ".join(m))
    if lead:
        nlines = len([p for p in lead.text_frame.text.split('\n') if p.strip()])
        cap = LEAD['w']*72/LEAD['sz']
        need2 = nlines>1 or wl(lead.text_frame.text.replace('\n',''))>cap
        h = 0.86 if need2 else 0.50
        m=place(lead, LEAD['x'], LEAD['y'], LEAD['w'], h)
        if m: msgs.append("リード文 "+" ".join(m))
    if src:
        m=place(src, SRC['x'], SRC['y'], SRC['w'], SRC['h'])
        if abs((size_of(src) or 0)-SRC['sz'])>0.01: set_size(src,SRC['sz']); m.append(f"sz→{SRC['sz']}")
        for pa in src.text_frame.paragraphs: pa.alignment=PP_ALIGN.RIGHT
        if m: msgs.append("出典 "+" ".join(m))
    if bar:
        m=place(bar, BAR['x'], BAR['y'], BAR['w'], BAR['h'])
        if m: msgs.append("結論バー(帯) "+" ".join(m))
    if bartx:
        solo = bartx is bar or (inch(bartx.width) or 0) > 11.4
        if solo:   # 塗り付きテキスト1個で帯を兼ねている型
            m=place(bartx, BAR['x'], BAR['y'], BAR['w'], BAR['h'])
        else:
            m=place(bartx, BARTX['x'], BARTX['y'], BARTX['w'], BARTX['h'])
        # 1行に収まる最大サイズ。26pt を下限とし、それ未満になる長文は26ptで2行にする
        want=max(26, min(BARTX['maxsz'], int((11.69-0.9)*72/max(1.0,wl(bartx.text_frame.text.strip())))))
        if size_of(bartx)!=want: set_size(bartx,want); m.append(f"sz{size_of(bartx)}→{want}")
        if m: msgs.append("結論バー(文字) "+" ".join(m))
    for d in divs:
        m=place(d, DIV['x'], DIV['y'], DIV['w'], DIV['h'])
        if size_of(d)!=DIV['sz']: set_size(d,DIV['sz']); m.append(f"sz→{DIV['sz']}")
        if m: msgs.append("章扉 "+" ".join(m))
    if msgs: log.append(f"P{n}: " + " ／ ".join(msgs))

# ---- P15：画像がリード文2行目に重なるため下げる ----
for sh in S[14].shapes:
    if sh.shape_type==13 and inch(sh.top) and inch(sh.top)<2.1:
        sh.top=Inches(2.06); log.append("P15: 画像をy2.06へ（リード文2行目との重なりを解消）")

prs.save('norm.pptx')
print("\n".join(log)); print(f"\n--- {len(log)}ページを調整 / norm.pptx")
