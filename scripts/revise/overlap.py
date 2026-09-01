# -*- coding: utf-8 -*-
"""あふれによって他要素と重なっているテキストだけを検出する"""
import sys
from pptx import Presentation
from measure import *

def eff_rect(sh):
    """テキストが実際に占める矩形（あふれた分を含む）"""
    L=sh.left/EMU; T=sh.top/EMU; W=sh.width/EMU; H=sh.height/EMU
    _,need,_=metrics(sh)
    h=max(H,need)
    anc=sh.text_frame.vertical_anchor
    if anc is not None and int(anc)==3:      # MIDDLE
        t=T+(H-h)/2
    elif anc is not None and int(anc)==4:    # BOTTOM
        t=T+H-h
    else:
        t=T
    return L,t,L+W,t+h,(h-H)

def inter(a,b):
    x=min(a[2],b[2])-max(a[0],b[0]); y=min(a[3],b[3])-max(a[1],b[1])
    return x,y

def run(path):
    prs=Presentation(path); hits=[]
    for n,s in enumerate(prs.slides,1):
        o=[]
        for sh in s.shapes: flat(sh,o)
        txt=[sh for sh in o if sh.has_text_frame and sh.text_frame.text.strip()
             and sh.width and sh.height and sh.width/EMU>0.6]
        rects={sh.shape_id:eff_rect(sh) for sh in txt}
        # 1) テキスト同士の重なり（元の枠では重なっていないもの）
        for i,a in enumerate(txt):
            for b in txt[i+1:]:
                ra,rb=rects[a.shape_id],rects[b.shape_id]
                if ra[4]<=0.02 and rb[4]<=0.02: continue
                xi,yi=inter(ra,rb)
                oa=(a.left/EMU,a.top/EMU,(a.left+a.width)/EMU,(a.top+a.height)/EMU)
                ob=(b.left/EMU,b.top/EMU,(b.left+b.width)/EMU,(b.top+b.height)/EMU)
                xo,yo=inter(oa,ob)
                if xi>0.15 and yi>0.06 and not (xo>0.15 and yo>0.06):
                    hits.append((n,'TXT×TXT',a,b,round(yi,2)))
        # 2) 画像との重なり
        pics=[sh for sh in o if sh.shape_type==13 and sh.width and sh.height]
        for a in txt:
            ra=rects[a.shape_id]
            if ra[4]<=0.02: continue
            for p in pics:
                rp=(p.left/EMU,p.top/EMU,(p.left+p.width)/EMU,(p.top+p.height)/EMU)
                oa=(a.left/EMU,a.top/EMU,(a.left+a.width)/EMU,(a.top+a.height)/EMU)
                xi,yi=inter(ra,rp); xo,yo=inter(oa,rp)
                if xi>0.2 and yi>0.08 and not (xo>0.2 and yo>0.08):
                    hits.append((n,'TXT×PIC',a,p,round(yi,2)))
        # 3) 収まっているべき枠（塗り枠 / 表 / 結論バー）からのはみ出し
        for a in txt:
            ra=rects[a.shape_id]
            if ra[4]<=0.02: continue
            for c in o:
                if c is a or not c.width or not c.height: continue
                if c.shape_type not in (1,9) : continue
                rc=(c.left/EMU,c.top/EMU,(c.left+c.width)/EMU,(c.top+c.height)/EMU)
                oa=(a.left/EMU,a.top/EMU,(a.left+a.width)/EMU,(a.top+a.height)/EMU)
                # a が c に内包されているのに、あふれて c の下端を越えた
                if oa[0]>=rc[0]-0.05 and oa[2]<=rc[2]+0.05 and oa[1]>=rc[1]-0.05 and oa[3]<=rc[3]+0.05:
                    if ra[3]>rc[3]+0.03 or ra[1]<rc[1]-0.03:
                        hits.append((n,'枠はみ出し',a,c,round(ra[3]-rc[3],2)))
    return hits

if __name__=='__main__':
    for n,kind,a,b,v in run(sys.argv[1]):
        ta=a.text_frame.text.replace('\n','/').replace('\x0b','/')[:30]
        tb=b.text_frame.text.replace('\n','/').replace('\x0b','/')[:22] if b.has_text_frame else str(b.shape_type)[:12]
        print(f'P{n:>2} {kind} id{a.shape_id}「{ta}」 ↔ id{b.shape_id}「{tb}」 {v}in')
