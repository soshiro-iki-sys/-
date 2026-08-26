# -*- coding: utf-8 -*-
"""第2回「聞く姿勢づくり＆必要性訴求①」研修資料ビルダー

使い方
------
1. アプローチブック「スマートハウスのある暮らし」のページ画像を用意する
     python <pptx-skill>/scripts/office/soffice.py --headless --convert-to pdf アプローチブック.pptx
     pdftoppm -jpeg -r 110 アプローチブック.pdf <ASSETS>/ab/ab
   → <ASSETS>/ab/ab-01.jpg 〜 ab-26.jpg ができる
2. 「契約までの4ステップ」ピラミッド図を <ASSETS>/pyramid.png に置く
   （合宿資料 ppt/media/image13.png）
3. ASSETS を書き換えて実行する
     python slides/build/build_session2.py

第3回以降はこのファイルをコピーし、スライド定義部分を差し替えて使う。
書式の定義は docs/研修資料フォーマット仕様.md、内容は docs/知識_*.md を参照。
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.opc.packuri import PackURI

REPO   = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
ASSETS = os.environ.get('TRAINING_ASSETS', '/tmp/training-assets')   # ← 手順1・2の出力先
AB     = os.path.join(ASSETS, 'ab')
TPL    = os.path.join(REPO, 'templates', '研修資料_フォーマット見本.pptx')
OUT    = os.path.join(REPO, 'slides', '2026営業研修_第2回_聞く姿勢づくり＆必要性訴求①.pptx')

GOTHIC='HGPｺﾞｼｯｸE'; MINCHO='HGP明朝E'; UD='BIZ UDPゴシック'; YU='游ゴシック'
RED='FF0000'; BLACK='000000'; WHITE='FFFFFF'; GREEN='9BBB59'; YELLOW='FFFF00'
GRAY='595959'; BLUE='0070C0'; LTGREEN='EAF1DD'; LTGRAY='F2F2F2'; LTYEL='FFF2CC'
LTBLUE='DEEBF7'; ORANGE='F79646'

# ---------------------------------------------------------------- helpers
def _rpr_font(run, name):
    rPr = run._r.get_or_add_rPr()
    rPr.get_or_add_latin().set('typeface', name)
    latin = rPr.find(qn('a:latin'))
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = parse_xml('<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        latin.addnext(ea)
    ea.set('typeface', name)

def _highlight(run, color):
    rPr = run._r.get_or_add_rPr()
    hl = parse_xml(
        '<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:srgbClr val="%s"/></a:highlight>' % color)
    latin = rPr.find(qn('a:latin'))
    if latin is not None:
        latin.addprevious(hl)
    else:
        rPr.append(hl)

def run(p, text, size=16, bold=False, color=BLACK, font=GOTHIC, hl=None):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)
    if hl: _highlight(r, hl)
    _rpr_font(r, font)
    return r

def tb(slide, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP, margin=0.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return box, tf

def para(tf, first=False, space_after=4, align=PP_ALIGN.LEFT, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after); p.alignment = align
    if line: p.line_spacing = line
    return p

def lines(tf, items, size=16, color=BLACK, bold=False, space=5, font=GOTHIC, line=None, align=PP_ALIGN.LEFT):
    """items: str | (str,dict)"""
    for i, it in enumerate(items):
        opts = {}
        if isinstance(it, tuple): it, opts = it
        p = para(tf, first=(i == 0), space_after=opts.pop('space', space),
                 align=opts.pop('align', align), line=opts.pop('line', line))
        run(p, it, size=opts.pop('size', size), bold=opts.pop('bold', bold),
            color=opts.pop('color', color), font=opts.pop('font', font),
            hl=opts.pop('hl', None))

def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=1.0,
         anchor=MSO_ANCHOR.TOP):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    else: s.fill.background()
    if line:
        s.line.color.rgb = RGBColor.from_string(line); s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    return s, tf

def pic(slide, path, x, y, w, border=True):
    p = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    if border:
        p.line.color.rgb = RGBColor.from_string('808080'); p.line.width = Pt(1)
    return p

def table(slide, x, y, w, col_w, rows, font_size=11, header=True,
          header_fill='000000', header_color=WHITE, row_h=0.30, head_h=0.30):
    n_r, n_c = len(rows), len(col_w)
    shp = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w),
                                 Inches(head_h + row_h * (n_r - 1)))
    t = shp.table
    t.first_row = header; t.horz_banding = False
    for i, cw in enumerate(col_w): t.columns[i].width = Inches(cw)
    t.rows[0].height = Inches(head_h)
    for i in range(1, n_r): t.rows[i].height = Inches(row_h)
    for ri, rowdata in enumerate(rows):
        for ci, cell in enumerate(rowdata):
            txt, opts = (cell, {}) if isinstance(cell, str) else cell
            c = t.cell(ri, ci)
            c.margin_left = c.margin_right = Inches(0.05)
            c.margin_top = c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = opts.get('fill', header_fill if (header and ri == 0) else None)
            if fill: c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string(fill)
            else: c.fill.background()
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = opts.get('align',
                PP_ALIGN.CENTER if (header and ri == 0) else PP_ALIGN.LEFT)
            run(p, txt, size=opts.get('size', font_size),
                bold=opts.get('bold', header and ri == 0),
                color=opts.get('color', header_color if (header and ri == 0) else BLACK),
                font=opts.get('font', YU))
    return t

# ---------------------------------------------------------------- slide frame
prs = Presentation(TPL)
BLANK = None
for l in prs.slide_masters[0].slide_layouts:
    if l.name == '白紙': BLANK = l
assert BLANK is not None

sldIdLst = prs.slides._sldIdLst
orig_ids = list(sldIdLst)
cover_id, colophon_id = orig_ids[0], orig_ids[6]
for sid in orig_ids[1:6]:                      # drop template samples 2-6
    rId = sid.get(qn('r:id')); sldIdLst.remove(sid)
    prs.part.drop_rel(rId)
# python-pptx names new slides by len(sldIdLst)+1 -> move colophon out of that range
prs.part.related_part(colophon_id.get(qn('r:id'))).partname = PackURI('/ppt/slides/slide900.xml')

page_no = [0]

def new_slide(heading=None):
    s = prs.slides.add_slide(BLANK)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    if heading is not None:
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.157), Inches(0.757), Inches(0.140), Inches(0.512))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor.from_string(BLACK)
        bar.line.color.rgb = RGBColor.from_string('D9D9D9'); bar.line.width = Pt(0.75)
        bar.shadow.inherit = False
        box, tf = tb(s, 0.348, 0.736, 9.9, 0.572)
        lines(tf, [heading], size=28, font=MINCHO)
    page_no[0] += 1
    box, tf = tb(s, 8.31, 7.18, 2.30, 0.33)
    lines(tf, [str(page_no[0])], size=12, color=WHITE, font=GOTHIC, align=PP_ALIGN.RIGHT)
    return s

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def ab_img(n):
    return os.path.join(AB, 'ab-%02d.jpg' % n)

# ---------------------------------------------------------------- 表紙
cover = prs.slides[0]
for shp in cover.shapes:
    if shp.has_text_frame and shp.name == 'テキスト ボックス 5':
        tf = shp.text_frame; tf.clear()
        lines(tf, [
            ('2026年度　営業研修', {'size': 32, 'bold': True, 'color': WHITE, 'font': UD, 'space': 8}),
            ('第2回　聞く姿勢づくり', {'size': 52, 'bold': True, 'color': WHITE, 'font': UD, 'space': 2}),
            ('＆必要性訴求①', {'size': 52, 'bold': True, 'color': WHITE, 'font': UD, 'space': 10}),
            ('〜太陽光＋蓄電池セット販売〜', {'size': 24, 'bold': True, 'color': WHITE, 'font': UD}),
        ])
    if shp.has_text_frame and shp.name == 'テキスト ボックス 6':
        tf = shp.text_frame; tf.clear()
        lines(tf, [('2026年8月　＠船井総研', {'size': 14, 'color': WHITE, 'font': UD})],
              align=PP_ALIGN.RIGHT)
page_no[0] = 1

# ================================================================ 2. 本日のゴール
s = new_slide('本日のゴール')
box, tf = tb(s, 0.45, 1.45, 9.9, 0.62)
lines(tf, [('アプローチブックを使って商談できるようになる',
            {'size': 28, 'bold': True, 'color': RED})])
_, tf = rect(s, 0.45, 2.35, 9.95, 3.05, fill=LTGREEN)
lines(tf, [
    ('研修が終わったときの「できる状態」', {'size': 18, 'bold': True, 'space': 10}),
    ('①　アプローチブックを見ながら、そのページで言うべきことが口から出る', {'size': 19, 'space': 9}),
    ('②　各ページの「結論の一文」を、資料を見ずに言える', {'size': 19, 'space': 9}),
    ('③　ページごとのゴール（お客様に言わせたい台詞）を引き出す質問が打てる', {'size': 19, 'space': 9}),
    ('④　ページをめくる順番と理由を説明できる', {'size': 19}),
])
_, tf = rect(s, 0.45, 5.65, 9.95, 0.85, fill=None, line=RED, lw=1.5)
lines(tf, [('本日の範囲：営業ルール ①〜⑨　／　アプローチブック p.2〜p.17',
            {'size': 22, 'bold': True, 'align': PP_ALIGN.CENTER})])
notes(s, '・知識を覚える会ではなく、アプローチブックが使えるようになる会だと最初に宣言する\n'
         '・「今日の終わりに暗唱チェックをします」と先に伝えておく')

# ================================================================ 3. アジェンダ
s = new_slide('本日のアジェンダ')
rows = [
    ['時間', '内容', '扱うもの'],
    ['0:00-0:10', '前回の振り返り', '第1回：業界・商品知識'],
    ['0:10-0:20', '全体像の確認', '4ステップ／20のルール／アプローチブックの地図'],
    ['0:20-1:00', ('聴く姿勢作り　ルール①〜⑥', {'bold': True}), 'アプローチブック p.2〜p.7'],
    ['1:00-1:10', '休憩', '−'],
    ['1:10-1:40', ('必要性訴求①　ルール⑦〜⑨', {'bold': True}), 'アプローチブック p.9〜p.17'],
    ['1:40-1:55', 'ロールプレイング', 'アプローチブック p.2〜p.17 通し'],
    ['1:55-2:00', '暗唱チェック／まとめ／宿題', '−'],
]
table(s, 0.45, 1.65, 9.95, [1.55, 4.3, 4.1], rows, font_size=14, row_h=0.55, head_h=0.40)
notes(s, '・ロープレの時間は必ず確保する。解説が延びたら解説側を削る')

# ================================================================ 4. 前回の振り返り
s = new_slide('前回の振り返り')
box, tf = tb(s, 0.45, 1.42, 9.9, 0.45)
lines(tf, [('第1回「業界商材の知識講座」で押さえたこと', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 2.00, 4.85, 2.45, fill=LTGRAY)
lines(tf, [
    ('エネルギー業界の3つの特徴', {'size': 17, 'bold': True, 'space': 8}),
    ('① 普及率が低い（戸建ての12〜13％）', {'size': 14, 'space': 4}),
    ('　→ 商品説明力が要る', {'size': 14, 'color': GRAY, 'space': 8}),
    ('② 単価が高く購買頻度が低い（1件200万円超）', {'size': 14, 'space': 4}),
    ('　→ 人柄力と知識量が要る', {'size': 14, 'color': GRAY, 'space': 8}),
    ('③ 耐用年数が長い（20年以上）', {'size': 14, 'space': 4}),
    ('　→ 20年後まで見据えた提案が要る', {'size': 14, 'color': GRAY}),
])
_, tf = rect(s, 5.55, 2.00, 4.85, 2.45, fill=LTGRAY)
lines(tf, [
    ('日本のエネルギー事情', {'size': 17, 'bold': True, 'space': 8}),
    ('・エネルギー自給率は 11.3％', {'size': 15, 'space': 6}),
    ('・火力発電が全体の 約7割', {'size': 15, 'space': 6}),
    ('・火力の資源は 99％が輸入', {'size': 15, 'space': 10}),
    ('だから燃料価格の高騰が、そのまま電気代に乗る', {'size': 15, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 0.45, 4.75, 9.95, 1.55, fill=LTYEL)
lines(tf, [
    ('第1回は「知識」＝自信の材料をつくる回でした。', {'size': 20, 'bold': True, 'space': 8}),
    ('第2回からは、その知識を', {'size': 20, 'bold': True, 'space': 0}),
])
p = tf.paragraphs[-1]
run(p, 'どの順番で、どのページで出すか', size=20, bold=True, color=RED)
run(p, ' を学びます。', size=20, bold=True)
notes(s, '・第1回の内容を覚えているか、口頭で2〜3問投げてから進める')

# ================================================================ 5. 全体像① 4ステップ
s = new_slide('全体像①：契約までの4ステップ')
pic(s, os.path.join(ASSETS, 'pyramid.png'), 0.55, 1.55, 6.0, border=False)
_, tf = rect(s, 6.95, 1.75, 3.45, 1.75, fill=None, line=RED, lw=2.0)
lines(tf, [
    ('本日の範囲', {'size': 20, 'bold': True, 'color': RED, 'align': PP_ALIGN.CENTER, 'space': 8}),
    ('STEP1　聴く姿勢づくり', {'size': 17, 'bold': True, 'space': 5}),
    ('STEP2　必要性訴求（前半）', {'size': 17, 'bold': True}),
])
_, tf = rect(s, 6.95, 3.75, 3.45, 2.35, fill=LTGRAY)
lines(tf, [
    ('積む順番が決まっている', {'size': 17, 'bold': True, 'space': 8}),
    ('下から4→3→2→1の順に積む。', {'size': 14, 'space': 6}),
    ('土台の「聴く姿勢」ができていない', {'size': 14, 'space': 0}),
    ('まま金額の話をしても、価格だけで', {'size': 14, 'space': 0}),
    ('判断されて終わります。', {'size': 14, 'space': 8}),
    ('だから今日は一番下の2段を', {'size': 14, 'bold': True, 'space': 0}),
    ('徹底的にやります。', {'size': 14, 'bold': True}),
])
notes(s, '・ピラミッドは毎回冒頭で出す。今日はどこをやるのかを必ず指差しで示す')

# ================================================================ 6. 全体像② 20のルール
s = new_slide('全体像②：太陽光＋蓄電池営業 20のルール')
HL = {'fill': LTYEL}
rows = [
    ['項目', 'ゴール', 'ルール', '詳細'],
    [('前準備', dict(HL)), ('必要な情報を揃える', dict(HL)), ('①商談テーブルの確認', dict(HL)),
     ('検針票・図面が手元にある状態で夫婦両主権者・1時間以上の時間を抑える', dict(HL))],
    [('聴く姿勢作り', dict(HL)), ('商談の趣旨を認識する', dict(HL)), ('②自社の概要・訪問趣旨の理解', dict(HL)),
     ('この地域で実績が豊富にあり、今回はその説明であることを理解してもらう', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('③未導入理由の確認', dict(HL)),
     ('これまで太陽光＋蓄電池を導入してこなかった理由を尋ねる', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('④前向きな検討に対する言質', dict(HL)),
     ('導入しなかった理由が解消できれば導入する予定であることに同意を得る', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('⑤太陽光が当たり前であることの理解', dict(HL)),
     ('10年の歴史を経て、蓄電池とのセット導入が標準になったことを理解してもらう', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('⑥直近の国策動向のおさらい', dict(HL)),
     ('国策を踏まえ、ますますスマートハウスが普及することを実感してもらう', dict(HL))],
    [('必要性訴求', dict(HL)), ('現状を正確に把握する', dict(HL)), ('⑦直近数か年の電気代推移の理解', dict(HL)),
     ('電力自由化でも電気代は下がっていないことを実感してもらう', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('⑧再エネ賦課金の構造の理解', dict(HL)),
     ('賦課金の上昇もあり、今後ますます光熱費が膨らむことを想像してもらう', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('⑨長期で支払う電気代総額の理解', dict(HL)),
     ('30年で400万円以上を支払う未来を創造してもらう', dict(HL))],
    ['', '', '⑩非常時に備える意義の理解', '第3回'],
    ['', '', '⑪〜⑫', '第3回'],
    ['時期訴求', '早期対策を進める', '⑬〜⑮', '第4回'],
    ['金額訴求', '費用対効果に納得する', '⑯〜⑳', '第4回'],
]
table(s, 0.30, 1.50, 10.25, [1.15, 1.70, 2.70, 4.70], rows, font_size=10, row_h=0.335, head_h=0.30)
_, tf = rect(s, 0.30, 6.42, 10.25, 0.42, fill=None, line=RED, lw=1.5)
lines(tf, [('黄色が本日の範囲。①〜⑨で「話を聞く姿勢」と「このままではまずい」までを作ります',
            {'size': 16, 'bold': True, 'align': PP_ALIGN.CENTER})])

# ================================================================ 7. アプローチブックの地図
s = new_slide('全体像③：アプローチブック26ページの地図')
box, tf = tb(s, 0.45, 1.38, 9.9, 0.4)
lines(tf, [('商談でお客様の前にめくるのは、この1冊だけです', {'size': 20, 'bold': True})])
blocks = [
    (1.85, 'Paragraph 01', 'これまでの10年と、これからの10年', 'p.1〜p.7',
     'ルール⑤⑥　→　本日', LTGREEN, '本日'),
    (3.20, 'Paragraph 02', 'なぜ必要とされているのか？', 'p.8〜p.21',
     'ルール⑦〜⑫　→　p.9〜17は本日／p.18〜21は第3回', LTGREEN, '本日＋第3回'),
    (4.55, 'Paragraph 03', 'お得に導入できるのは？', 'p.22〜p.26',
     'ルール⑬〜⑮　→　第4回', LTGRAY, '第4回'),
]
for y, tag, title, pages, rule, fill, badge in blocks:
    _, tf = rect(s, 0.45, y, 8.05, 1.18, fill=fill)
    lines(tf, [
        ('%s　%s' % (tag, title), {'size': 18, 'bold': True, 'space': 5}),
        ('%s　／　%s' % (pages, rule), {'size': 14, 'color': GRAY}),
    ])
    _, tf2 = rect(s, 8.70, y + 0.28, 1.70, 0.62,
                  fill=(RED if '本日' in badge else '808080'))
    lines(tf2, [(badge, {'size': 15, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER})])
_, tf = rect(s, 0.45, 5.90, 9.95, 0.95, fill=LTYEL)
lines(tf, [
    ('金額訴求（ルール⑯〜⑳）はアプローチブックの外。', {'size': 17, 'bold': True, 'space': 5}),
    ('シミュレーション・概算見積・カタログに持ち替えます（第4回）。', {'size': 17}),
])
notes(s, '・アプローチブックを実際に手に持たせ、3つのParagraphの区切りを指で確認させる')

# ================================================================ 8. なぜセット販売か
s = new_slide('なぜ「セット販売」なのか')
box, tf = tb(s, 0.45, 1.38, 9.9, 0.42)
lines(tf, [('太陽光と蓄電池は切り離さない。常にセットで話を進めます',
            {'size': 22, 'bold': True, 'color': RED})])
rows = [
    ['', '根拠', '中身'],
    ['①', 'シングル販売は「価格市場」に移行していく',
     '成約の決め手は「実績」→「価格」→「紹介」へ移る。D2Cの参入で価格崩壊が起きている地域もある'],
    ['②', 'セット販売は「空白市場」である',
     '太陽光設置済み約240万（総世帯の4％）に販売店が集中。持家は約3,000万（52％）、総世帯は約5,700万'],
    ['③', '経済合理性がセットにある',
     '売電単価15円 ＜ 買電単価30円。売るより貯めて使う方が得なので、太陽光だけでは最適化できない'],
    ['④', '設備面でもセットが前提になった',
     'ハイブリッドパワコンが普及（1台で制御でき変換効率も高い）。ZEH要件に蓄電池必須の基準案も提出'],
    ['⑤', 'セット販売の成功事例が出ている',
     '陽幸＝セット比率99％で月販110台／日本ライフサポート＝月販90台／プレジャーハウス＝7名で月販30台'],
]
table(s, 0.45, 1.92, 9.95, [0.50, 3.15, 6.30], rows, font_size=12, row_h=0.74, head_h=0.32)
_, tf = rect(s, 0.45, 6.30, 9.95, 0.58, fill=LTYEL)
lines(tf, [('「蓄電池を売る」ではなく「スマートハウスを導入いただく」と言い切る',
            {'size': 20, 'bold': True, 'align': PP_ALIGN.CENTER, 'hl': None})])
notes(s, '・ここが崩れると以降のトークが全部シングル販売になる。冒頭で必ず握る')

# ================================================================ 9. 章見出し 聴く姿勢作り
s = new_slide('聴く姿勢作り')
box, tf = tb(s, 0.45, 1.38, 9.9, 0.42)
lines(tf, [('ゴール：この営業マンの話を聞きたい、と思わせる', {'size': 22, 'bold': True})])
rows = [
    ['項目', 'ゴール', 'ルール', '詳細'],
    ['前準備', '必要な情報を揃える', '①商談テーブルの確認',
     '検針票・図面が手元にある状態で夫婦両主権者・1時間以上の時間を抑える'],
    ['聴く姿勢作り', '商談の趣旨を認識する', '②自社の概要・今回の訪問趣旨の理解',
     'この地域で実績が豊富にあり、今回はその説明であることを理解してもらう'],
    ['', '', '③太陽光＋蓄電池の未導入理由の確認', 'これまで導入してこなかった理由を尋ねる'],
    ['', '', '④前向きな検討に対する言質',
     '導入しなかった理由が解消できれば導入する予定であることに同意を得る'],
    ['', '', '⑤太陽光が当たり前であることの理解',
     '10年の歴史を経て、蓄電池とのセット導入が標準になったことを理解してもらう'],
    ['', '', '⑥直近の国策動向のおさらい',
     '国策を踏まえ、ますますスマートハウスが普及することを実感してもらう'],
]
for r in rows[1:]:
    for i in range(len(r)):
        r[i] = (r[i], {'fill': LTGREEN})
table(s, 0.30, 1.95, 10.25, [1.25, 1.75, 2.85, 4.40], rows, font_size=11, row_h=0.42, head_h=0.30)
_, tf = rect(s, 0.30, 4.88, 4.95, 2.00, fill=LTGRAY)
lines(tf, [
    ('使うアプローチブック', {'size': 17, 'bold': True, 'space': 8}),
    ('p.2　本日お話ししたいこと', {'size': 15, 'space': 5}),
    ('p.3　脱炭素社会に向け', {'size': 15, 'space': 5}),
    ('p.4　10年の普及年表', {'size': 15, 'space': 5}),
    ('p.5　大手ハウスメーカー', {'size': 15, 'space': 5}),
    ('p.6　お客様の声', {'size': 15, 'space': 5}),
    ('p.7　未導入のお家 vs 太陽光＋蓄電池', {'size': 15}),
])
_, tf = rect(s, 5.60, 4.88, 4.95, 2.00, fill=LTYEL)
lines(tf, [
    ('このパートの勝ち筋', {'size': 17, 'bold': True, 'space': 8}),
    ('①③④は資料を使わない「ヒアリング」。', {'size': 15, 'space': 6}),
    ('先に相手の話を聞くことが、', {'size': 15, 'space': 0}),
    ('自分の話を聞いてもらうことに直結します。', {'size': 15, 'space': 8}),
    ('②⑤⑥で資料を開き、', {'size': 15, 'space': 0}),
    ('「先生－生徒」の関係をつくります。', {'size': 15}),
])

# ---------------------------------------------------------------- rule helpers
def rule_head(s, no_title, point):
    box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
    lines(tf, [(no_title, {'size': 24, 'bold': True})])
    _, tf = rect(s, 0.45, 1.88, 9.95, 0.46, fill=None, line=RED, lw=1.25)
    p = para(tf, first=True, align=PP_ALIGN.LEFT)
    run(p, 'POINT　', size=15, bold=True, color=RED)
    run(p, point, size=16, bold=True)

def goal_box(s, x, y, w, h, quote):
    _, tf = rect(s, x, y, w, h, fill=LTBLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    lines(tf, [
        ('ゴール（お客様の反応）', {'size': 13, 'bold': True, 'color': BLUE, 'space': 6}),
    ])
    for q in quote:
        p = para(tf, space_after=3)
        run(p, '「%s」' % q, size=16, bold=True)

def points_box(s, x, y, w, h, items, title='説明のポイント'):
    _, tf = rect(s, x, y, w, h, fill=LTGRAY)
    lines(tf, [(title, {'size': 13, 'bold': True, 'color': GRAY, 'space': 6})])
    for it in items:
        opts = {}
        if isinstance(it, tuple): it, opts = it
        p = para(tf, space_after=opts.pop('space', 6))
        run(p, it, size=opts.pop('size', 12.5), bold=opts.pop('bold', False),
            color=opts.pop('color', BLACK))

def conclusion(s, text, y=6.36):
    _, tf = rect(s, 0.45, y, 9.95, 0.66, fill=LTGREEN)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, 'このページの結論　', size=12, bold=True, color=GRAY)
    run(p, text, size=15, bold=True, hl=YELLOW)

def talk_slide(s, no_title, talks, conc=None, note=None):
    box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
    lines(tf, [(no_title, {'size': 24, 'bold': True})])
    y = 1.92
    for label, body, h in talks:
        _, tf = rect(s, 0.45, y, 8.05, h, fill=LTGRAY, anchor=MSO_ANCHOR.MIDDLE)
        for i, ln in enumerate(body):
            p = para(tf, first=(i == 0), space_after=3)
            run(p, ln, size=14)
        _, tf2 = rect(s, 8.68, y, 1.72, h, fill=ORANGE, anchor=MSO_ANCHOR.MIDDLE)
        lines(tf2, [(label, {'size': 13, 'bold': True, 'color': WHITE,
                             'align': PP_ALIGN.CENTER})])
        y += h + 0.16
    if conc: conclusion(s, conc)
    if note: notes(s, note)

# ================================================================ 10-11 ルール①
s = new_slide('聴く姿勢作り')
rule_head(s, 'ルール①：商談テーブルの確認', '太陽光＋蓄電池商談の3条件＝検針票と図面／夫婦の両主権者／1時間')
conds = [('検針票・図面が\n手元にある', LTYEL), ('夫婦\n両主権者', LTYEL), ('1時間以上の\n時間を抑える', LTYEL)]
for i, (t, f) in enumerate(conds):
    _, tf = rect(s, 0.45 + i * 3.32, 2.55, 3.05, 1.15, fill=f, line='BF8F00')
    for j, ln in enumerate(t.split('\n')):
        p = para(tf, first=(j == 0), align=PP_ALIGN.CENTER, space_after=2)
        run(p, ln, size=19, bold=True)
_, tf = rect(s, 0.45, 3.95, 9.95, 2.35, fill=LTGRAY)
lines(tf, [
    ('考え方', {'size': 14, 'bold': True, 'color': GRAY, 'space': 8}),
    ('太陽光・蓄電池の失注理由は、「商談設定」か「商談内容」のどちらかしかありません。', {'size': 16, 'space': 6}),
    ('営業トークを磨けば「商談内容」の失注は防げますが、', {'size': 16, 'space': 6}),
    ('「商談設定」が適切に行われていなければ、決まるはずのものも決まらなくなります。', {'size': 16, 'space': 12}),
    ('だから、太陽光・蓄電池営業の第一歩は、この商談テーブルの設定です。', {'size': 18, 'bold': True, 'color': RED, 'space': 6}),
    ('確認は商談の前日に入れることを推奨します。', {'size': 16, 'bold': True}),
])
notes(s, '・「決まらなかった商談を思い出してください。3条件のどれが欠けていましたか？」と問いかける')

s = new_slide('聴く姿勢作り')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ルール①：よくあるネガと解消法', {'size': 24, 'bold': True})])
rows = [
    ['条件', 'よくあるネガ', '解消法'],
    ['①検針票・図面が手元にある', '図面が見つからない',
     '築20年以内のものであれば、住所のある役所に問い合わせれば入手可能です'],
    ['②夫婦両主権者', '片主権者が揃わない',
     '（1）双方それぞれ別の観点で同席してほしい、または（2）当社の規定上そうしている旨を伝えます'],
    ['③1時間以上の時間を抑える', '短くしてほしいと言われる',
     '「お家に合った話をしないと意味が無いので」ということを伝えます'],
]
table(s, 0.45, 1.95, 9.95, [2.45, 2.45, 5.05], rows, font_size=13, row_h=0.85, head_h=0.32)
_, tf = rect(s, 0.45, 4.90, 9.95, 1.35, fill=LTYEL)
lines(tf, [
    ('ワーク（3分）', {'size': 17, 'bold': True, 'space': 8}),
    ('自分が直近に失注した商談を1件思い出し、3条件のうち欠けていたものに印をつけてください。', {'size': 16, 'space': 5}),
    ('前日確認をしていれば防げたかどうか、ペアで共有します。', {'size': 16}),
])

# ================================================================ 12-14 ルール②
s = new_slide('聴く姿勢作り')
rule_head(s, 'ルール②：自社の概要・今回の訪問趣旨の理解', '自社に関する「ストーリー」は、商談冒頭に売る')
pic(s, ab_img(2), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.2　本日お話ししたいこと', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.05, ['商談の構成が分かる'])
points_box(s, 5.35, 3.70, 5.05, 2.60, [
    ('「で、金額は？」と言われないためにも、最初に商談の構成・内容をお伝えする', {'bold': True}),
    'ジャパネットHDは創業者の退任後、数千あった商品を600に絞った。話者が本当にその価値を語れる商品だけを残したという',
    'モノがあふれる時代、高価格商材の最大のポイントは「コトを売る」こと',
    ('自社の「ストーリー」の価値が、選ばれる理由になる', {'bold': True, 'color': RED, 'size': 13}),
])
notes(s, '・p.2を開いて、①なぜ増えているのか ②シミュレーション ③お申込みのご案内 の3点を指で追わせる')

s = new_slide('聴く姿勢作り')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ルール②：自社ストーリーの5段構成', {'size': 24, 'bold': True})])
rows = [
    ['', '目次', 'トーク例'],
    ['①', '企業理念', '当社は「太陽光を通じて幸せを届ける」ということを大事にさせて頂いています'],
    ['②', '創業ヒストリー',
     '元々、代表の田中が勤めていた訪販の会社では顧客目線で考えることが少なく、大クレームに発展して痛い目を見たことがありまして'],
    ['③', '当社のコダワリ', 'なので「目の前のお客様を大事にする」ということだけは譲らないようにしているのですよね'],
    ['④', 'サービス内容',
     '施工いただいたお客様には、必ず3年ごとに訪問させていただいています。形だけにならないよう報告書も書いて、モニターの写真も撮らせていただいています'],
    ['⑤', '営業マンの信条',
     'だから正直、私たちも末永いお付き合いだから嘘が付けなくって（笑）今回はぜひ、光熱費に関するリアルなお話が出来たらと思っています'],
]
table(s, 0.45, 1.95, 6.35, [0.40, 1.55, 4.40], rows, font_size=11, row_h=0.75, head_h=0.30)
_, tf = rect(s, 7.00, 1.95, 3.40, 4.05, fill=LTGRAY)
lines(tf, [
    ('効いている心理法則', {'size': 15, 'bold': True, 'space': 10}),
    ('物語効果', {'size': 14, 'bold': True, 'color': BLUE, 'space': 3}),
    ('事実の羅列より、物語の方が理解しやすく記憶に残る', {'size': 12, 'space': 10}),
    ('ラポール形成', {'size': 14, 'bold': True, 'color': BLUE, 'space': 3}),
    ('信条やコダワリが、個人的なつながりと共感を育む', {'size': 12, 'space': 10}),
    ('自己開示の法則', {'size': 14, 'bold': True, 'color': BLUE, 'space': 3}),
    ('一方が開示すると相手も開示しやすくなる', {'size': 12, 'space': 10}),
    ('独自の価値提案の明確化', {'size': 14, 'bold': True, 'color': BLUE, 'space': 3}),
    ('3年ごとの定期訪問など、他社にない強みが選ばれる理由になる', {'size': 12}),
])
conclusion(s, 'パンフレットを使いながら、鉄板で使える自社のストーリーを組み立てておく', y=6.20)

s = new_slide('聴く姿勢作り')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ワーク：自社のストーリーを書く（10分）', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.58)
lines(tf, [('実例をもとに「お客様から選ばれる自社（自分）のストーリー」を明確にしてみましょう。自社のパンフレットがあると尚良いです。',
            {'size': 14, 'color': GRAY})])
rows = [['', '目次', '自分のトーク（記入してください）']]
for n, t in [('①', '企業理念'), ('②', '創業ヒストリー'), ('③', '当社のコダワリ'),
             ('④', 'サービス内容'), ('⑤', '営業マンの信条')]:
    rows.append([n, t, ''])
table(s, 0.45, 2.52, 9.95, [0.45, 1.85, 7.65], rows, font_size=13, row_h=0.66, head_h=0.32)
_, tf = rect(s, 0.45, 6.32, 9.95, 0.52, fill=LTYEL)
lines(tf, [('書けたらペアで実演。聞き手は「この人から買ってもいいと思えたか」を答える',
            {'size': 17, 'bold': True, 'align': PP_ALIGN.CENTER})])

# ================================================================ 15-16 ルール③
s = new_slide('聴く姿勢作り')
rule_head(s, 'ルール③：太陽光＋蓄電池の未導入理由の確認', '3回聞くことで、お客様のNOに近づくヒアリング')
_, tf = rect(s, 0.45, 2.50, 4.55, 3.60, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('このルールは資料を使いません', {'size': 17, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 12}),
    ('アプローチブックは閉じたまま、', {'size': 15, 'space': 4}),
    ('お客様の顔を見て聞くパートです。', {'size': 15, 'space': 14}),
    ('手元に置くもの', {'size': 14, 'bold': True, 'color': GRAY, 'space': 6}),
    ('・検針票（お客様が出したもの）', {'size': 14, 'space': 4}),
    ('・ご自宅の図面', {'size': 14, 'space': 4}),
    ('・メモ（聞いた内容を必ず書き取る）', {'size': 14}),
])
points_box(s, 5.35, 2.50, 5.05, 3.60, [
    ('太陽光・蓄電池営業の基本は「お客様のNO（導入しない理由）を潰すこと」', {'bold': True, 'size': 15}),
    '商品力が高い（10数年で元を取れる）のに踏み切れないのは、お客様の中で何らかのNO（不便・不満・不安）が解消されなかったから',
    '訪問販売・催事販売・テレアポ等で、大多数の方は過去に太陽光の話を耳に入れている。基礎知識は持っている前提で、「導入しなかった理由」を明確にしてから商談に入る',
    ('冒頭に相手の話を「聞く」ことが、自分の話を「聞いてもらう」ことに直結する', {'bold': True, 'color': RED, 'size': 15}),
], title='考え方')

s = new_slide('聴く姿勢作り')
talk_slide(s, 'ルール③：3回のオープンクエスチョン', [
    ('オープンQ①', [
        'Q1「どうしてこれまでは太陽光を導入されてこなかったのですか？」',
        '＜よくある展開＞価格が高かったから／シミュレーションを見てピンとこなかったから／会社が信用しきれなかったから',
    ], 1.05),
    ('オープンQ②', [
        'Q2「何がキッカケで、今回は当社に聞いてみようと思われたのですか？」',
        '＜よくある展開＞たまたまチラシが入っていたのを見たので／買い物に来ていたついでに聞いた／安くなったから',
    ], 1.05),
    ('オープンQ③', [
        'Q3「蓄電池や太陽光には、どんなイメージをお持ちですか？」',
        '＜よくある展開＞まだ高そう／最近改めて国が推している／電気自動車の方がいいのでは',
    ], 1.05),
], conc='ここで出たNOを、⑤⑥と⑦〜⑨で1つずつ潰していく',
   note='・3問を声に出して言う練習まで行う。メモを取る姿勢もセットで確認する')

# ================================================================ 17-18 ルール④
s = new_slide('聴く姿勢作り')
rule_head(s, 'ルール④：前向きな検討に対する言質', '「やる／やらないを今日決める」ことの言質を取る')
points_box(s, 0.45, 2.50, 4.95, 3.60, [
    ('太陽光・蓄電池の営業は基本的に「即決」', {'bold': True, 'size': 16, 'color': RED}),
    '日にちをおくことで△（少し考える）が〇（成約）になることは多くありません。アポも商談も、その日に決めてもらうことを中心に考える必要があります',
    ('活用するのは「一貫性の法則」', {'bold': True, 'size': 16}),
    '一度Yesと首を縦に振ったことについて、人は否定することをためらうという原則です',
    'ルール③でヒアリングした「以前に太陽光を導入しなかった理由」が解消されるのであれば前向きに検討いただけそうか、について言質を取っておきます',
], title='考え方')
_, tf = rect(s, 5.75, 2.50, 4.65, 1.70, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('注意1　自然に行う', {'size': 16, 'bold': True, 'space': 6}),
    ('前段のトークより間を開けず、わざとらしさ・営業っぽさを一切なくすこと。', {'size': 13, 'space': 4}),
    ('間が空いたり言葉遣いが変だと一瞬で見抜かれ、主導権が逆転します。', {'size': 13}),
])
_, tf = rect(s, 5.75, 4.40, 4.65, 1.70, fill='FCE4E4', line=RED)
lines(tf, [
    ('注意2　特定商取引法への抵触', {'size': 16, 'bold': True, 'color': RED, 'space': 6}),
    ('訪問目的を秘匿して商品を売るのは違法行為で、通報されれば営業停止処分もありえます。', {'size': 13, 'space': 4}),
    ('アポ取得時の言い回しを必ず確認してください。', {'size': 13}),
])
notes(s, '・「今日決める」と言わせることが目的ではなく、決められる状態をつくるのが目的だと補足する')

s = new_slide('聴く姿勢作り')
talk_slide(s, '切り返しの鉄則（①〜④で使い倒す）', [
    ('Yes, but法', [
        'まず全力で同意する。相手は意見を受け入れてもらえたことで安心し、こちらの話も聞いてくれるようになる。',
        '「出費増えるのは大変ですよね！ただ、災害時の対策といった災害保険の観点もありますし」',
    ], 0.85),
    ('第三者話法', [
        '相手の意見に自分の意見を重ねるのは最大のNG。事例・事実で返す。',
        '「そうですよね、ご心配ですよね！B様というお客様も全く同じことをおっしゃってました！ただ、お話を聞いた後『やっぱり太陽光いいかも…』と言われたんですよね！」',
    ], 0.95),
    ('YES取り', [
        '「なんで、最初は出費が増えるのが嫌だっておっしゃってたB様が最終的に導入した理由、気になりませんか？」⇒「はい」',
    ], 0.62),
], conc='話を真に受けない。ずれたら「後でしっかりお伝えします！」で戻す',
   note='・聞く姿勢の途中で経済メリットや製品保証に話がずれると営業のリズムを失う。避け方まで練習させる')

# ================================================================ 19-20 ルール⑤
s = new_slide('聴く姿勢作り')
rule_head(s, 'ルール⑤：太陽光が当たり前であることの理解', '包み込みの法則①　太陽光の歴史を押さえる')
pic(s, ab_img(4), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.4　10年の普及年表', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.30,
         ['皆さん、太陽光・蓄電池をつけているのですね', '義務化・今後必須の流れになるの!?'])
points_box(s, 5.35, 3.92, 5.05, 2.38, [
    ('2012年 FIT開始 → 2021年 義務化の協議開始 → 2025年 東京都で新築に義務化 → 2030年 ZEH義務化', {'bold': True, 'size': 13}),
    '累計件数は 約120万件 → 約270万件 → 約350万件（2024年度末）と伸びている',
    'ZEH住宅の要件に蓄電池が必須になる基準案が提出されていることまで伝える',
    ('営業が優位に立てるのは、商品の「歴史」と「現状」を知っていること', {'bold': True, 'color': RED, 'size': 13}),
])
conclusion(s, '国の政策はすべて「電気を極力使わない生活へ」という方向に動いているのです')
notes(s, '・年表は年号を指差しながら1つずつめくる。累計件数の家アイコンが増えていくのを見せる')

s = new_slide('聴く姿勢作り')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ルール⑤：お客様のイメージ別・反論処理', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.35)
lines(tf, [('「そういう風に考えられてしまう時もありますよね」と包み込みながら、時系列で処理します',
            {'size': 15, 'color': GRAY})])
rows = [
    ['お持ちのイメージ', 'よくある誤解', '理解してもらうこと'],
    ['導入初期', '「FITありきで太陽光が増えた」',
     '① 制度施行以前にも50万件以上の実績があったこと　② 諸外国ではFIT制度で日本ほど爆増している所は少ないこと'],
    ['導入中期', '「補助金ありきの商品（設置しているのは一部）」',
     '① 安倍首相はじめ主要官僚のもと各ハウスメーカーに号令がかかったこと　② 設置対象外（築30年以上）を除けば設置率20％前後であること'],
    ['導入後期', '「もう機を逃しているのでは？」',
     '① 義務化が議論できるほど広く普及できるレベルになったこと　② 最もネックだったイニシャルコストが規定範囲内に収まってきたこと'],
]
table(s, 0.45, 2.35, 9.95, [1.55, 3.10, 5.30], rows, font_size=12, row_h=1.05, head_h=0.32)
_, tf = rect(s, 0.45, 6.00, 9.95, 0.80, fill=LTYEL)
lines(tf, [
    ('包み込みの法則＝お客様より優位に立てる範囲を明確にし、「先生－生徒」の関係をつくること',
     {'size': 16, 'bold': True, 'space': 4}),
    ('「私も当時その道は通りましたが、歴史を知ったことでその意気は脱しました」というトーンで。', {'size': 14}),
])

# ================================================================ 21-22 ルール⑥
s = new_slide('聴く姿勢作り')
rule_head(s, 'ルール⑥：直近の国策動向のおさらい', '包み込みの法則②　多数派の意見を押さえる')
pic(s, ab_img(5), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.5　大手ハウスメーカー', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.30,
         ['これからはスマートハウスが当たり前になるのですね'])
points_box(s, 5.35, 3.92, 5.05, 2.38, [
    ('積水ハウスの戸建住宅ZEH率は96％（ZEH累積世界No.1）', {'bold': True, 'size': 13}),
    '大和ハウスは2011年に堺市で全65区画をZEH化。1戸あたり太陽光5.2kW・蓄電池6.2kWh',
    '1922年の借家率は7％。「持ち家信仰」もたった100年の変化にすぎない',
    ('他の商品に例えながら「導入が当たり前になる」と強調する', {'bold': True, 'color': RED, 'size': 13}),
])
conclusion(s, '至るところでスマートハウスのお家が増えていっています')

s = new_slide('聴く姿勢作り')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ルール⑥：営業マンが持つべき「2つの当然意識」', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.32)
lines(tf, [('売り手側の自信の無さは、買い手側にすぐ分かると言われます', {'size': 15, 'color': GRAY})])
_, tf = rect(s, 0.45, 2.30, 4.95, 2.35, fill=LTGREEN)
lines(tf, [
    ('意識①　太陽光・蓄電池のあるお家が普通になっていく', {'size': 16, 'bold': True, 'space': 8}),
    ('・東京都での新築への太陽光パネル設置義務化', {'size': 14, 'space': 5}),
    ('・パリ協定の目標により、2050年までにネットゼロを達成するため国・地方自治体が設置を進めている', {'size': 14, 'space': 5}),
    ('・国交省「子育てグリーン住宅支援事業」で太陽光・蓄電池付き高性能住宅への補助金が増額し続けている', {'size': 14}),
])
_, tf = rect(s, 5.60, 2.30, 4.95, 2.35, fill=LTYEL)
lines(tf, [
    ('意識②　太陽光・蓄電池は基本的に一緒に導入する', {'size': 16, 'bold': True, 'space': 8}),
    ('・売電単価（15円）よりも買電単価（30円）が上回っている', {'size': 14, 'space': 5}),
    ('・ハイブリッドパワコンが広く普及している', {'size': 14, 'space': 12}),
    ('太陽光・蓄電池は切り離さず、基本的にセットで話を進めるのがポイント', {'size': 15, 'bold': True, 'color': RED}),
])
pic(s, ab_img(7), 6.90, 4.72, 3.10)
box, tf = tb(s, 6.90, 6.92, 3.10, 0.26)
lines(tf, [('アプローチブック p.7', {'size': 11, 'color': GRAY})])
_, tf = rect(s, 0.45, 4.80, 5.90, 1.75, fill=LTGRAY)
lines(tf, [
    ('p.7でテストクロージング（二者択一話法）', {'size': 16, 'bold': True, 'space': 8}),
    ('「未導入のお家と、太陽光＋蓄電池のお家、', {'size': 15, 'space': 2}),
    ('　どちらが良さそうですか？」', {'size': 15, 'space': 8}),
    ('→ ここでYESを取ってから必要性訴求へ進む', {'size': 15, 'bold': True, 'color': RED}),
])
notes(s, '・「じゃあ太陽光・蓄電池つけた方がいいかもねえ」を引き出せたら合格')

# ================================================================ 23 休憩
s = new_slide(None)
_, tf = rect(s, 2.40, 2.85, 6.05, 1.80, fill=LTGREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
lines(tf, [
    ('休憩　10分', {'size': 44, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 10}),
    ('後半は「必要性訴求」に入ります', {'size': 18, 'align': PP_ALIGN.CENTER}),
])

# ================================================================ 24 章見出し 必要性訴求
s = new_slide('必要性訴求')
box, tf = tb(s, 0.45, 1.38, 9.9, 0.42)
lines(tf, [('ゴール：現状を正確に把握してもらう', {'size': 22, 'bold': True})])
rows = [
    ['項目', 'ゴール', 'ルール', '詳細', '本日'],
    ['必要性訴求', '現状を正確に把握する', '⑦直近数か年の電気代推移の理解',
     '電力自由化でも電気代は下がっていないことを実感してもらう', '○'],
    ['', '', '⑧再エネ賦課金の構造の理解',
     '賦課金の上昇もあり、今後ますます光熱費が膨らむことを想像してもらう', '○'],
    ['', '', '⑨長期で支払う電気代総額の理解',
     '30年で400万円以上を支払う未来を創造してもらう', '○'],
    ['', '', '⑩非常時に備える意義の理解', '度重なる災害から、非常時に備える必要があることを理解してもらう', '第3回'],
    ['', '', '⑪非常時対策の効果の確認', '30年近く停電時の保証が得られることを理解してもらう', '第3回'],
    ['', '', '⑫太陽光・蓄電池の使い方の理解', 'シミュレーションに移る前に一般的な使い方を理解してもらう', '第3回'],
]
for i, r in enumerate(rows[1:], 1):
    fill = LTGREEN if i <= 3 else None
    for j in range(len(r)):
        r[j] = (r[j], {'fill': fill} if fill else {})
table(s, 0.30, 1.95, 10.25, [1.20, 1.65, 2.60, 3.90, 0.90], rows, font_size=11,
      row_h=0.44, head_h=0.30)
_, tf = rect(s, 0.30, 5.15, 10.25, 1.70, fill=LTYEL)
lines(tf, [
    ('このパートの基本設計は「下げてから上げる」', {'size': 20, 'bold': True, 'color': RED, 'space': 8}),
    ('ネガティブな現状を踏まえた上でポジティブな解決策を提示すると、商品の価値が高まります。', {'size': 16, 'space': 5}),
    ('お客様が「欲しい」と思う動機は、ネガ→ポジの「振れ幅」です。', {'size': 16, 'space': 5}),
    ('ポイントは、ネガティブな現状を「自分事」にすること。「ウチは関係ないしね」を無くします。', {'size': 16, 'bold': True}),
])

# ================================================================ 25-26 ルール⑦
s = new_slide('必要性訴求')
rule_head(s, 'ルール⑦：直近数か年の電気代推移の理解', '「しばらく電気代が下がっていない」事実に納得させる')
pic(s, ab_img(9), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.9　電気料金の推移', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.15, ['10年間で5.5万円も電気代は上がっているのですね'])
points_box(s, 5.35, 3.78, 5.05, 2.50, [
    ('電気代が具体的にいくら上がっているのかをヒアリングし、正解を伝えてあげる', {'bold': True, 'size': 14}),
    '2010年 ¥119,602 → 2025年 ¥180,675（4人家族・年5,500kWh想定）',
    'p.10「電力自由化でも下がらなかった」（21円→32円）、p.12「全国的な値上げ」まで続けて見せる',
    ('「全国平均の話で、ウチは関係ない」を無くすのが目的', {'bold': True, 'color': RED, 'size': 13}),
])
conclusion(s, '電気料金は2011年以降45％、約5.5万円も上昇している')
notes(s, '・お客様の検針票を手元に置き、実額と並べて話す。数字は必ず「〇〇様の場合は」に変換する')

s = new_slide('必要性訴求')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ルール⑦：よくあるネガと切り返し', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.32)
lines(tf, [('逃げ道を無くすことが目的です。お客様の細かなネガにも、なるべく対応できるようにしておきましょう',
            {'size': 15, 'color': GRAY})])
rows = [
    ['よくあるネガ', '切り返し方'],
    ['検針票が届かないので、いくらか把握できていない', '口座引き落としをご覧ください'],
    ['あまり費用が大きいと感じたことはない', '季節によっても変わります'],
    ['長く住んでいるので今さら…', '平均寿命を考えれば、定年からでも20年は払い続けますよね？'],
    ['節電対策はそれなりにやっている', 'こまめな節電は大事ですが、せいぜい1000円ぐらいです'],
    ['子供次第？', '大きくなって部屋も別々になれば、増えていきます'],
]
table(s, 0.45, 2.32, 9.95, [4.35, 5.60], rows, font_size=15, row_h=0.66, head_h=0.34)
conclusion(s, '節電で削減できるのは年13,646円が限界。買い続ける限り対策は「使用量を減らす」だけ', y=6.15)

# ================================================================ 27-28 ルール⑧
s = new_slide('必要性訴求')
rule_head(s, 'ルール⑧：再エネ賦課金の構造の理解', '賦課金は「13年で18倍」に上がり、その賦課金は「わざわざ払っている」')
pic(s, ab_img(14), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.14　再エネ賦課金の推移', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.15, ['再エネ賦課金が上昇することで、電気代が上昇するのですね'])
points_box(s, 5.35, 3.78, 5.05, 2.50, [
    ('太陽光を設置している人ですら、今でも知らない人が大多数', {'bold': True, 'size': 14}),
    '検針票も電子化され、上昇を把握しているケースは稀。丁寧に説明して理解してもらう',
    '2012年 0.22円/kWh（世帯負担1,210円）→ 2026年 4.18円/kWh（世帯負担22,990円）',
    ('⑦の上にもう一押し「今のままじゃダメだ」を突きつける', {'bold': True, 'color': RED, 'size': 13}),
])
conclusion(s, '裏返せば、太陽光を導入している方は続々とお得になっています！')
notes(s, '・2023年だけ下がって見えるのは化石燃料費高騰による例外措置。聞かれたら必ず補足する')

s = new_slide('必要性訴求')
talk_slide(s, 'ルール⑧：3段で説明し、確実にYESを取る', [
    ('①賦課金とは', [
        '・使用量に応じて10年近く払い続けていること（年間で2万円程度にもなること）',
        '・検針票にも小さくしか書いておらず、気づかないうちに支払っていること',
    ], 0.78),
    ('②徴収構造', [
        '・基本料金とは別に課されているお金であること',
        '・裏返すと、太陽光・蓄電池を導入している方は続々とお得になってきていること',
    ], 0.78),
    ('③今後の上昇幅', [
        '・（少なくとも）5円／kWh程度まで上昇する見込みがあること',
        '・制度が先行して施行されている国では、更に上昇している地域もあること',
    ], 0.78),
    ('YES取り', [
        '「これからも電気代、上がりそうですよね？」⇒ ここでYESを取ってから⑨へ進む',
    ], 0.52),
], conc='電気を買い続ける＝賦課金を払い続けること')

# ================================================================ 29-30 ルール⑨
s = new_slide('必要性訴求')
rule_head(s, 'ルール⑨：長期で支払う電気代総額の理解', '第三者話法で「だから太陽光・蓄電池が導入されている」事実を伝える')
pic(s, ab_img(17), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.17　生涯に支払う電気代', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.15, ['30年間でそんなにも電気代を払わないといけないのですね'])
points_box(s, 5.35, 3.78, 5.05, 2.50, [
    ('月12,000円の方は、30年間で432万円。住宅ローンの一部に匹敵する金額', {'bold': True, 'size': 14}),
    '月2万円なら30年で700万円以上。毎月の支出としては認識していても、生涯でこれほどの金額とは想像していない',
    'この「見過ごされてきた巨額の支出」を明示し、「このままではいけない」という危機感を喚起する',
    ('煽るだけで終わらせず、解決策があることを示唆する', {'bold': True, 'color': RED, 'size': 13}),
])
conclusion(s, '月々2万円なら、30年間の負担額は700万円以上にもなる見通しです')

s = new_slide('必要性訴求')
talk_slide(s, 'ルール⑨：4ステップで自分事にする', [
    ('①自分事化', [
        '「これはあくまで一般的なご家庭の例ですが…」と前置きすることで、',
        '「自分の家もこれくらいか、あるいはもっとかかっているかもしれない」と捉えやすくさせる。',
    ], 0.78),
    ('②実感させる', [
        '教育費など、他の大きな支出や人生の重要なイベントと比較し、金額の大きさをリアルに感じさせる。',
    ], 0.55),
    ('③損失回避', [
        '「これは、今の生活を続けていれば『必ず』支払うことになるお金、とも言えますよね。」',
        '「しかも、先ほどお話ししたように、電気代は上がることはあっても、大幅に下がることは考えにくい状況です。」',
    ], 0.78),
    ('④次への誘導', [
        '「もし、この432万円という支出を、賢く活用する方法があるとしたら、少しお話を聞いてみたいと思われませんか？」',
    ], 0.55),
], conc='静かに、しかし確実に迫るリスクとして伝える',
   note='・声を大きくしない。淡々と数字を置いていくほうが効く')

# ================================================================ 31 ロープレ指示
s = new_slide('ロールプレイング')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('アプローチブックを手に持って、p.2〜p.17を通す（15分）', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 1.95, 4.95, 2.55, fill=LTGREEN)
lines(tf, [
    ('やること', {'size': 17, 'bold': True, 'space': 8}),
    ('① 3人1組（営業／お客様役ご主人／お客様役奥様）', {'size': 15, 'space': 6}),
    ('② 営業役はアプローチブックを実際にめくる', {'size': 15, 'space': 6}),
    ('③ ルール①〜⑨の順にひと通り進める', {'size': 15, 'space': 6}),
    ('④ 5分で交代。全員が営業役をやる', {'size': 15}),
])
_, tf = rect(s, 5.60, 1.95, 4.95, 2.55, fill=LTYEL)
lines(tf, [
    ('見られているポイント', {'size': 17, 'bold': True, 'space': 8}),
    ('・発言比率は 営業6：お客様4 になっているか', {'size': 15, 'space': 6}),
    ('・お客様側の4は、ご主人2：奥様2 になっているか', {'size': 15, 'space': 6}),
    ('・ラリー数は30回に届いているか', {'size': 15, 'space': 6}),
    ('・最後までセット提案を崩していないか', {'size': 15, 'bold': True, 'color': RED}),
])
rows = [
    ['進行', 'アプローチブック', '狙う反応'],
    ['ルール②', 'p.2', '商談の構成が分かる'],
    ['ルール③④', '（閉じたまま）', '導入しなかった理由が出る／今日決めることに同意'],
    ['ルール⑤⑥', 'p.3〜p.7', 'じゃあつけた方がいいかもねえ'],
    ['ルール⑦⑧⑨', 'p.9〜p.17', '30年でそんなにも払わないといけないのですね'],
]
table(s, 0.45, 4.62, 9.95, [1.60, 2.60, 5.75], rows, font_size=13, row_h=0.42, head_h=0.30)

# ================================================================ 32 顧客設定
s = new_slide('ロールプレイング')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('顧客設定（お客様役はこの設定で演じてください）', {'size': 24, 'bold': True})])
rows = [
    ['項目', '設定'],
    ['ご家族', '田中様ご夫婦（ご主人48歳・会社員／奥様45歳・パート）＋ 高校生・中学生のお子様2人'],
    ['お住まい', '築14年の戸建て（4LDK）。南向き切妻屋根。オール電化ではない'],
    ['導入状況', ('太陽光・蓄電池ともに未導入', {'bold': True, 'color': RED})],
    ['電気代', '月平均 18,000円（夏場は24,000円）。検針票はご主人が保管している'],
    ['きっかけ', '折込チラシを見て、奥様が問い合わせた'],
    ['温度感', 'ご主人＝慎重。「元が取れるのか」が最大の関心。奥様＝前向き。停電と電気代が心配'],
    ['断り文句', '「以前も訪問販売が来て、高かったからやめた」「一度持ち帰って考えたい」'],
]
table(s, 0.45, 1.92, 9.95, [1.60, 8.35], rows, font_size=12.5, row_h=0.54, head_h=0.32)
_, tf = rect(s, 0.45, 6.35, 9.95, 0.56, fill=LTYEL)
lines(tf, [('この設定は「太陽光も蓄電池も未導入」＝セット販売のケースです',
            {'size': 18, 'bold': True, 'align': PP_ALIGN.CENTER})])
notes(s, '・卒FIT客（太陽光あり・蓄電池のみ追加）はシングル販売のケース。今回は扱わない')

# ================================================================ 33 チェックシート
s = new_slide('チェックシート')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.42)
lines(tf, [('ロープレのあと、自分で○△×をつけてください', {'size': 22, 'bold': True})])
rows = [
    ['ルール', '該当ページ', '狙った台詞を引き出せたか', '結論を言えたか'],
    ['①商談テーブルの確認', '−', '3条件を確認できた', '−'],
    ['②自社の概要・訪問趣旨', 'p.2', '商談の構成が分かる', '−'],
    ['③未導入理由の確認', '−', '導入しなかった理由が3つ出た', '−'],
    ['④前向きな検討への言質', '−', '今日決めることに同意を得た', '−'],
    ['⑤太陽光が当たり前', 'p.3〜4', '義務化・必須の流れになるの!?', '電気を極力使わない生活へ'],
    ['⑥直近の国策動向', 'p.5〜7', 'スマートハウスが当たり前に', 'スマートハウスが増えている'],
    ['⑦電気代推移の理解', 'p.9〜12', '10年で5.5万円も上がっている', '2011年以降45％上昇'],
    ['⑧再エネ賦課金の構造', 'p.11・13〜16', '賦課金で電気代が上がる', '太陽光の方は続々とお得に'],
    ['⑨電気代総額の理解', 'p.17', '30年でそんなにも払うのですね', '30年で700万円以上'],
    [('商談運営', {'bold': True}), '−', '発言比率6：4／ラリー数30回', 'セット提案を崩さなかった'],
]
table(s, 0.30, 1.90, 10.25, [2.55, 1.55, 3.35, 2.80], rows, font_size=11.5, row_h=0.40, head_h=0.32)
_, tf = rect(s, 0.30, 6.35, 10.25, 0.50, fill=LTGREEN)
lines(tf, [('最後に「結論を言えたか」の欄だけ、資料を閉じて声に出して確認します（暗唱チェック）',
            {'size': 17, 'bold': True, 'align': PP_ALIGN.CENTER})])

# ================================================================ 34 まとめ・宿題
s = new_slide('本日のまとめ')
_, tf = rect(s, 0.45, 1.45, 9.95, 2.30, fill=LTGREEN)
lines(tf, [
    ('本日持ち帰ること', {'size': 18, 'bold': True, 'space': 10}),
    ('① 失注の原因は「商談設定」か「商談内容」しかない。前日確認で設定を潰す', {'size': 17, 'space': 8}),
    ('② 先に相手の話を聞くことが、自分の話を聞いてもらうことに直結する', {'size': 17, 'space': 8}),
    ('③ 必要性訴求は「下げてから上げる」。ネガを自分事にしてから解決策を出す', {'size': 17, 'space': 8}),
    ('④ 太陽光と蓄電池は切り離さない。常にセットで話す', {'size': 17, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 0.45, 4.00, 9.95, 2.30, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('宿題（次回までに）', {'size': 18, 'bold': True, 'space': 10}),
    ('① 自社ストーリー5段構成を清書し、20秒で言えるようにしてくること', {'size': 17, 'space': 8}),
    ('② アプローチブック p.2〜p.17 の結論の一文を、資料を見ずに言えるようにしてくること', {'size': 17, 'space': 8}),
    ('③ 実商談を1件、ルール①〜⑨のチェックシートをつけて振り返ってくること', {'size': 17}),
])
_, tf = rect(s, 0.45, 6.55, 9.95, 0.42, fill=None, line=RED, lw=1.5)
lines(tf, [('次回・第3回は「必要性訴求②（ルール⑩〜⑫）」／アプローチブック p.18〜p.21',
            {'size': 16, 'bold': True, 'align': PP_ALIGN.CENTER})])
notes(s, '・忘却曲線の話（翌日には50％忘れる）を添えて、その日のうちの復習を促す')

# ---------------------------------------------------------------- 並べ替え
ids = list(sldIdLst)
colophon = [x for x in ids if x is colophon_id][0]
sldIdLst.remove(colophon)
sldIdLst.append(colophon)
page_no[0] += 1

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print('saved:', OUT, '/ slides:', len(prs.slides.__iter__.__self__._sldIdLst))

# ---------------------------------------------------------------- content-type fix
def ensure_jpg_content_type(path):
    """テンプレートが jpg を image/png と誤宣言しているため、python-pptx が
    保存時に Default を落とすことがある。表紙画像が壊れるので明示的に直す。"""
    import zipfile, shutil, re, tempfile
    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        ct = z.read('[Content_Types].xml').decode('utf-8')
        needs = any(n.lower().endswith('.jpg') for n in names)
        if not needs or 'Extension="jpg"' in ct:
            if 'Extension="jpg" ContentType="image/jpeg"' in ct or not needs:
                return False
        data = {n: z.read(n) for n in names}
    ct = re.sub(r'<Default Extension="jpg"[^/]*/>', '', ct)
    ct = ct.replace('<Default Extension="png"',
                    '<Default Extension="jpg" ContentType="image/jpeg"/><Default Extension="png"', 1)
    data['[Content_Types].xml'] = ct.encode('utf-8')
    tmp = path + '.tmp'
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as z:
        for n in names:
            z.writestr(n, data[n])
    shutil.move(tmp, path)
    return True

if ensure_jpg_content_type(OUT):
    print('fixed [Content_Types].xml (jpg -> image/jpeg)')
