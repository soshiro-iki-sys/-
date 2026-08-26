# -*- coding: utf-8 -*-
"""第3回「必要性訴求② ─ 災害対策・商品説明（FAB）・経済メリット」研修資料ビルダー

使い方は slides/build/build_session2.py と同じ。加えて本回は以下の画像が必要。

  <ASSETS>/ab/ab-18.jpg 〜 ab-21.jpg  アプローチブックのページ画像
  <ASSETS>/pyramid.png                契約までの4ステップ（本回では未使用）
  <ASSETS>/sim_trust.jpg              シミュレーションの信憑性グラフ
                                      （「20のポイント解説」slide30 の円グラフを切り出したもの）

    TRAINING_ASSETS=/path/to/assets python slides/build/build_session3.py

内容の出典：docs/知識_*.md ＋「20のポイント解説」
（契約とは／7つの阻害要因／3つの切り口／FAB／シミュレーションの2つの伝え方）
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
MSO_SHAPE_TYPE_AUTO = MSO_SHAPE_TYPE.AUTO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.opc.packuri import PackURI

REPO   = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
ASSETS = os.environ.get('TRAINING_ASSETS', '/tmp/training-assets')
AB     = os.path.join(ASSETS, 'ab')
TPL    = os.path.join(REPO, 'templates', '研修資料_フォーマット見本.pptx')
OUT    = os.path.join(REPO, 'slides', '2026営業研修_第3回_必要性訴求②＆ロープレ.pptx')

GOTHIC='HGPｺﾞｼｯｸE'; MINCHO='HGP明朝E'; UD='BIZ UDPゴシック'; YU='游ゴシック'
RED='FF0000'; BLACK='000000'; WHITE='FFFFFF'; GREEN='9BBB59'; YELLOW='FFFF00'
GRAY='595959'; BLUE='0070C0'; LTGREEN='EAF1DD'; LTGRAY='F2F2F2'; LTYEL='FFF2CC'
LTBLUE='DEEBF7'; ORANGE='F79646'

# ---------------------------------------------------------------- helpers
def _rpr_font(run, name):
    rPr = run._r.get_or_add_rPr()
    rPr.get_or_add_latin().set('typeface', name)
    latin = rPr.find(qn('a:latin'))
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = parse_xml('<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        latin.addnext(ea)
    ea.set('typeface', name)

def _highlight(run, color):
    rPr = run._r.get_or_add_rPr()
    hl = parse_xml(
        '<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:srgbClr val="%s"/></a:highlight>' % color)
    latin = rPr.find(qn('a:latin'))
    if latin is not None:
        latin.addprevious(hl)
    else:
        rPr.append(hl)

def run(p, text, size=16, bold=False, color=BLACK, font=GOTHIC, hl=None):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)
    if hl: _highlight(r, hl)
    _rpr_font(r, font)
    return r

def tb(slide, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP, margin=0.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return box, tf

def para(tf, first=False, space_after=4, align=PP_ALIGN.LEFT, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after); p.alignment = align
    if line: p.line_spacing = line
    return p

def lines(tf, items, size=16, color=BLACK, bold=False, space=5, font=GOTHIC, line=None, align=PP_ALIGN.LEFT):
    """items: str | (str,dict)"""
    for i, it in enumerate(items):
        opts = {}
        if isinstance(it, tuple): it, opts = it
        p = para(tf, first=(i == 0), space_after=opts.pop('space', space),
                 align=opts.pop('align', align), line=opts.pop('line', line))
        run(p, it, size=opts.pop('size', size), bold=opts.pop('bold', bold),
            color=opts.pop('color', color), font=opts.pop('font', font),
            hl=opts.pop('hl', None))

def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=1.0,
         anchor=MSO_ANCHOR.TOP):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    else: s.fill.background()
    if line:
        s.line.color.rgb = RGBColor.from_string(line); s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    return s, tf

def pic(slide, path, x, y, w, border=True):
    p = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    if border:
        p.line.color.rgb = RGBColor.from_string('808080'); p.line.width = Pt(1)
    return p

def table(slide, x, y, w, col_w, rows, font_size=11, header=True,
          header_fill='000000', header_color=WHITE, row_h=0.30, head_h=0.30):
    n_r, n_c = len(rows), len(col_w)
    shp = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w),
                                 Inches(head_h + row_h * (n_r - 1)))
    t = shp.table
    t.first_row = header; t.horz_banding = False
    for i, cw in enumerate(col_w): t.columns[i].width = Inches(cw)
    t.rows[0].height = Inches(head_h)
    for i in range(1, n_r): t.rows[i].height = Inches(row_h)
    for ri, rowdata in enumerate(rows):
        for ci, cell in enumerate(rowdata):
            txt, opts = (cell, {}) if isinstance(cell, str) else cell
            c = t.cell(ri, ci)
            c.margin_left = c.margin_right = Inches(0.05)
            c.margin_top = c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = opts.get('fill', header_fill if (header and ri == 0) else None)
            if fill: c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string(fill)
            else: c.fill.background()
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = opts.get('align',
                PP_ALIGN.CENTER if (header and ri == 0) else PP_ALIGN.LEFT)
            run(p, txt, size=opts.get('size', font_size),
                bold=opts.get('bold', header and ri == 0),
                color=opts.get('color', header_color if (header and ri == 0) else BLACK),
                font=opts.get('font', YU))
    return t

# ---------------------------------------------------------------- slide frame
prs = Presentation(TPL)
BLANK = None
for l in prs.slide_masters[0].slide_layouts:
    if l.name == '白紙': BLANK = l
assert BLANK is not None

sldIdLst = prs.slides._sldIdLst
orig_ids = list(sldIdLst)
cover_id, colophon_id = orig_ids[0], orig_ids[6]
for sid in orig_ids[1:6]:                      # drop template samples 2-6
    rId = sid.get(qn('r:id')); sldIdLst.remove(sid)
    prs.part.drop_rel(rId)
# python-pptx names new slides by len(sldIdLst)+1 -> move colophon out of that range
prs.part.related_part(colophon_id.get(qn('r:id'))).partname = PackURI('/ppt/slides/slide900.xml')

page_no = [0]

def new_slide(heading=None):
    s = prs.slides.add_slide(BLANK)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    if heading is not None:
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.157), Inches(0.757), Inches(0.140), Inches(0.512))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor.from_string(BLACK)
        bar.line.color.rgb = RGBColor.from_string('D9D9D9'); bar.line.width = Pt(0.75)
        bar.shadow.inherit = False
        box, tf = tb(s, 0.348, 0.736, 9.9, 0.572)
        lines(tf, [heading], size=28, font=MINCHO)
    page_no[0] += 1
    box, tf = tb(s, 8.31, 7.18, 2.30, 0.33)
    lines(tf, [str(page_no[0])], size=12, color=WHITE, font=GOTHIC, align=PP_ALIGN.RIGHT)
    return s

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def ab_img(n):
    return os.path.join(AB, 'ab-%02d.jpg' % n)

# ---------------------------------------------------------------- 表紙
cover = prs.slides[0]
for shp in cover.shapes:
    if shp.shape_type == MSO_SHAPE_TYPE_AUTO and shp.name == '正方形/長方形 4':
        sp = shp._element.spPr
        fill = sp.find(qn('a:solidFill'))
        for ch in list(fill): fill.remove(ch)
        fill.append(parse_xml(
            '<a:srgbClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'val="F79646"><a:alpha val="80000"/></a:srgbClr>'))
    if shp.has_text_frame and shp.name == 'テキスト ボックス 5':
        tf = shp.text_frame; tf.clear()
        lines(tf, [
            ('2026年度　営業研修', {'size': 30, 'bold': True, 'color': WHITE, 'font': UD, 'space': 8}),
            ('第3回　必要性訴求②', {'size': 48, 'bold': True, 'color': WHITE, 'font': UD, 'space': 8}),
            ('災害対策・商品説明・経済メリット', {'size': 30, 'bold': True, 'color': WHITE, 'font': UD, 'space': 10}),
            ('〜太陽光＋蓄電池セット販売〜', {'size': 22, 'bold': True, 'color': WHITE, 'font': UD}),
        ])
    if shp.has_text_frame and shp.name == 'テキスト ボックス 6':
        tf = shp.text_frame; tf.clear()
        lines(tf, [('2026年8月下旬　＠船井総研', {'size': 14, 'color': WHITE, 'font': UD})],
              align=PP_ALIGN.RIGHT)
page_no[0] = 1

# ---------------------------------------------------------------- rule helpers
def rule_head(s, no_title, point):
    box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
    lines(tf, [(no_title, {'size': 24, 'bold': True})])
    _, tf = rect(s, 0.45, 1.88, 9.95, 0.46, fill=None, line=RED, lw=1.25)
    p = para(tf, first=True, align=PP_ALIGN.LEFT)
    run(p, 'POINT　', size=15, bold=True, color=RED)
    run(p, point, size=16, bold=True)

def goal_box(s, x, y, w, h, quote):
    _, tf = rect(s, x, y, w, h, fill=LTBLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    lines(tf, [
        ('ゴール（お客様の反応）', {'size': 13, 'bold': True, 'color': BLUE, 'space': 6}),
    ])
    for q in quote:
        p = para(tf, space_after=3)
        run(p, '「%s」' % q, size=16, bold=True)

def points_box(s, x, y, w, h, items, title='説明のポイント'):
    _, tf = rect(s, x, y, w, h, fill=LTGRAY)
    lines(tf, [(title, {'size': 13, 'bold': True, 'color': GRAY, 'space': 6})])
    for it in items:
        opts = {}
        if isinstance(it, tuple): it, opts = it
        p = para(tf, space_after=opts.pop('space', 6))
        run(p, it, size=opts.pop('size', 12.5), bold=opts.pop('bold', False),
            color=opts.pop('color', BLACK))

def conclusion(s, text, y=6.36):
    _, tf = rect(s, 0.45, y, 9.95, 0.66, fill=LTGREEN)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, 'このページの結論　', size=12, bold=True, color=GRAY)
    run(p, text, size=15, bold=True, hl=YELLOW)

def talk_slide(s, no_title, talks, conc=None, note=None):
    box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
    lines(tf, [(no_title, {'size': 24, 'bold': True})])
    y = 1.92
    for label, body, h in talks:
        _, tf = rect(s, 0.45, y, 8.05, h, fill=LTGRAY)
        for i, ln in enumerate(body):
            p = para(tf, first=(i == 0), space_after=3)
            run(p, ln, size=14)
        _, tf2 = rect(s, 8.68, y, 1.72, 0.38, fill=ORANGE)
        lines(tf2, [(label, {'size': 12, 'bold': True, 'color': WHITE,
                             'align': PP_ALIGN.CENTER})])
        y += h + 0.16
    if conc: conclusion(s, conc)
    if note: notes(s, note)

# ================================================================ 2. 本日のゴール
s = new_slide('本日のゴール')
box, tf = tb(s, 0.45, 1.45, 9.9, 0.62)
lines(tf, [('お客様に合わせて商品を説明し、経済メリットを納得させる',
            {'size': 26, 'bold': True, 'color': RED})])
_, tf = rect(s, 0.45, 2.35, 9.95, 3.05, fill=LTGREEN)
lines(tf, [
    ('研修が終わったときの「できる状態」', {'size': 18, 'bold': True, 'space': 10}),
    ('①　災害の話を第三者話法だけで、3分以内に終えられる', {'size': 19, 'space': 9}),
    ('②　FABの3要素を使い分け、お客様ごとの「B（利益）」を言える', {'size': 19, 'space': 9}),
    ('③　シミュレーションを2通りの伝え方で説明できる', {'size': 19, 'space': 9}),
    ('④　金額を出さずに、経済メリットだけで納得をつくれる', {'size': 19}),
])
_, tf = rect(s, 0.45, 5.65, 9.95, 0.85, fill=None, line=RED, lw=1.5,
             anchor=MSO_ANCHOR.MIDDLE)
lines(tf, [
    ('本日の比重', {'size': 14, 'bold': True, 'color': GRAY, 'align': PP_ALIGN.CENTER, 'space': 3}),
    ('災害対策15分　／　商品説明43分　／　経済メリット30分',
     {'size': 20, 'bold': True, 'align': PP_ALIGN.CENTER}),
])
notes(s, '・災害は「短く終える技術」、FABと経済メリットは「難しいから時間をかける」と最初に説明する\n'
         '・金額（見積り）は第4回。今日は出さない')

# ================================================================ 3. アジェンダ
s = new_slide('本日のアジェンダ')
rows = [
    ['時間', '内容', '扱うもの'],
    ['0:00-0:12', '前回の振り返り／契約とは／3つの切り口', '第2回：ルール①〜⑨'],
    ['0:12-0:27', ('② 災害対策　ルール⑩⑪', {'bold': True}), 'アプローチブック p.18〜p.19'],
    ['0:27-0:37', '休憩', '−'],
    ['0:37-1:20', ('商品説明＝FAB式営業　ルール⑫', {'bold': True, 'color': RED}),
     'FAB／アプローチブック p.20・p.21'],
    ['1:20-1:50', ('③ 経済メリットの伝え方', {'bold': True, 'color': RED}), 'シミュレーションの2つの見せ方'],
    ['1:50-2:00', 'ロープレ／チェック／まとめ', '−'],
]
table(s, 0.45, 1.68, 9.95, [1.45, 4.30, 4.20], rows, font_size=13.5, row_h=0.58, head_h=0.40)
_, tf = rect(s, 0.45, 5.90, 9.95, 0.95, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('後半2つ（FAB・経済メリット）が本日の主役です', {'size': 18, 'bold': True, 'color': RED, 'space': 5}),
    ('この2つは「知っている」だけでは使えません。今日は手を動かす時間を長く取ります。', {'size': 15}),
])

# ================================================================ 4. 前回の振り返り
s = new_slide('前回の振り返り')
box, tf = tb(s, 0.45, 1.42, 9.9, 0.45)
lines(tf, [('第2回「聞く姿勢づくり＆必要性訴求①」で押さえたこと', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 2.00, 4.85, 2.30, fill=LTGRAY)
lines(tf, [
    ('聴く姿勢作り（①〜⑥）', {'size': 17, 'bold': True, 'space': 8}),
    ('・失注理由は「商談設定」か「商談内容」しかない', {'size': 14, 'space': 5}),
    ('・先に聞くから、自分の話を聞いてもらえる', {'size': 14, 'space': 5}),
    ('・包み込みの法則で「先生－生徒」の関係をつくる', {'size': 14, 'space': 5}),
    ('・太陽光と蓄電池は切り離さず、常にセットで話す', {'size': 14, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 5.55, 2.00, 4.85, 2.30, fill=LTGRAY)
lines(tf, [
    ('必要性訴求①（⑦〜⑨）', {'size': 17, 'bold': True, 'space': 8}),
    ('・電気代は2011年以降45％、約5.5万円上昇', {'size': 14, 'space': 5}),
    ('・再エネ賦課金は13年で18倍（4.18円/kWh）', {'size': 14, 'space': 5}),
    ('・30年で432万円〜700万円以上を支払う', {'size': 14, 'space': 5}),
    ('・基本設計は「下げてから上げる」', {'size': 14, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 0.45, 4.60, 9.95, 1.75, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('宿題の確認（5分）', {'size': 18, 'bold': True, 'space': 8}),
    ('① 自社ストーリー5段構成を、20秒で言ってみてください（2名指名）', {'size': 16, 'space': 6}),
    ('② アプローチブック p.9・p.14・p.17 の結論の一文を、資料を見ずに言えますか', {'size': 16, 'space': 6}),
    ('③ 実商談のチェックシートで、○が付かなかったルールはどれでしたか', {'size': 16}),
])

# ================================================================ 5. 契約とは
s = new_slide('契約とは')
box, tf = tb(s, 0.45, 1.42, 9.95, 0.45)
lines(tf, [('契約は「取る」ものではありません', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 1.98, 9.95, 1.10, fill=LTGREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
lines(tf, [
    ('契約とは', {'size': 15, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 4}),
    ('契約を阻害する要因（ネック）が、すべて解消された状態のこと',
     {'size': 22, 'bold': True, 'color': RED, 'align': PP_ALIGN.CENTER}),
])
box, tf = tb(s, 0.45, 3.20, 9.95, 0.35)
lines(tf, [('押しの強さで獲得するのではなく、お客様が抱える不安や疑問を取り除くことが契約に繋がります',
            {'size': 15, 'color': GRAY})])
rows = [
    ['', '7つの阻害要因', '解消の仕方'],
    ['①', '会社', '会社紹介をしっかり行い、信用を築く。来店商談も有効'],
    ['②', '自分', '自己紹介をしっかり行い、信用を築く。ビジネスマナーの徹底も重要'],
    [('③', {'fill': LTYEL}), ('プラン', {'fill': LTYEL, 'bold': True}),
     ('要望をしっかり聴き、優先順位に基づいて複数のプランを作成・提案する', {'fill': LTYEL})],
    [('④', {'fill': LTYEL}), ('金額', {'fill': LTYEL, 'bold': True}),
     ('予算と上限額を確認し、逸脱しないプランを提案する', {'fill': LTYEL})],
    ['⑤', 'キーマン', 'キーマンだけでなく、可能であれば家族全員に商談に参加してもらう'],
    ['⑥', '時期', '今後のスケジュールをお客様とともに確認・設定し、購買意欲を高める'],
    ['⑦', '競合', 'スピード営業で参入障壁を高めると同時にアポを取り続け、競合を待つ'],
]
table(s, 0.45, 3.56, 9.95, [0.50, 1.75, 7.70], rows, font_size=12, row_h=0.36, head_h=0.30)
_, tf = rect(s, 0.45, 6.46, 9.95, 0.45, fill=None, line=RED, lw=1.5)
lines(tf, [('本日つぶすのは③プランと④金額。FAB＝プランのネック、経済メリット＝金額のネックです',
            {'size': 15, 'bold': True, 'align': PP_ALIGN.CENTER})])

# ================================================================ 6. 3つの切り口
s = new_slide('必要性訴求')
box, tf = tb(s, 0.45, 1.40, 9.95, 0.45)
lines(tf, [('必要性を訴求する3つの切り口', {'size': 26, 'bold': True})])
cuts = [
    (1.98, '①', '電気代が高騰している ⇒ 払わなくていい',
     'プロスペクト理論（損失回避）\n人は、利益よりも損失の方が大きく感じられる\n例：5%off! ＜ 消費税5%増税',
     LTGRAY, '第2回で実施済み'),
    (3.48, '②', '災害対策',
     '災害時も安心。ただし長く話すほど失敗する。\n主な災害時の停電発生件数と日数を「事実」で見せる',
     LTGREEN, '本日・前半15分'),
    (4.98, '③', '経済メリット',
     '一番重要な要素ではあるが、メリットがあまり出ないお宅にも訴求するためには、\n'
     '経済メリット訴求のみでは安定感がなくなってしまう',
     LTYEL, '本日・後半30分'),
]
for y, n, t, d, fill, badge in cuts:
    _, tf = rect(s, 0.45, y, 0.62, 1.32, fill=fill)
    lines(tf, [(n, {'size': 26, 'bold': True, 'align': PP_ALIGN.CENTER})])
    _, tf = rect(s, 1.20, y, 7.15, 1.32, fill=fill)
    lines(tf, [(t, {'size': 19, 'bold': True, 'space': 6})])
    for i, ln in enumerate(d.split('\n')):
        p = para(tf, space_after=2); run(p, ln, size=13)
    _, tf = rect(s, 8.50, y + 0.35, 1.90, 0.62,
                 fill=(GRAY if '済' in badge else RED))
    lines(tf, [(badge, {'size': 12, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER})])
_, tf = rect(s, 0.45, 6.50, 9.95, 0.50, fill=None, line=RED, lw=1.5)
lines(tf, [('②③に加えて、本日は「商品説明（FAB式営業）」を挟みます。ここが一番難しいパートです',
            {'size': 16, 'bold': True, 'align': PP_ALIGN.CENTER})])

# ================================================================ 7. 本日の流れ
s = new_slide('本日の流れ')
box, tf = tb(s, 0.45, 1.40, 9.95, 0.45)
lines(tf, [('3つのパートを、この順番でつなぎます', {'size': 24, 'bold': True})])
parts = [
    (0.45, 'PART 1', '② 災害対策', 'ルール⑩⑪', 'p.18・p.19', '15分', LTGREEN,
     ['「もしも」の不安をつくる', '長く話さない。3分で抜ける']),
    (3.85, 'PART 2', '商品説明（FAB）', 'ルール⑫', 'p.20・p.21', '43分', LTYEL,
     ['お客様ごとの「B」を言う', '手を動かす時間を長く取る']),
    (7.25, 'PART 3', '③ 経済メリット', 'シミュレーション', '（資料の外）', '30分', LTBLUE,
     ['2通りの伝え方を使い分ける', '金額（見積り）はまだ出さない']),
]
for x, tag, title, rule, page, mins, fill, notes_ in parts:
    _, tf = rect(s, x, 2.00, 3.15, 3.55, fill=fill)
    lines(tf, [
        (tag, {'size': 14, 'bold': True, 'color': GRAY, 'align': PP_ALIGN.CENTER, 'space': 4}),
        (title, {'size': 21, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 10}),
        (mins, {'size': 26, 'bold': True, 'color': RED, 'align': PP_ALIGN.CENTER, 'space': 12}),
        ('ルール：%s' % rule, {'size': 13, 'space': 4}),
        ('資料：%s' % page, {'size': 13, 'space': 10}),
    ])
    for ln in notes_:
        p = para(tf, space_after=4); run(p, '・' + ln, size=13, bold=True)
for x in (3.55, 6.95):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(3.55), Inches(0.24), Inches(0.32))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor.from_string('808080')
    a.line.fill.background(); a.shadow.inherit = False
_, tf = rect(s, 0.45, 5.80, 9.95, 1.05, fill=LTGRAY)
lines(tf, [
    ('なぜこの順番か', {'size': 16, 'bold': True, 'space': 6}),
    ('不安（PART1）→ 解決策の中身（PART2）→ それがいくら得か（PART3）。'
     '中身を説明する前に「お得です」と言っても、何がお得なのか伝わりません。', {'size': 15}),
])

# ================================================================ 8. PART1 ルール⑩-1
s = new_slide('PART 1　災害対策')
rule_head(s, 'ルール⑩：非常時に備える意義の理解', '「いつ身の回りで起きてもおかしくない」に腹落ちさせる')
pic(s, ab_img(18), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.18　予測が困難な地震情報', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.15, ['いつ災害が起こるか分からないですね'])
points_box(s, 5.35, 3.78, 5.05, 2.50, [
    ('地震リスクが最低とされた札幌で、大規模停電が起きた', {'bold': True, 'size': 13}),
    '地震調査委員会が2018年6月に発表した「今後30年以内に震度6弱以上」の確率で札幌が最下位',
    'その3か月後に北海道胆振東部地震が発生。気象庁も「確度の高い地震の予測は難しい」と明言',
    ('「うちの地域は地震来ないから」を潰すためのページ', {'bold': True, 'color': RED, 'size': 13}),
])
conclusion(s, '専門家ですら、いつ「もしものこと」が起きるのか分からない')

# ================================================================ 9. PART1 ルール⑩-2
s = new_slide('PART 1　災害対策')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ルール⑩：3段構成で逃げ道を無くす', {'size': 24, 'bold': True})])
talks = [
    ('①権威', '「実は、日本で最も権威のある『地震調査委員会』という機関が、定期的に地震のデータを出しているんです」'),
    ('②事実', '「そのデータの2018年6月版では、札幌市の地震発生確率が全国で最も低かったんですね」'),
    ('③事実', '「ところが、その3か月後に胆振東部地震が起きてしまいました。気象庁も『確度の高い地震の予測は難しい』と明言しています」'),
    ('着地', '「どこでどのくらい大きな地震が来るかなんて、分からないですよね？」'),
]
y = 1.92
for label, body in talks:
    _, tf = rect(s, 0.45, y, 8.05, 0.68, fill=LTGRAY, anchor=MSO_ANCHOR.MIDDLE)
    lines(tf, [(body, {'size': 14})])
    _, tf2 = rect(s, 8.68, y, 1.72, 0.68, fill=ORANGE, anchor=MSO_ANCHOR.MIDDLE)
    lines(tf2, [(label, {'size': 13, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER})])
    y += 0.80
rows = [
    ['よくあるネガ', '切り返し方'],
    ['この地域では災害なんて起きていない', 'そうですよね。実は札幌の方も同じことを（→p.18へ）'],
    ['備蓄はしてあるから大丈夫', '食料は備えられますが、電気だけは備蓄できないんですよね'],
    ['停電なんてすぐ復旧するでしょう', '胆振東部地震では、復旧まで2週間かかった地域もありました'],
]
table(s, 0.45, 5.20, 9.95, [3.65, 6.30], rows, font_size=13, row_h=0.44, head_h=0.30)
conclusion(s, '否定せず、事実と事例で包み込む。議論にしない', y=6.75)

# ================================================================ 10. PART1 ルール⑩-3
s = new_slide('PART 1　災害対策')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ルール⑩：停電被害の実績で「規模」を見せる', {'size': 24, 'bold': True})])
pic(s, ab_img(19), 0.45, 1.95, 4.55)
box, tf = tb(s, 0.45, 5.30, 4.55, 0.3)
lines(tf, [('アプローチブック p.19　自然災害と停電被害', {'size': 12, 'color': GRAY})])
rows = [
    ['災害', '停電戸数', '備考'],
    ['2016年4月　熊本地震', '約50万戸', '−'],
    ['2018年9月　台風21号（関東）', '240万戸', '−'],
    ['2018年9月　台風24号', '約180万戸', '−'],
    ['2018年9月　北海道胆振東部地震', '約300万戸', ('日本初のブラックアウト', {'bold': True, 'color': RED})],
    ['2019年9月　台風15号（関東）', '93万戸', '5日連続で停電'],
    ['今後10年以内　南海トラフ地震？', '？', '−'],
]
table(s, 5.35, 1.95, 5.05, [2.55, 1.05, 1.45], rows, font_size=10, row_h=0.48, head_h=0.30)
_, tf = rect(s, 5.35, 5.30, 5.05, 1.00, fill=LTGRAY)
lines(tf, [
    ('胆振東部地震は平成最大規模。', {'size': 14, 'space': 4}),
    ('復旧まで2週間かかった地域もありました。', {'size': 14}),
])
conclusion(s, '災害大国だからこそ、もしもの時には備える必要があるのです')

# ================================================================ 11. PART1 ルール⑪
s = new_slide('PART 1　災害対策')
rule_head(s, 'ルール⑪：非常時対策の効果の確認', '災害トークは「第三者話法」でしか伝えない')
_, tf = rect(s, 0.45, 2.50, 4.95, 2.05, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('このルールに対応するページは', {'size': 17, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 3}),
    ('アプローチブックにありません', {'size': 17, 'bold': True, 'color': RED,
                        'align': PP_ALIGN.CENTER, 'space': 10}),
    ('「事例」でないと意味が無いからです。', {'size': 14, 'space': 5}),
    ('印刷された一般論では「うちの場合は？」に答えられません。', {'size': 14}),
])
_, tf = rect(s, 5.60, 2.50, 4.95, 2.05, fill=LTGREEN)
lines(tf, [
    ('第三者話法の型', {'size': 17, 'bold': True, 'space': 8}),
    ('✕「停電しても安心ですよ」', {'size': 15, 'space': 4}),
    ('　← 自分の意見。刺さらない', {'size': 13, 'color': GRAY, 'space': 8}),
    ('○「実際に○○市のA様が、台風の停電で3日間、冷蔵庫と照明を動かせたそうです」',
     {'size': 15, 'bold': True}),
])
rows = [
    ['メーカー', '特徴', 'URL'],
    ['Panasonic', '動画＋テキストの記事があり使いやすい', 'https://sumai.panasonic.jp/chikuden/'],
    ['SmartStar', '災害時の活用に特化した動画', 'https://www.smartstar.jp/voice/'],
    ['ニチコン', '販売店にフィットしたコンテンツ', 'https://www.nichicon.co.jp/products/ess/about/voice.html'],
]
table(s, 0.45, 4.75, 9.95, [1.45, 3.20, 5.30], rows, font_size=11.5, row_h=0.52, head_h=0.32)
conclusion(s, '一般論ではなく「事例」。「私が〜」ではなく「実際に〜あります」', y=6.55)

# ================================================================ 12. PART1 落とし穴
s = new_slide('PART 1　災害対策')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('災害トークの落とし穴：長く話すほど失注する', {'size': 24, 'bold': True, 'color': RED})])
_, tf = rect(s, 0.45, 1.98, 4.95, 3.55, fill='FCE4E4', line=RED)
lines(tf, [
    ('✕　やりすぎ例', {'size': 20, 'bold': True, 'color': RED, 'space': 10}),
    ('「南海トラフが来たら、この地域は震度7です。', {'size': 14, 'space': 3}),
    ('　ライフラインは1か月止まると言われています。', {'size': 14, 'space': 3}),
    ('　お子さんもいらっしゃいますよね。', {'size': 14, 'space': 3}),
    ('　真っ暗な中で、水も出ない状態で……」', {'size': 14, 'space': 12}),
    ('何が起きるか', {'size': 15, 'bold': True, 'space': 5}),
    ('・商談の雰囲気が暗くなる', {'size': 14, 'space': 4}),
    ('・「不安を煽られた」と警戒される', {'size': 14, 'space': 4}),
    ('・お客様が黙り、ラリーが止まる', {'size': 14, 'space': 4}),
    ('・そのあとの商品説明が入らなくなる', {'size': 14}),
])
_, tf = rect(s, 5.60, 1.98, 4.95, 3.55, fill=LTGREEN, line='70A040')
lines(tf, [
    ('○　適切な例', {'size': 20, 'bold': True, 'space': 10}),
    ('「札幌が地震の確率が最も低いと出ていたのに、', {'size': 14, 'space': 3}),
    ('　3か月後に大きな地震が来てしまったんですね。', {'size': 14, 'space': 3}),
    ('　気象庁も予測は難しいと言っています。', {'size': 14, 'space': 3}),
    ('　どこで起きるかは分からないですよね？」', {'size': 14, 'space': 12}),
    ('何が違うか', {'size': 15, 'bold': True, 'space': 5}),
    ('・主語が「事実」と「第三者」', {'size': 14, 'space': 4}),
    ('・お客様の家族を想像させていない', {'size': 14, 'space': 4}),
    ('・最後が質問で終わり、ラリーが続く', {'size': 14, 'space': 4}),
    ('・3分で次のページへ抜けられる', {'size': 14}),
])
conclusion(s, '災害は「事実 → 事例 → サラッと次へ」。3分で抜ける', y=5.70)
notes(s, '・講師はまず「やりすぎ例」を実演して見せる。受講者に違和感を体験させてから正解を示す')

# ================================================================ 13. 休憩
s = new_slide(None)
_, tf = rect(s, 2.40, 2.85, 6.05, 1.80, fill=LTGREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
lines(tf, [
    ('休憩　10分', {'size': 44, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 10}),
    ('後半は商品説明（FAB）と経済メリットです', {'size': 18, 'align': PP_ALIGN.CENTER}),
])

# ================================================================ 14. PART2 なぜ商品説明が難しいか
s = new_slide('PART 2　商品説明')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('なぜ「商品説明」が必要なのか', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.90, 1.95, 4.15, 1.55, fill=LTBLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
lines(tf, [
    ('購入 ＝ 価値 ÷ 価格', {'size': 22, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 8}),
    ('この値が1を超えたとき、', {'size': 14, 'align': PP_ALIGN.CENTER, 'space': 2}),
    ('人は購入に至る', {'size': 14, 'align': PP_ALIGN.CENTER}),
])
_, tf = rect(s, 5.35, 1.95, 5.05, 1.55, fill=LTGRAY)
lines(tf, [
    ('100均で考えると', {'size': 15, 'bold': True, 'space': 6}),
    ('100円の商品を買うのは、100円以上の価値を感じているから。', {'size': 14, 'space': 5}),
    ('つまり X > 100円 ÷ 100円 の状態です。', {'size': 14}),
])
_, tf = rect(s, 0.45, 3.72, 4.95, 2.20, fill=LTGREEN)
lines(tf, [
    ('やり方①　価値を上げる', {'size': 18, 'bold': True, 'space': 8}),
    ('①電気代のネガ（第2回）', {'size': 15, 'space': 4}),
    ('②災害対策（PART1）', {'size': 15, 'space': 4}),
    ('＋ 商品そのものの価値（FAB）', {'size': 15, 'bold': True, 'color': RED, 'space': 10}),
    ('PART2はここ', {'size': 16, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 5.60, 3.72, 4.95, 2.20, fill=LTYEL)
lines(tf, [
    ('やり方②　価格を下げる', {'size': 18, 'bold': True, 'space': 8}),
    ('価格 −（経済メリット③）を下げる', {'size': 15, 'space': 4}),
    ('＝ 実質いくらになるか', {'size': 15, 'space': 10}),
    ('PART3はここ', {'size': 16, 'bold': True, 'color': RED, 'space': 4}),
    ('※ 見積り金額そのものは第4回', {'size': 13, 'color': GRAY}),
])
conclusion(s, '価値を上げるか、価格を下げるか。この2つしかない', y=6.10)

# ================================================================ 15. PART2 FABとは
s = new_slide('PART 2　商品説明')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('FAB式営業：価値の伝え方は3階建て', {'size': 24, 'bold': True})])
rows = [
    ['', '要素', '意味', '固定／変動', 'できる営業の割合'],
    ['F', 'Feature（特徴）', 'その製品の特長・長所', '固定', 'ほとんどの営業マンができる'],
    ['A', 'Advantage（メリット）', 'その特徴が「全ての人に」もたらす利点', '固定', '約8割ができる'],
    ['B', ('Benefit（ベネフィット）', {'bold': True}),
     ('「目の前のお客様が」受ける利益', {'bold': True}),
     ('変動', {'bold': True, 'color': RED}), ('トップ2割しかできない', {'bold': True, 'color': RED})],
]
table(s, 0.45, 1.95, 9.95, [0.45, 2.50, 3.35, 1.15, 2.50], rows, font_size=12, row_h=0.72, head_h=0.34)
_, tf = rect(s, 0.45, 4.55, 9.95, 1.70, fill=LTGRAY)
lines(tf, [
    ('例）コーヒーマシン', {'size': 16, 'bold': True, 'space': 6}),
    ('F：専用カプセルを入れてボタン一つで本格的なコーヒーが淹れられる', {'size': 15, 'space': 4}),
    ('A：誰でも安定した味を再現でき、手間も汚れも少ない', {'size': 15, 'space': 4}),
    ('B：忙しい朝でもボタン1つで一息つける／毎月カプセルが届くので買い忘れがない',
     {'size': 15, 'bold': True, 'color': RED}),
])
conclusion(s, 'お客様に対してのB（利益）を、いかに具体的に話せるかが最重要', y=6.40)

# ================================================================ 16. PART2 太陽光のFAB
s = new_slide('PART 2　商品説明')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('太陽光＋蓄電池のFAB（模範例）', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 1.95, 9.95, 0.85, fill=LTGRAY)
lines(tf, [
    ('F　特徴（固定）', {'size': 15, 'bold': True, 'color': GRAY, 'space': 4}),
    ('太陽光で発電した電気を蓄電池に貯めて、夜も自宅で使えます', {'size': 16}),
])
_, tf = rect(s, 0.45, 2.95, 9.95, 0.85, fill=LTGRAY)
lines(tf, [
    ('A　メリット（固定・誰にでも同じ）', {'size': 15, 'bold': True, 'color': GRAY, 'space': 4}),
    ('電力会社から買う電気が減り、停電のときも電気が使えます', {'size': 16}),
])
_, tf = rect(s, 0.45, 3.95, 9.95, 2.35, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('B　ベネフィット（変動・お客様ごとに変える）', {'size': 15, 'bold': True, 'color': RED, 'space': 6}),
    ('「お子様が受験生でしたね。夜遅くまでエアコンをつけても電気代を気にせずに済みます」',
     {'size': 15, 'space': 6}),
    ('「奥様が日中パートに出られない日も、太陽光の電気で洗濯乾燥機を回せます」', {'size': 15, 'space': 6}),
    ('「ご主人は在宅勤務が週2日でしたね。在宅の日ほどお得になります」', {'size': 15, 'space': 6}),
    ('「停電のとき、お母様が避難所に行かずに済みます」', {'size': 15}),
])
conclusion(s, 'Bの材料は、ルール③のヒアリングで聞いた話の中にある', y=6.45)
notes(s, '・Bは想像で作らない。ルール③で聞いた「導入しなかった理由」「今の困りごと」から作る')

# ================================================================ 17. PART2 Bを引き出す
s = new_slide('PART 2　商品説明')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('Bの材料は、ヒアリングで手に入れる', {'size': 24, 'bold': True})])
rows = [
    ['ルール③で聞いたこと', 'そこから作るB（ベネフィット）'],
    ['「価格が高かったからやめた」', '「月々のご負担は電気代の範囲に収まります。今と変わらない支出で、資産が残ります」'],
    ['「シミュレーションがピンとこなかった」', '「今日は数字だけでなく、〇〇様のお家で1日どう電気が動くかでご説明します」'],
    ['「会社が信用しきれなかった」', '「当社は3年ごとに訪問して報告書をお出ししています。設置して終わりにしません」'],
    ['「子供が受験で電気代が心配」', '「夜遅くまでエアコンをつけても、電気代を気にせずに済みます」'],
    ['「停電の経験がなくて実感がない」', '「〇〇市のA様も同じでしたが、台風の停電で3日間、冷蔵庫が止まらなかったそうです」'],
]
table(s, 0.45, 1.95, 9.95, [3.55, 6.40], rows, font_size=13, row_h=0.72, head_h=0.34)
_, tf = rect(s, 0.45, 6.00, 9.95, 0.85, fill=LTGREEN)
lines(tf, [
    ('聞いていないことは、Bにできません。', {'size': 17, 'bold': True, 'space': 4}),
    ('だからルール③のオープンクエスチョン3回が効いてきます。', {'size': 16, 'bold': True, 'color': RED}),
])

# ================================================================ 18. PART2 ルール⑫
s = new_slide('PART 2　商品説明')
rule_head(s, 'ルール⑫：太陽光・蓄電池の使い方の理解', 'シミュレーションに移る前に基本概念をおさらいする')
pic(s, ab_img(21), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.21　太陽光発電・蓄電池の使い方', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.15, ['太陽光・蓄電池の活用方法が分かる'])
points_box(s, 5.35, 3.78, 5.05, 2.50, [
    ('売電価格が電気代より安い今は「自家消費」が主流', {'bold': True, 'size': 13}),
    'p.20で「買うより自給自足」を示し、p.21で具体的な使い方を見せる',
    '売電単価15円 ＜ 買電単価30円。貯めて使う方が得',
    ('この説明がFのベース。曖昧だとBも作れない', {'bold': True, 'color': RED, 'size': 13}),
])
conclusion(s, '具体的なシミュレーションを見ていきましょう！')

# ================================================================ 19. PART2 1日の流れ
s = new_slide('PART 2　商品説明')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('1日の電気の流れ（これは図なしで言えるように）', {'size': 24, 'bold': True})])
blocks = [
    (0.45, '朝', '発電量が少ない時間帯', ['電力会社から電気を購入する', '（買電消費）'], LTBLUE),
    (3.75, '昼', '太陽光発電が活発な時間帯', ['① 太陽光の電気をまず家庭で使う（自家消費）',
                                     '② 余った電気は蓄電池に充電する', '③ さらに余った電気は売電する'], LTYEL),
    (7.05, '夕方〜夜間', '発電できない時間帯', ['蓄電池に貯めた電気を使う'], LTGREEN),
]
for x, t, sub, items, fill in blocks:
    _, tf = rect(s, x, 1.98, 3.10, 2.45, fill=fill)
    lines(tf, [
        (t, {'size': 22, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 4}),
        (sub, {'size': 12, 'color': GRAY, 'align': PP_ALIGN.CENTER, 'space': 10}),
    ])
    for it in items:
        p = para(tf, space_after=6); run(p, it, size=13)
for x in (3.58, 6.88):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(2.95), Inches(0.30), Inches(0.35))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor.from_string('808080')
    a.line.fill.background(); a.shadow.inherit = False
_, tf = rect(s, 0.45, 4.60, 9.95, 1.65, fill=LTGRAY)
lines(tf, [
    ('このサイクルは、太陽光と蓄電池が「両方」あって初めて成立します',
     {'size': 18, 'bold': True, 'color': RED, 'space': 8}),
    ('・太陽光だけ　→　昼に余った電気は売るしかない。売電15円＜買電30円なので損', {'size': 15, 'space': 5}),
    ('・蓄電池だけ　→　貯める電気を電力会社から買うことになる', {'size': 15, 'space': 5}),
    ('・セット　　　→　自分で創った電気を、自分で使い切れる', {'size': 15, 'bold': True}),
])
conclusion(s, '創る（太陽光）→ 貯める（蓄電池）→ 夜に使う', y=6.40)

# ================================================================ 20. PART2 ワーク①
s = new_slide('PART 2　商品説明')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ワーク①：お手元の顧客設定集からFABを考える（15分）', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.35)
lines(tf, [('F・Aは全員同じで構いません。Bだけは、その顧客設定に合わせて3本書いてください',
            {'size': 15, 'color': GRAY})])
rows = [['', '要素', '記入してください']]
for n, t in [('F', '特徴（固定）'), ('A', 'メリット（固定）'),
             ('B①', 'ベネフィット（変動）'), ('B②', 'ベネフィット（変動）'), ('B③', 'ベネフィット（変動）')]:
    rows.append([n, t, ''])
table(s, 0.45, 2.32, 9.95, [0.60, 2.10, 7.25], rows, font_size=13, row_h=0.72, head_h=0.32)
_, tf = rect(s, 0.45, 6.20, 9.95, 0.68, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('B①②③は、必ず「〇〇様は〜でしたよね」から書き始めること。', {'size': 16, 'bold': True, 'space': 4}),
    ('主語がお客様になっていなければ、それはAです。', {'size': 15, 'bold': True, 'color': RED}),
])

# ================================================================ 21. PART2 ワーク②
s = new_slide('PART 2　商品説明')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ワーク②：ペアで実演する（10分）', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 1.95, 4.95, 2.10, fill=LTGREEN)
lines(tf, [
    ('やること', {'size': 17, 'bold': True, 'space': 8}),
    ('① 2人1組。営業役／お客様役', {'size': 15, 'space': 6}),
    ('② 書いたFABを、そのまま口で言う', {'size': 15, 'space': 6}),
    ('③ 2分で交代', {'size': 15}),
])
_, tf = rect(s, 5.60, 1.95, 4.95, 2.10, fill=LTYEL)
lines(tf, [
    ('聞き手のチェック項目', {'size': 17, 'bold': True, 'space': 8}),
    ('・Bの主語が「お客様」になっていたか', {'size': 15, 'space': 6}),
    ('・Bが3本とも違う切り口だったか', {'size': 15, 'space': 6}),
    ('・スペック（kWh・サイクル数）で終わっていないか', {'size': 15}),
])
box, tf = tb(s, 0.45, 4.30, 9.95, 0.40)
lines(tf, [('FABでやりがちな失敗', {'size': 20, 'bold': True, 'color': RED})])
rows = [
    ['やりがちな失敗', 'なぜダメか', '直し方'],
    ['Fだけ並べる（kWh・変換効率・サイクル数）', 'スペックはお客様の生活と結びつかない', 'Aに翻訳してからBへ'],
    ['Aで止まる（「電気代が安くなります」）', '全員に同じことを言っている', '「〇〇様は」で始める'],
    ['Bを想像で作る', 'ヒアリングしていないので外れる', 'ルール③で聞いた話から作る'],
]
table(s, 0.45, 4.78, 9.95, [3.75, 3.20, 3.00], rows, font_size=12.5, row_h=0.58, head_h=0.32)

# ================================================================ 22. PART3 経済メリットとは
s = new_slide('PART 3　経済メリット')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('③ 経済メリット：一番重要だが、それだけでは危ない', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 1.98, 9.95, 1.25, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('経済メリットは一番重要な要素。ただしメリットがあまり出ないお宅にも訴求するには、',
     {'size': 17, 'bold': True, 'space': 4}),
    ('経済メリット訴求のみでは安定感がなくなってしまう', {'size': 17, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 0.45, 3.45, 4.95, 1.65, fill='FCE4E4', line=RED)
lines(tf, [
    ('✕　経済メリットだけで押す', {'size': 17, 'bold': True, 'color': RED, 'space': 8}),
    ('狭小住宅・北向き屋根・電気使用量が少ないお宅では、数字が出ません。', {'size': 14, 'space': 5}),
    ('数字が出ない＝売れない、になってしまいます。', {'size': 14}),
])
_, tf = rect(s, 5.60, 3.45, 4.95, 1.65, fill=LTGREEN, line='70A040')
lines(tf, [
    ('○　①②③を積んだ上で③を出す', {'size': 17, 'bold': True, 'space': 8}),
    ('①電気代のネガ ②災害対策 ＋ 商品の価値（FAB）を積んでから、', {'size': 14, 'space': 5}),
    ('最後に③経済メリットで背中を押す。', {'size': 14, 'bold': True}),
])
_, tf = rect(s, 0.45, 5.35, 9.95, 1.05, fill=LTGRAY)
lines(tf, [
    ('金額ばかり褒めてはいけない理由も同じです', {'size': 16, 'bold': True, 'space': 5}),
    ('「元が取れる／取れない」で判断させると、狭小住宅への提案が困難になります。', {'size': 15, 'space': 3}),
    ('金額の褒め → 体験の褒めへ変換しましょう。', {'size': 15}),
])
conclusion(s, '③は「決め手」であって「入り口」ではない', y=6.50)

# ================================================================ 23. PART3 シミュレーションの読み方
s = new_slide('PART 3　経済メリット')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('シミュレーションの読み方（例）', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.35)
lines(tf, [('FIT期間中と、FIT期間が終わったあとで、毎月の実質負担額が変わります',
            {'size': 15, 'color': GRAY})])
for x, ttl, a, b, c, d in (
    (0.45, 'FIT期間', '24,002 円/月', '14,429 円/月', '9,573', '約 319 円'),
    (5.60, 'FIT期間終了後', '24,002 円/月', '12,686 円/月', '11,316', '約 377 円'),
):
    _, tf = rect(s, x, 2.32, 4.95, 0.42, fill='DAEEF3')
    lines(tf, [(ttl, {'size': 16, 'bold': True, 'align': PP_ALIGN.CENTER})])
    rows = [
        ['毎月の分割支払額', a],
        ['− 実質光熱費削減額', b],
    ]
    table(s, x, 2.80, 4.95, [3.05, 1.90], rows, font_size=13, row_h=0.44, head_h=0.44, header=False)
    _, tf = rect(s, x, 3.72, 4.95, 0.95, fill=LTGREEN)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, '＝ 毎月の実質負担額　', size=14, bold=True)
    run(p, c, size=26, bold=True, color=RED)
    run(p, ' 円/月', size=14, bold=True)
    p = para(tf, align=PP_ALIGN.CENTER)
    run(p, '（1日あたり %s）' % d, size=13)
_, tf = rect(s, 0.45, 4.90, 9.95, 1.35, fill=LTYEL)
lines(tf, [
    ('この表のどこを読むか', {'size': 16, 'bold': True, 'space': 6}),
    ('お客様が見るのは「毎月の実質負担額」と「1日あたり」だけです。', {'size': 15, 'space': 4}),
    ('分割支払額（24,002円）を先に見せると、そこで思考が止まります。', {'size': 15, 'bold': True, 'color': RED}),
])
conclusion(s, '1日あたりの金額まで落として初めて、判断できる数字になる', y=6.40)

# ================================================================ 24. PART3 伝え方①
s = new_slide('PART 3　経済メリット')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('伝え方①：ローン返済額を含めたトータル金額でお伝えする', {'size': 22, 'bold': True})])
_, tf = rect(s, 0.45, 1.92, 9.95, 0.72, fill=LTGRAY)
lines(tf, [
    ('＜ポイント＞　分割払いの場合、返済中は合計の支出が増える。', {'size': 14, 'bold': True, 'space': 3}),
    ('　　　　　　　完済後は返済額が無くなるため、支出が減少する', {'size': 14, 'bold': True}),
])
rows = [
    ['（例）', 'ローン支払い中', '完済後'],
    ['もともとの電気代', '15,000', '15,000'],
    ['＋ ローンの分割額', '22,000', '0'],
    ['− 電気代削減額', '−14,000', '−14,000'],
    [('支出合計', {'bold': True}), ('23,000', {'bold': True, 'fill': LTYEL}),
     ('1,000', {'bold': True, 'fill': LTYEL, 'color': RED})],
]
table(s, 0.45, 2.80, 9.95, [3.55, 3.20, 3.20], rows, font_size=15, row_h=0.52, head_h=0.36)
_, tf = rect(s, 0.45, 5.30, 4.95, 1.05, fill=LTGREEN)
lines(tf, [
    ('向いているお客様', {'size': 16, 'bold': True, 'space': 4}),
    ('・長期で考えられる方', {'size': 14, 'space': 2}),
    ('・完済後の姿に価値を感じる方', {'size': 14}),
])
_, tf = rect(s, 5.60, 5.30, 4.95, 1.05, fill='FCE4E4', line=RED)
lines(tf, [
    ('注意', {'size': 16, 'bold': True, 'color': RED, 'space': 4}),
    ('返済中は支出が増えます。', {'size': 14, 'space': 2}),
    ('ここを隠すと後で必ず不信になります。', {'size': 14}),
])
conclusion(s, '完済後に支出が「1,000円」まで下がる絵を見せる', y=6.50)

# ================================================================ 25. PART3 伝え方②
s = new_slide('PART 3　経済メリット')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('伝え方②：電気代の減収分のみをお伝えする', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.35)
lines(tf, [('ローンの話をいったん外し、「電気代がいくら減るか」だけを見せる方法です',
            {'size': 15, 'color': GRAY})])
rows = [
    ['（例）', '金額', '期間', '削減額'],
    ['もともとの電気代', '15,000', '1年間での削減額', '¥168,000'],
    ['− 電気代削減額', '14,000', '10年間での削減額', '¥1,680,000'],
    [('支出合計', {'bold': True}), ('1,000', {'bold': True, 'fill': LTYEL}),
     ('30年間での削減額', {'bold': True}), ('¥5,040,000', {'bold': True, 'fill': LTYEL, 'color': RED})],
]
table(s, 0.45, 2.32, 9.95, [2.75, 1.80, 2.60, 2.80], rows, font_size=15, row_h=0.58, head_h=0.36)
_, tf = rect(s, 0.45, 4.65, 4.95, 1.60, fill=LTGREEN)
lines(tf, [
    ('向いているお客様', {'size': 16, 'bold': True, 'space': 6}),
    ('・現金一括で検討される方', {'size': 14, 'space': 4}),
    ('・「月々の支出が増える」ことに強い抵抗がある方', {'size': 14, 'space': 4}),
    ('・数字が多いと混乱される方', {'size': 14}),
])
_, tf = rect(s, 5.60, 4.65, 4.95, 1.60, fill=LTYEL)
lines(tf, [
    ('使い分けの原則', {'size': 16, 'bold': True, 'space': 6}),
    ('①で反応が鈍い（＝返済中の支出増で止まった）と', {'size': 14, 'space': 2}),
    ('感じたら、②に切り替える。', {'size': 14, 'space': 6}),
    ('どちらも同じ事実を、別の枠組みで', {'size': 14, 'bold': True, 'space': 2}),
    ('見せているだけです。', {'size': 14, 'bold': True}),
])
conclusion(s, '30年で500万円以上の削減。これが「価格を下げる」の中身', y=6.42)

# ================================================================ 26. PART3 疑われている
s = new_slide('PART 3　経済メリット')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('前提：シミュレーションは疑われています', {'size': 24, 'bold': True, 'color': RED})])
pic(s, os.path.join(ASSETS, 'sim_trust.jpg'), 0.45, 1.95, 5.55)
box, tf = tb(s, 0.45, 5.45, 5.55, 0.28)
lines(tf, [('出典：エネがえる運営事務局調べ（国際航業株式会社）', {'size': 10, 'color': GRAY})])
_, tf = rect(s, 6.25, 1.95, 4.15, 1.55, fill='FCE4E4', line=RED)
lines(tf, [
    ('経済効果シミュレーションの信憑性を', {'size': 15, 'bold': True, 'space': 3}),
    ('疑ったことがある人', {'size': 15, 'bold': True, 'space': 6}),
    ('75.4％', {'size': 34, 'bold': True, 'color': RED, 'align': PP_ALIGN.CENTER}),
])
_, tf = rect(s, 6.25, 3.70, 4.15, 2.55, fill=LTGRAY)
lines(tf, [
    ('だからどうするか', {'size': 16, 'bold': True, 'space': 8}),
    ('・金額の大きさで勝とうとしない', {'size': 14, 'space': 6}),
    ('・前提条件を先に口で言う', {'size': 14, 'space': 2}),
    ('　（電気使用量・単価・設置容量）', {'size': 14, 'color': GRAY, 'space': 6}),
    ('・お客様の検針票の実額から出発する', {'size': 14, 'space': 6}),
    ('・「盛っていない」ことが伝わる方が', {'size': 14, 'bold': True, 'color': RED, 'space': 2}),
    ('　金額より効く', {'size': 14, 'bold': True, 'color': RED}),
])
conclusion(s, '✕ シミュレーション金額　　◎ 安心感', y=6.45)

# ================================================================ 27. PART3 信頼×比較
s = new_slide('PART 3　経済メリット')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('相見積もりになったときのポイント', {'size': 24, 'bold': True})])
_, tf = rect(s, 2.75, 1.95, 5.35, 1.20, fill=LTGREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
lines(tf, [('信頼　×　比較', {'size': 34, 'bold': True, 'align': PP_ALIGN.CENTER})])
rows = [
    ['', 'やること', '言い方'],
    ['信頼', '前提条件を隠さない。数字の出どころを言う',
     '「この数字は〇〇様の検針票から計算しています」'],
    ['比較', '他社下げをせず、お客様に判断軸を与える',
     '「設置後の点検体制で比べてみてください」'],
    ['−', ('他社の強みで戦わず、自社の強みで戦う', {'bold': True}),
     ('会社の安定性／地域密着性／メーカーとしての価値', {'bold': True})],
]
table(s, 0.45, 3.35, 9.95, [0.95, 4.20, 4.80], rows, font_size=13, row_h=0.70, head_h=0.34)
conclusion(s, '金額で勝とうとせず、安心感で選ばれる', y=6.45)

# ================================================================ 28. PART3 今日は金額を出さない
s = new_slide('PART 3　経済メリット')
box, tf = tb(s, 0.45, 1.45, 9.95, 0.62)
lines(tf, [('今日は「金額」を出しません', {'size': 30, 'bold': True, 'color': RED})])
rows = [
    ['', '今日やったこと（経済メリット）', '第4回でやること（金額訴求）'],
    ['目的', '「電気代がいくら減るか」を納得してもらう', '「毎月いくら払うか」を払える形にする'],
    ['出す数字', '削減額（月14,000円／30年504万円）', '月々の支払額・概算見積・総額'],
    ['使う道具', 'シミュレーション', 'シミュレーション＋概算見積＋カタログ'],
    ['ルール', '（必要性訴求の一部）', '⑯〜⑳'],
]
table(s, 0.45, 2.30, 9.95, [1.30, 4.30, 4.35], rows, font_size=13.5, row_h=0.72, head_h=0.36)
_, tf = rect(s, 0.45, 5.60, 9.95, 1.25, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('なぜ分けるのか', {'size': 16, 'bold': True, 'space': 6}),
    ('「いくら得か」が腹落ちしていない状態で「いくら払うか」を出すと、支払額だけが記憶に残ります。',
     {'size': 15, 'space': 4}),
    ('今日は「得の大きさ」を作りきるところまでです。', {'size': 15, 'bold': True, 'color': RED}),
])

# ================================================================ 29. ロープレ
s = new_slide('ロールプレイング')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('FABと経済メリットを通す（10分）', {'size': 24, 'bold': True})])
rows = [
    ['進行', '使うもの', '時間', '狙う反応'],
    ['災害（⑩⑪）', 'アプローチブック p.18・p.19＋事例', ('1分', {'bold': True, 'color': RED}),
     'いつ災害が起こるか分からないですね'],
    ['使い方（⑫）', 'p.20・p.21', '1分', '買うより自給自足した方が良いですね'],
    [('商品説明（FAB）', {'bold': True}), ('ワーク①で書いたB3本', {'bold': True}),
     ('2分', {'bold': True}), ('うちの場合はそうなるんですね', {'bold': True})],
    [('経済メリット', {'bold': True}), ('シミュレーション（伝え方①または②）', {'bold': True}),
     ('2分', {'bold': True}), ('思ったより減るんですね', {'bold': True})],
]
table(s, 0.45, 1.95, 9.95, [1.90, 3.55, 0.90, 3.60], rows, font_size=13, row_h=0.62, head_h=0.32)
_, tf = rect(s, 0.45, 4.65, 4.95, 2.15, fill=LTGRAY)
lines(tf, [
    ('顧客設定（第2回と同じ田中様ご夫婦）', {'size': 15, 'bold': True, 'space': 6}),
    ('・ご主人48歳／奥様45歳／高校生・中学生のお子様', {'size': 13, 'space': 4}),
    ('・築14年・南向き切妻屋根・オール電化ではない', {'size': 13, 'space': 4}),
    ('・太陽光・蓄電池ともに未導入', {'size': 13, 'bold': True, 'color': RED, 'space': 4}),
    ('・電気代 月18,000円　・住宅ローン返済中', {'size': 13, 'space': 4}),
    ('・お子様が来年受験／ご主人は在宅勤務が週2日', {'size': 13, 'bold': True}),
])
_, tf = rect(s, 5.60, 4.65, 4.95, 2.15, fill=LTYEL)
lines(tf, [
    ('この周で見るポイント', {'size': 15, 'bold': True, 'space': 6}),
    ('・災害が1分で終わったか', {'size': 14, 'space': 5}),
    ('・Bの主語が「田中様」になっていたか', {'size': 14, 'bold': True, 'color': RED, 'space': 5}),
    ('・受験・在宅勤務をBに使えたか', {'size': 14, 'space': 5}),
    ('・伝え方①②のどちらを選んだか、理由を言えるか', {'size': 14, 'space': 5}),
    ('・金額（見積り）を口にしていないか', {'size': 14, 'bold': True}),
])

# ================================================================ 30. チェックシート
s = new_slide('チェックシート')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.42)
lines(tf, [('ロープレのあと、自分で○△×をつけてください', {'size': 22, 'bold': True})])
rows = [
    ['パート', '項目', 'できたか'],
    ['① 災害対策', '第三者話法だけで語れた（自分の意見を混ぜない）', ''],
    ['', '3分以内（ロープレでは1分）で終えられた', ''],
    ['', '「うちの地域は災害が来ない」を切り返せた', ''],
    [('② 商品説明（FAB）', {'bold': True, 'fill': LTYEL}), ('F・A・Bを区別して言えた', {'fill': LTYEL}), ('', {'fill': LTYEL})],
    [('', {'fill': LTYEL}), ('Bの主語が「お客様」になっていた', {'fill': LTYEL, 'bold': True}), ('', {'fill': LTYEL})],
    [('', {'fill': LTYEL}), ('Bを3本、違う切り口で言えた', {'fill': LTYEL}), ('', {'fill': LTYEL})],
    [('', {'fill': LTYEL}), ('ヒアリングで聞いた話からBを作れた', {'fill': LTYEL}), ('', {'fill': LTYEL})],
    [('③ 経済メリット', {'bold': True, 'fill': LTYEL}), ('伝え方①②のどちらを使うか選べた', {'fill': LTYEL}), ('', {'fill': LTYEL})],
    [('', {'fill': LTYEL}), ('前提条件（検針票の実額）から出発できた', {'fill': LTYEL, 'bold': True}), ('', {'fill': LTYEL})],
    [('', {'fill': LTYEL}), ('削減額を「1日あたり」まで落とせた', {'fill': LTYEL}), ('', {'fill': LTYEL})],
    [('', {'fill': LTYEL}), ('金額（見積り）を口にしなかった', {'fill': LTYEL}), ('', {'fill': LTYEL})],
]
table(s, 0.45, 1.86, 9.95, [2.30, 6.45, 1.20], rows, font_size=12, row_h=0.34, head_h=0.30)
_, tf = rect(s, 0.45, 6.30, 9.95, 0.58, fill=LTGREEN)
lines(tf, [
    ('商談運営：発言比率6：4／最後までセット提案を崩さなかったか', {'size': 15, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 3}),
    ('黄色が本日の主役。×が付いた項目は、次の商談までに必ず直してください', {'size': 14, 'bold': True, 'align': PP_ALIGN.CENTER}),
])

# ================================================================ 31. まとめ・宿題
s = new_slide('本日のまとめ')
_, tf = rect(s, 0.45, 1.45, 9.95, 2.30, fill=LTGREEN)
lines(tf, [
    ('本日持ち帰ること', {'size': 18, 'bold': True, 'space': 10}),
    ('① 契約とは「阻害要因がすべて解消された状態」。今日つぶしたのは③プランと④金額', {'size': 17, 'space': 8}),
    ('② 災害は「事実 → 事例 → サラッと次へ」。3分で抜ける', {'size': 17, 'space': 8}),
    ('③ FABのBは「お客様が主語」。材料はヒアリングの中にある', {'size': 17, 'bold': True, 'color': RED, 'space': 8}),
    ('④ 経済メリットは2通りの伝え方を使い分ける。金額で勝とうとせず、安心感で選ばれる',
     {'size': 17, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 0.45, 4.00, 9.95, 2.30, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('宿題（次回までに）', {'size': 18, 'bold': True, 'space': 10}),
    ('① 実際の担当顧客1件について、FABのBを3本書いてくること', {'size': 17, 'space': 8}),
    ('② 伝え方①②の両方を、シミュレーションを見ずに口で説明できるようにしてくること', {'size': 17, 'space': 8}),
    ('③ 実商談を1件、災害パートに何分使ったか計測してくること', {'size': 17}),
])
_, tf = rect(s, 0.45, 6.55, 9.95, 0.42, fill=None, line=RED, lw=1.5)
lines(tf, [('次回・第4回は「時期訴求・料金訴求（ルール⑬〜⑳）」／ここで初めて金額を出します',
            {'size': 16, 'bold': True, 'align': PP_ALIGN.CENTER})])
notes(s, '・その日のうちに復習すると定着する（翌日には50％忘れる）ことを添える')

# ---------------------------------------------------------------- 並べ替え
ids = list(sldIdLst)
colophon = [x for x in ids if x is colophon_id][0]
sldIdLst.remove(colophon)
sldIdLst.append(colophon)
page_no[0] += 1

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print('saved:', OUT, '/ slides:', len(prs.slides.__iter__.__self__._sldIdLst))

# ---------------------------------------------------------------- content-type fix
def ensure_jpg_content_type(path):
    """テンプレートが jpg を image/png と誤宣言しているため、python-pptx が
    保存時に Default を落とすことがある。表紙画像が壊れるので明示的に直す。"""
    import zipfile, shutil, re, tempfile
    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        ct = z.read('[Content_Types].xml').decode('utf-8')
        needs = any(n.lower().endswith('.jpg') for n in names)
        if not needs or 'Extension="jpg"' in ct:
            if 'Extension="jpg" ContentType="image/jpeg"' in ct or not needs:
                return False
        data = {n: z.read(n) for n in names}
    ct = re.sub(r'<Default Extension="jpg"[^/]*/>', '', ct)
    ct = ct.replace('<Default Extension="png"',
                    '<Default Extension="jpg" ContentType="image/jpeg"/><Default Extension="png"', 1)
    data['[Content_Types].xml'] = ct.encode('utf-8')
    tmp = path + '.tmp'
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as z:
        for n in names:
            z.writestr(n, data[n])
    shutil.move(tmp, path)
    return True

if ensure_jpg_content_type(OUT):
    print('fixed [Content_Types].xml (jpg -> image/jpeg)')
