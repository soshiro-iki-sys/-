# -*- coding: utf-8 -*-
"""第4回「時期・料金訴求＆ロープレ」研修資料ビルダー

使い方は slides/build/build_session2.py と同じ。
アプローチブックのページ画像（<ASSETS>/ab/ab-01.jpg 〜 ab-26.jpg）と
ピラミッド図（<ASSETS>/pyramid.png）を用意してから実行する。

    TRAINING_ASSETS=/path/to/assets python slides/build/build_session4.py

書式は docs/研修資料フォーマット仕様.md、内容は docs/知識_*.md を参照。
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
MSO_SHAPE_TYPE_AUTO = MSO_SHAPE_TYPE.AUTO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.opc.packuri import PackURI

REPO   = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
ASSETS = os.environ.get('TRAINING_ASSETS', '/tmp/training-assets')
AB     = os.path.join(ASSETS, 'ab')
TPL    = os.path.join(REPO, 'templates', '研修資料_フォーマット見本.pptx')
OUT    = os.path.join(REPO, 'slides', '2026営業研修_第4回_時期・料金訴求＆ロープレ.pptx')

GOTHIC='HGPｺﾞｼｯｸE'; MINCHO='HGP明朝E'; UD='BIZ UDPゴシック'; YU='游ゴシック'
RED='FF0000'; BLACK='000000'; WHITE='FFFFFF'; GREEN='9BBB59'; YELLOW='FFFF00'
GRAY='595959'; BLUE='0070C0'; LTGREEN='EAF1DD'; LTGRAY='F2F2F2'; LTYEL='FFF2CC'
LTBLUE='DEEBF7'; ORANGE='F79646'

# ---------------------------------------------------------------- helpers
def _rpr_font(run, name):
    rPr = run._r.get_or_add_rPr()
    rPr.get_or_add_latin().set('typeface', name)
    latin = rPr.find(qn('a:latin'))
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = parse_xml('<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        latin.addnext(ea)
    ea.set('typeface', name)

def _highlight(run, color):
    rPr = run._r.get_or_add_rPr()
    hl = parse_xml(
        '<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:srgbClr val="%s"/></a:highlight>' % color)
    latin = rPr.find(qn('a:latin'))
    if latin is not None:
        latin.addprevious(hl)
    else:
        rPr.append(hl)

def run(p, text, size=16, bold=False, color=BLACK, font=GOTHIC, hl=None):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)
    if hl: _highlight(r, hl)
    _rpr_font(r, font)
    return r

def tb(slide, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP, margin=0.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return box, tf

def para(tf, first=False, space_after=4, align=PP_ALIGN.LEFT, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after); p.alignment = align
    if line: p.line_spacing = line
    return p

def lines(tf, items, size=16, color=BLACK, bold=False, space=5, font=GOTHIC, line=None, align=PP_ALIGN.LEFT):
    """items: str | (str,dict)"""
    for i, it in enumerate(items):
        opts = {}
        if isinstance(it, tuple): it, opts = it
        p = para(tf, first=(i == 0), space_after=opts.pop('space', space),
                 align=opts.pop('align', align), line=opts.pop('line', line))
        run(p, it, size=opts.pop('size', size), bold=opts.pop('bold', bold),
            color=opts.pop('color', color), font=opts.pop('font', font),
            hl=opts.pop('hl', None))

def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=1.0):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    else: s.fill.background()
    if line:
        s.line.color.rgb = RGBColor.from_string(line); s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    return s, tf

def pic(slide, path, x, y, w, border=True):
    p = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    if border:
        p.line.color.rgb = RGBColor.from_string('808080'); p.line.width = Pt(1)
    return p

def table(slide, x, y, w, col_w, rows, font_size=11, header=True,
          header_fill='000000', header_color=WHITE, row_h=0.30, head_h=0.30):
    n_r, n_c = len(rows), len(col_w)
    shp = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w),
                                 Inches(head_h + row_h * (n_r - 1)))
    t = shp.table
    t.first_row = header; t.horz_banding = False
    for i, cw in enumerate(col_w): t.columns[i].width = Inches(cw)
    t.rows[0].height = Inches(head_h)
    for i in range(1, n_r): t.rows[i].height = Inches(row_h)
    for ri, rowdata in enumerate(rows):
        for ci, cell in enumerate(rowdata):
            txt, opts = (cell, {}) if isinstance(cell, str) else cell
            c = t.cell(ri, ci)
            c.margin_left = c.margin_right = Inches(0.05)
            c.margin_top = c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = opts.get('fill', header_fill if (header and ri == 0) else None)
            if fill: c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string(fill)
            else: c.fill.background()
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = opts.get('align',
                PP_ALIGN.CENTER if (header and ri == 0) else PP_ALIGN.LEFT)
            run(p, txt, size=opts.get('size', font_size),
                bold=opts.get('bold', header and ri == 0),
                color=opts.get('color', header_color if (header and ri == 0) else BLACK),
                font=opts.get('font', YU))
    return t

# ---------------------------------------------------------------- slide frame
prs = Presentation(TPL)
BLANK = None
for l in prs.slide_masters[0].slide_layouts:
    if l.name == '白紙': BLANK = l
assert BLANK is not None

sldIdLst = prs.slides._sldIdLst
orig_ids = list(sldIdLst)
cover_id, colophon_id = orig_ids[0], orig_ids[6]
for sid in orig_ids[1:6]:                      # drop template samples 2-6
    rId = sid.get(qn('r:id')); sldIdLst.remove(sid)
    prs.part.drop_rel(rId)
# python-pptx names new slides by len(sldIdLst)+1 -> move colophon out of that range
prs.part.related_part(colophon_id.get(qn('r:id'))).partname = PackURI('/ppt/slides/slide900.xml')

page_no = [0]

def new_slide(heading=None):
    s = prs.slides.add_slide(BLANK)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    if heading is not None:
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.157), Inches(0.757), Inches(0.140), Inches(0.512))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor.from_string(BLACK)
        bar.line.color.rgb = RGBColor.from_string('D9D9D9'); bar.line.width = Pt(0.75)
        bar.shadow.inherit = False
        box, tf = tb(s, 0.348, 0.736, 9.9, 0.572)
        lines(tf, [heading], size=28, font=MINCHO)
    page_no[0] += 1
    box, tf = tb(s, 8.31, 7.18, 2.30, 0.33)
    lines(tf, [str(page_no[0])], size=12, color=WHITE, font=GOTHIC, align=PP_ALIGN.RIGHT)
    return s

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def ab_img(n):
    return os.path.join(AB, 'ab-%02d.jpg' % n)

# ---------------------------------------------------------------- 表紙
cover = prs.slides[0]
for shp in cover.shapes:
    if shp.shape_type == MSO_SHAPE_TYPE_AUTO and shp.name == '正方形/長方形 4':
        sp = shp._element.spPr
        fill = sp.find(qn('a:solidFill'))
        for ch in list(fill): fill.remove(ch)
        fill.append(parse_xml(
            '<a:srgbClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'val="4BACC6"><a:alpha val="80000"/></a:srgbClr>'))
    if shp.has_text_frame and shp.name == 'テキスト ボックス 5':
        tf = shp.text_frame; tf.clear()
        lines(tf, [
            ('2026年度　営業研修', {'size': 32, 'bold': True, 'color': WHITE, 'font': UD, 'space': 8}),
            ('第4回　時期・料金訴求', {'size': 52, 'bold': True, 'color': WHITE, 'font': UD, 'space': 2}),
            ('＆ロープレ', {'size': 52, 'bold': True, 'color': WHITE, 'font': UD, 'space': 10}),
            ('〜太陽光＋蓄電池セット販売〜', {'size': 24, 'bold': True, 'color': WHITE, 'font': UD}),
        ])
    if shp.has_text_frame and shp.name == 'テキスト ボックス 6':
        tf = shp.text_frame; tf.clear()
        lines(tf, [('2026年9月上旬　＠船井総研', {'size': 14, 'color': WHITE, 'font': UD})],
              align=PP_ALIGN.RIGHT)
page_no[0] = 1

# ---------------------------------------------------------------- rule helpers
def rule_head(s, no_title, point):
    box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
    lines(tf, [(no_title, {'size': 24, 'bold': True})])
    _, tf = rect(s, 0.45, 1.88, 9.95, 0.46, fill=None, line=RED, lw=1.25)
    p = para(tf, first=True, align=PP_ALIGN.LEFT)
    run(p, 'POINT　', size=15, bold=True, color=RED)
    run(p, point, size=16, bold=True)

def goal_box(s, x, y, w, h, quote):
    _, tf = rect(s, x, y, w, h, fill=LTBLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    lines(tf, [
        ('ゴール（お客様の反応）', {'size': 13, 'bold': True, 'color': BLUE, 'space': 6}),
    ])
    for q in quote:
        p = para(tf, space_after=3)
        run(p, '「%s」' % q, size=16, bold=True)

def points_box(s, x, y, w, h, items, title='説明のポイント'):
    _, tf = rect(s, x, y, w, h, fill=LTGRAY)
    lines(tf, [(title, {'size': 13, 'bold': True, 'color': GRAY, 'space': 6})])
    for it in items:
        opts = {}
        if isinstance(it, tuple): it, opts = it
        p = para(tf, space_after=opts.pop('space', 6))
        run(p, it, size=opts.pop('size', 12.5), bold=opts.pop('bold', False),
            color=opts.pop('color', BLACK))

def conclusion(s, text, y=6.36):
    _, tf = rect(s, 0.45, y, 9.95, 0.66, fill=LTGREEN)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, 'このページの結論　', size=12, bold=True, color=GRAY)
    run(p, text, size=15, bold=True, hl=YELLOW)

def talk_slide(s, no_title, talks, conc=None, note=None):
    box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
    lines(tf, [(no_title, {'size': 24, 'bold': True})])
    y = 1.92
    for label, body, h in talks:
        _, tf = rect(s, 0.45, y, 8.05, h, fill=LTGRAY)
        for i, ln in enumerate(body):
            p = para(tf, first=(i == 0), space_after=3)
            run(p, ln, size=14)
        _, tf2 = rect(s, 8.68, y, 1.72, 0.38, fill=ORANGE)
        lines(tf2, [(label, {'size': 12, 'bold': True, 'color': WHITE,
                             'align': PP_ALIGN.CENTER})])
        y += h + 0.16
    if conc: conclusion(s, conc)
    if note: notes(s, note)

# ================================================================ 2. 本日のゴール
s = new_slide('本日のゴール')
box, tf = tb(s, 0.45, 1.45, 9.9, 0.62)
lines(tf, [('「今が一番お得だ」と思わせ、払える金額として着地させる',
            {'size': 28, 'bold': True, 'color': RED})])
_, tf = rect(s, 0.45, 2.35, 9.95, 3.05, fill=LTGREEN)
lines(tf, [
    ('研修が終わったときの「できる状態」', {'size': 18, 'bold': True, 'space': 10}),
    ('①　「待てば安くなる」を、データと事実で断ち切れる', {'size': 19, 'space': 9}),
    ('②　太陽光＋蓄電池を「出費」ではなく「投資」として説明できる', {'size': 19, 'space': 9}),
    ('③　金額を 1日 → 1か月 → 総額 の順で出せる', {'size': 19, 'space': 9}),
    ('④　アプローチブックを閉じて、見積りへ自然に受け渡せる', {'size': 19}),
])
_, tf = rect(s, 0.45, 5.65, 9.95, 0.85, fill=None, line=RED, lw=1.5)
lines(tf, [('本日の範囲：営業ルール ⑬〜⑳　／　アプローチブック p.23〜p.26 ＋ 見積り',
            {'size': 21, 'bold': True, 'align': PP_ALIGN.CENTER})])
notes(s, '・本日で20のルールが全部つながる。次回は通しロープレだけだと予告しておく\n'
         '・金額の話は「安く見せる」ではなく「払える形にする」ことだと最初に握る')

# ================================================================ 3. アジェンダ
s = new_slide('本日のアジェンダ')
rows = [
    ['時間', '内容', '扱うもの'],
    ['0:00-0:10', '前回の振り返り／宿題の確認', '第3回：ルール⑩〜⑫'],
    ['0:10-0:50', ('時期訴求　ルール⑬〜⑮', {'bold': True}), 'アプローチブック p.23〜p.26'],
    ['0:50-1:00', '休憩', '−'],
    ['1:00-1:15', ('金額訴求　ルール⑯', {'bold': True}), '月々支払い額のみ／逃げ方3つ'],
    ['1:15-1:30', ('分割して伝える（フレーミング効果）', {'bold': True}), '1日 → 1か月 → 総額'],
    ['1:30-1:45', ('金額訴求　ルール⑰〜⑳', {'bold': True}), '容量・見積り・後始末'],
    ['1:45-1:57', 'ロールプレイング（2周）', 'p.23〜26　→　金額提示まで'],
    ['1:57-2:00', 'チェック／まとめ／宿題', '−'],
]
table(s, 0.45, 1.62, 9.95, [1.55, 4.3, 4.1], rows, font_size=14, row_h=0.50, head_h=0.38)
notes(s, '・8ルールと多い。⑲⑳は厳密なルールではなく推奨事項なので巻いてよい')

# ================================================================ 4. 前回の振り返り
s = new_slide('前回の振り返り')
box, tf = tb(s, 0.45, 1.42, 9.9, 0.45)
lines(tf, [('第3回「必要性訴求②＆ロープレ」で押さえたこと', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 2.00, 4.85, 2.30, fill=LTGRAY)
lines(tf, [
    ('必要性訴求②（⑩〜⑫）', {'size': 17, 'bold': True, 'space': 8}),
    ('・災害は「事実 → 事例 → サラッと次へ」。3分で抜ける', {'size': 14, 'space': 5}),
    ('・⑪にページが無いのは、一般論では意味が無いから', {'size': 14, 'space': 5}),
    ('・創る → 貯める → 夜に使う', {'size': 14, 'space': 5}),
    ('・このサイクルはセットでしか成立しない', {'size': 14, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 5.55, 2.00, 4.85, 2.30, fill=LTYEL)
lines(tf, [
    ('ここまでで作れているもの', {'size': 17, 'bold': True, 'space': 8}),
    ('聴く姿勢（①〜⑥）', {'size': 15, 'space': 5}),
    ('＋ 電気代のネガ（⑦⑧⑨）', {'size': 15, 'space': 5}),
    ('＋ もしもの備え（⑩⑪）', {'size': 15, 'space': 5}),
    ('＋ 使い方のイメージ（⑫）', {'size': 15, 'space': 8}),
    ('＝ お客様は「欲しい」状態', {'size': 16, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 0.45, 4.60, 9.95, 1.75, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('宿題の確認（10分）', {'size': 18, 'bold': True, 'space': 8}),
    ('① 災害事例を90秒で言ってみてください（2名指名）', {'size': 16, 'space': 6}),
    ('② アプローチブック p.9〜p.21 の結論の一文を、通しで言えますか', {'size': 16, 'space': 6}),
    ('③ 実商談で、災害パートに何分使いましたか', {'size': 16}),
])

# ================================================================ 5. 全体像① 4ステップ
s = new_slide('全体像①：契約までの4ステップ')
pic(s, os.path.join(ASSETS, 'pyramid.png'), 0.55, 1.55, 6.0, border=False)
_, tf = rect(s, 6.95, 1.75, 3.45, 1.55, fill=None, line=RED, lw=2.0)
lines(tf, [
    ('本日の範囲', {'size': 20, 'bold': True, 'color': RED, 'align': PP_ALIGN.CENTER, 'space': 8}),
    ('STEP3　時期訴求（⑬〜⑮）', {'size': 17, 'bold': True, 'space': 5}),
    ('STEP4　金額訴求（⑯〜⑳）', {'size': 17, 'bold': True}),
])
_, tf = rect(s, 6.95, 3.55, 3.45, 2.55, fill=LTGRAY)
lines(tf, [
    ('土台はもう積んである', {'size': 17, 'bold': True, 'space': 8}),
    ('①〜⑫で「聴く姿勢」と', {'size': 14, 'space': 0}),
    ('「必要性」は作れました。', {'size': 14, 'space': 8}),
    ('残るはピラミッドの上2段。', {'size': 14, 'space': 8}),
    ('ここで初めて', {'size': 14, 'space': 0}),
    ('「いつ」と「いくら」', {'size': 14, 'bold': True, 'color': RED, 'space': 0}),
    ('の話をします。', {'size': 14, 'space': 8}),
    ('順番を飛ばすと、価格だけで', {'size': 14, 'space': 0}),
    ('判断されて終わります。', {'size': 14}),
])

# ================================================================ 6. 全体像② 20のルール
s = new_slide('全体像②：太陽光＋蓄電池営業 20のルール')
DONE = {'fill': LTGRAY, 'color': GRAY}
HL = {'fill': LTYEL}
rows = [
    ['項目', 'ゴール', 'ルール', '詳細'],
    [('前準備・聴く姿勢作り', dict(DONE)), ('商談の趣旨を認識する', dict(DONE)), ('①〜⑥', dict(DONE)), ('第2回で実施済み', dict(DONE))],
    [('必要性訴求', dict(DONE)), ('現状を正確に把握する', dict(DONE)), ('⑦〜⑨', dict(DONE)), ('第2回で実施済み', dict(DONE))],
    [('', dict(DONE)), ('', dict(DONE)), ('⑩〜⑫', dict(DONE)), ('第3回で実施済み', dict(DONE))],
    [('時期訴求', dict(HL)), ('早期対策を進める', dict(HL)), ('⑬導入に向けた考え方の転換', dict(HL)),
     ('「導入するかどうか」から「いつ導入するのが最も賢明か」へ転換させる', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('⑭物価指数とコストトレンドの説明', dict(HL)),
     ('消費者物価指数が上昇中で、待つほど導入費用が上がる理解を得る', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('⑮投資対効果の視点', dict(HL)),
     ('「出費」ではなく将来の電気代削減に繋がる「投資」と捉えてもらう', dict(HL))],
    [('金額訴求', dict(HL)), ('費用対効果に納得する', dict(HL)), ('⑯概算の機器導入費用のおさらい', dict(HL)),
     ('電気代出費の範囲内に収まれば導入に問題ないかどうか尋ねる', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('⑰太陽光設置容量の確認', dict(HL)),
     ('屋根に載せられるパネルの量と発電量を算出して伝える', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('⑱概算見積りの提示', dict(HL)),
     ('最適な見積書を提示し、その場で導入することに前向きな答えをもらう', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('⑲正式見積りの提示', dict(HL)),
     ('1週間以内に工事スタッフと共に現場調査の上、正確な見積もりを提示する', dict(HL))],
    [('', dict(HL)), ('', dict(HL)), ('⑳後始末の実施', dict(HL)),
     ('契約書を交わし、その日中に御礼のハガキをポストに投函する', dict(HL))],
]
table(s, 0.30, 1.55, 10.25, [1.30, 1.60, 2.70, 4.65], rows, font_size=10.5, row_h=0.40, head_h=0.30)
_, tf = rect(s, 0.30, 6.35, 10.25, 0.50, fill=None, line=RED, lw=1.5)
lines(tf, [('本日で20のルールが全部つながります。次回は通しロープレだけです',
            {'size': 17, 'bold': True, 'align': PP_ALIGN.CENTER})])

# ================================================================ 7. 章見出し 時期訴求
s = new_slide('時期訴求')
box, tf = tb(s, 0.45, 1.38, 9.9, 0.42)
lines(tf, [('ゴール：早期対策を進める', {'size': 22, 'bold': True})])
rows = [
    ['項目', 'ゴール', 'ルール', 'POINT'],
    ['時期訴求', '早期対策を進める', '⑬導入に向けた考え方の転換',
     '「導入するか」の迷いを断ち、「いつ導入するのが最も合理的か」へ移行させる'],
    ['', '', '⑭物価指数とコストトレンドの説明',
     '「待てば安くなる」に、客観データと「実質負担額」の視点で反論する'],
    ['', '', '⑮投資対効果の視点',
     '単なる「節約」ではなく「資産形成」であることを、持ち家との比較で納得させる'],
]
for r in rows[1:]:
    for i in range(len(r)):
        r[i] = (r[i], {'fill': LTGREEN})
table(s, 0.30, 1.95, 10.25, [1.25, 1.75, 2.85, 4.40], rows, font_size=12, row_h=0.62, head_h=0.32)
_, tf = rect(s, 0.30, 4.25, 4.95, 2.20, fill=LTGRAY)
lines(tf, [
    ('使うアプローチブック', {'size': 17, 'bold': True, 'space': 8}),
    ('p.23　10年後に設置する vs すぐ設置する', {'size': 15, 'space': 6}),
    ('p.24　待っても安くならないの？', {'size': 15, 'space': 6}),
    ('p.25　まとまった費用をかけずに導入する方法', {'size': 15, 'space': 6}),
    ('p.26　かしこいスマートハウスの導入時期', {'size': 15}),
])
_, tf = rect(s, 5.60, 4.25, 4.95, 2.20, fill=LTYEL)
lines(tf, [
    ('このパートの勝ち筋', {'size': 17, 'bold': True, 'space': 8}),
    ('お客様の頭の中を', {'size': 15, 'space': 4}),
    ('「導入するかどうか」から', {'size': 16, 'bold': True, 'space': 2}),
    ('「いつ導入するのが賢いか」へ', {'size': 16, 'bold': True, 'color': RED, 'space': 4}),
    ('切り替えるのが目的です。', {'size': 15, 'space': 8}),
    ('やるかどうかの議論に戻らせない。', {'size': 15, 'bold': True}),
])

# ================================================================ 8. ルール⑬-1
s = new_slide('時期訴求')
rule_head(s, 'ルール⑬：導入に向けた考え方の転換', '「導入するか」の迷いを断ち、「いつ導入するのが最も合理的か」へ移行させる')
pic(s, ab_img(23), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.23　10年後に設置 vs すぐ設置', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.15, ['今、設置した方が将来的にお得なのですね'])
points_box(s, 5.35, 3.78, 5.05, 2.50, [
    ('10年後に設置すると、その10年で144万円の電気代を払うだけ（月12,000円の場合）', {'bold': True, 'size': 13}),
    'その間、売電収入も自家消費による電気代削減の恩恵も受けられない',
    '設置までの10年で、電気代の支払いと売電単価の減少というデメリットしかない',
    ('この「機会損失」を具体的に認識させ、待つことのデメリットを明確にする', {'bold': True, 'color': RED, 'size': 13}),
])
conclusion(s, '設置を遅くするにつれ、経済メリットは少なくなります！')

# ================================================================ 9. ルール⑬-2
s = new_slide('時期訴求')
talk_slide(s, 'ルール⑬：「機会損失」を自分事にする3手', [
    ('①言葉を置く', [
        '「導入を待つということは、この経済効果を得られるチャンスを先送りにする、つまり『機会損失』が生じているのと同じなんです」',
        '（損失回避の心理を刺激する。「損」という言葉を必ず使う）',
    ], 0.85),
    ('②数字を置換', [
        '「〇〇様の場合は月々約18,000円ですので、10年間ですと216万円の機会損失になりますね」',
        '（検針票の実額に必ず置き換える。一般論のままにしない）',
    ], 0.85),
    ('③質問で促す', [
        '「この『10年間何もしない場合の支出』と『今すぐ始めて得られるメリット』、どちらが大きいと思われますか？」',
        '「もし、この216万円を別のことに使えるとしたら、何に使いたいですか？」',
    ], 0.85),
], conc='お客様自身に考えさせ、損失の大きさを実感させる',
   note='・数字は必ずお客様の検針票の額に置き換える。144万円のまま話さない')

# ================================================================ 10. 時期訴求 スマホの例え
s = new_slide('時期訴求')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('「今じゃなくていい」を潰す：スマホで考える', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.35)
lines(tf, [('最新のスマホが欲しいと思っているのに、まだ買っていない人の言い分と同じ構造です',
            {'size': 15, 'color': GRAY})])
_, tf = rect(s, 0.45, 2.32, 4.95, 2.15, fill='FCE4E4', line=RED)
lines(tf, [
    ('買わない人の言い分', {'size': 18, 'bold': True, 'color': RED, 'space': 10}),
    ('① 今じゃなくていいや', {'size': 17, 'space': 8}),
    ('② 待てば安くなる', {'size': 17, 'space': 8}),
    ('③ 最新のものが出た後、型落ちを買おう', {'size': 17}),
])
_, tf = rect(s, 5.60, 2.32, 4.95, 2.15, fill=LTGREEN, line='70A040')
lines(tf, [
    ('潰し方（3点）', {'size': 18, 'bold': True, 'space': 10}),
    ('① 電気代は毎年のように上がり続けている', {'size': 17, 'space': 8}),
    ('② 待っている間の機会損失（プロスペクト理論）', {'size': 17, 'space': 8}),
    ('③ 数年で性能は変わらない', {'size': 17}),
])
_, tf = rect(s, 0.45, 4.70, 9.95, 1.55, fill=LTGRAY)
lines(tf, [
    ('【必要性】は、もう作れているはずです', {'size': 17, 'bold': True, 'space': 8}),
    ('最新のものが欲しい／バッテリーの持ちをよくしたい／いい画質で写真を撮りたい ── これがスマホの必要性。',
     {'size': 15, 'space': 5}),
    ('太陽光＋蓄電池で言えば、①〜⑫でここまで作ってきました。あとは「いつ」の問題だけです。',
     {'size': 15, 'bold': True}),
])
conclusion(s, '「今が一番お得だ」と感じてもらうのがゴール', y=6.40)

# ================================================================ 11. ルール⑭-1
s = new_slide('時期訴求')
rule_head(s, 'ルール⑭：物価指数とコストトレンドの説明', '「待てば安くなる」に、客観データと「実質負担額」の視点で反論する')
pic(s, ab_img(24), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.24　待っても安くならないの？', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.15, ['待っても物価は安くなるどころか、上がっていく一方なのですね'])
points_box(s, 5.35, 3.78, 5.05, 2.50, [
    ('消費者物価指数は2020年1月を100として、2026年3月時点で112.7', {'bold': True, 'size': 13}),
    '6年で物価は12.7％上昇。電気代だけでなく世の中全体が上がっている',
    '「太陽光・蓄電池だけが将来大幅に安くなる」という期待は現実的ではない',
    ('待って価格がわずかに下がるのを期待するより、今導入して高騰リスクを回避する方が賢明', {'bold': True, 'color': RED, 'size': 13}),
])
conclusion(s, '待っていても必ずしも安くなるとは限らないのです')

# ================================================================ 12. ルール⑭-2
s = new_slide('時期訴求')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ルール⑭：なぜ蓄電池は安くならないのか（根拠）', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.35)
lines(tf, [('「本当に安くならないの？」と突っ込まれたときの材料。原材料の供給構造が変わらないためです',
            {'size': 15, 'color': GRAY})])
_, tf = rect(s, 0.45, 2.32, 4.95, 2.35, fill=LTGRAY)
lines(tf, [
    ('炭酸リチウム', {'size': 18, 'bold': True, 'space': 8}),
    ('・世界最大の生産国はチリ', {'size': 15, 'space': 5}),
    ('・日本は100％輸入、うち80％がチリから', {'size': 15, 'space': 5}),
    ('・アタカマ塩湖でかん水を1年間天日干しにして作る', {'size': 15, 'space': 5}),
    ('・世界シェアNo.1の企業ですら半自動化どまり', {'size': 15, 'space': 8}),
    ('生産量を急に増やせない', {'size': 15, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 5.60, 2.32, 4.95, 2.35, fill=LTGRAY)
lines(tf, [
    ('コバルト', {'size': 18, 'bold': True, 'space': 8}),
    ('・レアメタル。コンゴ民主共和国に生産が集中', {'size': 15, 'space': 5}),
    ('・コンゴ産の20％が手掘り', {'size': 15, 'space': 5}),
    ('・約4万人の子供が働いていると推計（アムネスティ）', {'size': 15, 'space': 5}),
    ('・児童労働問題により生産量を増やせない', {'size': 15, 'space': 8}),
    ('特定の国に依存し、増産できない', {'size': 15, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 0.45, 4.90, 9.95, 1.35, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('言い方の型', {'size': 16, 'bold': True, 'space': 8}),
    ('「もちろん、将来の価格を正確に予測することは誰にもできません。ただ、現在のデータやトレンドから合理的に判断すると、'
     '大幅な価格低下を期待して好機を逃すよりも、確実なメリットを早く享受する方が賢明ではないでしょうか」',
     {'size': 14}),
])
notes(s, '・断定しない。「予測はできない、ただしデータでは」という冷静な分析のトーンで話す')

# ================================================================ 13. ルール⑮-1
s = new_slide('時期訴求')
rule_head(s, 'ルール⑮：投資対効果の視点', '単なる「節約」ではなく「資産形成」であることを、持ち家との比較で納得させる')
pic(s, ab_img(25), 0.45, 2.50, 4.55)
box, tf = tb(s, 0.45, 5.85, 4.55, 0.3)
lines(tf, [('アプローチブック p.25　まとまった費用をかけずに導入', {'size': 12, 'color': GRAY})])
goal_box(s, 5.35, 2.50, 5.05, 1.15, ['長期的には、財産になるのですね'])
points_box(s, 5.35, 3.78, 5.05, 2.50, [
    ('実質負担額 ＝ 設置費用 − メリット', {'bold': True, 'size': 14}),
    '太陽光が生み出す経済メリットが、導入時の実質的な負担を軽くする（設置費用は太陽光が稼いでくれる）',
    '家賃を払いたくないから持ち家を買う。電気代を払いたくないから創蓄を導入する',
    ('創蓄は支払っている最中も経済メリットが出るのが、住宅ローンとの違い', {'bold': True, 'color': RED, 'size': 13}),
])
conclusion(s, '持ち家と同じように、長期的な財産になるものなのです')

# ================================================================ 14. ルール⑮-2
s = new_slide('時期訴求')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ルール⑮：「賃貸 vs 持ち家」で言い換える', {'size': 24, 'bold': True})])
rows = [
    ['', '払い続けるもの', '結果'],
    ['賃貸住宅', '家賃', '払い続けても自分のものにはならない'],
    [('電気を買う生活', {'bold': True, 'color': RED}), ('電気代', {'bold': True, 'color': RED}),
     ('払い続けても自分のものにはならない', {'bold': True, 'color': RED})],
    ['持ち家', 'ローン（初期費用）', '払い終われば資産になる'],
    [('太陽光＋蓄電池', {'bold': True, 'color': RED}), ('ローン（初期費用）', {'bold': True, 'color': RED}),
     ('払い終われば資産。しかも電気代が大幅に減る', {'bold': True, 'color': RED})],
]
table(s, 0.45, 1.95, 9.95, [2.45, 3.10, 4.40], rows, font_size=14, row_h=0.52, head_h=0.32)
pic(s, ab_img(26), 0.45, 4.45, 3.40)
box, tf = tb(s, 0.45, 6.85, 3.40, 0.26)
lines(tf, [('アプローチブック p.26', {'size': 11, 'color': GRAY})])
_, tf = rect(s, 4.20, 4.45, 6.20, 1.30, fill=LTGRAY)
lines(tf, [
    ('①「財産（資産）」という言葉を明確に使う', {'size': 15, 'bold': True, 'space': 5}),
    ('「単に電気代が安くなるだけでなく、ご自宅に新たな『価値』や『財産』が生まれるとお考えください」',
     {'size': 13}),
])
_, tf = rect(s, 4.20, 5.90, 6.20, 1.25, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('②「賢いお金の使い方」の視点', {'size': 15, 'bold': True, 'space': 5}),
    ('「同じお金を支払うなら、消費してなくなるものより、将来の支出を減らしてくれるものに使う方が賢いと言えませんか」',
     {'size': 13}),
])

# ================================================================ 15. 休憩
s = new_slide(None)
_, tf = rect(s, 2.40, 2.85, 6.05, 1.80, fill=LTGREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
lines(tf, [
    ('休憩　10分', {'size': 44, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 10}),
    ('後半は「金額訴求」に入ります', {'size': 18, 'align': PP_ALIGN.CENTER}),
])

# ================================================================ 16. 章見出し 金額訴求
s = new_slide('金額訴求')
box, tf = tb(s, 0.45, 1.38, 9.9, 0.42)
lines(tf, [('ゴール：費用対効果に納得してもらう', {'size': 22, 'bold': True})])
rows = [
    ['項目', 'ルール', 'POINT'],
    ['金額訴求', '⑯概算の機器導入費用のおさらい', '支払っている光熱費で太陽光＋蓄電池を導入できることを伝える'],
    ['', '⑰太陽光設置容量の確認', '設置容量を確認することで設置後のイメージを膨らませる'],
    ['', '⑱概算見積りの提示', '月々1万円程度の予算で提案し、金額面でのハードルを大きく下げる'],
    ['', '⑲正式見積りの提示', '成約率最大化のために、正式見積りの提出で即決を促す'],
    ['', '⑳後始末の実施', '契約後のアフターフォローをしっかり行うことでお客様の不安を払拭する'],
]
for r in rows[1:]:
    for i in range(len(r)):
        r[i] = (r[i], {'fill': LTGREEN})
table(s, 0.30, 1.95, 10.25, [1.30, 3.05, 5.90], rows, font_size=12.5, row_h=0.52, head_h=0.32)
_, tf = rect(s, 0.30, 5.00, 4.95, 1.45, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('ここからアプローチブックを離れます', {'size': 17, 'bold': True, 'color': RED, 'space': 6}),
    ('使うのはシミュレーション・概算見積・カタログ・契約書・お礼状です。', {'size': 14}),
])
_, tf = rect(s, 5.60, 5.00, 4.95, 1.45, fill=LTGRAY)
lines(tf, [
    ('⑲⑳は「推奨事項」', {'size': 17, 'bold': True, 'space': 6}),
    ('厳密な営業ルールではなく、成約率を最大化するために守ることを推奨する内容です。',
     {'size': 14}),
])

# ================================================================ 17. ルール⑯-1
s = new_slide('金額訴求')
rule_head(s, 'ルール⑯：概算の機器導入費用（月額）のおさらい', '支払っている光熱費で太陽光＋蓄電池を導入できることを伝える')
_, tf = rect(s, 0.45, 2.50, 4.95, 1.55, fill=LTBLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
lines(tf, [
    ('光熱費 ＝ 生涯ローン', {'size': 26, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 8}),
    ('いま払っている光熱費が、どこまで太陽光＋蓄電池に置き換わるか', {'size': 13, 'align': PP_ALIGN.CENTER}),
])
points_box(s, 5.60, 2.50, 4.95, 1.55, [
    ('機器の費用だけでなく、既にかかっている電気代と常に天秤にかけながら進める', {'bold': True, 'size': 13}),
    '「電気代の範囲内に収まれば、導入には問題ないですか？」と尋ねる',
], title='考え方')
_, tf = rect(s, 0.45, 4.30, 9.95, 1.95, fill='FCE4E4', line=RED)
lines(tf, [
    ('まずは「月々支払い額」のみを理解いただくことに終始する', {'size': 20, 'bold': True, 'color': RED, 'space': 10}),
    ('総額費用の具体的な提示は、ここでは禁物です。', {'size': 16, 'space': 6}),
    ('総額から入ると、価値ではなく価格だけで判断されてしまいます。', {'size': 16}),
])
conclusion(s, '電気代の範囲内に収まるかどうか。それだけをここで握る', y=6.40)

# ================================================================ 18. ルール⑯-2
s = new_slide('金額訴求')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('ルール⑯：「総額はいくら？」と聞かれたときの逃げ方3つ', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.35)
lines(tf, [('ことあるごとに聞かれます。使いやすい「かわし方」を先に持っておきましょう',
            {'size': 15, 'color': GRAY})])
rows = [
    ['', 'かわし方', '評価'],
    ['①', '「改めてご自宅にあったお見積りをお出ししたいので…」',
     '最も理にかなっているが、「ザックリでいいから教えて」と言われやすいのが難点'],
    ['②', '「上に確認を取らないと分からなくて…」',
     '使い勝手は良いが、皆さんのお立場とのバランスになる'],
    ['③', '「この後のお話で、一通り出てきますので…」',
     'やや強引な「ルールでの包み込み」パターン。軽いツッコミなら十分交わせる'],
]
table(s, 0.45, 2.32, 9.95, [0.50, 4.35, 5.10], rows, font_size=14, row_h=0.85, head_h=0.34)
_, tf = rect(s, 0.45, 5.30, 9.95, 0.95, fill=LTYEL)
lines(tf, [
    ('ワーク（3分）', {'size': 16, 'bold': True, 'space': 6}),
    ('自分が一番言いやすいものを1つ選び、声に出して言ってみてください。棒読みだと見抜かれます。',
     {'size': 15}),
])
conclusion(s, '総額は最後。ここでは月々の額だけを握る', y=6.40)

# ================================================================ 19. 分割して伝える①
s = new_slide('分割して伝える')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('①　フレーミング効果を使う', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 1.95, 9.95, 0.80, fill='FCE4E4', line=RED)
lines(tf, [('✕　「結論、材工込みで総額200万円です」',
            {'size': 24, 'bold': True, 'color': RED, 'align': PP_ALIGN.CENTER})])
_, tf = rect(s, 0.45, 3.00, 9.95, 1.30, fill=LTGRAY)
lines(tf, [
    ('フレーミング効果とは', {'size': 16, 'bold': True, 'space': 6}),
    ('物事を表現する枠組み（フレーム）を変えることで、与える印象が変わり、購入意欲が高まる効果。',
     {'size': 15}),
])
_, tf = rect(s, 0.45, 4.55, 4.95, 1.70, fill=LTGREEN)
lines(tf, [
    ('例）10万円のバッグ', {'size': 17, 'bold': True, 'space': 8}),
    ('「10年間使ったら1万円/年だし、安いかも…」', {'size': 16, 'space': 6}),
    ('分割して伝えると、印象が変わる', {'size': 15, 'bold': True}),
])
_, tf = rect(s, 5.60, 4.55, 4.95, 1.70, fill=LTGREEN)
lines(tf, [
    ('例）飲み会に置き換える', {'size': 17, 'bold': True, 'space': 8}),
    ('「会社の飲み会を2回我慢すれば手に入ります」', {'size': 16, 'space': 6}),
    ('身近な支出に置き換えると、判断しやすくなる', {'size': 15, 'bold': True}),
])
conclusion(s, '金額は「いくらか」ではなく「どう見せるか」', y=6.40)

# ================================================================ 20. 分割して伝える②
s = new_slide('分割して伝える')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('②　伝える順番は 1日 → 1か月 → 総額', {'size': 24, 'bold': True})])
steps = [
    (1.98, '1日あたり', '1日500円 ≒ コーヒー一杯分ぐらい',
     '【YES取り】「旦那様、1日500円くらいコンビニで何か買われてますよね？」', LTGREEN),
    (3.55, '1か月あたり', '1か月1万円 ≒ 電気代の支払額より安い（＋停電保険）',
     '【YES取り】「奥さんも、旦那さんの飲み会2回くらい削ったっていいですよね？」\n'
     '【第三者話法】「皆さん、停電保険って言われますね」', LTYEL),
    (5.35, '総支払額', '「300万円」',
     '総支払額は最後に伝えるのがベター', LTGRAY),
]
for y, t, mid, talk, fill in steps:
    _, tf = rect(s, 0.45, y, 2.30, 1.35, fill=fill)
    lines(tf, [(t, {'size': 20, 'bold': True, 'align': PP_ALIGN.CENTER, 'anchor': None})])
    _, tf = rect(s, 2.95, y, 3.15, 1.35, fill=None, line='BFBFBF')
    lines(tf, [(mid, {'size': 14, 'bold': True})])
    _, tf = rect(s, 6.30, y, 4.10, 1.35, fill=None, line='BFBFBF')
    for i, ln in enumerate(talk.split('\n')):
        p = para(tf, first=(i == 0), space_after=4)
        run(p, ln, size=12.5)
for y in (3.38, 5.18):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.45), Inches(y - 0.06), Inches(0.30), Inches(0.26))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor.from_string('808080')
    a.line.fill.background(); a.shadow.inherit = False
conclusion(s, '小さい単位から出す。総額はいちばん最後', y=6.85)

# ================================================================ 21. ルール⑰
s = new_slide('金額訴求')
rule_head(s, 'ルール⑰：太陽光設置容量の確認', '設置容量を確認することで設置後のイメージを膨らませる')
_, tf = rect(s, 0.45, 2.48, 9.95, 1.25, fill=LTBLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, '① 屋根に載る太陽光の容量を決める', size=19, bold=True)
run(p, '　→　', size=19, bold=True, color=GRAY)
run(p, '② 蓄電池は 太陽光容量 × 1.5 で自動的に決まる', size=19, bold=True, color=RED)
p = para(tf, align=PP_ALIGN.CENTER, space_after=0)
run(p, '蓄電池の容量は据え置き。お客様に選択肢を与えない', size=14, color=GRAY)
points_box(s, 0.45, 3.90, 4.95, 2.36, [
    ('⑰は「テストクロージング」にあたる', {'bold': True, 'size': 14}),
    '即決営業の技術は「マインドセット」「言質」「反論処理」「テストクロージング」など',
    '⑯の概算費用を踏まえ、「パネル枚数」を使ってお客様の出費の限度額（≒温度感）を測る',
    ('ここで温度感を測ってから⑱の見積り提示に進む', {'bold': True, 'color': RED, 'size': 13}),
], title='考え方')
_, tf = rect(s, 5.60, 3.90, 4.95, 2.36, fill=LTYEL)
lines(tf, [
    ('実践するときのコツ', {'size': 15, 'bold': True, 'space': 8}),
    ('・容量と価格帯を一覧表にまとめておくと便利', {'size': 14, 'space': 6}),
    ('・容量によってパワコン価格も変わるので注意', {'size': 14, 'space': 6}),
    ('・補助金がもらえるエリアを除き、蓄電池は適正容量で提案する', {'size': 14, 'space': 8}),
    ('パネル枚数で話すと、お客様が設置後の家を想像しやすくなる', {'size': 14, 'bold': True}),
])
conclusion(s, '太陽光の容量が決まれば、蓄電池は自動的に決まる', y=6.40)

# ================================================================ 22. ルール⑱
s = new_slide('金額訴求')
rule_head(s, 'ルール⑱：概算見積りの提示', '月々1万円程度の予算で提案し、金額面でのハードルを大きく下げる')
box, tf = tb(s, 0.45, 2.48, 9.95, 0.32)
lines(tf, [('価格の「根源的分岐点」── 人は予算を「1, 2, 3, 5」、購入価格を「4, 8, 18（, 27）」の循環で考えます',
            {'size': 14, 'color': GRAY})])
rows = [
    ['お客様の予算', '実際に許容する価格帯'],
    [('1万円予算', {'bold': True}), ('8,000円　〜　18,000円', {'bold': True, 'color': RED})],
    ['2万円予算', '18,000円　〜　27,000円'],
    ['3万円予算', '27,000円　〜　40,000円'],
]
table(s, 0.45, 2.85, 4.95, [2.15, 2.80], rows, font_size=14, row_h=0.52, head_h=0.32)
_, tf = rect(s, 5.60, 2.85, 4.95, 1.88, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('クリアすべきこと', {'size': 16, 'bold': True, 'space': 6}),
    ('お客様の「予算の中での割安感」を得ること。', {'size': 14, 'space': 6}),
    ('月額1万円予算内（上限18,000円内）で納められる提案を入れておくことがポイントです。',
     {'size': 14, 'bold': True}),
])
_, tf = rect(s, 0.45, 4.85, 9.95, 1.40, fill=LTGRAY)
lines(tf, [
    ('精度を上げるツール', {'size': 16, 'bold': True, 'space': 6}),
    ('① SolarMaster … 航空写真を用いた診断ツール。屋根面積と方角から搭載可能容量を算出'
     '（新型・一部非対応パネル、屋根材は考慮外）', {'size': 14, 'space': 5}),
    ('② 国際航業／エネがえる … 簡易診断ツール。費用や商談フロー上①が難しい場合に', {'size': 14}),
])
conclusion(s, '月1万円台に収める提案を、必ず1本用意しておく', y=6.40)

# ================================================================ 23. 受け渡し
s = new_slide('金額訴求')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('アプローチブックを閉じて、見積りへ受け渡す', {'size': 24, 'bold': True})])
pic(s, ab_img(26), 0.45, 1.95, 4.10)
box, tf = tb(s, 0.45, 4.80, 4.10, 0.28)
lines(tf, [('アプローチブック p.26（最終ページ）', {'size': 11, 'color': GRAY})])
_, tf = rect(s, 4.90, 1.95, 5.50, 1.50, fill=LTGREEN)
lines(tf, [
    ('①　p.26を閉じる前に握る', {'size': 16, 'bold': True, 'space': 6}),
    ('【実演】「費用対効果でみれば、早いうちにセット導入するのがお得です。ここまではよろしいでしょうか？」',
     {'size': 13}),
])
_, tf = rect(s, 4.90, 3.60, 5.50, 1.50, fill=LTGREEN)
lines(tf, [
    ('②　電気代と天秤にかける', {'size': 16, 'bold': True, 'space': 6}),
    ('【実演】「いま月18,000円お支払いですよね。その範囲内に収まれば、導入に問題はないですか？」',
     {'size': 13}),
])
_, tf = rect(s, 4.90, 5.25, 5.50, 1.50, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('③　シミュレーションを開く', {'size': 16, 'bold': True, 'space': 6}),
    ('【実演】「では、〇〇様のお家の屋根で実際にどうなるか、見ていきましょう」',
     {'size': 13, 'bold': True}),
])
notes(s, '・ここで資料が切り替わる。手元でアプローチブックを閉じ、シミュレーションを出す動作まで練習させる')

# ================================================================ 24. ルール⑲
s = new_slide('金額訴求')
rule_head(s, 'ルール⑲：正式見積りの提示', '成約率最大化のために、正式見積りの提出で即決を促す')
box, tf = tb(s, 0.45, 2.46, 9.95, 0.52)
lines(tf, [('セオリーは商談時にそのまま提示。図面が無い・条件が変わる場合は別日提示もやむなし（1〜2割は仕方ありません）',
            {'size': 14, 'color': GRAY})])
rows = [
    ['', 'ポイント', 'なぜそうするか'],
    ['①', '「ケーブル単価」まで記載する',
     '「その他部材」とまとめない。「高い」と感じても「やらない」ではなく「どこか削れないか」という交渉に変わる'],
    ['②', '「お客様の声」ご協力お値引を実施する',
     'お客様が感じる「安さ」は実際の価格ではなく「基準値との差 × その為に要した労力」。'
     '単なる値引きはありがたみを感じにくい'],
    ['③', '（定価記載のある）メーカーカタログを同封する', '本商材は「比較される前提」で考える'],
]
table(s, 0.45, 3.02, 9.95, [0.50, 3.35, 6.10], rows, font_size=13, row_h=0.88, head_h=0.32)
conclusion(s, '安さとは「基準値との差 × その為に要した労力」', y=6.40)

# ================================================================ 25. ルール⑳
s = new_slide('金額訴求')
rule_head(s, 'ルール⑳：後始末の実施', '契約後のアフターフォローをしっかり行うことでお客様の不安を払拭する')
box, tf = tb(s, 0.45, 2.46, 9.95, 0.52)
lines(tf, [('できる営業マンほど「後始末」まで気を抜きません。クーリング・オフ期間（8日間）はキャンセルを見越して対応します',
            {'size': 14, 'color': GRAY})])
rows = [
    ['時期', 'やること'],
    ['契約時', '正式見積りの提示が前後することになっても、必ず一筆いただけるようにする'],
    ['契約直後〜5日以内',
     'なるべく翌日までに、お礼状（ハガキで構いません）を投函する。直筆のコメントが添えられていると尚良い'],
    ['8日間〜', '期間終了後も取り消しができないわけではない。油断せず、完工まで2週間ごとに連絡を取る'],
]
table(s, 0.45, 2.98, 9.95, [2.15, 7.80], rows, font_size=14, row_h=0.62, head_h=0.32)
_, tf = rect(s, 0.45, 5.32, 9.95, 0.95, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('取り組み事例：エコプラスワン様', {'size': 16, 'bold': True, 'space': 6}),
    ('お礼状に「訪問時の話題」と「夫婦両者のお名前」を入れる　→　クーリングオフほぼなし', {'size': 15, 'bold': True}),
])
conclusion(s, '契約時は良い買い物と思っていても、その後に不安が強くなる', y=6.40)

# ================================================================ 26. ロープレ①
s = new_slide('ロールプレイング①')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('1周目：Paragraph 03（p.23〜p.26）を通す（6分）', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 1.95, 4.95, 2.30, fill=LTGREEN)
lines(tf, [
    ('やること', {'size': 17, 'bold': True, 'space': 8}),
    ('① 2人1組（営業／お客様役）', {'size': 15, 'space': 6}),
    ('② p.23から始めて、p.26の「早いうちにセット導入がお得」で終わる', {'size': 15, 'space': 6}),
    ('③ 3分で交代', {'size': 15}),
])
_, tf = rect(s, 5.60, 1.95, 4.95, 2.30, fill=LTYEL)
lines(tf, [
    ('この周で見るポイント', {'size': 17, 'bold': True, 'space': 8}),
    ('・「機会損失」という言葉を使えたか', {'size': 15, 'space': 6}),
    ('・数字をお客様の実額（月18,000円）に置き換えたか', {'size': 15, 'space': 6}),
    ('・「財産」「投資」という言葉で締めたか', {'size': 15, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 0.45, 4.50, 9.95, 1.75, fill=LTGRAY)
lines(tf, [
    ('お客様役はこの断り文句を必ず1回入れてください', {'size': 17, 'bold': True, 'space': 8}),
    ('「もう少し待てば、蓄電池も安くなるんじゃないですか？」', {'size': 17, 'bold': True, 'color': RED, 'space': 6}),
    ('→ 営業役は物価指数（6年で12.7％上昇）と原材料の話で返す', {'size': 15}),
])
conclusion(s, 'やるかどうかの議論に戻らせない', y=6.40)

# ================================================================ 27. 顧客設定
s = new_slide('ロールプレイング')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('顧客設定（第2回・第3回と同じ田中様ご夫婦）', {'size': 24, 'bold': True})])
rows = [
    ['項目', '設定'],
    ['ご家族', '田中様ご夫婦（ご主人48歳・会社員／奥様45歳・パート）＋ 高校生・中学生のお子様2人'],
    ['お住まい', '築14年の戸建て（4LDK）。南向き切妻屋根。オール電化ではない'],
    ['導入状況', ('太陽光・蓄電池ともに未導入', {'bold': True, 'color': RED})],
    ['電気代', ('月平均 18,000円（夏場は24,000円）　→　30年で648万円', {'bold': True})],
    ['温度感', 'ご主人＝慎重。「元が取れるのか」が最大の関心。奥様＝前向き'],
    ['本日の追加設定', ('住宅ローンを月9万円返済中。「これ以上ローンを増やしたくない」という意識がある',
                  {'bold': True, 'color': RED})],
    ['断り文句', '「もう少し待てば、蓄電池も安くなるんじゃないですか？」「総額でいくらになるんですか？」'],
]
table(s, 0.45, 1.92, 9.95, [1.85, 8.10], rows, font_size=12.5, row_h=0.54, head_h=0.32)
_, tf = rect(s, 0.45, 6.35, 9.95, 0.56, fill=LTYEL)
lines(tf, [('ローン返済中の設定です。「追加の出費」ではなく「置き換え」だと伝えられるかを見ています',
            {'size': 16, 'bold': True, 'align': PP_ALIGN.CENTER})])
notes(s, '・「電力会社への捨て金を、同額のままご自宅の資産にスライドさせるだけ」という言い方を思い出させる')

# ================================================================ 28. ロープレ②
s = new_slide('ロールプレイング②')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('2周目：金額提示まで通す（6分）', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 1.95, 4.95, 2.20, fill=LTGREEN)
lines(tf, [
    ('やること', {'size': 17, 'bold': True, 'space': 8}),
    ('① 3人1組（営業／ご主人役／奥様役）', {'size': 15, 'space': 6}),
    ('② p.23から始め、アプローチブックを閉じて金額提示まで', {'size': 15, 'space': 6}),
    ('③ 3分で交代', {'size': 15}),
])
_, tf = rect(s, 5.60, 1.95, 4.95, 2.20, fill=LTYEL)
lines(tf, [
    ('この周で見るポイント', {'size': 17, 'bold': True, 'space': 8}),
    ('・総額を先に言ってしまっていないか', {'size': 15, 'bold': True, 'color': RED, 'space': 6}),
    ('・1日 → 1か月 → 総額 の順を守れたか', {'size': 15, 'space': 6}),
    ('・最後までセット提案を崩していないか', {'size': 15}),
])
rows = [
    ['進行', '使うもの', '時間の目安'],
    ['時期訴求（⑬⑭⑮）', 'アプローチブック p.23〜p.26', '1分30秒'],
    ['受け渡し', 'アプローチブックを閉じる', '15秒'],
    [('金額訴求（⑯⑰⑱）', {'bold': True, 'color': RED}), ('シミュレーション・概算見積', {'bold': True}),
     ('1分15秒', {'bold': True, 'color': RED})],
]
table(s, 0.45, 4.40, 9.95, [3.55, 3.20, 3.20], rows, font_size=14, row_h=0.58, head_h=0.34)
conclusion(s, '総額を口にするのは最後の15秒だけ', y=6.35)

# ================================================================ 29. チェックシート
s = new_slide('チェックシート')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.42)
lines(tf, [('ロープレのあと、自分で○△×をつけてください', {'size': 22, 'bold': True})])
rows = [
    ['ルール', '該当ページ／道具', '狙った台詞を引き出せたか', '結論を言えたか'],
    ['⑬導入に向けた考え方の転換', 'p.23', '今設置した方がお得なのですね', '設置を遅くするほどメリットが減る'],
    ['⑭物価指数とコストトレンド', 'p.24', '待っても上がる一方なのですね', '待っても安くなるとは限らない'],
    ['⑮投資対効果の視点', 'p.25・26', '長期的には財産になるのですね', '持ち家と同じく長期的な財産になる'],
    ['⑯概算の機器導入費用', 'シミュレーション', '電気代の範囲なら問題ない', '月々の額だけを握れた'],
    ['⑰太陽光設置容量の確認', '図面・診断ツール', '（枚数のイメージが湧いた反応）', '蓄電池は容量×1.5で提案した'],
    ['⑱概算見積りの提示', '概算見積', '月1万円台なら、と前向きな答え', '−'],
    [('伝える順番', {'bold': True}), '−', ('1日 → 1か月 → 総額 の順を守れた', {'bold': True}), '−'],
    [('商談運営', {'bold': True}), '−', '発言比率6：4／ご主人2：奥様2', 'セット提案を崩さなかった'],
]
table(s, 0.30, 1.90, 10.25, [2.55, 1.95, 3.15, 2.60], rows, font_size=11, row_h=0.50, head_h=0.32)
_, tf = rect(s, 0.30, 6.25, 10.25, 0.58, fill=LTGREEN)
lines(tf, [('「結論を言えたか」の欄は、資料を閉じて声に出して確認します（暗唱チェック）',
            {'size': 17, 'bold': True, 'align': PP_ALIGN.CENTER})])

# ================================================================ 30. まとめ・宿題
s = new_slide('本日のまとめ')
_, tf = rect(s, 0.45, 1.45, 9.95, 2.30, fill=LTGREEN)
lines(tf, [
    ('本日持ち帰ること', {'size': 18, 'bold': True, 'space': 10}),
    ('① お客様の頭を「導入するか」から「いつ導入するのが賢いか」へ切り替える', {'size': 17, 'space': 8}),
    ('② 待っても安くならない。物価は6年で12.7％上がり、原材料の構造も変わらない', {'size': 17, 'space': 8}),
    ('③ 電気代は「賃貸の家賃」。創蓄は払い終われば「持ち家」になる', {'size': 17, 'space': 8}),
    ('④ 金額は 1日 → 1か月 → 総額 の順。総額は最後', {'size': 17, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 0.45, 4.00, 9.95, 2.30, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('宿題（次回までに）', {'size': 18, 'bold': True, 'space': 10}),
    ('① 自社の「容量と価格帯の一覧表」をつくり、月1万円台の提案を1本用意してくること', {'size': 17, 'space': 8}),
    ('② 20のルール全部を、アプローチブックのページと対応させて言えるようにしてくること', {'size': 17, 'space': 8}),
    ('③ 実商談を1件、総額を口にしたタイミングを記録してくること', {'size': 17}),
])
_, tf = rect(s, 0.45, 6.55, 9.95, 0.42, fill=None, line=RED, lw=1.5)
lines(tf, [('次回・第5回は「総仕上げ通しロープレ」／アプローチブック p.2〜p.26 を頭から通します',
            {'size': 16, 'bold': True, 'align': PP_ALIGN.CENTER})])
notes(s, '・次回は解説をほとんどしない。20のルールを全部持ってくるよう強く伝える')

# ---------------------------------------------------------------- 並べ替え
ids = list(sldIdLst)
colophon = [x for x in ids if x is colophon_id][0]
sldIdLst.remove(colophon)
sldIdLst.append(colophon)
page_no[0] += 1

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print('saved:', OUT, '/ slides:', len(prs.slides.__iter__.__self__._sldIdLst))
