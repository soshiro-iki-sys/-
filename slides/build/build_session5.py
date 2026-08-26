# -*- coding: utf-8 -*-
"""第5回「総仕上げ通しロープレ」研修資料ビルダー

使い方は slides/build/build_session2.py と同じ。
ピラミッド図（<ASSETS>/pyramid.png）を用意してから実行する。
（第5回はアプローチブックのページ画像を使わない）

    TRAINING_ASSETS=/path/to/assets python slides/build/build_session5.py

書式は docs/研修資料フォーマット仕様.md、内容は docs/知識_*.md を参照。
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
MSO_SHAPE_TYPE_AUTO = MSO_SHAPE_TYPE.AUTO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.opc.packuri import PackURI

REPO   = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
ASSETS = os.environ.get('TRAINING_ASSETS', '/tmp/training-assets')
AB     = os.path.join(ASSETS, 'ab')
TPL    = os.path.join(REPO, 'templates', '研修資料_フォーマット見本.pptx')
OUT    = os.path.join(REPO, 'slides', '2026営業研修_第5回_総仕上げ通しロープレ.pptx')

GOTHIC='HGPｺﾞｼｯｸE'; MINCHO='HGP明朝E'; UD='BIZ UDPゴシック'; YU='游ゴシック'
RED='FF0000'; BLACK='000000'; WHITE='FFFFFF'; GREEN='9BBB59'; YELLOW='FFFF00'
GRAY='595959'; BLUE='0070C0'; LTGREEN='EAF1DD'; LTGRAY='F2F2F2'; LTYEL='FFF2CC'
LTBLUE='DEEBF7'; ORANGE='F79646'

# ---------------------------------------------------------------- helpers
def _rpr_font(run, name):
    rPr = run._r.get_or_add_rPr()
    rPr.get_or_add_latin().set('typeface', name)
    latin = rPr.find(qn('a:latin'))
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = parse_xml('<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
        latin.addnext(ea)
    ea.set('typeface', name)

def _highlight(run, color):
    rPr = run._r.get_or_add_rPr()
    hl = parse_xml(
        '<a:highlight xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:srgbClr val="%s"/></a:highlight>' % color)
    latin = rPr.find(qn('a:latin'))
    if latin is not None:
        latin.addprevious(hl)
    else:
        rPr.append(hl)

def run(p, text, size=16, bold=False, color=BLACK, font=GOTHIC, hl=None):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)
    if hl: _highlight(r, hl)
    _rpr_font(r, font)
    return r

def tb(slide, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP, margin=0.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return box, tf

def para(tf, first=False, space_after=4, align=PP_ALIGN.LEFT, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after); p.alignment = align
    if line: p.line_spacing = line
    return p

def lines(tf, items, size=16, color=BLACK, bold=False, space=5, font=GOTHIC, line=None, align=PP_ALIGN.LEFT):
    """items: str | (str,dict)"""
    for i, it in enumerate(items):
        opts = {}
        if isinstance(it, tuple): it, opts = it
        p = para(tf, first=(i == 0), space_after=opts.pop('space', space),
                 align=opts.pop('align', align), line=opts.pop('line', line))
        run(p, it, size=opts.pop('size', size), bold=opts.pop('bold', bold),
            color=opts.pop('color', color), font=opts.pop('font', font),
            hl=opts.pop('hl', None))

def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=1.0,
         anchor=MSO_ANCHOR.TOP):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    else: s.fill.background()
    if line:
        s.line.color.rgb = RGBColor.from_string(line); s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    return s, tf

def pic(slide, path, x, y, w, border=True):
    p = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    if border:
        p.line.color.rgb = RGBColor.from_string('808080'); p.line.width = Pt(1)
    return p

def table(slide, x, y, w, col_w, rows, font_size=11, header=True,
          header_fill='000000', header_color=WHITE, row_h=0.30, head_h=0.30):
    n_r, n_c = len(rows), len(col_w)
    shp = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w),
                                 Inches(head_h + row_h * (n_r - 1)))
    t = shp.table
    t.first_row = header; t.horz_banding = False
    for i, cw in enumerate(col_w): t.columns[i].width = Inches(cw)
    t.rows[0].height = Inches(head_h)
    for i in range(1, n_r): t.rows[i].height = Inches(row_h)
    for ri, rowdata in enumerate(rows):
        for ci, cell in enumerate(rowdata):
            txt, opts = (cell, {}) if isinstance(cell, str) else cell
            c = t.cell(ri, ci)
            c.margin_left = c.margin_right = Inches(0.05)
            c.margin_top = c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = opts.get('fill', header_fill if (header and ri == 0) else None)
            if fill: c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string(fill)
            else: c.fill.background()
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = opts.get('align',
                PP_ALIGN.CENTER if (header and ri == 0) else PP_ALIGN.LEFT)
            run(p, txt, size=opts.get('size', font_size),
                bold=opts.get('bold', header and ri == 0),
                color=opts.get('color', header_color if (header and ri == 0) else BLACK),
                font=opts.get('font', YU))
    return t

# ---------------------------------------------------------------- slide frame
prs = Presentation(TPL)
BLANK = None
for l in prs.slide_masters[0].slide_layouts:
    if l.name == '白紙': BLANK = l
assert BLANK is not None

sldIdLst = prs.slides._sldIdLst
orig_ids = list(sldIdLst)
cover_id, colophon_id = orig_ids[0], orig_ids[6]
for sid in orig_ids[1:6]:                      # drop template samples 2-6
    rId = sid.get(qn('r:id')); sldIdLst.remove(sid)
    prs.part.drop_rel(rId)
# python-pptx names new slides by len(sldIdLst)+1 -> move colophon out of that range
prs.part.related_part(colophon_id.get(qn('r:id'))).partname = PackURI('/ppt/slides/slide900.xml')

page_no = [0]

def new_slide(heading=None):
    s = prs.slides.add_slide(BLANK)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    if heading is not None:
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.157), Inches(0.757), Inches(0.140), Inches(0.512))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor.from_string(BLACK)
        bar.line.color.rgb = RGBColor.from_string('D9D9D9'); bar.line.width = Pt(0.75)
        bar.shadow.inherit = False
        box, tf = tb(s, 0.348, 0.736, 9.9, 0.572)
        lines(tf, [heading], size=28, font=MINCHO)
    page_no[0] += 1
    box, tf = tb(s, 8.31, 7.18, 2.30, 0.33)
    lines(tf, [str(page_no[0])], size=12, color=WHITE, font=GOTHIC, align=PP_ALIGN.RIGHT)
    return s

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def ab_img(n):
    return os.path.join(AB, 'ab-%02d.jpg' % n)

# ---------------------------------------------------------------- 表紙
cover = prs.slides[0]
for shp in cover.shapes:
    if shp.shape_type == MSO_SHAPE_TYPE_AUTO and shp.name == '正方形/長方形 4':
        sp = shp._element.spPr
        fill = sp.find(qn('a:solidFill'))
        for ch in list(fill): fill.remove(ch)
        fill.append(parse_xml(
            '<a:srgbClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'val="8064A2"><a:alpha val="80000"/></a:srgbClr>'))
    if shp.has_text_frame and shp.name == 'テキスト ボックス 5':
        tf = shp.text_frame; tf.clear()
        lines(tf, [
            ('2026年度　営業研修', {'size': 32, 'bold': True, 'color': WHITE, 'font': UD, 'space': 8}),
            ('第5回　総仕上げ', {'size': 52, 'bold': True, 'color': WHITE, 'font': UD, 'space': 2}),
            ('通しロープレ', {'size': 52, 'bold': True, 'color': WHITE, 'font': UD, 'space': 10}),
            ('〜太陽光＋蓄電池セット販売〜', {'size': 24, 'bold': True, 'color': WHITE, 'font': UD}),
        ])
    if shp.has_text_frame and shp.name == 'テキスト ボックス 6':
        tf = shp.text_frame; tf.clear()
        lines(tf, [('2026年9月下旬　＠船井総研', {'size': 14, 'color': WHITE, 'font': UD})],
              align=PP_ALIGN.RIGHT)
page_no[0] = 1

# ---------------------------------------------------------------- rule helpers
def rule_head(s, no_title, point):
    box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
    lines(tf, [(no_title, {'size': 24, 'bold': True})])
    _, tf = rect(s, 0.45, 1.88, 9.95, 0.46, fill=None, line=RED, lw=1.25)
    p = para(tf, first=True, align=PP_ALIGN.LEFT)
    run(p, 'POINT　', size=15, bold=True, color=RED)
    run(p, point, size=16, bold=True)

def goal_box(s, x, y, w, h, quote):
    _, tf = rect(s, x, y, w, h, fill=LTBLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    lines(tf, [
        ('ゴール（お客様の反応）', {'size': 13, 'bold': True, 'color': BLUE, 'space': 6}),
    ])
    for q in quote:
        p = para(tf, space_after=3)
        run(p, '「%s」' % q, size=16, bold=True)

def points_box(s, x, y, w, h, items, title='説明のポイント'):
    _, tf = rect(s, x, y, w, h, fill=LTGRAY)
    lines(tf, [(title, {'size': 13, 'bold': True, 'color': GRAY, 'space': 6})])
    for it in items:
        opts = {}
        if isinstance(it, tuple): it, opts = it
        p = para(tf, space_after=opts.pop('space', 6))
        run(p, it, size=opts.pop('size', 12.5), bold=opts.pop('bold', False),
            color=opts.pop('color', BLACK))

def conclusion(s, text, y=6.36):
    _, tf = rect(s, 0.45, y, 9.95, 0.66, fill=LTGREEN)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, 'このページの結論　', size=12, bold=True, color=GRAY)
    run(p, text, size=15, bold=True, hl=YELLOW)

def talk_slide(s, no_title, talks, conc=None, note=None):
    box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
    lines(tf, [(no_title, {'size': 24, 'bold': True})])
    y = 1.92
    for label, body, h in talks:
        _, tf = rect(s, 0.45, y, 8.05, h, fill=LTGRAY, anchor=MSO_ANCHOR.MIDDLE)
        for i, ln in enumerate(body):
            p = para(tf, first=(i == 0), space_after=3)
            run(p, ln, size=14)
        _, tf2 = rect(s, 8.68, y, 1.72, h, fill=ORANGE, anchor=MSO_ANCHOR.MIDDLE)
        lines(tf2, [(label, {'size': 13, 'bold': True, 'color': WHITE,
                             'align': PP_ALIGN.CENTER})])
        y += h + 0.16
    if conc: conclusion(s, conc)
    if note: notes(s, note)

# ================================================================ 2. 本日のゴール
s = new_slide('本日のゴール')
box, tf = tb(s, 0.45, 1.45, 9.9, 0.62)
lines(tf, [('アプローチブック26ページを、止まらずに1人で通せる',
            {'size': 25, 'bold': True, 'color': RED})])
_, tf = rect(s, 0.45, 2.35, 9.95, 3.05, fill=LTGREEN)
lines(tf, [
    ('研修が終わったときの「できる状態」', {'size': 18, 'bold': True, 'space': 10}),
    ('①　p.2からp.26まで、資料を見ながら止まらずに話せる', {'size': 19, 'space': 9}),
    ('②　20のルールを、アプローチブックのページと対応させて言える', {'size': 19, 'space': 9}),
    ('③　よくある断り文句5つに、その場で切り返せる', {'size': 19, 'space': 9}),
    ('④　最後までセット提案を崩さずに、仮契約まで運べる', {'size': 19}),
])
_, tf = rect(s, 0.45, 5.65, 9.95, 0.85, fill=None, line=RED, lw=1.5)
lines(tf, [('本日は解説をほとんどしません。時間の大半をロープレに使います',
            {'size': 22, 'bold': True, 'align': PP_ALIGN.CENTER})])
notes(s, '・本日は「教わる場」ではなく「試す場」だと最初に宣言する\n'
         '・全員が必ず1回は最後まで通す。途中で止まっても止めない')

# ================================================================ 3. アジェンダ
s = new_slide('本日のアジェンダ')
rows = [
    ['時間', '内容', '扱うもの'],
    ['0:00-0:10', '全5回の振り返り', '第1回〜第4回'],
    ['0:10-0:25', ('20のルール ↔ アプローチブックの総復習', {'bold': True}), '対応表／商談の設計図'],
    ['0:25-0:40', '講師による通しデモ', 'アプローチブック p.2〜p.26'],
    ['0:40-0:50', '休憩', '−'],
    ['0:50-1:15', ('ロールプレイング①　前半', {'bold': True}), 'p.2〜p.17（ルール①〜⑨）'],
    ['1:15-1:40', ('ロールプレイング②　後半', {'bold': True}), 'p.18〜p.26＋見積り（ルール⑩〜⑳）'],
    ['1:40-1:55', '個別フィードバック', '評価シート'],
    ['1:55-2:00', '全5回の総括／明日からの行動計画', '−'],
]
table(s, 0.45, 1.62, 9.95, [1.55, 4.3, 4.1], rows, font_size=14, row_h=0.50, head_h=0.38)
notes(s, '・解説パートが延びたら削る。ロープレの50分は死守する')

# ================================================================ 4. 全5回の振り返り
s = new_slide('全5回でやってきたこと')
rows = [
    ['回', 'テーマ', '扱ったルール', 'アプローチブック', '核になる考え方'],
    ['第1回', '業界商材の知識講座', '−', '−', '知識は「自信」の材料'],
    ['第2回', '聞く姿勢づくり＆必要性訴求①', '①〜⑨', 'p.2〜p.17',
     '先に聞くから、聞いてもらえる／下げてから上げる'],
    ['第3回', '災害対策・商品説明・経済メリット', '⑩〜⑫', 'p.18〜p.21',
     '災害は3分で抜ける／FABのBは「お客様が主語」／経済メリットは2通り'],
    ['第4回', '時期・料金訴求＆ロープレ', '⑬〜⑳', 'p.23〜p.26＋見積り',
     '「いつ導入するのが賢いか」へ転換／総額は最後'],
    [('第5回', {'bold': True, 'color': RED}), ('総仕上げ通しロープレ', {'bold': True, 'color': RED}),
     ('①〜⑳', {'bold': True, 'color': RED}), ('p.2〜p.26 通し', {'bold': True, 'color': RED}),
     ('全部つなげて、1人で通す', {'bold': True, 'color': RED})],
]
table(s, 0.30, 1.55, 10.25, [0.80, 2.55, 1.15, 1.90, 3.85], rows, font_size=12, row_h=0.72, head_h=0.34)
_, tf = rect(s, 0.30, 5.95, 10.25, 0.90, fill=LTYEL)
lines(tf, [
    ('ここまでで、商談で使う材料はすべて揃っています。', {'size': 18, 'bold': True, 'space': 5}),
    ('あとは「順番どおりに、止まらずに出せるか」だけです。', {'size': 18, 'bold': True, 'color': RED}),
])

# ================================================================ 5. 4ステップ総復習
s = new_slide('総復習：契約までの4ステップ')
pic(s, os.path.join(ASSETS, 'pyramid.png'), 0.55, 1.55, 6.0, border=False)
_, tf = rect(s, 6.95, 1.60, 3.45, 4.55, fill=LTGRAY)
lines(tf, [
    ('積む順番を絶対に飛ばさない', {'size': 17, 'bold': True, 'space': 10}),
    ('STEP1　聴く姿勢（①〜⑥）', {'size': 15, 'bold': True, 'space': 3}),
    ('この営業の話を聞きたいと思わせる', {'size': 13, 'color': GRAY, 'space': 9}),
    ('STEP2　必要性訴求（⑦〜⑫）', {'size': 15, 'bold': True, 'space': 3}),
    ('太陽光・蓄電池は必要だと思ってもらう', {'size': 13, 'color': GRAY, 'space': 9}),
    ('STEP3　時期訴求（⑬〜⑮）', {'size': 15, 'bold': True, 'space': 3}),
    ('今導入するのがお得だと認識させる', {'size': 13, 'color': GRAY, 'space': 9}),
    ('STEP4　金額訴求（⑯〜⑳）', {'size': 15, 'bold': True, 'space': 3}),
    ('自分でも払えそうな金額だと理解させる', {'size': 13, 'color': GRAY, 'space': 10}),
    ('土台が無いまま金額の話をすると、価格だけで判断されて終わります。',
     {'size': 13, 'bold': True, 'color': RED}),
])

# ================================================================ 6. 20のルール ↔ ページ対応表（保存版）
s = new_slide('保存版：20のルール ↔ アプローチブック')
box, tf = tb(s, 0.45, 1.30, 9.95, 0.32)
lines(tf, [('この1枚を手元に置いてロープレしてください', {'size': 16, 'bold': True})])
R = RED
rows = [
    ['ルール', 'ページ', '狙う反応（お客様の台詞）', '使う話法'],
    ['①商談テーブルの確認', '（商談前）', '検針票・図面／夫婦両主権者／1時間を押さえる', '−'],
    ['②自社の概要・訪問趣旨の理解', 'p.2', '商談の構成が分かる', '物語効果・ラポール形成'],
    ['③未導入理由の確認', '（資料なし）', '導入しなかった理由が3つ出る', 'オープンクエスチョン×3'],
    ['④前向きな検討に対する言質', '（資料なし）', '解消できれば導入する予定です', '一貫性の法則'],
    ['⑤太陽光が当たり前であることの理解', 'p.3・p.4', '義務化・今後必須の流れになるの!?', '包み込みの法則①'],
    ['⑥直近の国策動向のおさらい', 'p.5〜p.7', 'じゃあつけた方がいいかもねえ', '包み込みの法則②・二者択一話法'],
    ['⑦直近数か年の電気代推移の理解', 'p.9・10・12', '10年で5.5万円も上がっているのですね', '−'],
    ['⑧再エネ賦課金の構造の理解', 'p.11・13〜16', '賦課金が上昇して電気代が上がるのですね', 'YES取り'],
    ['⑨長期で支払う電気代総額の理解', 'p.17', '30年でそんなにも払わないといけないのですね', '第三者話法・損失回避'],
    ['⑩非常時に備える意義の理解', 'p.18・p.19', 'いつ災害が起こるか分からないですね', '−'],
    ['⑪非常時対策の効果の確認', ('（ページなし）', {'color': R}), '停電時も電気が使えた方が安心です', '第三者話法'],
    ['⑫太陽光・蓄電池の使い方の理解', 'p.20・p.21', '買うより自給自足した方が良いですね', '−'],
    ['⑬導入に向けた考え方の転換', 'p.23', '今設置した方が将来的にお得なのですね', 'プロスペクト理論（機会損失）'],
    ['⑭物価指数とコストトレンドの説明', 'p.24', '待っても上がる一方なのですね', '−'],
    ['⑮投資対効果の視点', 'p.25・p.26', '長期的には財産になるのですね', '二者択一話法'],
    ['⑯概算の機器導入費用のおさらい', 'シミュレーション', '電気代の範囲なら問題ないです', '−'],
    ['⑰太陽光設置容量の確認', '図面・診断ツール', '（枚数のイメージが湧いた反応）', 'テストクロージング'],
    ['⑱概算見積りの提示', '概算見積', '月1万円台なら、と前向きな答え', 'フレーミング効果'],
    ['⑲正式見積りの提示', '正式見積・カタログ', '（その場で即決）', '−'],
    ['⑳後始末の実施', '契約書・お礼状', '（クーリングオフさせない）', '−'],
]
table(s, 0.30, 1.66, 10.25, [2.90, 1.50, 3.60, 2.25], rows, font_size=9, row_h=0.255, head_h=0.28)

# ================================================================ 7. 商談の設計図
s = new_slide('商談の設計図')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('60分の商談を、どこに何分使うか', {'size': 24, 'bold': True})])
rows = [
    ['パート', 'ルール', 'アプローチブック', '時間', '出てはいけない話'],
    ['聴く姿勢作り', '①〜⑥', 'p.2〜p.7', '15分', '金額・製品スペック'],
    ['必要性訴求', '⑦〜⑫', 'p.9〜p.21', '20分', '金額'],
    [('　うち災害（⑩⑪）', {'color': RED}), ('⑩⑪', {'color': RED}), ('p.18・p.19', {'color': RED}),
     ('3分以内', {'bold': True, 'color': RED}), ('お客様の家族を想像させる話', {'color': RED})],
    ['時期訴求', '⑬〜⑮', 'p.23〜p.26', '10分', '総額'],
    ['金額訴求', '⑯〜⑱', 'シミュレーション・概算見積', '12分', '−'],
    ['後始末', '⑲⑳', '正式見積・契約書', '3分', '−'],
]
table(s, 0.45, 1.95, 9.95, [2.30, 1.10, 3.10, 1.05, 2.40], rows, font_size=12, row_h=0.50, head_h=0.32)
_, tf = rect(s, 0.45, 5.20, 4.95, 1.62, fill=LTYEL)
lines(tf, [
    ('発言比率とラリー数', {'size': 17, 'bold': True, 'space': 8}),
    ('・営業マン：お客様 ＝ 6：4', {'size': 15, 'space': 6}),
    ('・お客様側の4は ご主人2：奥様2', {'size': 15, 'space': 6}),
    ('・ラリー数は30回（「うん」を除く）', {'size': 15}),
])
_, tf = rect(s, 5.60, 5.20, 4.95, 1.62, fill=LTGREEN)
lines(tf, [
    ('順番を守る理由', {'size': 17, 'bold': True, 'space': 8}),
    ('価値が伝わる前に価格を出すと、', {'size': 15, 'space': 4}),
    ('お客様は価格だけで判断します。', {'size': 15, 'space': 8}),
    ('聞かれても「この後のお話で一通り出てきますので」でかわす。', {'size': 15, 'bold': True}),
])

# ================================================================ 8. 講師デモの見どころ
s = new_slide('講師による通しデモ')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('p.2〜p.26を15分で通します。ここを見てください', {'size': 24, 'bold': True})])
items = [
    ('①', 'ページをめくる速さ', '1ページに何秒かけているか。止まっているページはどこか'),
    ('②', '質問を投げる位置', 'どこで話を振り、どこで黙っているか'),
    ('③', '結論の一文の言い方', '各ページの緑帯を、どのタイミングで口にしているか'),
    ('④', '災害パートの長さ', '時計を見て測ってください。3分を超えていないか'),
    ('⑤', 'セット提案の一貫性', '「蓄電池が」ではなく「太陽光と蓄電池で」と言い続けているか'),
]
y = 1.95
for n, t, d in items:
    _, tf = rect(s, 0.45, y, 0.60, 0.80, fill=LTGREEN)
    lines(tf, [(n, {'size': 22, 'bold': True, 'align': PP_ALIGN.CENTER})])
    _, tf = rect(s, 1.20, y, 9.20, 0.80, fill=LTGRAY)
    lines(tf, [
        (t, {'size': 17, 'bold': True, 'space': 3}),
        (d, {'size': 14, 'color': GRAY}),
    ])
    y += 0.92
_, tf = rect(s, 0.45, 6.60, 9.95, 0.42, fill=LTYEL)
lines(tf, [('デモ中に気づいたことは、手元の評価シートにメモしてください',
            {'size': 16, 'bold': True, 'align': PP_ALIGN.CENTER})])
notes(s, '・講師は完璧にやらない。あえて1か所詰まってみせ、そこからの立て直し方も見せる')

# ================================================================ 9. 休憩
s = new_slide(None)
_, tf = rect(s, 2.40, 2.85, 6.05, 1.80, fill=LTGREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
lines(tf, [
    ('休憩　10分', {'size': 44, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 10}),
    ('後半はロールプレイングです', {'size': 18, 'align': PP_ALIGN.CENTER}),
])

# ================================================================ 10. ロープレのルール
s = new_slide('ロールプレイングの進め方')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('3人1組。全員が必ず1回は営業役をやります', {'size': 24, 'bold': True})])
rows = [
    ['役割', 'やること'],
    ['営業役', 'アプローチブックを実際にめくりながら進める。詰まっても止まらない'],
    ['お客様役（ご主人）', '慎重派。「元が取れるのか」を必ず一度は聞く。指定の断り文句を1回入れる'],
    ['お客様役（奥様）', '前向き。ただし黙り込まない。ご主人との温度差を出す'],
    ['（観察役）', '4人組の場合。評価シートを見ながら、発言比率とラリー数を数える'],
]
table(s, 0.45, 1.95, 9.95, [2.45, 7.50], rows, font_size=13.5, row_h=0.62, head_h=0.34)
_, tf = rect(s, 0.45, 4.60, 4.95, 2.20, fill=LTGREEN)
lines(tf, [
    ('進め方', {'size': 17, 'bold': True, 'space': 8}),
    ('① 合図で開始。時間は講師が計ります', {'size': 15, 'space': 6}),
    ('② 詰まっても止めない。飛ばして先に進む', {'size': 15, 'space': 6}),
    ('③ 終わったら、すぐ評価シートに○△×', {'size': 15, 'space': 6}),
    ('④ 交代。全員が営業役をやるまで回す', {'size': 15}),
])
_, tf = rect(s, 5.60, 4.60, 4.95, 2.20, fill='FCE4E4', line=RED)
lines(tf, [
    ('やってはいけないこと', {'size': 17, 'bold': True, 'color': RED, 'space': 8}),
    ('✕ 途中で「今のところもう一回」と巻き戻す', {'size': 15, 'space': 6}),
    ('✕ お客様役が意地悪をしすぎる', {'size': 15, 'space': 6}),
    ('✕ 資料を読み上げるだけで質問を投げない', {'size': 15, 'space': 6}),
    ('✕ 蓄電池だけ／太陽光だけの提案になる', {'size': 15, 'bold': True}),
])

# ================================================================ 11. 顧客設定A
s = new_slide('顧客設定A')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('田中様ご夫婦（第2〜4回と同じ・標準ケース）', {'size': 24, 'bold': True})])
rows = [
    ['項目', '設定'],
    ['ご家族', '田中様ご夫婦（ご主人48歳・会社員／奥様45歳・パート）＋ 高校生・中学生のお子様2人'],
    ['お住まい', '築14年の戸建て（4LDK）。南向き切妻屋根。オール電化ではない'],
    ['導入状況', ('太陽光・蓄電池ともに未導入', {'bold': True, 'color': RED})],
    ['電気代', '月平均 18,000円（夏場は24,000円）。30年で648万円'],
    ['温度感', 'ご主人＝慎重。「元が取れるのか」が最大の関心。奥様＝前向き'],
    ['その他', '住宅ローンを月9万円返済中。停電の経験はない'],
    ['断り文句', '「以前も訪問販売が来て、高かったからやめた」／「総額でいくらになるんですか？」'],
]
table(s, 0.45, 1.92, 9.95, [1.60, 8.35], rows, font_size=12.5, row_h=0.54, head_h=0.32)
_, tf = rect(s, 0.45, 6.20, 9.95, 0.68, fill=LTYEL)
lines(tf, [
    ('ロールプレイング①（前半 p.2〜p.17）で使います', {'size': 17, 'bold': True,
                                       'align': PP_ALIGN.CENTER, 'space': 4}),
    ('太陽光も蓄電池も未導入＝セット販売のケースです', {'size': 14, 'align': PP_ALIGN.CENTER}),
])

# ================================================================ 12. 顧客設定B
s = new_slide('顧客設定B')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('佐藤様ご夫婦（即決を渋るケース）', {'size': 24, 'bold': True})])
rows = [
    ['項目', '設定'],
    ['ご家族', '佐藤様ご夫婦（ご主人52歳・自営業／奥様50歳・会社員）＋ 大学生のお子様1人'],
    ['お住まい', '築8年の戸建て（3LDK）。東西面の寄棟屋根。オール電化'],
    ['導入状況', ('太陽光・蓄電池ともに未導入', {'bold': True, 'color': RED})],
    ['電気代', '月平均 24,000円（オール電化のため冬が高い）。30年で864万円'],
    ['温度感', ('ご主人＝比較検討したがる。奥様＝決められない。二人とも「その場では決めない」主義',
            {'bold': True, 'color': RED})],
    ['その他', '知人が別会社で見積りを取っている。補助金の情報をネットで見ている'],
    ['断り文句', ('「補助金が出るまで待ちたい」／「他の会社の話も聞いてみたい」／「妻と相談して決めます」',
             {'bold': True, 'color': RED})],
]
table(s, 0.45, 1.92, 9.95, [1.60, 8.35], rows, font_size=12.5, row_h=0.54, head_h=0.32)
_, tf = rect(s, 0.45, 6.20, 9.95, 0.68, fill=LTYEL)
lines(tf, [
    ('ロールプレイング②（後半 p.18〜p.26＋見積り）で使います', {'size': 17, 'bold': True,
                                              'align': PP_ALIGN.CENTER, 'space': 4}),
    ('断り文句3つを必ず出してください。切り返しの練習が目的です', {'size': 14, 'align': PP_ALIGN.CENTER}),
])

# ================================================================ 13. 評価シート
s = new_slide('評価シート')
box, tf = tb(s, 0.45, 1.30, 9.95, 0.32)
lines(tf, [('ロープレのたびに○△×をつけてください（①＝前半／②＝後半）', {'size': 16, 'bold': True})])
LEFT = [
    ['ルール', 'ページ', '①', '②'],
    ['①商談テーブルの確認', '（商談前）', ' ', ' '],
    ['②自社の概要・訪問趣旨', 'p.2', ' ', ' '],
    ['③未導入理由の確認', '（資料なし）', ' ', ' '],
    ['④前向きな検討への言質', '（資料なし）', ' ', ' '],
    ['⑤太陽光が当たり前であること', 'p.3・p.4', ' ', ' '],
    ['⑥直近の国策動向のおさらい', 'p.5〜p.7', ' ', ' '],
    ['⑦電気代推移の理解', 'p.9・10・12', ' ', ' '],
    ['⑧再エネ賦課金の構造の理解', 'p.11・13〜16', ' ', ' '],
    ['⑨電気代総額の理解', 'p.17', ' ', ' '],
    ['⑩非常時に備える意義', 'p.18・p.19', ' ', ' '],
]
RIGHT = [
    ['ルール', 'ページ／道具', '①', '②'],
    ['⑪非常時対策の効果の確認', '（ページなし）', ' ', ' '],
    ['⑫太陽光・蓄電池の使い方', 'p.20・p.21', ' ', ' '],
    ['⑬導入に向けた考え方の転換', 'p.23', ' ', ' '],
    ['⑭物価指数とコストトレンド', 'p.24', ' ', ' '],
    ['⑮投資対効果の視点', 'p.25・p.26', ' ', ' '],
    ['⑯概算の機器導入費用', 'シミュレーション', ' ', ' '],
    ['⑰太陽光設置容量の確認', '図面・診断', ' ', ' '],
    ['⑱概算見積りの提示', '概算見積', ' ', ' '],
    ['⑲⑳正式見積り・後始末', '見積・お礼状', ' ', ' '],
    [('商談運営（6：4／30ラリー）', {'bold': True}), '−', ' ', ' '],
]
table(s, 0.30, 1.70, 5.00, [2.60, 1.25, 0.57, 0.58], LEFT, font_size=10, row_h=0.30, head_h=0.30)
table(s, 5.55, 1.70, 5.00, [2.60, 1.25, 0.57, 0.58], RIGHT, font_size=10, row_h=0.30, head_h=0.30)
_, tf = rect(s, 0.30, 5.35, 10.25, 0.62, fill='FCE4E4', line=RED)
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, '最後までセット提案を崩さなかったか　', size=18, bold=True, color=RED)
run(p, '①　　　　②', size=18, bold=True)
_, tf = rect(s, 0.30, 6.12, 10.25, 0.75, fill=None, line='BFBFBF')
lines(tf, [('気づいたこと（自由記入）', {'size': 13, 'bold': True, 'color': GRAY})])

# ================================================================ 14. フィードバックの型
s = new_slide('フィードバックの型')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('お互いにフィードバックするときの言い方', {'size': 24, 'bold': True})])
box, tf = tb(s, 0.45, 1.88, 9.95, 0.35)
lines(tf, [('褒め方は商談でもそのまま使えます。「質の高い褒め」を「数」多く言うのが基本です',
            {'size': 15, 'color': GRAY})])
_, tf = rect(s, 0.45, 2.32, 4.95, 1.95, fill=LTGREEN)
lines(tf, [
    ('型①　比較して褒める', {'size': 18, 'bold': True, 'space': 8}),
    ('✕「良かったです」', {'size': 15, 'space': 5}),
    ('○「①では災害に5分使っていましたが、②では2分半に収まっていました」', {'size': 15, 'bold': True}),
])
_, tf = rect(s, 5.60, 2.32, 4.95, 1.95, fill=LTGREEN)
lines(tf, [
    ('型②　人を褒める', {'size': 18, 'bold': True, 'space': 8}),
    ('✕「切り返しが良かったです」', {'size': 15, 'space': 5}),
    ('○「あの場面でとっさに事例を出せた○○さんが凄いですね」', {'size': 15, 'bold': True}),
])
_, tf = rect(s, 0.45, 4.50, 9.95, 1.75, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('直す点は「1人1つだけ」', {'size': 18, 'bold': True, 'space': 8}),
    ('一度に3つも4つも指摘すると、どれも直りません。', {'size': 15, 'space': 5}),
    ('その人が次の商談で一番効く1つに絞って伝えてください。', {'size': 15, 'space': 5}),
    ('順番は必ず「良い点 → 直す点1つ」。逆にしない。', {'size': 16, 'bold': True, 'color': RED}),
])
conclusion(s, '“質”の高い褒めを“数”多く言うことが重要である', y=6.40)

# ================================================================ 15. ロープレ①
s = new_slide('ロールプレイング①')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('前半：p.2〜p.17を通す（顧客設定A・25分）', {'size': 24, 'bold': True})])
rows = [
    ['進行', 'アプローチブック', '時間', '狙う反応'],
    ['ルール①', '（商談前の確認）', '1分', '3条件が揃っている'],
    ['ルール②', 'p.2', '3分', '商談の構成が分かる'],
    ['ルール③④', '（閉じたまま）', '4分', '導入しなかった理由／今日決めることに同意'],
    ['ルール⑤⑥', 'p.3〜p.7', '5分', 'じゃあつけた方がいいかもねえ'],
    ['ルール⑦⑧⑨', 'p.9〜p.17', '7分', '30年でそんなにも払わないといけないのですね'],
]
table(s, 0.45, 1.95, 9.95, [1.35, 2.55, 0.90, 5.15], rows, font_size=13, row_h=0.52, head_h=0.32)
_, tf = rect(s, 0.45, 4.85, 4.95, 1.95, fill=LTYEL)
lines(tf, [
    ('この周で見るポイント', {'size': 17, 'bold': True, 'space': 8}),
    ('・自社ストーリーを20秒で言えたか', {'size': 15, 'space': 6}),
    ('・オープンクエスチョンを3回投げたか', {'size': 15, 'space': 6}),
    ('・p.9〜p.17を7分で抜けられたか', {'size': 15}),
])
_, tf = rect(s, 5.60, 4.85, 4.95, 1.95, fill=LTGRAY)
lines(tf, [
    ('お客様役が必ず入れる断り文句', {'size': 17, 'bold': True, 'space': 8}),
    ('「以前も訪問販売が来て、高かったからやめた」', {'size': 16, 'bold': True, 'color': RED, 'space': 6}),
    ('→ Yes,but法 → 第三者話法 → YES取り で返す', {'size': 15}),
])

# ================================================================ 16. ロープレ②
s = new_slide('ロールプレイング②')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('後半：p.18〜p.26＋見積り（顧客設定B・25分）', {'size': 24, 'bold': True})])
rows = [
    ['進行', 'アプローチブック／道具', '時間', '狙う反応'],
    ['ルール⑩⑪', 'p.18・p.19＋事例', ('3分以内', {'bold': True, 'color': RED}), 'いつ災害が起こるか分からないですね'],
    ['ルール⑫', 'p.20・p.21', '3分', '買うより自給自足した方が良いですね'],
    ['ルール⑬⑭⑮', 'p.23〜p.26', '6分', '長期的には財産になるのですね'],
    ['受け渡し', 'アプローチブックを閉じる', '30秒', '−'],
    ['ルール⑯⑰⑱', 'シミュレーション・概算見積', '7分', '月1万円台なら、と前向きな答え'],
]
table(s, 0.45, 1.95, 9.95, [1.35, 3.20, 1.00, 4.40], rows, font_size=13, row_h=0.52, head_h=0.32)
_, tf = rect(s, 0.45, 4.85, 4.95, 1.95, fill=LTYEL)
lines(tf, [
    ('この周で見るポイント', {'size': 17, 'bold': True, 'space': 8}),
    ('・災害が3分以内に終わったか', {'size': 15, 'space': 6}),
    ('・1日 → 1か月 → 総額 の順を守れたか', {'size': 15, 'space': 6}),
    ('・断り文句3つを全部さばけたか', {'size': 15, 'bold': True, 'color': RED}),
])
_, tf = rect(s, 5.60, 4.85, 4.95, 1.95, fill=LTGRAY)
lines(tf, [
    ('お客様役が必ず入れる断り文句（3つ）', {'size': 17, 'bold': True, 'space': 8}),
    ('①「補助金が出るまで待ちたい」', {'size': 15, 'space': 5}),
    ('②「他の会社の話も聞いてみたい」', {'size': 15, 'space': 5}),
    ('③「妻と相談して決めます」', {'size': 15}),
])

# ================================================================ 17. 断り文句 早見表①
s = new_slide('よく出る断り文句 早見表①')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.40)
lines(tf, [('切り返しの原則：① 同意 → ②「だからこそ」→ ③ 第三者事例', {'size': 20, 'bold': True})])
_, tf = rect(s, 0.45, 1.88, 9.95, 1.75, fill=LTGREEN)
lines(tf, [
    ('「補助金が出るまで待ちたい」', {'size': 18, 'bold': True, 'space': 8}),
    ('「確かに補助金があれば嬉しいですよね。実は今、『補助金待ちで損をする』方が非常に多いんです。'
     '待っている間に値上がりし続ける電気代を払い続けると、おそらく年間15〜20万円以上になります。'
     'しかも補助金は抽選式ですので、必ずしももらえるとは限りません。'
     'そう考えると、最も安く導入できるタイミングは、間違いなく『今』なんです」', {'size': 14}),
])
_, tf = rect(s, 0.45, 3.78, 9.95, 1.50, fill=LTGREEN)
lines(tf, [
    ('「また考えて連絡します」（その場で解決できる場合）', {'size': 18, 'bold': True, 'space': 8}),
    ('「お気持ちは痛いほど分かります。一生に何度もない大きな買い物ですから。'
     'よければ、検討したいとお考えのことをお教えいただけませんか？」'
     '→ 出てきた不安に、シミュレーションと現地調査で答える', {'size': 14}),
])
_, tf = rect(s, 0.45, 5.43, 9.95, 1.45, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('「また考えて連絡します」（即決に至らなかった場合）', {'size': 18, 'bold': True, 'space': 8}),
    ('「時期が空くと導入しない言い訳を探してしまいがちで、絶好のタイミングを逃してしまうかもしれません。'
     '○○日11時までにご検討いただけますか？ 〇日に私からご連絡します。'
     'それまでは、このキャンペーン枠を私の責任で確保しておきますね」', {'size': 14, 'bold': True}),
])

# ================================================================ 18. 断り文句 早見表②
s = new_slide('よく出る断り文句 早見表②')
_, tf = rect(s, 0.45, 1.40, 9.95, 1.55, fill=LTGREEN)
lines(tf, [
    ('「妻と相談して決めます」', {'size': 18, 'bold': True, 'space': 8}),
    ('「ごもっともです。奥様のご意見は大切ですよね。では仮のお話ですが、もしこの後奥様にお話しされて、'
     '『家計も助かるし、災害時も安心だからやりましょう』と快諾いただけたとしたら、'
     '〇〇様ご自身としては『進めたい』というお気持ちでよろしいでしょうか？」', {'size': 14, 'space': 5}),
    ('（NO／曖昧な場合）→「ということは、〇〇様ご自身の中に、まだ引っかかる点があるということですね？」',
     {'size': 14, 'bold': True}),
])
_, tf = rect(s, 0.45, 3.10, 9.95, 1.55, fill=LTGREEN)
lines(tf, [
    ('「他の会社の話も聞いてみたい」', {'size': 18, 'bold': True, 'space': 8}),
    ('「おっしゃる通りです。ただ、太陽光や蓄電池選びで一番怖いことをご存じでしょうか。'
     'それは『表面上の価格だけで選んでしまい、施工やアフターフォローで失敗すること』なんです。'
     '設置後の定期点検、トラブル時の即日対応、15年20年先まで責任を持つ企業体制。'
     'これら『設置後の安心』まで含めた総合力では、絶対に負けない自信がございます」', {'size': 14}),
])
_, tf = rect(s, 0.45, 4.80, 9.95, 1.55, fill=LTGREEN)
lines(tf, [
    ('「今は資金的に余裕がない」', {'size': 18, 'bold': True, 'space': 8}),
    ('「一つだけ整理させてください。今のままでも、電気代という『終わりのないローン』は払い続けますよね？'
     '今回のお話はただの追加の出費ではありません。『電力会社への捨て金』を、'
     '同額のまま『ご自宅の資産』にスライドさせるだけです」', {'size': 14}),
])
_, tf = rect(s, 0.45, 6.50, 9.95, 0.42, fill=LTYEL)
lines(tf, [('相手の意見に自分の意見を重ねるのは最大のNG。「私が〜」ではなく「実際に〜あります」で返す',
            {'size': 15, 'bold': True, 'align': PP_ALIGN.CENTER})])

# ================================================================ 19. クロージングとは
s = new_slide('クロージング')
_, tf = rect(s, 0.45, 1.55, 9.95, 1.20, fill=None, line='BFBFBF')
p = para(tf, first=True, align=PP_ALIGN.CENTER)
run(p, 'クロージングとは、買うべき理由を説明・説得して', size=20)
p = para(tf, align=PP_ALIGN.CENTER)
run(p, '契約を迫る行為', size=20)
run(p, 'ではない', size=20, bold=True, color=RED)
_, tf = rect(s, 0.45, 3.00, 9.95, 1.35, fill=LTGREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
lines(tf, [
    ('クロージングとは', {'size': 18, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 4}),
    ('お客様の意思決定をサポートすることである', {'size': 30, 'bold': True, 'color': RED,
                                 'align': PP_ALIGN.CENTER}),
])
_, tf = rect(s, 0.45, 4.60, 9.95, 1.05, fill=None, line='BFBFBF')
lines(tf, [
    ('つまり「買わない理由を解消し、お客様自身に買うべき理由を発見していただく作業」のことである',
     {'size': 19, 'align': PP_ALIGN.CENTER}),
])
_, tf = rect(s, 0.45, 5.90, 9.95, 0.95, fill=LTYEL)
lines(tf, [
    ('断り文句が出るのは「必要でないから」ではなく、単に「迷っているから」。', {'size': 16, 'bold': True, 'space': 5}),
    ('営業マンがすべきことは「背中を押してあげること」だけです。', {'size': 16, 'bold': True, 'color': RED}),
])

# ================================================================ 20. 商談の正しい進め方 総復習
s = new_slide('総復習：商談の正しい進め方')
rows = [
    ['', '原則', '中身'],
    ['（1）', '蓄電池の営業スタイル', '「説明」だけで50万円超は売れない。①買わない理由をなくす ②買う理由を確認する'],
    ['（2）', 'ラリー数と発言比率', '営業：お客様＝6：4（お客様側は夫婦で2：2）。ラリー数は30回'],
    ['（3）', '話の振り方', 'ネガティブな話＝クローズド／ポジティブな話＝オープン。A（言わせたいこと）ありきでQを投げる'],
    ['（4）', '切り返し', '① 同意（Yes,but法）→ ②「だからこそ」（オウム返し法）→ ③ 第三者事例'],
    ['（5）', '主権者商談', '本主権者がいない状態での価格提示は厳禁。価値が伝わらないまま価格に引っ張られる'],
]
table(s, 0.45, 1.50, 9.95, [0.75, 2.60, 6.60], rows, font_size=12.5, row_h=0.72, head_h=0.32)
_, tf = rect(s, 0.45, 5.80, 4.95, 1.05, fill=LTGRAY)
lines(tf, [
    ('話を振る目的', {'size': 16, 'bold': True, 'space': 5}),
    ('本来こちらが説明する内容を、相手の口から言わせること。', {'size': 14}),
])
_, tf = rect(s, 5.60, 5.80, 4.95, 1.05, fill=LTGRAY)
lines(tf, [
    ('「だからこそ」の使い方', {'size': 16, 'bold': True, 'space': 5}),
    ('前後は論理的に意味をなしている必要は全くありません。', {'size': 14}),
])

# ================================================================ 21. 後始末と追客
s = new_slide('後始末と追客')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('契約したあとが、次の紹介につながります', {'size': 24, 'bold': True})])
rows = [
    ['時期', 'やること'],
    ['契約時', '正式見積りの提示が前後することになっても、必ず一筆いただけるようにする'],
    ['契約直後〜5日以内',
     'なるべく翌日までにお礼状（ハガキで可）を投函。直筆コメント・訪問時の話題・夫婦両者のお名前を入れる'],
    ['8日間〜', '期間終了後も取り消しができないわけではない。油断せず、完工まで2週間ごとに連絡を取る'],
]
table(s, 0.45, 1.95, 9.95, [2.15, 7.80], rows, font_size=13.5, row_h=0.68, head_h=0.32)
_, tf = rect(s, 0.45, 4.45, 4.95, 2.35, fill=LTGRAY)
lines(tf, [
    ('現地調査までに熱が冷めるとき', {'size': 16, 'bold': True, 'space': 8}),
    ('・お礼手紙の送付', {'size': 14, 'space': 4}),
    ('　→ 判断は「論理」、決断は「感情」', {'size': 14, 'color': GRAY, 'space': 8}),
    ('・裏側の近況報告（いま図面を確認しています）', {'size': 14, 'space': 4}),
    ('　→ 放置されている感を与えない', {'size': 14, 'color': GRAY, 'space': 8}),
    ('・ドローンで屋根の状態まで調査すると伝える', {'size': 14, 'space': 4}),
    ('　→ プロとしての価値を高める', {'size': 14, 'color': GRAY}),
])
_, tf = rect(s, 5.60, 4.45, 4.95, 2.35, fill=LTYEL)
lines(tf, [
    ('連絡が途絶えたとき', {'size': 16, 'bold': True, 'space': 8}),
    ('・期日には必ず営業マン側から連絡する', {'size': 14, 'space': 4}),
    ('　→ 待つと音信不通になる確率が高まる', {'size': 14, 'color': GRAY, 'space': 8}),
    ('・商談終了時に懸念点へのアプローチをしておくと、', {'size': 14, 'space': 2}),
    ('　次の連絡が自然なトークの運びになる', {'size': 14, 'space': 10}),
    ('取り組み事例：エコプラスワン様', {'size': 15, 'bold': True, 'space': 4}),
    ('お礼状の工夫でクーリングオフほぼなし', {'size': 14, 'bold': True}),
])

# ================================================================ 22. 個別フィードバック
s = new_slide('個別フィードバック')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('評価シートを見ながら、1人ずつ確認します（15分）', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 1.95, 9.95, 1.30, fill=LTGREEN)
lines(tf, [
    ('全員に共通で確認すること', {'size': 17, 'bold': True, 'space': 6}),
    ('①〜⑳のうち、○が付かなかったルールはどれか　／　その原因は「知らない」か「言えない」か',
     {'size': 16, 'bold': True}),
])
rows = [
    ['原因', '対処'],
    ['知らない（内容が頭に入っていない）', '該当ページの結論の一文を暗唱する。docs/知識_営業ルール20.md を読み直す'],
    ['言えない（知っているが口から出ない）', '声に出す回数が足りない。次の商談まで毎日3回、そのページだけ音読する'],
    ['飛ばした（時間が足りなかった）', '手前のパートが長すぎる。商談の設計図の時間配分に戻る'],
]
table(s, 0.45, 3.45, 9.95, [3.55, 6.40], rows, font_size=13.5, row_h=0.72, head_h=0.34)
_, tf = rect(s, 0.45, 5.95, 9.95, 0.90, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('直す点は1人1つだけ。次の商談で一番効く1つに絞って伝えます。', {'size': 16, 'bold': True, 'space': 4}),
    ('順番は必ず「良い点 → 直す点1つ」。', {'size': 16, 'bold': True, 'color': RED}),
])

# ================================================================ 23. 明日からの行動計画
s = new_slide('明日からの行動計画')
box, tf = tb(s, 0.45, 1.38, 9.95, 0.45)
lines(tf, [('研修は今日で終わりですが、定着はここからです', {'size': 24, 'bold': True})])
_, tf = rect(s, 0.45, 1.92, 9.95, 1.05, fill=LTGRAY)
lines(tf, [
    ('忘却曲線', {'size': 16, 'bold': True, 'space': 5}),
    ('翌日には50％、1週間後にはほとんど忘れます。その日のうちに復習することで記憶が定着します。',
     {'size': 15}),
])
rows = [
    ['いつ', 'やること', '記入してください'],
    ['今日中', 'アプローチブック p.2〜p.26 の結論の一文を通しで音読する', ''],
    ['今週中', '実商談を1件、評価シートをつけて振り返る', ''],
    ['今週中', '○が付かなかったルールを3つ書き出し、対処を決める', ''],
    ['今月中', '通しロープレを社内で2回やる（相手：　　　　　　　）', ''],
]
table(s, 0.45, 3.15, 9.95, [1.30, 5.20, 3.45], rows, font_size=13, row_h=0.68, head_h=0.34)
_, tf = rect(s, 0.45, 6.15, 9.95, 0.72, fill=LTYEL, line='BF8F00')
lines(tf, [
    ('成長の7割は「経験」から生まれます（ロミンガーの法則：経験70／薫陶20／研修10）。',
     {'size': 16, 'bold': True, 'space': 4}),
    ('今日学んだ10を、明日の現場で70に変えてください。', {'size': 16, 'bold': True, 'color': RED}),
])

# ================================================================ 24. 修了メッセージ
s = new_slide('おわりに')
_, tf = rect(s, 0.45, 1.70, 9.95, 1.55, fill=LTGREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
lines(tf, [
    ('太陽光と蓄電池は、切り離さない。', {'size': 30, 'bold': True, 'align': PP_ALIGN.CENTER, 'space': 6}),
    ('スマートハウスを、お届けする。', {'size': 30, 'bold': True, 'color': RED, 'align': PP_ALIGN.CENTER}),
])
_, tf = rect(s, 0.45, 3.55, 9.95, 2.20, fill=None, line='BFBFBF')
lines(tf, [
    ('全5回でお伝えしてきたことは、突き詰めればこの2つです。', {'size': 18, 'space': 10,
                                            'align': PP_ALIGN.CENTER}),
    ('① お客様の「買わない理由」を、順番どおりに1つずつ潰していくこと', {'size': 18, 'bold': True, 'space': 8}),
    ('② お客様自身に「買うべき理由」を発見していただくこと', {'size': 18, 'bold': True, 'space': 10}),
    ('アプローチブックは、そのための地図です。何度も開いてください。', {'size': 18,
                                                'align': PP_ALIGN.CENTER}),
])
_, tf = rect(s, 0.45, 6.05, 9.95, 0.80, fill=LTYEL)
lines(tf, [('全5回、お疲れさまでした。現場でお会いしましょう。',
            {'size': 22, 'bold': True, 'align': PP_ALIGN.CENTER})])
notes(s, '・最後は必ず全員に一言ずつ言ってもらってから締める')

# ---------------------------------------------------------------- 並べ替え
ids = list(sldIdLst)
colophon = [x for x in ids if x is colophon_id][0]
sldIdLst.remove(colophon)
sldIdLst.append(colophon)
page_no[0] += 1

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print('saved:', OUT, '/ slides:', len(prs.slides.__iter__.__self__._sldIdLst))

# ---------------------------------------------------------------- content-type fix
def ensure_jpg_content_type(path):
    """テンプレートが jpg を image/png と誤宣言しているため、python-pptx が
    保存時に Default を落とすことがある。表紙画像が壊れるので明示的に直す。"""
    import zipfile, shutil, re, tempfile
    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        ct = z.read('[Content_Types].xml').decode('utf-8')
        needs = any(n.lower().endswith('.jpg') for n in names)
        if not needs or 'Extension="jpg"' in ct:
            if 'Extension="jpg" ContentType="image/jpeg"' in ct or not needs:
                return False
        data = {n: z.read(n) for n in names}
    ct = re.sub(r'<Default Extension="jpg"[^/]*/>', '', ct)
    ct = ct.replace('<Default Extension="png"',
                    '<Default Extension="jpg" ContentType="image/jpeg"/><Default Extension="png"', 1)
    data['[Content_Types].xml'] = ct.encode('utf-8')
    tmp = path + '.tmp'
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as z:
        for n in names:
            z.writestr(n, data[n])
    shutil.move(tmp, path)
    return True

if ensure_jpg_content_type(OUT):
    print('fixed [Content_Types].xml (jpg -> image/jpeg)')
