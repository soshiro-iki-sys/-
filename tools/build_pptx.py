#!/usr/bin/env python3
"""
工務店アライアンス DL資料 を「編集可能なPowerPoint」として生成する。

Marp の --pptx は各ページを画像1枚として貼るだけなので編集できない。
このスクリプトは支給テンプレート（2026年8月版）を土台に、
すべてネイティブの図形・テキストボックス・表で組み立てる。

  python3 tools/build_pptx.py <テンプレート.pptx> <出力.pptx>
"""
import copy
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

# ---------------------------------------------------------------- 定数
NAVY = RGBColor(0x00, 0x20, 0x60)
RED = RGBColor(0xFF, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x40, 0x40, 0x40)      # 表ヘッダー
MUTED = RGBColor(0x80, 0x80, 0x80)     # 出典・注記
HL = RGBColor(0xFF, 0xF9, 0xD6)        # 表の強調行
FONT = "Meiryo UI"

# レイアウト座標（cm）— docs/資料フォーマット定義.md に準拠
L = 0.64                # 左マージン
W = 36.05               # 本文幅
CHAPTER = (L, 0.27, 30.0, 1.45)
MIDTITLE = (L, 2.35, 36.0, 1.28)
SUBHEAD = (L, 3.60, 36.0, 1.00)
BAND = (0.0, 17.78, 37.33, 2.34)
SRC_Y = 16.90
BODY_TOP = 4.60         # 本文の開始
BODY_BOTTOM = 17.40     # 本文の下限

GRID_STYLE = "{5940675A-B579-460E-94D1-54222C63F5DA}"   # スタイルなし、表のグリッド線あり


# ---------------------------------------------------------------- 低レベル補助
def _ea(run, name=FONT):
    """python-pptx は latin しか設定しないので、日本語（ea）も明示する。"""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def tf_setup(shape, margin=0.0, anchor=MSO_ANCHOR.TOP, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Cm(margin)
    tf.margin_top = tf.margin_bottom = Cm(0)
    return tf


def put(slide, x, y, w, h, lines, size=16, bold=True, color=BLACK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, margin=0.0):
    """テキストボックスを置く。lines は str か (文字列, {上書き}) のリスト。"""
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tf_setup(box, margin, anchor)
    if isinstance(lines, str):
        lines = [lines]
    for i, item in enumerate(lines):
        txt, ov = (item, {}) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("spacing", spacing)
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = FONT
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)
        _ea(r)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    tf_setup(sp, 0.15, MSO_ANCHOR.MIDDLE)
    sp.text_frame.text = ""
    return sp


def label(sp, lines, size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER,
          spacing=1.0):
    """図形に文字を入れる。"""
    tf = sp.text_frame
    if isinstance(lines, str):
        lines = [lines]
    for i, item in enumerate(lines):
        txt, ov = (item, {}) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("spacing", spacing)
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = FONT
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)
        _ea(r)
    return sp


def table(slide, x, y, w, h, rows, widths=None, size=11, header=True,
          row_h=None, hl_rows=()):
    """テンプレートの表スタイル（グレーヘッダー＋黒グリッド）で表を置く。"""
    nr, nc = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nr, nc, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = gf.table
    # 表スタイルを「スタイルなし、表のグリッド線あり」に固定
    el = tbl._tbl.find(qn("a:tblPr"))
    if el is not None:
        for sid in el.findall(qn("a:tableStyleId")):
            el.remove(sid)
        sid = el.makeelement(qn("a:tableStyleId"), {})
        sid.text = GRID_STYLE
        el.append(sid)
        el.set("firstRow", "1" if header else "0")
        el.set("bandRow", "0")
    if widths:
        for i, cw in enumerate(widths):
            tbl.columns[i].width = Cm(cw)
    if row_h:
        for i, rh in enumerate(row_h):
            tbl.rows[i].height = Cm(rh)
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Cm(0.18)
            cell.margin_top = cell.margin_bottom = Cm(0.08)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_head = header and ri == 0
            cell.fill.solid()
            if is_head:
                cell.fill.fore_color.rgb = GRAY
            elif ri in hl_rows:
                cell.fill.fore_color.rgb = HL
            else:
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            parts = cell_text.split("\n") if isinstance(cell_text, str) else cell_text
            for pi, part in enumerate(parts):
                p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
                p.line_spacing = 1.15
                sz, col = size, (WHITE if is_head else BLACK)
                if part.startswith("§"):        # §で始まる行は小さめのグレー
                    part, sz, col = part[1:], size - 2, MUTED
                if part.startswith("!"):        # !で始まる行は赤
                    part, col = part[1:], RED
                r = p.add_run()
                r.text = part
                f = r.font
                f.name = FONT
                f.size = Pt(sz)
                f.bold = True
                f.color.rgb = col
                _ea(r)
    return tbl


# ---------------------------------------------------------------- ページ組み立て
def page(prs, layout, no, chapter, midtitle, sub=None, band=None, src=None):
    """中面を1枚つくる。"""
    s = prs.slides.add_slide(layout)
    put(s, *CHAPTER, chapter, size=28, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    put(s, *MIDTITLE, midtitle, size=24, color=RED, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        box = put(s, *SUBHEAD, "　" + sub, size=20, color=BLACK,
                  anchor=MSO_ANCHOR.MIDDLE)
        # 行頭記号はテンプレートと同じ Wingdings の "Ø"（＝➢）を使う
        para = box.text_frame.paragraphs[0]
        bullet = copy.deepcopy(para.runs[0]._r)
        para._p.insert(list(para._p).index(para.runs[0]._r), bullet)
        from pptx.oxml.ns import qn as _q
        bullet.find(_q("a:t")).text = "Ø"
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = bullet.get_or_add_rPr().find(_q(tag))
            if el is not None:
                el.set("typeface", "Wingdings")
    if band:
        b = rect(s, *BAND, fill=NAVY)
        lines = band if isinstance(band, list) else [band]
        label(b, lines, size=24 if len(lines) == 1 else 20, color=WHITE,
              spacing=1.2)
    if src:
        put(s, L, SRC_Y, W, 0.7, src, size=9, bold=False, color=MUTED)
    if no is not None:
        put(s, 28.93, 20.38, 8.4, 0.6, str(no), size=11, bold=False,
            color=WHITE, align=PP_ALIGN.CENTER)
    return s


def card(slide, x, y, w, h, badge, title, body, accent=NAVY):
    """左に太い縦帯を持つカード。"""
    rect(slide, x, y, w, h, fill=None, line=BLACK, line_w=0.75)
    rect(slide, x, y, 0.28, h, fill=accent)
    put(slide, x + 0.55, y + 0.28, w - 0.9, 0.6, badge, size=11,
        color=GRAY, align=PP_ALIGN.LEFT)
    put(slide, x + 0.55, y + 0.92, w - 0.9, 0.9, title, size=15,
        color=accent, align=PP_ALIGN.LEFT)
    put(slide, x + 0.55, y + 1.90, w - 0.9, h - 2.2, body, size=11,
        bold=True, align=PP_ALIGN.LEFT, spacing=1.25)


def main(tpl_path, out_path):
    prs = Presentation(tpl_path)
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)

    # 「中面」レイアウトを取得
    layout = None
    for m in prs.slide_masters:
        for lo in m.slide_layouts:
            if lo.name == "中面":
                layout = lo
    assert layout is not None, "レイアウト「中面」が見つかりません"

    # テンプレートのサンプル中面（2〜5枚目）を削除。表紙とCTAは残す。
    # sldIdLst から外すだけではパートが残り ZIP が壊れるので、関連付けも落とす。
    for i in (4, 3, 2, 1):
        prs.part.drop_rel(ids[i].rId)
        sldIdLst.remove(ids[i])

    # CTAページのパート名を退避しておく。空いた slide2〜6 に新規スライドが
    # 割り当てられると slide6.xml が衝突してZIPが壊れるため。
    from pptx.opc.packuri import PackURI
    cta_part = prs.part.related_part(list(sldIdLst)[1].rId)
    cta_part.partname = PackURI("/ppt/slides/slide900.xml")

    # ---------- P1 表紙：文字だけ差し替える ----------
    cover = prs.slides[0]
    for sh in cover.shapes:
        if sh.name == "タイトル 1":
            tf = sh.text_frame
            tf.clear()
            for i, t in enumerate(["工務店アライアンス", "成功のポイントとは？"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                r = p.add_run()
                r.text = t
                r.font.name = FONT
                r.font.size = Pt(60)
                r.font.bold = True
                r.font.color.rgb = BLACK
                _ea(r)
        elif sh.name == "object 17":
            tf = sh.text_frame
            tf.clear()
            for i, t in enumerate(["受注件数の方程式と実践KPI",
                                   "地域密着・営業5名規模の販売店向け"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p.add_run()
                r.text = t
                r.font.name = FONT
                r.font.size = Pt(22)
                r.font.bold = True
                r.font.color.rgb = BLACK
                _ea(r)

    CH1 = "1. アライアンスとは？"
    CH2 = "2. アライアンスをやるべき理由"
    CH3 = "3. アライアンス成功のポイント"

    # ---------- P2 はじめに ----------
    s = page(prs, layout, 2, "はじめに", "この資料でわかること",
             "地域密着型・営業人員5名程度の販売店が、工務店との提携で受注を伸ばすための実務書です",
             "提携の「やり方」ではなく、どの変数を改善すべきかが分かる資料です")
    cw, gap = 11.55, 0.62
    cards = [
        ("01", "アライアンスの\n座組と使う手法",
         "工務店から新築施主の太陽光・蓄電池を発注いただく「元請け開拓」の座組を整理します。"
         "開拓手法は3つあり、自社の規模と商圏でどれを選ぶべきかを示します。"),
        ("02", "工務店が\n動けていない理由",
         "必要性を理解していても対応できない工務店が数多く存在します。その理由を"
         "人材・知識・収益・運用の4つに分解し、提携の余地がどこにあるかを示します。"),
        ("03", "受注を分解する\n方程式と打ち手",
         "受注件数を5つの変数の掛け算として捉え、変数ごとのKPIと具体施策、"
         "そして自社の弱点を特定する診断チェックリストをご用意しました。"),
    ]
    for i, (no, ttl, body) in enumerate(cards):
        x = L + i * (cw + gap)
        rect(s, x, 5.0, cw, 9.6, fill=None, line=BLACK, line_w=0.75)
        b = rect(s, x + 0.6, 5.6, 1.9, 0.75, fill=GRAY)
        label(b, no, size=12, color=WHITE)
        put(s, x + 0.6, 6.7, cw - 1.2, 2.0, ttl.split("\n"), size=17,
            color=NAVY, align=PP_ALIGN.LEFT, spacing=1.2)
        put(s, x + 0.6, 9.0, cw - 1.2, 5.2, body, size=12,
            align=PP_ALIGN.LEFT, spacing=1.35)

    # ---------- P3 目次 ----------
    s = page(prs, layout, 3, "目次", "")
    put(s, L, 2.98, 36.19, 7.78,
        ["① アライアンスとは？",
         "② アライアンスをやるべき理由",
         "③ アライアンス成功のポイント"],
        size=40, color=BLACK, spacing=1.6)

    # ---------- P4 定義 ----------
    s = page(prs, layout, 4, CH1, "定義と本資料が扱う範囲",
             "工務店から新築施主の太陽光・蓄電池導入を発注いただく「元請け開拓」を扱います",
             "目指すのは案件の獲得ではなく、案件が生まれ続ける関係をつくること")
    put(s, L, 4.75, W, 1.8,
        "工務店が持つ「顧客との接点と信頼」と、販売店が持つ「商品・提案・施工・アフターの実行力」を"
        "組み合わせ、双方の顧客に価値を提供する継続的な業務提携です。"
        "一度きりの案件紹介でも、元請けからの工事受注でもありません。",
        size=13, spacing=1.35)
    table(s, L, 7.0, W, 9.3, [
        ["", "下請け", "単発の紹介", "アライアンス"],
        ["関係性", "発注者と受注者の上下関係", "都度のスポット取引", "対等な継続的パートナー"],
        ["顧客との接点", "元請け経由。自社は表に出ない", "紹介された案件のみ", "顧客名簿単位で継続的に接触"],
        ["価格の決定権", "元請けが決定", "案件ごとに交渉", "提携時に条件を取り決め"],
        ["収益の形", "工事単価のみ", "単発の紹介フィー", "継続的な受注＋紹介手数料"],
        ["再現性", "元請けの受注量に依存", "低い（属人的）", "高い（仕組みで積み上がる）"],
    ], widths=[6.0, 10.0, 8.5, 11.55], size=12)

    # ---------- P5 座組 ----------
    s = page(prs, layout, 5, CH1, "座組と役割分担",
             "工務店は「紹介だけ」から「自社で営業し施工だけ依頼」まで関与度を選べる",
             "工務店の不安は「自社の顧客を任せること」。そこを消す設計が成否を決める")
    lw = 18.6
    b = rect(s, 5.4, 5.0, 8.0, 1.2, fill=NAVY)
    label(b, "オーナー顧客", size=15, color=WHITE)
    put(s, 0.9, 6.6, 8.2, 1.9,
        ["② 訪問・関心の確認",
         ("（自社で営業する場合はそのまま営業）", {"size": 9, "color": GRAY})],
        size=11, align=PP_ALIGN.CENTER, spacing=1.2)
    put(s, 10.0, 6.6, 8.2, 1.9,
        ["④ 営業代行＋施工", ("④' 施工のみ", {"size": 9, "color": GRAY})],
        size=11, align=PP_ALIGN.CENTER, spacing=1.2)
    n1 = rect(s, 0.9, 8.7, 8.2, 1.9, fill=NAVY)
    label(n1, ["工務店", ("顧客接点・地域での信頼", {"size": 9, "color": RGBColor(0xC9, 0xD6, 0xEA)})],
          size=14, color=WHITE, spacing=1.2)
    n2 = rect(s, 10.0, 8.7, 8.2, 1.9, fill=NAVY)
    label(n2, ["販売店", ("提案・施工・アフター", {"size": 9, "color": RGBColor(0xC9, 0xD6, 0xEA)})],
          size=14, color=WHITE, spacing=1.2)
    a1 = rect(s, 0.9, 11.1, 17.3, 1.4, fill=None, line=BLACK, line_w=0.75)
    rect(s, 0.9, 11.1, 0.18, 1.4, fill=RED)
    label(a1, "工務店 ▶ 販売店　③ 案件の紹介／③' 成約顧客リストの提供",
          size=11, align=PP_ALIGN.LEFT)
    a2 = rect(s, 0.9, 12.8, 17.3, 1.4, fill=None, line=BLACK, line_w=0.75)
    rect(s, 0.9, 12.8, 0.18, 1.4, fill=NAVY)
    label(a2, "販売店 ▶ 工務店　① 商材の卸・営業ツール／⑤ バックマージン",
          size=11, align=PP_ALIGN.LEFT)
    table(s, 19.7, 5.0, 17.0, 9.2, [
        ["", "工務店", "販売店（自社）"],
        ["提供する\nもの", "顧客接点\n地域での信頼\n顧客名簿",
         "商品調達／提案・見積\n施工・アフター\n補助金申請"],
        ["収益", "紹介手数料\n§販売金額の5〜7％が目安", "工事売上\n継続的な提案機会"],
        ["リスク", "顧客満足度の低下\n§施工品質は販売店に依存",
         "提携先の紹介停滞\n個人情報の管理責任"],
    ], widths=[3.4, 6.6, 7.0], size=12)

    # ---------- P6 3手法 ----------
    s = page(prs, layout, 6, CH1, "工務店を開拓する3つの手法｜自社が使うのはどれか",
             "地域密着・営業5名以下なら「セミナー（勉強会）」を選ぶ",
             "本資料は以降、セミナー（勉強会）による開拓を前提に解説します")
    table(s, L, 4.9, W, 10.5, [
        ["施策", "メリット", "デメリット", "商圏", "対象とする\n販売店", "出口"],
        ["① 電話\n§TELマーケティング",
         "自社のリソースで手間をかけずにアポが取れる／確実にアプローチできる",
         "決裁者までたどり着くのに時間がかかる", "関東圏／全国", "営業10名以上",
         "新築＋工務店リストの提供を受ける"],
        ["② ネットワーク活用",
         "上手くいけば一気に広げやすい／手っ取り早く決裁者に会いやすい",
         "団体ごとの制約あり（有償、営業NGなど）", "—", "誰でも",
         "新築＋工務店リストの提供を受ける"],
        ["③ セミナー\n§勉強会",
         "意欲の高い会社を集められる／外部講師などを含めてフックを使える",
         "着手〜セミナー開催に時間を要する（3〜4か月）", "地域密着型", "5名以下",
         "新築＋工務店リストの提供を受け、太陽光・蓄電池の営業を行う"],
    ], widths=[5.2, 9.0, 7.2, 3.6, 3.9, 7.15], size=11, hl_rows=(3,))

    # ---------- P7 金利 ----------
    s = page(prs, layout, 7, CH2, "金利上昇が、工務店を「追加受注」に向かわせている",
             "金利が1％上がると施主の予算は約600万円縮む。棟数だけでなく1棟単価も落ちている",
             ["金利上昇で工務店の粗利は縮み、施主は「月々いくらか」に敏感になっている",
              "その両方に効くのが太陽光・蓄電池。だが工務店にはそれを売る力がない"],
             src="※借入可能額・返済額は4,000万円・35年・元利均等返済での試算。"
                 "電気代の削減額は設置容量・地域・電気料金単価により異なります。")
    bw, bx = 11.55, [L, L + 12.17, L + 24.34]
    heads = [("① 施主に起きていること", "借りられる額が縮む", NAVY, GRAY),
             ("② 工務店に起きていること", "棟数も単価も落ちる", NAVY, GRAY),
             ("③ ではなぜ太陽光なのか", "月々の支払いはむしろ下がる", RED, RED)]
    for i, (head, ttl, col, hcol) in enumerate(heads):
        x = bx[i]
        rect(s, x, 4.9, bw, 10.9, fill=None, line=BLACK, line_w=0.75)
        hb = rect(s, x, 4.9, bw, 0.9, fill=hcol)
        label(hb, head, size=12, color=WHITE)
        put(s, x + 0.4, 6.1, bw - 0.8, 0.9, ttl, size=15, color=col,
            align=PP_ALIGN.CENTER)
    # ① 表と数値
    table(s, bx[0] + 0.5, 7.3, bw - 1.0, 3.6, [
        ["金利", "月々返済", "借入可能額"],
        ["0.5％", "103,834円", "4,000万円"],
        ["1.0％", "112,914円", "3,678万円"],
        ["1.5％", "122,474円", "3,391万円"],
    ], widths=[2.9, 4.0, 3.65], size=10)
    put(s, bx[0] + 0.5, 11.3, bw - 1.0, 1.0, "▲609万円", size=20, color=NAVY,
        align=PP_ALIGN.CENTER)
    put(s, bx[0] + 0.5, 12.4, bw - 1.0, 2.0,
        "同じ月10.4万円の返済で借りられる額", size=11, align=PP_ALIGN.CENTER,
        spacing=1.25)
    # ②
    put(s, bx[1] + 0.5, 7.3, bw - 1.0, 8.0,
        ["着工棟数が減るうえに、施主の予算縮小で1棟あたりの単価も下がる。掛け算で粗利額が縮みます。",
         "",
         "建物本体で挽回するのは難しく、1顧客あたりの取引額を上げるしか手がありません。",
         "",
         "だから工務店は、いま追加受注の商材を探しています。"],
        size=11.5, align=PP_ALIGN.LEFT, spacing=1.35)
    # ③
    put(s, bx[2] + 0.5, 7.3, bw - 1.0, 2.4,
        "「予算が縮むなら設備は削られるのでは？」——逆です。",
        size=11.5, align=PP_ALIGN.LEFT, spacing=1.35)
    put(s, bx[2] + 0.5, 9.5, bw - 1.0, 1.4,
        "太陽光150万円をローンに乗せても返済増は", size=11.5,
        align=PP_ALIGN.LEFT, spacing=1.35)
    put(s, bx[2] + 0.5, 11.0, bw - 1.0, 1.0, "月4,593円", size=20, color=RED,
        align=PP_ALIGN.CENTER)
    put(s, bx[2] + 0.5, 12.2, bw - 1.0, 3.2,
        "電気代の削減がこれを上回れば、ローン＋光熱費の総額はむしろ下がる。"
        "金利上昇局面ほど刺さる提案です。",
        size=11.5, align=PP_ALIGN.LEFT, spacing=1.35)

    # ---------- P8 4つの課題 ----------
    s = page(prs, layout, 8, CH2, "工務店が抱える4つの課題",
             "追加受注を取りたい。しかし工務店は動けていない。理由は4つの構造的な壁にある",
             "動けない理由は「やる気」ではなく構造。だから外部との組み合わせで解ける")
    quad = [
        ("課題① 人　材", "営業リソースが足りない",
         "新築の受注活動で手一杯。既存顧客に追加提案する担当を置けない。"
         "採用も難しく、有効求人倍率は高止まりしています。"),
        ("課題② 知　識", "提案の仕方が分からない",
         "容量設計、経済効果シミュレーション、見積の組み立てに知見がなく、"
         "聞かれても答えられないため提案自体を避けてしまう。"),
        ("課題③ 収　益", "粗利を確保できない",
         "仕入ルートがなく単価が高い。相見積で価格勝負になり、"
         "手間の割に利益が残らない。だから優先度が下がる。"),
        ("課題④ 運　用", "補助金・法対応の負荷",
         "補助金は年度ごとに要件が変わり、申請期限も短い。"
         "系統連系の申請も含め、事務負担が読めない。"),
    ]
    qw, qh = 17.72, 5.3
    for i, (badge, ttl, body) in enumerate(quad):
        x = L + (i % 2) * (qw + 0.61)
        y = 5.0 + (i // 2) * (qh + 0.6)
        card(s, x, y, qw, qh, badge, ttl, body)

    # ---------- P9 課題の深掘り ----------
    s = page(prs, layout, 9, CH2, "課題の深掘り", None,
             "4つの課題はすべて「自社だけで抱えるから解けない」という共通点を持つ")
    deep = [
        ("課題① 営業リソース不足", "接点はあるのに提案されない",
         "・オーナー顧客は増えているが、そのメリットを生かしきれていない\n"
         "・点検やアフター訪問の機会があっても提案に踏み込めない\n"
         "・担当者が現場管理と兼務で追客の時間が取れない",
         "▶ 放置すると、顧客は他社の情報だけで判断してしまう"),
        ("課題② 見積と提案の難しさ", "自信がないから切り出せない",
         "・新規の売り方は分かるが、リフォーム商材の提案方法が分からない\n"
         "・効果を語れないため会話が続かない\n"
         "・切り出すタイミングが分からず点検で終わる",
         "▶ 放置すると、受け身になり案件が生まれない"),
        ("課題③ 粗利確保", "手間の割に利益が残らない",
         "・仕入ルートがなく、単発発注のため仕入単価が高い\n"
         "・訪販業者と相見積になり価格で負ける\n"
         "・工事管理の手間に対して利益が見合わない",
         "▶ 放置すると、社内で優先順位が下がり定着しない"),
        ("課題④ 補助金・法対応", "事務負担が読めない",
         "・補助金は年度ごとに要件・予算・締切が変わる\n"
         "・予算到達で早期終了することがあり社内で追えない\n"
         "・系統連系や各種申請の手順が確立していない",
         "▶ 放置すると、申請漏れで施主の期待に応えられない"),
    ]
    qh2 = 5.6
    for i, (badge, ttl, body, risk) in enumerate(deep):
        x = L + (i % 2) * (qw + 0.61)
        y = 4.3 + (i // 2) * (qh2 + 0.6)
        rect(s, x, y, qw, qh2, fill=None, line=BLACK, line_w=0.75)
        rect(s, x, y, 0.28, qh2, fill=NAVY)
        put(s, x + 0.55, y + 0.25, qw - 0.9, 0.6, badge, size=11, color=GRAY,
            align=PP_ALIGN.LEFT)
        put(s, x + 0.55, y + 0.88, qw - 0.9, 0.8, ttl, size=15, color=NAVY,
            align=PP_ALIGN.LEFT)
        put(s, x + 0.55, y + 1.80, qw - 0.9, 2.6, body.split("\n"), size=11,
            align=PP_ALIGN.LEFT, spacing=1.3)
        put(s, x + 0.55, y + 4.55, qw - 0.9, 0.8, risk, size=10.5, color=RED,
            align=PP_ALIGN.LEFT)

    # ---------- P10 だからこそ ----------
    s = page(prs, layout, 10, CH2, "だからこそのアライアンス",
             "4つの課題は、販売店が既に持っている機能をそのまま当てれば解消できる",
             ["工務店は「手間なく満足度と売上が上がる」、販売店は「広告費ゼロで良質な案件が入る」",
              "実績や規模は関係ない。営業マン1名の販売店でも提携は成立している"])
    table(s, L, 5.1, W, 10.8, [
        ["", "工務店の課題", "アライアンスによる解決"],
        ["① 人　材", "提案する人がいない。採用もできない",
         "営業代行：販売店の営業が顧客訪問から提案・クロージングまで実施。工務店は紹介するだけでよい"],
        ["② 知　識", "容量設計・経済効果・見積が分からない",
         "ツールと勉強会：アプローチブック・シミュレーション資料の提供と、営業マン向けの個社別勉強会を定期開催"],
        ["③ 収　益", "仕入が高く粗利が残らない",
         "商材卸＋紹介手数料：まとめ仕入れによる原価低減と、販売金額の5〜7％を目安とした紹介フィー"],
        ["④ 運　用", "補助金・申請の事務負担が読めない",
         "申請代行と施工：補助金申請・系統連系・施工・アフターまで販売店が一括対応"],
    ], widths=[4.2, 11.0, 20.85], size=12)

    # ---------- P11 章の締め ----------
    s = page(prs, layout, 11, CH2, "第2章まとめ")
    put(s, 0.0, 7.4, 37.33, 5.0,
        [("必要性は感じている。", {"color": NAVY}),
         ("しかし、動けていない。", {"color": RED})],
        size=46, align=PP_ALIGN.CENTER, spacing=1.4)

    # ---------- P12 方程式 ----------
    s = page(prs, layout, 12, CH3, "受注件数の方程式",
             "アライアンスの受注件数は、5つの変数の掛け算で決まる",
             ["1つでもゼロなら成果はゼロ。提携社数を増やしても紹介率0％なら受注は生まれない",
              "各変数1.5倍で受注は約7.6倍。一点突破より全変数の底上げが効く"])
    eq = rect(s, L, 5.0, W, 2.1, fill=None, line=BLACK, line_w=1.0)
    label(eq, "受注件数 ＝ ①提携社数 × ②顧客名簿数 × ③紹介率 × ④商談率 × ⑤成約率",
          size=22, color=NAVY)
    bw2 = 8.4
    b1 = rect(s, 6.2, 8.2, bw2, 5.6, fill=None, line=BLACK, line_w=0.75)
    h1 = rect(s, 6.2, 8.2, bw2, 1.0, fill=GRAY)
    label(h1, "改善前", size=13, color=WHITE)
    put(s, 6.2, 9.6, bw2, 2.6,
        ["提携10社 × 名簿100件", "× 紹介率5％ × 商談率60％", "× 成約率30％"],
        size=11, align=PP_ALIGN.CENTER, spacing=1.3)
    put(s, 6.2, 12.3, bw2, 1.2, "年9件", size=22, color=NAVY,
        align=PP_ALIGN.CENTER)
    put(s, 15.0, 9.6, 3.4, 2.6, ["各変数を", "1.5倍に", "↓"], size=13,
        color=NAVY, align=PP_ALIGN.CENTER, spacing=1.3)
    b2 = rect(s, 18.8, 8.2, bw2, 5.6, fill=None, line=RED, line_w=1.5)
    h2 = rect(s, 18.8, 8.2, bw2, 1.0, fill=RED)
    label(h2, "改善後", size=13, color=WHITE)
    put(s, 18.8, 9.6, bw2, 2.6,
        ["提携15社 × 名簿150件", "× 紹介率7.5％ × 商談率90％", "× 成約率45％"],
        size=11, align=PP_ALIGN.CENTER, spacing=1.3)
    put(s, 18.8, 12.3, bw2, 1.2, "年68件", size=22, color=RED,
        align=PP_ALIGN.CENTER)
    put(s, 27.6, 9.9, 5.6, 2.0, ["約7.6倍", "（1.5の5乗）"], size=16, color=NAVY,
        align=PP_ALIGN.CENTER, spacing=1.3)

    # ---------- P13 5変数 ----------
    s = page(prs, layout, 13, CH3, "5変数の打ち手一覧", None,
             "「背中を見て覚えろ」からの脱却。売れる営業の育成は仕組み")
    table(s, L, 4.4, W, 12.0, [
        ["変数", "見るべきKPI", "具体施策"],
        ["① 提携社数", "アタックリスト数\n面談率／提携率／継続率",
         "地域密着の工務店を第一ターゲットに、SEO・MEO／SUUMO／工務店団体／商工会議所・銀行から"
         "リストアップ。訪問のゴールを「提携」ではなく「勉強会への招待」に下げる（詳細はP14）"],
        ["② 顧客名簿数", "提携先のオーナー数\n名簿の共有可否（0 or 1）\n年間棟数",
         "提案の幅×案件の透明化。1年目PV・2年目蓄電池・3年目塗装・4年目床下・5年目内装と"
         "5年サイクルの提案設計を示す。顧客管理ツールで営業履歴を開示し、預ける不安を消す"],
        ["③ 紹介率", "提携先別の紹介率\n§（年間棟数に対する紹介件数）",
         "紹介が止まる阻害要因を特定する。社長のトップダウン不足／営業の理解不足／提案資料がない／"
         "上下の熱量差。対策は個社別勉強会、営業マンへの直接説明、商談への同席、個別チラシの作成"],
        ["④ 商談率", "紹介アポ→商談の到達率\n§目標：100％",
         "紹介案件は見積提出のみで済むことが多く100％を目指せる。案件管理を紙からデジタルへ移行"
         "（スプレッドシート／kintone等）し、管理表を全社で共有して抜け漏れを防ぐ"],
        ["⑤ 成約率", "商談→成約の到達率\n営業担当者別のばらつき",
         "エース営業の「心構え・アポトーク・営業ツール・クロージングトーク」を言語化しマニュアル化。"
         "営業同行・1on1・ロープレ・営業動画で組織に展開する"],
    ], widths=[4.4, 8.4, 23.25], size=11)

    # ---------- P14 勉強会 ----------
    s = page(prs, layout, 14, CH3, "①提携社数を増やす｜勉強会の設計とKPI",
             "訪問120社から締結3社。数字で追うから改善できる",
             "訪問のゴールは「提携」ではなく「勉強会への招待」。ハードルを下げるから続く")
    fw = 8.0
    fx = [3.6, 14.6, 25.6]
    ftx = [("訪問", "120社", ""), ("勉強会 参加", "15社", "20名ほど"),
           ("提携 締結", "3社", "")]
    for i, (t1, t2, t3) in enumerate(ftx):
        rect(s, fx[i], 4.9, fw, 3.0, fill=None, line=BLACK, line_w=0.75)
        put(s, fx[i], 5.3, fw, 0.8, t1, size=13, align=PP_ALIGN.CENTER)
        put(s, fx[i], 6.1, fw, 1.1, t2, size=22, color=NAVY,
            align=PP_ALIGN.CENTER)
        if t3:
            put(s, fx[i], 7.2, fw, 0.6, t3, size=9, color=GRAY,
                align=PP_ALIGN.CENTER)
    for i, rt in enumerate(["12.5％", "20％"]):
        put(s, fx[i] + fw + 0.3, 5.7, 2.4, 1.4, [rt, "▶"], size=13, color=RED,
            align=PP_ALIGN.CENTER, spacing=1.2)
    table(s, L, 8.6, 17.72, 6.4, [
        ["企　画", "やること"],
        ["頻度", "3ヵ月に1回。年間の講座構成を先に決める"],
        ["ゲスト", "メーカー・自治体を招き、集客のフックにする"],
        ["会場", "商工会議所など、地域で信頼される場所"],
        ["集客", "遅くとも開催の1か月半前から着手する"],
    ], widths=[4.0, 13.72], size=12)
    table(s, 18.97, 8.6, 17.72, 6.4, [
        ["運　営", "やること"],
        ["当日オペ", "分単位のタイムテーブルを組む"],
        ["訪問", "2週間に1回の訪問×電話で接点を保つ"],
        ["面談率", "役職者が訪問／雨の日を狙う／事務員と関係を作る"],
        ["締結", "社長同席のもと、先方の社長に条件を伝える"],
    ], widths=[4.0, 13.72], size=12)

    # ---------- P15 チェックリスト ----------
    s = page(prs, layout, 15, CH3, "自社診断チェックリスト",
             "「いいえ」が多い変数が、いま最優先で手を打つべき箇所",
             "13個以上「はい」なら仕組みは完成間近。7個以下なら伸びしろは大きい")
    col = [
        [("① 提携社数", None),
         (None, "□ 1　アタックリストが50社以上ある"),
         (None, "□ 2　ターゲットを地域密着の工務店に絞れている"),
         (None, "□ 3　訪問後に2回目以降の接触をしている"),
         (None, "□ 4　勉強会など提携以外の出口を用意している"),
         ("② 顧客名簿数", None),
         (None, "□ 5　提携先ごとのオーナー数を把握している"),
         (None, "□ 6　太陽光以外の商材も提案できる"),
         (None, "□ 7　5年先までの提案サイクルを設計している"),
         (None, "□ 8　個人情報の管理体制を説明できる")],
        [("③ 紹介率", None),
         (None, "□ 9　提携先ごとに紹介率を数値で把握している"),
         (None, "□ 10　紹介が止まる理由を相手に確認している"),
         (None, "□ 11　提携先の営業マンに直接説明する機会がある"),
         ("④ 商談率", None),
         (None, "□ 12　案件管理をデジタルで一元化している"),
         (None, "□ 13　追客・見積の抜け漏れが起きない仕組みがある"),
         ("⑤ 成約率", None),
         (None, "□ 14　営業トークをマニュアル化している"),
         (None, "□ 15　ロープレ・同行の機会を定例化している")],
    ]
    for ci, items in enumerate(col):
        x = L + ci * 18.33
        y = 4.9
        for grp, item in items:
            if grp:
                g = rect(s, x, y, 17.4, 0.75, fill=GRAY)
                label(g, grp, size=12, color=WHITE, align=PP_ALIGN.LEFT)
                y += 0.95
            else:
                put(s, x + 0.15, y, 17.2, 0.7, item, size=12,
                    align=PP_ALIGN.LEFT)
                y += 0.78

    # ---------- P16 まとめ ----------
    s = page(prs, layout, 16, "まとめ", "次の一歩",
             "自社の「売りたい」ではなく、工務店の「困っている」から始める",
             "第二の集客軸の策定。そしてさらなる業績拡大へ")
    eq = rect(s, L, 5.0, W, 2.1, fill=None, line=BLACK, line_w=1.0)
    label(eq, "受注件数 ＝ ①提携社数 × ②顧客名簿数 × ③紹介率 × ④商談率 × ⑤成約率",
          size=22, color=NAVY)
    steps = [
        ("STEP 1", "自社の数値を方程式に入れる",
         "P.15のチェックリストで弱い変数を特定してください。"),
        ("STEP 2", "勉強会を1本企画する",
         "メーカーをゲストに、地域の工務店を招く形から始めるのが取り組みやすい方法です。"),
        ("STEP 3", "紹介率を計測する",
         "提携先ごとに年間棟数と紹介件数を並べ、阻害要因を特定します。"),
    ]
    sw = 11.55
    for i, (hd, ttl, body) in enumerate(steps):
        x = L + i * (sw + 0.62)
        rect(s, x, 8.2, sw, 6.6, fill=None, line=BLACK, line_w=0.75)
        h = rect(s, x + 0.6, 8.8, 3.6, 0.8, fill=NAVY)
        label(h, hd, size=12, color=WHITE)
        put(s, x + 0.6, 10.0, sw - 1.2, 1.4, ttl, size=14, color=NAVY,
            align=PP_ALIGN.LEFT, spacing=1.25)
        put(s, x + 0.6, 11.6, sw - 1.2, 2.8, body, size=12,
            align=PP_ALIGN.LEFT, spacing=1.35)

    # ---------- CTA を最後へ ----------
    cta = list(sldIdLst)[1]          # テンプレート由来の最終ページ
    sldIdLst.remove(cta)
    sldIdLst.append(cta)

    prs.save(out_path)
    print(f"生成しました: {out_path}（{len(sldIdLst)}枚）")


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
