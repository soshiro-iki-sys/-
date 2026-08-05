#!/usr/bin/env python3
"""
PPTX のレイアウト確認用の簡易プレビュー。

このコンテナには LibreOffice も Meiryo UI も無いため、
python-pptx で読んだ座標・色・文字を PIL で描き起こして
「はみ出し」「重なり」だけを目視できるようにする。
（字形は代替フォントなので、実際の見た目とは異なる）

  python3 tools/preview_pptx.py <入力.pptx> <出力ディレクトリ> [ページ番号...]
"""
import glob
import io
import os
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

SCALE = 1280 / 37.33      # cm -> px
FONTS = sorted(glob.glob("/usr/share/fonts/**/ipag*.ttf", recursive=True)) or \
        sorted(glob.glob("/usr/share/fonts/**/*.ttc", recursive=True))
_cache = {}


def font(pt):
    px = max(8, int(pt * 1.20953))
    if px not in _cache:
        try:
            _cache[px] = ImageFont.truetype(FONTS[0], px)
        except Exception:
            _cache[px] = ImageFont.load_default()
    return _cache[px]


def cm(v):
    return Emu(v).cm if v is not None else 0


def px(v):
    return int(round(cm(v) * SCALE))


def rgb(color, default=(0, 0, 0)):
    try:
        if color is None or color.type is None:
            return default
        return tuple(bytes.fromhex(str(color.rgb)))
    except Exception:
        return default


def wrap(draw, text, fnt, maxw):
    out, cur = [], ""
    for ch in text:
        if draw.textlength(cur + ch, font=fnt) > maxw and cur:
            out.append(cur)
            cur = ch
        else:
            cur += ch
    out.append(cur)
    return out


def draw_tf(draw, tf, x, y, w, h, anchor="t"):
    lines = []
    for p in tf.paragraphs:
        runs = p.runs
        if not runs:
            lines.append(("", 12, (0, 0, 0), "l", 1.0))
            continue
        txt = "".join(r.text for r in runs)
        r0 = runs[0]
        sz = r0.font.size.pt if r0.font.size else 12
        col = rgb(r0.font.color)
        al = {1: "l", 2: "c", 3: "r"}.get(
            int(p.alignment) if p.alignment is not None else 1, "l")
        ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
        lines.append((txt, sz, col, al, ls))
    rendered = []
    for txt, sz, col, al, ls in lines:
        fnt = font(sz)
        for seg in (wrap(draw, txt, fnt, w) if txt else [""]):
            rendered.append((seg, fnt, col, al, int(sz * 1.20953 * ls * 1.25)))
    total = sum(r[4] for r in rendered)
    cy = y + (h - total) / 2 if anchor == "m" else y
    for seg, fnt, col, al, lh in rendered:
        tw = draw.textlength(seg, font=fnt)
        cx = x + (w - tw) / 2 if al == "c" else (x + w - tw if al == "r" else x)
        draw.text((cx, cy), seg, font=fnt, fill=col)
        cy += lh


def render(slide, W=1280, H=720):
    img = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(img)
    # マスターの帯（テンプレート由来）
    d.rectangle([0, 0, W, px(713232)], fill=(0, 0x20, 0x60))          # 1.98cm
    d.rectangle([0, px(7311390), W, H], fill=(0, 0x20, 0x60))         # 20.31cm

    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(sh.shapes)
                continue
            x, y = px(sh.left), px(sh.top)
            w, h = px(sh.width), px(sh.height)
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
                    img.paste(im.resize((max(w, 1), max(h, 1))), (x, y))
                except Exception:
                    d.rectangle([x, y, x + w, y + h], fill=(220, 226, 236))
                continue
            if sh.has_table:
                tbl = sh.table
                cy = y
                for row in tbl.rows:
                    cx = x
                    rh = px(row.height)
                    for ci, cell in enumerate(row.cells):
                        cw = px(tbl.columns[ci].width)
                        fill = (255, 255, 255)
                        try:
                            if cell.fill.type is not None and str(cell.fill.type).startswith("SOLID"):
                                fill = rgb(cell.fill.fore_color, (255, 255, 255))
                        except Exception:
                            pass
                        d.rectangle([cx, cy, cx + cw, cy + rh], fill=fill,
                                    outline=(0, 0, 0))
                        draw_tf(d, cell.text_frame, cx + 6, cy + 3, cw - 12,
                                rh - 6, "m")
                        cx += cw
                    cy += rh
                continue
            fill = None
            try:
                if sh.fill.type is not None and str(sh.fill.type).startswith("SOLID"):
                    fill = rgb(sh.fill.fore_color, None)
            except Exception:
                pass
            line = None
            try:
                if sh.line.fill.type is not None and str(sh.line.fill.type).startswith("SOLID"):
                    line = rgb(sh.line.color, (0, 0, 0))
            except Exception:
                pass
            if fill or line:
                d.rectangle([x, y, x + w, y + h], fill=fill, outline=line)
            if sh.has_text_frame and sh.text_frame.text.strip():
                anchor = "m" if str(sh.text_frame.vertical_anchor or "") .startswith("MIDDLE") else "t"
                draw_tf(d, sh.text_frame, x + 4, y + 2, w - 8, h - 4, anchor)

    walk(slide.shapes)
    return img


def main(src, outdir, pages=None):
    os.makedirs(outdir, exist_ok=True)
    prs = Presentation(src)
    for i, s in enumerate(prs.slides, 1):
        if pages and i not in pages:
            continue
        render(s).save(os.path.join(outdir, f"p{i:02d}.png"))
    print("プレビューを書き出しました:", outdir)


if __name__ == "__main__":
    pg = [int(a) for a in sys.argv[3:]] or None
    main(sys.argv[1], sys.argv[2], pg)
