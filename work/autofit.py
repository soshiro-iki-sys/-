# -*- coding: utf-8 -*-
"""ボックスからあふれているテキストを、収まる最大サイズまで自動縮小する"""
from pptx import Presentation
from pptx.util import Pt
import unicodedata, math
EMU=914400
prs=Presentation('final_norm.pptx'); S=prs.slides; log=[]

def wide(c): return unicodedata.east_asian_width(c) in ('W','F','A')
def flat(sh,o):
    if sh.shape_type==6:
        for c in sh.shapes: flat(c,o)
    else: o.append(sh)

def need_height(sh, scale=1.0):
    """段落ごとに、ラン単位の実寸で必要高さを推定"""
    W=(sh.width or 0)/EMU
    if W<=0.2: return 0
    inner = W - 0.12
    total=0.0
    for pa in sh.text_frame.paragraphs:
        runs=[(r.text, (r.font.size.pt if r.font.size else 18)*scale) for r in pa.runs]
        if not runs:
            total += 18*scale/72*1.2; continue
        w=0.0; mx=0.0
        for t,s in runs:
            em=s/72.0; mx=max(mx,s)
            w += sum(em if wide(c) else em*0.55 for c in t)
        lines=max(1, math.ceil(w/inner - 1e-6))
        total += lines * mx/72.0 * 1.22
    return total

def cur_sizes(sh):
    return [r.font.size.pt for pa in sh.text_frame.paragraphs for r in pa.runs if r.font.size]

# 正規化済みの共通要素は触らない
def is_grid(sh):
    L=(sh.left or 0)/EMU; T=(sh.top or 0)/EMU; W=(sh.width or 0)/EMU
    if T<0.2 and W>4: return True                    # 章タイトル
    if 0.6<T<0.95 and L<0.4 and W<8: return True     # チップ
    if 1.0<T<1.3 and L<0.4 and W>10.5: return True   # リード文
    if 6.85<T<7.25 and W>9: return True              # 出典
    if 7.0<T<7.6 and W>9: return True                # 結論バー
    return False

plan={}
for n,s in enumerate(S,1):
    o=[]
    for sh in s.shapes: flat(sh,o)
    for sh in o:
        if not sh.has_text_frame or not sh.text_frame.text.strip(): continue
        if sh.left is None or sh.height is None: continue
        if is_grid(sh): continue
        H=sh.height/EMU; W=sh.width/EMU
        if H<=0.2 or W<2.2: continue   # 幅の狭いラベルは折り返さない前提なので対象外
        sizes=cur_sizes(sh)
        if not sizes: continue
        base=max(sizes)
        need=need_height(sh)
        if need <= H+0.05: continue
        # 収まるまで 0.5pt 刻みで縮小（下限12pt、元サイズの65%まで）
        scale=1.0; floor=max(12.0/base, 0.65)
        while scale>floor:
            scale-=0.02
            if need_height(sh,scale) <= H+0.03: break
        if scale<1.0:
            plan.setdefault((n, base, round(W*2)/2), []).append((sh, scale))

# 同じ役割の並列ボックス（同サイズ・ほぼ同幅）は、最も小さくなるものに揃える
for (n, base, w), items in plan.items():
    sc=min(x[1] for x in items)
    # 同グループで縮小対象外だった兄弟も拾って揃える
    o=[]
    for sh in S[n-1].shapes: flat(sh,o)
    for sh in o:
        if not sh.has_text_frame or not sh.text_frame.text.strip(): continue
        if sh.left is None or sh.height is None or sh.width is None: continue
        if is_grid(sh) or sh.width/EMU<2.2: continue
        cs=cur_sizes(sh)
        if not cs or abs(max(cs)-base)>0.01: continue
        if abs(round(sh.width/EMU*2)/2 - w)>0.01: continue
        for pa in sh.text_frame.paragraphs:
            for r in pa.runs:
                if r.font.size: r.font.size=Pt(max(14.0, round(r.font.size.pt*sc*2)/2))
    log.append(f"P{n}: {base:.0f}pt→{base*sc:.0f}pt ×{len(items)}箇所（同幅{w}inのグループを一括）")

prs.save('fit.pptx')
print("\n".join(log)); print(f"\n--- {len(log)}箇所を縮小 / fit.pptx")
