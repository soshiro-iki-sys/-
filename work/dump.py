from pptx import Presentation
import sys, unicodedata, math
EMU=914400
prs=Presentation(sys.argv[1])
def flat(sh,o,pre=""):
    if sh.shape_type==6:
        for c in sh.shapes: flat(c,o,pre+"  ")
    else: o.append((pre,sh))
for n in [int(x) for x in sys.argv[2].split(',')]:
    s=prs.slides[n-1]; o=[]
    for sh in s.shapes: flat(sh,o)
    print(f"===== P{n} =====")
    for pre,sh in o:
        L=(sh.left or 0)/EMU; T=(sh.top or 0)/EMU; W=(sh.width or 0)/EMU; H=(sh.height or 0)/EMU
        t=""
        if sh.has_text_frame:
            t=sh.text_frame.text.replace('\n','|').replace('\x0b','|')[:70]
            szs=sorted({r.font.size.pt for pa in sh.text_frame.paragraphs for r in pa.runs if r.font.size})
            anc=sh.text_frame.vertical_anchor
            wr=sh.text_frame.word_wrap
            t=f"[{szs}] anc={anc} wrap={wr} {t}"
        print(f"{pre}{sh.shape_id:>4} {str(sh.shape_type)[:12]:12} x{L:6.2f} y{T:6.2f} w{W:6.2f} h{H:6.2f} {t}")
