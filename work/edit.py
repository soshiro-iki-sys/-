# -*- coding: utf-8 -*-
"""9月5日ヤシロ勉強会資料 修正スクリプト（添付44p版 → 45p）"""
from pptx import Presentation
from pptx.dml.color import RGBColor
import copy, sys

NAVY = RGBColor(0x00,0x20,0x60); RED = RGBColor(0xC0,0x00,0x00)
GRAY = RGBColor(0xA6,0xA6,0xA6); INK = RGBColor(0x40,0x40,0x40)
BLACK = RGBColor(0,0,0); WHITE = RGBColor(0xFF,0xFF,0xFF)

prs = Presentation('step1.pptx')
S = prs.slides
log = []

def flat(sh, out):
    if sh.shape_type == 6:
        for c in sh.shapes: flat(c, out)
    else: out.append(sh)

def shapes(idx):
    out = []
    for sh in S[idx-1].shapes: flat(sh, out)
    return out

def norm(t):
    return t.replace('\x0b', '').replace('\n', '').replace(' ', '')

def find(idx, needle, nth=0):
    hits = [sh for sh in shapes(idx)
            if sh.has_text_frame and norm(needle) in norm(sh.text_frame.text)]
    if len(hits) <= nth:
        raise SystemExit(f"NOT FOUND on P{idx}: {needle!r}")
    return hits[nth]

def set_lines(sh, lines, page='', label=''):
    tf = sh.text_frame
    paras = list(tf.paragraphs)
    for i, para in enumerate(paras):
        want = lines[i] if i < len(lines) else ""
        for br in para._p.findall(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}br'):
            para._p.remove(br)          # 段落内の強制改行を除去
        if para.runs:
            para.runs[0].text = want
            for r in para.runs[1:]:
                r._r.getparent().remove(r._r)
        elif want:
            para.text = want
    if len(lines) > len(paras):
        raise SystemExit(f"too many lines for {label} on P{page}")
    log.append(f"P{page} {label}: {' / '.join(lines)[:70]}")

def set_text(idx, needle, new, label='', nth=0):
    sh = find(idx, needle, nth)
    set_lines(sh, [new], idx, label or needle[:16])

def set_multi(idx, needle, lines, label='', nth=0):
    sh = find(idx, needle, nth)
    set_lines(sh, lines, idx, label or needle[:16])

def cell_text(tbl, r, c, new):
    cell = tbl.rows[r].cells[c]
    para = cell.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = new
        for x in para.runs[1:]: x._r.getparent().remove(x._r)
    else:
        para.text = new

def table_of(idx, nth=0):
    tbls = [sh for sh in shapes(idx) if sh.has_table]
    return tbls[nth].table

def strike(sh):
    for para in sh.text_frame.paragraphs:
        for r in para.runs:
            r._r.get_or_add_rPr().set('strike', 'sngStrike')

def color(sh, rgb):
    for para in sh.text_frame.paragraphs:
        for r in para.runs:
            r.font.color.rgb = rgb

# ============ 章タイトルの表記統一 ============
for pg in (10, 11):
    set_text(pg, "1.新潟県のエネルギー事情", "1.長岡市のエネルギー事情と災害対策", "章タイトル")

# ============ P7 北陸地方に戻す（A案） ============
set_text(7, "新潟県を含む新潟地方も例外",
         "電気代は全国で上がり続けており、新潟県を含む北陸地方も例外ではない。", "リード")
set_text(7, "新潟地方の平均電気代推移",
         "北陸地方（新潟県を含む）の平均電気代推移（二人以上の世帯・月額）", "グラフ見出し")
set_text(7, "※新潟地方・二人以上の世帯ベース",
         "※北陸地方・二人以上の世帯ベース（家計調査の地方区分では新潟県は北陸地方に含まれる）　出典：総務省「家計調査」を基に作成", "出典")

# ============ P11 他社名の残存を削除 ============
set_text(11, "出典：山岸調べ",
         "出典：内閣府・各電力会社の公表資料をもとに作成", "出典")

# ============ P14 結論バーを言い切り型に ============
set_text(14, "生活環境に大きく影響しますよね",
         "雪国の停電は、命に直結します", "結論バー")

# ============ P20 ハイブ長岡に統一・降雪量に修正 ============
set_text(20, "自社発電所（柏崎・長岡）の実測値",
         "出典：長岡市「雪国対応の太陽光発電実証実験」（ハイブ長岡）の実測値", "出典")
t = table_of(20)
for r in range(len(t.rows)):
    if "除雪量" in t.rows[r].cells[2].text:
        cell_text(t, r, 2, "降雪量（cm）"); log.append("P20 表ヘッダ: 降雪量（cm）")

# ============ P22 経済メリットを2段構成に ============
set_text(22, "5kWシステムで年間どれだけ得になるか",
         "5kWシステムの経済メリットを、5年目以降まで含めて試算する。", "リード")
t = table_of(22)
rows = [
    ("項目", "太陽光のみ", "太陽光＋蓄電池"),
    ("自家消費率", "約30％", "約65％"),
    ("当初4年（売電24円）", "約141,000円/年", "約153,000円/年"),
    ("5年目以降（売電8.3円）", "約81,000円/年", "約123,000円/年"),
    ("5年目以降の差", "―", "＋約42,000円/年"),
    ("10年間の累計", "約106万円", "約135万円"),
]
for r, vals in enumerate(rows):
    for c, v in enumerate(vals):
        cell_text(t, r, c, v)
log.append("P22 表: 当初4年／5年目以降／10年累計の3段に変更")
set_text(22, "単価はご契約プランで変わる", "5年目から差が大きく開きます", "右カード見出し")
set_multi(22, "・契約プラン（従量電灯／時間帯別）", [
    "売電単価は当初4年が24円、",
    "5年目以降は8.3円に下がります。",
    "",
    "売った電気の価値が下がるぶん、",
    "自分で使う（自家消費）ほど得になる。",
    "",
    "蓄電池は自家消費率を約30％→約65％に引き上げる設備です。",
], "右カード本文")
set_text(22, "前提：発電5,500kWh",
         "前提：発電5,500kWh／買電単価30円/kWh（新潟県の一般的な家庭用単価）／売電単価 当初4年24円・5年目以降8.3円（2026年度FIT）", "前提")
set_text(22, "蓄電池を足すと、経済メリットはさらに伸びます",
         "5年目以降、蓄電池の有無で差が大きく開きます", "結論バー")

# ============ P24 ZEH住宅を受付終了に ============
zeh_head = find(24, "ZEH住宅")
zeh_head.fill.solid(); zeh_head.fill.fore_color.rgb = GRAY
color(zeh_head, WHITE)
zeh_amt = find(24, "55万円")
color(zeh_amt, GRAY); strike(zeh_amt)
zeh_note = [sh for sh in shapes(24) if sh.has_text_frame and sh.text_frame.text.strip() == "一律（新築）"]
set_lines(zeh_note[0], ["受付終了"], 24, "ZEH受付終了")
color(zeh_note[0], RED)
set_text(24, "※ZEH／ZEH＋の対象要件・金額は要確認",
         "出典：長岡市「令和8年度 雪国長岡での再エネ導入促進補助金」　※ZEH住宅（55万円）は受付終了。ZEH＋住宅は継続。残予算は要確認", "出典")

# ============ P25 受付終了を実例として使う ============
set_text(25, "年度内でも、早期に締め切られることがあります",
         "実際にZEH住宅枠は、すでに受付を終了しています", "③の説明")

# ============ P31 表記ゆれ ============
set_text(31, "太陽光5kw＋蓄電池のセット", "太陽光5kW＋蓄電池のセット：200～300万円", "kW表記")

# ============ P34 文体統一・店舗表現の削除 ============
set_text(34, "太陽光・蓄電池は30年以上使えるものですが",
         "太陽光・蓄電池は長く使える設備だが、定期的なメンテナンスが必要になる。", "リード")
set_multi(34, "各店舗から30分圏内に",
          ["地域密着だからこそ、", "不具合があればすぐに駆けつけられる"], "対応②")

# ============ P38 ①地域密着（基準→参考の構成に） ============
set_text(38, "設立から62年、盤石な経営基盤で雪国長岡に",
         "業者選びでは、その会社が地域にどれだけ根を張っているかを確認したい。", "リード")
set_text(38, "信頼と実績", "地域密着で確認する3つのポイント", "タブ")
set_text(38, "悪質な訪問販売のような",
         "その地域で何年営業しているか（＝売り逃げができない立場か）", "①")
set_text(38, "長岡市周辺に対応エリアを絞った",
         "何分で駆けつけられるエリアか（対応エリアの広さ）", "②")
set_text(38, "設立から62年の盤石な経営基盤",
         "地元での施工実績がどれだけあるか", "③")

# ============ P39 ②適正価格（誤字＋基準化） ============
set_text(39, "高すぎるもつもりはもちろん",
         "高すぎるのはもちろん、安すぎる見積りにも理由がある。", "リード")
set_text(39, "【安すぎる見積りの危うさ】", "【見積りで確認する3点】", "左タブ")
set_text(39, "【ヤシロの考え方】", "【参考：ヤシロの場合】", "右タブ")
set_text(39, "工事費を削るために施工の質が落ちる",
         "内訳が開示されるか。「一式」だけの見積書は比較できない", "左①")
set_text(39, "売ったあとの不具合に誠実な対応をしない",
         "相場から極端に安くないか。工事の質が落ちている場合がある", "左②")
set_text(39, "「安すぎる価格と押し売り」は、比較させない",
         "その場で決めさせようとしないか。「今日だけ特別価格」は要注意", "左③")
set_text(39, "理由を添えて、適正な価格で販売いたします",
         "理由を添えて、適正な価格でご提案します", "右②")
set_text(39, "じっくり検討していただして",
         "他社と比較する時間をお取りください。即決はお願いしません", "右③")

# ============ P40 ③メーカーと保証 ============
set_text(40, "③国産メーカーの取り扱い", "③メーカーと保証", "チップ")
set_text(40, "弊社が国産メーカーの長州産業を扱う理由",
         "設備そのものより、不具合が起きたときに誰が動くかを確認したい。", "リード")
set_text(40, "長州産業を取り扱う理由", "保証で確認する3つのポイント", "左カード見出し")
set_text(40, "国内生産のメーカーである", "保証年数（機器・出力・自然災害）", "左①")
set_text(40, "生産から供給まで国内で完結している", "年数だけでなく、無償かどうかも確認する", "左①補足")
set_text(40, "不具合時の対応がしっかりしている", "雪害が自然災害補償の対象か", "左②")
set_text(40, "連絡がつき、話が前に進む。これがいちばん大きい", "雪国では、ここが対象外だと意味がない", "左②補足")
set_text(40, "雪害が自然災害補償の対象", "不具合のとき、誰が動くか", "左③", nth=1)
set_text(40, "雪国で使ううえで、これは大きな差になる", "メーカー・販売店・施工店のどこが窓口かを確認する", "左③補足")
set_text(40, "保証内容", "参考：弊社取り扱い（長州産業）の保証内容", "表見出し")
t = table_of(40)
for r in range(len(t.rows)):
    if "施工保証" in t.rows[r].cells[0].text:
        cell_text(t, r, 0, "自然災害補償")
        cell_text(t, r, 1, "10年")
        log.append("P40 表: 自然災害補償10年に差し替え（施工保証の行を廃止）")

# ============ P41 ④現地調査 ============
set_text(41, "④現地調査の丁寧さ", "④現地調査でどこを見るか", "チップ")
set_text(41, "工務店として住宅を扱ってきたからこそ",
         "屋根に載せる工事である以上、屋根そのものを見てもらう必要がある。", "リード")
set_text(41, "屋根の状態によっては「今はやめておきましょう」",
         "参考：ヤシロは工務店として屋根の下地まで確認し、向かない場合は「今はやめておきましょう」とお伝えします。", "下部")

# ============ P42 新規：業者選びチェックリスト ============
set_text(42, "④現地調査の丁寧さ", "業者選びチェックリスト", "チップ")
set_text(42, "工務店として住宅を扱ってきたからこそ",
         "他社と比較するときに、そのまま使える確認項目をまとめた。", "リード")
set_text(42, "現地調査で確認するべき項目", "どの会社にも、この6つを聞いてください", "見出し")
CHK = [
    ("屋根材の種類と状態", "地域で何年営業していますか",
     "瓦・スレート・金属。材質ごとに固定方法が変わる", "会社の設立年と、太陽光を始めた年を聞く"),
    ("下地・野地板の健全性", "何分で駆けつけられますか",
     "表面ではなく、留め付ける下地が持つかどうか", "対応エリアと、緊急時の体制を聞く"),
    ("雨仕舞い", "見積りの内訳を出せますか",
     "穴を開ける以上、水の処理をどうするか", "「一式」ではなく、項目ごとの金額を出せるか"),
    ("積雪荷重", "他社と比較する時間をくれますか",
     "パネルと雪の重さに、構造が耐えられるか", "即決を迫る会社は、その時点で除外してよい"),
    ("落雪の方向と着地点", "保証は何年で、雪害は対象ですか",
     "隣家・道路・カーポートに落とさない配置か", "機器・出力・自然災害の3つを確認する"),
    ("電気容量・分電盤", "屋根の下地まで見てくれますか",
     "蓄電池・V2Hを入れられる余地があるか", "表面だけの調査では、後の雨漏りを防げない"),
]
for old_t, new_t, old_d, new_d in CHK:
    set_text(42, old_t, new_t, "チェック項目")
    set_text(42, old_d, new_d, "チェック補足")
set_text(42, "屋根の状態によっては「今はやめておきましょう」",
         "この6つに即答できない会社は、避けたほうが安全です。", "下部")
set_text(42, "屋根に載せる前に、屋根そのものを見ます",
         "迷ったら、この6つを聞いてみてください", "結論バー")

# ============ P43 まとめ（ZEH終了を反映） ============
set_text(43, "補助金は早い者勝ち。動くなら早めに",
         "補助金は早い者勝ち。ZEH住宅枠はすでに終了しています", "まとめ④")

prs.save('out.pptx')
print("\n".join(log))
print(f"\n--- {len(log)} 箇所を修正 / 保存: out.pptx")
