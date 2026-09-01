# -*- coding: utf-8 -*-
"""2次修正：文字量をボックス幅に合わせて調整＋P40出典の復旧"""
from pptx import Presentation
from pptx.dml.color import RGBColor
BLACK = RGBColor(0,0,0)
prs = Presentation('out.pptx'); S = prs.slides; log=[]
NS='{http://schemas.openxmlformats.org/drawingml/2006/main}'

def flat(sh,out):
    if sh.shape_type==6:
        for c in sh.shapes: flat(c,out)
    else: out.append(sh)
def shapes(i):
    o=[]
    for sh in S[i-1].shapes: flat(sh,o)
    return o
def norm(t): return t.replace('\x0b','').replace('\n','').replace(' ','')
def find(i,needle,nth=0,exact=False):
    if exact:
        h=[s for s in shapes(i) if s.has_text_frame and norm(s.text_frame.text)==norm(needle)]
    else:
        h=[s for s in shapes(i) if s.has_text_frame and norm(needle) in norm(s.text_frame.text)]
    if len(h)<=nth: raise SystemExit(f"NOT FOUND P{i}: {needle!r}")
    return h[nth]
def set_lines(sh,lines,page,label):
    tf=sh.text_frame; paras=list(tf.paragraphs)
    if len(lines)>len(paras): raise SystemExit(f"too many lines P{page} {label}")
    for i,pa in enumerate(paras):
        want = lines[i] if i<len(lines) else ""
        for br in pa._p.findall(NS+'br'): pa._p.remove(br)
        if pa.runs:
            pa.runs[0].text=want
            for r in pa.runs[1:]: r._r.getparent().remove(r._r)
        elif want: pa.text=want
    log.append(f"P{page} {label}: {' / '.join(lines)[:60]}")
def st(i,needle,new,label,nth=0,exact=False):
    set_lines(find(i,needle,nth,exact),[new],i,label)
def sm(i,needle,lines,label,nth=0):
    set_lines(find(i,needle,nth),lines,i,label)
def recolor(sh,rgb):
    for pa in sh.text_frame.paragraphs:
        for r in pa.runs: r.font.color.rgb=rgb

# ---- P7 見出し・サマリーを枠内に ----
st(7,"北陸地方（新潟県を含む）の平均電気代推移","北陸地方の平均電気代推移（二人以上の世帯・月額）","グラフ見出し")
sm(7,"2025年平均でも月13,951円",
   ["2025年平均でも月13,951円。",
    "直近12ヵ月平均は月17,166円と、負担は高止まりしています。"],"サマリー")

# ---- P22 右カード本文を枠内に ----
sm(22,"売電単価は当初4年が24円",[
    "売電単価は当初4年が24円、",
    "5年目以降は8.3円に下がります。",
    "",
    "売った電気の価値が下がるぶん、",
    "自分で使うほど得になります。",
    "",
    "蓄電池は自家消費率を高める設備です。",
],"右カード本文")

# ---- P38 1行に収める ----
st(38,"地域密着で確認する3つのポイント","地域密着で見る3点","タブ")
st(38,"その地域で何年営業しているか","その地域で何年営業しているか","①")
st(38,"何分で駆けつけられるエリアか","何分で駆けつけられるエリアか","②")
sh=find(38,"地元での施工実績がどれだけあるか")
set_lines(sh,["地元での施工実績はどれだけあるか"],38,"③"); recolor(sh,BLACK)

# ---- P39 2行に収める ----
st(39,"内訳が開示されるか","内訳が開示されるか。「一式」だけでは比較できない","左①")
st(39,"相場から極端に安くないか","相場から極端に安くないか。工事の質が落ちる","左②")
st(39,"その場で決めさせようとしないか","その場で決めさせようとしないか","左③")
st(39,"見積りの内訳を開示し","見積りの内訳を開示し、何にいくらかかるかを説明します","右①")
st(39,"他社と比較する時間をお取りください","他社と比較する時間をお取りください","右③")
st(39,"価格内訳とその理由を説明できる会社","価格の理由を説明できる会社かどうかで判断を","結論バー")

# ---- P40 出典の復旧＋左カードを枠内に ----
st(40,"参考：弊社取り扱い（長州産業）の保証内容",
   "出典：長州産業の公表保証内容（2026年時点）　※適用条件・最新内容はメーカー資料をご確認ください","出典（復旧）")
st(40,"保証内容","参考：長州産業（弊社取り扱い）の保証内容","表見出し",exact=True)
st(40,"保証で確認する3つのポイント","保証で確認する3点","左カード見出し")
st(40,"保証年数（機器・出力・自然災害）","保証の年数と無償かどうか","左①")
st(40,"年数だけでなく、無償かどうかも確認する","機器・出力・自然災害の3つを確認する","左①補足")
st(40,"雪害が自然災害補償の対象か","雪害は補償の対象か","左②")

prs.save('out2.pptx')
print("\n".join(log)); print(f"\n--- {len(log)} 箇所を調整 / out2.pptx")
