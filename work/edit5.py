# -*- coding: utf-8 -*-
"""5次修正：ZEH＋（100万円）も受付終了を反映"""
from pptx import Presentation
from pptx.dml.color import RGBColor
GRAY=RGBColor(0xA6,0xA6,0xA6); RED=RGBColor(0xC0,0x00,0x00); WHITE=RGBColor(0xFF,0xFF,0xFF)
NS='{http://schemas.openxmlformats.org/drawingml/2006/main}'
prs=Presentation('done.pptx'); S=prs.slides; log=[]
def flat(sh,o):
    if sh.shape_type==6:
        for c in sh.shapes: flat(c,o)
    else: o.append(sh)
def shapes(i):
    o=[]
    for sh in S[i-1].shapes: flat(sh,o)
    return o
def norm(t): return t.replace('\x0b','').replace('\n','').replace(' ','')
def find(i,needle,nth=0,exact=False):
    if exact:
        h=[s for s in shapes(i) if s.has_text_frame and norm(s.text_frame.text)==norm(needle)]
    else:
        h=[s for s in shapes(i) if s.has_text_frame and norm(needle) in norm(s.text_frame.text)]
    if len(h)<=nth: raise SystemExit(f"NOT FOUND P{i}: {needle!r}")
    return h[nth]
def setp(sh,lines):
    for k,pa in enumerate(list(sh.text_frame.paragraphs)):
        want=lines[k] if k<len(lines) else ""
        for br in pa._p.findall(NS+'br'): pa._p.remove(br)
        if pa.runs:
            pa.runs[0].text=want
            for r in pa.runs[1:]: r._r.getparent().remove(r._r)
        elif want: pa.text=want
def st(i,needle,lines,label,nth=0,exact=False):
    if isinstance(lines,str): lines=[lines]
    setp(find(i,needle,nth,exact),lines); log.append(f"P{i} {label}: {' / '.join(lines)[:56]}")
def color(sh,rgb):
    for pa in sh.text_frame.paragraphs:
        for r in pa.runs: r.font.color.rgb=rgb
def strike(sh):
    for pa in sh.text_frame.paragraphs:
        for r in pa.runs: r._r.get_or_add_rPr().set('strike','sngStrike')

# ---- P24 ZEH＋も受付終了に ----
h=find(24,"ZEH＋住宅",exact=True)
h.fill.solid(); h.fill.fore_color.rgb=GRAY; color(h,WHITE)
a=find(24,"100万円",exact=True); color(a,GRAY); strike(a)
n=find(24,"一律（新築）",exact=True)          # 残っているのは ZEH＋ 側のみ
setp(n,["受付終了"]); color(n,RED)
log.append("P24 ZEH＋住宅：グレー＋取り消し線＋『受付終了』")
st(24,"太陽光＋蓄電池（既築もOK）",["太陽光＋蓄電池の枠は継続中","最大91.4万円"],"左まとめカード")
st(24,"ZEH＋と併用できる新築なら",["新築向けのZEH・ZEH＋枠は","受付終了しています"],"右まとめカード")
st(24,"※ZEH住宅（55万円）は受付終了",
   "出典：長岡市「令和8年度 雪国長岡での再エネ導入促進補助金」　※ZEH（55万円）・ZEH＋（100万円）はいずれも受付終了。太陽光・蓄電池の枠は継続。残予算は要確認","出典")
st(24,"補助金をフル活用すれば","太陽光＋蓄電池なら、最大91.4万円の補助が使えます","結論バー")

# ---- P25 / P43 ----
st(25,"実際にZEH住宅枠は","実際にZEH・ZEH＋の枠は、すでに受付を終了しています","③の説明")
st(43,"補助金は早い者勝ち。ZEH住宅枠はすでに終了しています",
   "補助金は早い者勝ち。ZEH・ZEH＋の枠はすでに終了","まとめ④")

prs.save('done2.pptx')
print("\n".join(log)); print(f"\n--- {len(log)} 件 / done2.pptx")
