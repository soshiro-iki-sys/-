# -*- coding: utf-8 -*-
"""6次修正：191.4万円の記述を残し、受付終了の注記を足す"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
RED=RGBColor(0xC0,0x00,0x00); WHITE=RGBColor(0xFF,0xFF,0xFF); INK=RGBColor(0x40,0x40,0x40)
NS='{http://schemas.openxmlformats.org/drawingml/2006/main}'
prs=Presentation('final_v2.pptx'); S=prs.slides; log=[]
def flat(sh,o):
    if sh.shape_type==6:
        for c in sh.shapes: flat(c,o)
    else: o.append(sh)
def shapes(i):
    o=[]
    for sh in S[i-1].shapes: flat(sh,o)
    return o
def norm(t): return t.replace('\x0b','').replace('\n','').replace(' ','')
def find(i,needle,nth=0):
    h=[s for s in shapes(i) if s.has_text_frame and norm(needle) in norm(s.text_frame.text)]
    if len(h)<=nth: raise SystemExit(f"NOT FOUND P{i}: {needle!r}")
    return h[nth]
def setp(sh,lines):
    for k,pa in enumerate(list(sh.text_frame.paragraphs)):
        want=lines[k] if k<len(lines) else ""
        for br in pa._p.findall(NS+'br'): pa._p.remove(br)
        if pa.runs:
            pa.runs[0].text=want
            for r in pa.runs[1:]: r._r.getparent().remove(r._r)
        elif want: pa.text=want
def st(i,needle,lines,label,nth=0):
    if isinstance(lines,str): lines=[lines]
    setp(find(i,needle,nth),lines); log.append(f"P{i} {label}: {' / '.join(lines)[:56]}")

# ---- 右まとめカードに191.4万円を戻し、取り消し線を掛ける ----
card = find(24,"新築向けのZEH・ZEH＋枠は")
setp(card, ["ZEH＋と併用できる新築なら", "最大191.4万円"])
paras = list(card.text_frame.paragraphs)
for r in paras[1].runs:                       # 金額側だけ取り消し線
    r._r.get_or_add_rPr().set('strike','sngStrike')
log.append("P24 右まとめカード: ZEH＋と併用できる新築なら / 最大191.4万円（取り消し線）")

# ---- 「受付終了」バッジを重ねる ----
sl = S[23]
badge = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.62), Inches(5.36), Inches(1.4), Inches(0.42))
badge.fill.solid(); badge.fill.fore_color.rgb = RED
badge.line.fill.background(); badge.shadow.inherit = False
tf = badge.text_frame; tf.word_wrap = False; tf.margin_left = tf.margin_right = 0
tf.margin_top = tf.margin_bottom = 0; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
run = p0.add_run(); run.text = "受付終了"
run.font.name = "Meiryo UI"; run.font.size = Pt(15); run.font.bold = True; run.font.color.rgb = WHITE
log.append("P24 『受付終了』バッジを右カードに追加")

# ---- 2枚のカードの下に注記を1行 ----
note = sl.shapes.add_textbox(Inches(0.6), Inches(6.62), Inches(10.5), Inches(0.34))
tf = note.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
run = p0.add_run()
run.text = "※ZEH（55万円）・ZEH＋（100万円）は予算上限に達し、すでに受付終了。太陽光・蓄電池の枠は継続中です。"
run.font.name = "Meiryo UI"; run.font.size = Pt(14); run.font.bold = True; run.font.color.rgb = RED
log.append("P24 注記行を追加（ZEH・ZEH＋は受付終了／太陽光・蓄電池は継続）")

# ---- 出典と結論バー ----
st(24,"出典：長岡市","出典：長岡市「令和8年度 雪国長岡での再エネ導入促進補助金」　※受付状況・残予算は変動します","出典")
st(24,"太陽光＋蓄電池なら、最大91.4万円の補助が使えます",
   "太陽光＋蓄電池なら、今も最大91.4万円が使えます","結論バー")

prs.save('final_v3.pptx')
print("\n".join(log)); print(f"\n--- {len(log)} 件 / final_v3.pptx")
