# -*- coding: utf-8 -*-
"""レンダリング検証で見つかった体裁の乱れ（あふれ・重なり・サイズ不揃い）をまとめて是正する"""
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from measure import flat, EMU, NS

prs = Presentation('final_norm.pptx')
S = prs.slides
log = []

def shapes(n):
    o = []
    for sh in S[n-1].shapes: flat(sh, o)
    return {sh.shape_id: sh for sh in o}

def set_sz(sh, pt):
    for pa in sh.text_frame.paragraphs:
        for r in pa.runs:
            r.font.size = Pt(pt)

def box(sh, x=None, y=None, w=None, h=None):
    if x is not None: sh.left = Emu(int(x*EMU))
    if y is not None: sh.top = Emu(int(y*EMU))
    if w is not None: sh.width = Emu(int(w*EMU))
    if h is not None: sh.height = Emu(int(h*EMU))

def margins(sh, l=None, r=None, t=None, b=None):
    tf = sh.text_frame
    if l is not None: tf.margin_left = Emu(int(l*EMU))
    if r is not None: tf.margin_right = Emu(int(r*EMU))
    if t is not None: tf.margin_top = Emu(int(t*EMU))
    if b is not None: tf.margin_bottom = Emu(int(b*EMU))

def relines(sh, lines):
    """1段落＋<a:br> で lines を書き直す（先頭ランの書式を保つ）"""
    tf = sh.text_frame
    pa = tf.paragraphs[0]
    for br in pa._p.findall(NS+'br'): pa._p.remove(br)
    runs = pa.runs
    if not runs: return
    keep = runs[0]
    for r in runs[1:]: r._r.getparent().remove(r._r)
    keep.text = lines[0]
    for extra in lines[1:]:
        br = keep._r.makeelement(NS+'br', {}); pa._p.append(br)
        new = keep._r.makeelement(NS+'r', {})
        import copy
        new = copy.deepcopy(keep._r); pa._p.append(new)
        for t in new.findall(NS+'t'): t.text = extra
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)

# ── P1 表紙タイトル：1行に収める ──────────────────────────
d = shapes(1)
box(d[6], x=0.15, w=11.39)
for pa in d[6].text_frame.paragraphs:
    for r in pa.runs:
        if r.font.size and r.font.size.pt > 50: r.font.size = Pt(56)
log.append("P1 表紙タイトル 60→56pt・枠を 11.39in に拡張（折り返し解消）")

# ── P3 6つの吹き出し：サイズ 24/28pt 混在 → 22pt に統一し、2行で揃える ──
d = shapes(3)
CH = {11: ["「長岡の雪じゃ", "発電しないでしょ？」"],
      13: ["「屋根が雪で", "傷むのでは？」"],
      15: ["「訪問販売が", "怖い」"],
      17: ["「業者によって", "言うことが違う」"],
      19: ["「どのメーカーが", "いいの？」"],
      21: ["「いつ導入するのが", "お得なの？」"]}
for sid, lines in CH.items():
    sh = d[sid]
    card_x = round((sh.left/EMU - 0.12), 2)
    box(sh, x=card_x+0.05, w=3.35)
    margins(sh, l=0.05, r=0.05)
    relines(sh, lines)
    set_sz(sh, 22)
log.append("P3 吹き出し6枚 24/28pt混在→22pt統一・全て2行で揃え・枠幅を3.35inに統一")

# ── P15 リード文：改行位置を再配分して2行に収める ─────────────
d = shapes(15)
tf = d[10].text_frame
ps = tf.paragraphs
ps[0].runs[0].text = "このように復旧まで時間がかかる場合もある中、太陽光発電や"
ps[1].runs[0].text = "蓄電池のない家庭で停電が起こると生活に関わる深刻な問題が多発します。"
log.append("P15 リード文の改行位置を再配分（3行あふれ→2行）")

# ── P23 リード文が2行のため、表を 2.02in に下げて重なりを解消 ─────
for sh in S[22].shapes:
    if sh.has_table:
        box(sh, y=2.02)
        for r in sh.table.rows: r.height = Emu(int(0.855*EMU))
log.append("P23 表を y2.02in へ移動・行高0.855inに調整（2行リード文との重なり解消）")

# ── P31〜P35 デメリット／対応カード：見出し28pt・本文22ptに統一 ────
for n in (31, 32, 33, 34, 35):
    d = shapes(n)
    for sh in d.values():
        if not sh.has_text_frame or not sh.text_frame.text.strip(): continue
        szs = [r.font.size.pt for pa in sh.text_frame.paragraphs for r in pa.runs if r.font.size]
        if not szs: continue
        if abs(max(szs)-32) < 0.01 and sh.text_frame.text.startswith("【"):
            set_sz(sh, 28)
        elif abs(max(szs)-28) < 0.01:
            set_sz(sh, 22)
    # 本文枠の幅を左右で揃える
    for sh in d.values():
        if sh.has_text_frame and sh.width and 3.9 < sh.width/EMU < 4.3 and sh.top and sh.top/EMU > 2.5:
            box(sh, w=4.38)
log.append("P31〜P35 カード見出し32→28pt／本文28→22pt に統一・本文枠幅を4.38inに統一")

# ── P37 相談トラブル4件の見出し：20→18pt（折り返し解消） ─────────
d = shapes(37)
for sid in (13, 16, 19, 22):
    set_sz(d[sid], 18); box(d[sid], w=4.15)
log.append("P37 トラブル見出し4件 20→18pt・枠幅4.15inに統一（折り返し解消）")

# ── P38 ①の文言：改行位置を指定して語中改行を解消・3項目の左端を揃える ──
d = shapes(38)
relines(d[40], ["悪質な訪問販売のような", "「売り逃げ」などがないかどうか"])
for sid in (40, 38, 26):
    box(d[sid], x=1.16, w=10.33)
log.append("P38 ①の改行位置を指定（語中改行を解消）・3項目の左端と幅を統一")

# ── P39 見積り3点／ヤシロ3点：24→20pt に統一し、手動改行を解除 ────
d = shapes(39)
for sid in (16, 19, 22, 25, 28, 31):
    sh = d[sid]
    txt = sh.text_frame.text.replace('\x0b', '').replace('\n', '')
    relines(sh, [txt]); set_sz(sh, 20); box(sh, w=4.20)
log.append("P39 説明文6件 24→20pt に統一・手動改行を解除・枠幅4.20inに統一")

# ── P41 現地調査6項目の見出し：枠を広げて折り返しを解消 ───────────
d = shapes(41)
for sid in (14, 19, 24, 29, 34, 39):
    box(d[sid], w=2.81); margins(d[sid], l=0.02, r=0.02); set_sz(d[sid], 20)
log.append("P41 項目見出し6件 24→20pt・枠幅2.81inに拡張（折り返し解消）")

# ── P43 参加者特典カード：文字サイズを枠に合わせる ─────────────
d = shapes(43)
set_sz(d[15], 28)
for sid in (17, 19):
    for pa in d[sid].text_frame.paragraphs:
        for r in pa.runs: r.font.size = Pt(16)
log.append("P43 小冊子名 32→28pt・説明/連絡先 18→16pt（折り返し解消）")

prs.save('polished.pptx')
print("\n".join("・"+x for x in log))
print(f"\n--- {len(log)}件を是正 → polished.pptx")
