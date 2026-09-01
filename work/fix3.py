# -*- coding: utf-8 -*-
"""3巡目：残りの折り返し・重なりを是正する"""
import copy
from pptx import Presentation
from pptx.util import Pt, Emu
from measure import flat, em_len, seg_runs, EMU, NS

prs = Presentation('polished2.pptx')
S = prs.slides
log = []

def shapes(n):
    o = []
    for sh in S[n-1].shapes: flat(sh, o)
    return {sh.shape_id: sh for sh in o}

def set_sz(sh, pt):
    for pa in sh.text_frame.paragraphs:
        for r in pa.runs: r.font.size = Pt(pt)

def box(sh, x=None, y=None, w=None, h=None):
    if x is not None: sh.left = Emu(int(x*EMU))
    if y is not None: sh.top = Emu(int(y*EMU))
    if w is not None: sh.width = Emu(int(w*EMU))
    if h is not None: sh.height = Emu(int(h*EMU))

def widest(sh):
    mx = 0.0
    for pa in sh.text_frame.paragraphs:
        base = max([r.font.size.pt for r in pa.runs if r.font.size] or [18])
        for seg in seg_runs(pa):
            mx = max(mx, sum(em_len(t)*((sz or base)/72.0) for t, sz in seg))
    return mx

def fit_line(sh, ratio=0.94, pad=0.10):
    avail = (sh.width/EMU - pad) * ratio
    w = widest(sh)
    if w <= avail: return False
    f = avail / w
    for pa in sh.text_frame.paragraphs:
        for r in pa.runs:
            if r.font.size: r.font.size = Pt(round(r.font.size.pt*f*2)/2)
    return True

# ── P8 年間負担増の1行／注記の1行 ──────────────────────────
d = shapes(8)
fit_line(d[24]); set_sz(d[25], 15)
log.append("P8 「年間 約25,000円 の負担増」を1行に収まるサイズへ再調整・注記16→15pt")

# ── P18 リード文 24→22pt（1行に収める） ───────────────────
set_sz(shapes(18)[9], 22)
log.append("P18 リード文 24→22pt（1行に収める）")

# ── P19 発電量カード：単位を改行で送り 36→34pt・枠高2.30in ────────
d = shapes(19); sh = d[15]
pa = sh.text_frame.paragraphs[0]
runs = pa.runs
if len(runs) >= 3 and not pa._p.findall(NS+'br'):
    br = runs[0]._r.makeelement(NS+'br', {})
    runs[2]._r.addprevious(br)
for r in pa.runs:
    if r.font.size and r.font.size.pt > 30: r.font.size = Pt(34)
box(sh, h=2.30)
log.append("P19 発電量カード：単位を改行で送り 36→34pt・枠高2.30inへ（注記との重なり解消）")

prs.save('polished3.pptx')
print("\n".join("・"+x for x in log))
