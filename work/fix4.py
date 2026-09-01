# -*- coding: utf-8 -*-
"""4巡目：リード文の体裁を全ページで統一する（1行=h0.50 / 2行=h0.86・改行位置を均等配分）"""
import copy
from pptx import Presentation
from pptx.util import Emu
from measure import flat, em_len, EMU, NS

prs = Presentation('polished3.pptx')
PUNC = '、。」』）'
log = []

def best_split(text, avail, pt):
    """2行に均等配分する分割位置（句読点・閉じ括弧の直後を優先）"""
    cands = []
    for i in range(1, len(text)):
        w1 = em_len(text[:i]) * pt / 72.0
        w2 = em_len(text[i:]) * pt / 72.0
        if w1 > avail or w2 > avail: continue
        cands.append((text[i-1] in PUNC, -max(w1, w2), i))
    if not cands: return None
    return max(cands)[2]

for n, s in enumerate(prs.slides, 1):
    o = []
    for sh in s.shapes: flat(sh, o)
    for sh in o:
        if not sh.has_text_frame or sh.left is None: continue
        if not (abs(sh.left/EMU-0.16) < 0.03 and abs(sh.top/EMU-1.12) < 0.03 and sh.width/EMU > 10.5):
            continue
        tf = sh.text_frame
        pas = tf.paragraphs
        text = ''.join(p.text for p in pas).replace('\x0b', '')
        pt = max([r.font.size.pt for p in pas for r in p.runs if r.font.size] or [22])
        avail = sh.width/EMU - 0.20
        has_br = any(p._p.findall(NS+'br') for p in pas) or len(pas) > 1
        if em_len(text)*pt/72.0 <= avail:                      # 1行
            if sh.height/EMU != 0.50: sh.height = Emu(int(0.50*EMU)); log.append(f"P{n} リード文 枠高→0.50in（1行）")
            continue
        # 2行に整える
        sh.height = Emu(int(0.86*EMU))
        if has_br: continue                                    # すでに改行指定あり
        i = best_split(text, avail, pt)
        if not i: continue
        pa = pas[0]
        runs = pa.runs
        if not runs: continue
        keep = runs[0]
        for r in runs[1:]: r._r.getparent().remove(r._r)
        keep.text = text[:i]
        pa._p.append(keep._r.makeelement(NS+'br', {}))
        nr = copy.deepcopy(keep._r)
        for t in nr.findall(NS+'t'): t.text = text[i:]
        pa._p.append(nr)
        for p in list(tf.paragraphs)[1:]: p._p.getparent().remove(p._p)
        log.append(f"P{n} リード文 枠高→0.86in・改行位置を均等配分「{text[:i]}／{text[i:]}」")

prs.save('final_fix.pptx')
print("\n".join("・"+x for x in log))
print(f"--- {len(log)}件 → final_fix.pptx")
