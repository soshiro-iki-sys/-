# -*- coding: utf-8 -*-
"""P22 の累計行を 10年間 → 20年間 に変更する"""
from pptx import Presentation
from measure import flat

prs = Presentation('final_fix2.pptx')
NEW = ["20年間の累計", "約186万円", "約258万円"]

def set_cell(cell, text):
    pa = cell.text_frame.paragraphs[0]
    runs = pa.runs
    if not runs:
        pa.text = text; return
    runs[0].text = text
    for r in runs[1:]: r._r.getparent().remove(r._r)
    for p in list(cell.text_frame.paragraphs)[1:]:
        p._p.getparent().remove(p._p)

o = []
for sh in prs.slides[21].shapes: flat(sh, o)
for sh in o:
    if sh.has_table:
        row = sh.table.rows[5]
        before = [c.text for c in row.cells]
        for c, t in zip(row.cells, NEW): set_cell(c, t)
        print("before:", " | ".join(before))
        print("after :", " | ".join(c.text for c in row.cells))

prs.save('final_fix3.pptx')

# ── P22 吹き出しの見出しが2行に割れていたので枠を広げて18ptに ──────
from pptx.util import Pt, Emu
from measure import EMU as _E
prs2 = Presentation('final_fix3.pptx')
o2 = []
for sh in prs2.slides[21].shapes: flat(sh, o2)
for sh in o2:
    if sh.shape_id in (14, 15, 16):
        sh.left = Emu(int(7.70*_E)); sh.width = Emu(int(3.75*_E))
    if sh.shape_id == 14:
        for pa in sh.text_frame.paragraphs:
            for r in pa.runs: r.font.size = Pt(18)
prs2.save('final_fix3.pptx')
print("P22 吹き出し見出し 20→18pt・枠幅3.75inに拡張（折り返し解消）")
