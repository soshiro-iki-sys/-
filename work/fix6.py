# -*- coding: utf-8 -*-
"""長岡市アンケート（実測データ）ページの追加と、発電量ページの実測ベース化"""
import copy
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from measure import flat, em_len, EMU, NS

NAVY = RGBColor(0x00, 0x20, 0x60)
INK  = RGBColor(0x40, 0x40, 0x40)
FONT = 'Meiryo UI'
prs = Presentation('withpage.pptx')
log = []

def shapes(n):
    o = []
    for sh in prs.slides[n-1].shapes: flat(sh, o)
    return {sh.shape_id: sh for sh in o}

def set_runs(sh, texts, para=0):
    """段落内のランに1対1で流し込む（各ランの書式をそのまま保つ）"""
    pa = sh.text_frame.paragraphs[para]
    runs = pa.runs
    for r, t in zip(runs, texts): r.text = t
    for r in runs[len(texts):]: r._r.getparent().remove(r._r)

def set_lines(sh, lines):
    """段落1つ＋<a:br> で書き換える（先頭ランの書式を保つ）"""
    tf = sh.text_frame
    pa = tf.paragraphs[0]
    for br in pa._p.findall(NS+'br'): pa._p.remove(br)
    runs = pa.runs
    if not runs:
        tf.text = "\n".join(lines); return
    keep = runs[0]
    for r in runs[1:]: r._r.getparent().remove(r._r)
    keep.text = lines[0]
    for extra in lines[1:]:
        pa._p.append(keep._r.makeelement(NS+'br', {}))
        nr = copy.deepcopy(keep._r)
        for t in nr.findall(NS+'t'): t.text = extra
        pa._p.append(nr)
    for p in list(tf.paragraphs)[1:]: p._p.getparent().remove(p._p)

# ══ 新P20：長岡市アンケートの実測データ ═══════════════════════
s20 = prs.slides[19]
for sh in list(s20.shapes):                     # 複製で共有されたグラフを外す
    if sh.has_chart:
        el = sh._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/chart}chart')
        if el is not None:
            rId = el.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            s20.part.rels._rels.pop(rId, None)
        sh._element.getparent().remove(sh._element)

cd = CategoryChartData()
cd.categories = ['1月','2月','3月','4月','5月','6月','7月','8月','9月','10月','11月','12月']
cd.add_series('1kW当たりの発電量（kWh／月）',
              (14.2, 21.5, 49.2, 104.6, 118.0, 106.4, 97.6, 86.1, 90.7, 78.4, 45.0, 15.7))
ch = s20.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                          Inches(0.42), Inches(2.10), Inches(7.30), Inches(4.55), cd).chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = '長岡市 1kWあたりの月別発電量（kWh／月・回答平均）'
r = ch.chart_title.text_frame.paragraphs[0].runs[0]
r.font.size, r.font.bold, r.font.name = Pt(14), True, FONT
r.font.color.rgb = NAVY
pl = ch.plots[0]
pl.gap_width = 60
pl.has_data_labels = True
pl.data_labels.font.size = Pt(9)
pl.data_labels.font.name = FONT
pl.data_labels.font.color.rgb = INK
pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
pl.series[0].format.fill.solid()
pl.series[0].format.fill.fore_color.rgb = NAVY
for ax in (ch.category_axis, ch.value_axis):
    ax.tick_labels.font.size = Pt(11)
    ax.tick_labels.font.name = FONT
    ax.tick_labels.font.color.rgb = INK
ch.value_axis.minimum_scale = 0.0
ch.value_axis.maximum_scale = 130.0

d = shapes(20)
set_runs(d[6], ['長岡市の実測データ'])
set_lines(d[7], ['長岡市が、補助金を使って太陽光を設置した市民に行ったアンケートの結果である。'])
set_lines(d[8], ['出典：長岡市「太陽光発電設備設置に関するアンケート調査の結果について」'
                 '（平成27年1〜2月実施・回答65件／回収率56.0％）'])
set_runs(d[10], ['長岡で実際に', '設置した方の実測データです'])
set_runs(d[13], ['アンケートで分かったこと'])
set_runs(d[15], ['1kW当たり　', '827.4', ' kWh／年'])
for i, t in enumerate(['平均出力　4.88kW', '年間発電量　約4,130kWh']):
    set_runs(d[15], [t], para=2+i)
    for rr in d[15].text_frame.paragraphs[2+i].runs: rr.font.size = Pt(18)
set_lines(d[16], ['雪による被害は「ない」が87.7％。',
                  '1月でも5月（ピーク）の',
                  '約12％を発電しています。'])
log.append('新P20「長岡市の実測データ」を追加（月別発電量グラフ／平均出力4.88kW・827.4kWh/kW・年／雪害「ない」87.7％）')

# ══ P19：社内メモが残っていた出典を整える ══════════════════════
set_lines(shapes(19)[8], ['出典：NEDO日射量データベースをもとにした都道府県別の年間発電量試算値より作成'
                          '（標準的な設置条件での試算値）'])
log.append('P19 出典から社内メモ「※前提条件は要確認」を削除し、試算値である旨を明記')

# ══ P22（旧P21）：長岡市アンケートの実測値ベースに ═══════════════
d = shapes(22)
set_lines(d[7], ['長岡市の実測データでは、平均的な設備で家庭の電気使用量の多くを賄えている。'])
set_lines(d[8], ['※発電量は長岡市アンケート（平均出力4.88kW・回答65件）の実測平均／'
                 '電気使用量は資源エネルギー庁「平成24年度エネルギー使用合理化基盤整備事業」報告書'])
set_runs(d[10], ['消費電力の93％', 'をカバーすることが可能に'])
set_runs(d[11], ['長岡市での実測：平均4.88kWシステムの場合'])
set_runs(d[14], ['4.88kW × 827kWh'])
set_runs(d[15], ['約4,130', ' kWh／年'])
set_runs(d[17], ['1年間の電気使用量（全国1世帯当たり）'])
set_runs(d[18], ['4,432', ' kWh／年'])
set_runs(d[22], ['93', '％'])
set_runs(d[23], ['4,130kWh ÷ 4,432kWh'])
log.append('P22 発電量試算を長岡市アンケートの実測値ベースに変更'
           '（4.88kW×827kWh＝約4,130kWh／全国1世帯4,432kWhに対しカバー率84％→93％）')

# ══ P23（旧P22）：経済試算は据え置き、前提の違いを明記 ══════════
set_lines(shapes(23)[8],
          ['前提：5kW・発電5,500kWh（JIS C 8907の試算値。前ページの実測平均とは前提が異なる）'
           '／買電30円/kWh／売電 当初4年24円・5年目以降8.3円'])
log.append('P23 経済試算の前提行に、前ページの実測平均とは前提が異なる旨を明記')

prs.save('final_v4.pptx')
print("\n".join("・"+x for x in log))
