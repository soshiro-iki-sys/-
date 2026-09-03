# -*- coding: utf-8 -*-
"""新P20の右カードの収まりを整え、リード文を1行に収める"""
from pptx import Presentation
from pptx.util import Pt, Emu
from measure import flat, em_len, EMU, NS

prs = Presentation('final_v4.pptx')
log = []

def shapes(n):
    o = []
    for sh in prs.slides[n-1].shapes: flat(sh, o)
    return {sh.shape_id: sh for sh in o}

d = shapes(20)
# リード文を1行に収まる長さに
pa = d[7].text_frame.paragraphs[0]
for br in pa._p.findall(NS+'br'): pa._p.remove(br)
pa.runs[0].text = '補助金を使って太陽光を設置した市民に、長岡市がアンケートを実施した。'
for r in pa.runs[1:]: r._r.getparent().remove(r._r)
for p in list(d[7].text_frame.paragraphs)[1:]: p._p.getparent().remove(p._p)
d[7].height = Emu(int(0.50*EMU))
log.append('新P20 リード文を1行に収まる文言に調整')

# 数値2行を16ptにして折り返しを解消し、注記を下げて重なりを解消
for i in (2, 3):
    for r in d[15].text_frame.paragraphs[i].runs: r.font.size = Pt(16)
d[15].height = Emu(int(2.10*EMU))
d[16].top = Emu(int(5.22*EMU))
d[16].height = Emu(int(1.36*EMU))
log.append('新P20 右カードの数値2行を18→16pt・注記をy5.22inへ（折り返しと重なりを解消）')

prs.save('final_v5.pptx')
print("\n".join("・"+x for x in log))

# ── 注記の位置をさらに下げ、空行を詰める（重なりの最終調整） ──────
prs2 = Presentation('final_v5.pptx')
o2 = []
for sh in prs2.slides[19].shapes: flat(sh, o2)
for sh in o2:
    if sh.shape_id == 15:
        blank = sh.text_frame.paragraphs[1]
        blank.font.size = Pt(8)
    if sh.shape_id == 16:
        sh.top = Emu(int(5.30*EMU)); sh.height = Emu(int(1.25*EMU))
        pa2 = sh.text_frame.paragraphs[0]
        for br in pa2._p.findall(NS+'br'): pa2._p.remove(br)
        keep = pa2.runs[0]
        for r in pa2.runs[1:]: r._r.getparent().remove(r._r)
        keep.text = '雪による被害は'
        keep.font.size = Pt(15)
        import copy as _c
        for extra in ['「ない」が87.7％。', '1月でも5月（ピーク）の', '約12％を発電しています。']:
            pa2._p.append(keep._r.makeelement(NS+'br', {}))
            nr = _c.deepcopy(keep._r)
            for t in nr.findall(NS+'t'): t.text = extra
            pa2._p.append(nr)
prs2.save('final_v5.pptx')
print('・新P20 右カードの空行を8ptに詰め、注記を15pt・4行に整えてy5.30inへ（重なりを解消）')
